"""
Modive MarTech Slide Modifier
Applies comment-based changes to modive_martech_v1.pptx
"""
import os, io, urllib.request, copy, sys
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

PPTX_IN  = "C:/Agent/pepper/output/modive_martech_v1.pptx"
PPTX_OUT = "C:/Agent/pepper/output/modive_martech_v2.pptx"

IMG_SLIDE5 = ("C:/Agent/pepper/output/img_s5.jpg",
    "https://images.unsplash.com/photo-1714846201700-35b42d937158"
    "?ixid=M3w5MzA1MTF8MHwxfGFsbHx8fHx8fHx8fDE3NzY4MTY5NzR8"
    "&ixlib=rb-4.1.0&w=1200&h=1400&fit=crop&q=90")

IMG_SLIDE6 = ("C:/Agent/pepper/output/img_s6.jpg",
    "https://images.unsplash.com/photo-1558655146-6c222b05fce4"
    "?ixid=M3w5MzA1MTF8MHwxfGFsbHx8fHx8fHx8fDE3NzY4MTY5NzV8"
    "&ixlib=rb-4.1.0&w=1200&h=1400&fit=crop&q=90")

IMG_SLIDE7 = ("C:/Agent/pepper/output/img_s7.jpg",
    "https://images.unsplash.com/photo-1581090690925-3898802525e2"
    "?ixid=M3w5MzA1MTF8MHwxfGFsbHx8fHx8fHx8fDE3NzY4MTY5NzZ8"
    "&ixlib=rb-4.1.0&w=1260&h=1400&fit=crop&q=90")


# ── helpers ──────────────────────────────────────────────────────────────────

def dl(url, path):
    req = urllib.request.Request(url, headers={"User-Agent": "Mozilla/5.0"})
    with urllib.request.urlopen(req, timeout=45) as r:
        data = r.read()
    with open(path, "wb") as f:
        f.write(data)
    print(f"  ✓ downloaded {os.path.basename(path)} ({len(data)//1024} KB)")
    return path


def set_run_text(run_elem, text):
    """Set text of an <a:r> element."""
    t = run_elem.find(qn("a:t"))
    if t is None:
        t = etree.SubElement(run_elem, qn("a:t"))
    t.text = text


def replace_shape_text(shape, new_text):
    """
    Replace all text in a shape with new_text, keeping the first run's
    character formatting (font/color/size).  All paragraphs after the first
    are removed; line-breaks inside the paragraph are removed.
    """
    if not shape.has_text_frame:
        return
    txBody = shape.text_frame._txBody
    paras  = txBody.findall(qn("a:p"))

    # collect first run element (for formatting reference)
    first_run_elem = None
    for p in paras:
        for r in p.findall(qn("a:r")):
            first_run_elem = r
            break
        if first_run_elem is not None:
            break

    # remove extra paragraphs
    for p in paras[1:]:
        txBody.remove(p)

    # clean first paragraph: remove all <a:r> and <a:br>
    first_p = paras[0]
    for child in list(first_p):
        if child.tag in (qn("a:r"), qn("a:br")):
            first_p.remove(child)

    # build new run
    if first_run_elem is not None:
        new_r = copy.deepcopy(first_run_elem)
        set_run_text(new_r, new_text)
    else:
        new_r = etree.SubElement(first_p, qn("a:r"))
        t = etree.SubElement(new_r, qn("a:t"))
        t.text = new_text
        new_r = None  # already appended via SubElement

    if new_r is not None:
        first_p.append(new_r)


def replace_run_substring(shape, old, new):
    """Find and replace a substring inside any run in the shape."""
    if not shape.has_text_frame:
        return False
    changed = False
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            if old in run.text:
                run.text = run.text.replace(old, new)
                changed = True
    return changed


def replace_picture(slide, shape_idx, img_path):
    """
    Replace the picture at shape_idx with img_path,
    preserving position, size, and z-order.
    """
    shapes  = slide.shapes
    old_shp = shapes[shape_idx]
    left, top, width, height = old_shp.left, old_shp.top, old_shp.width, old_shp.height
    spTree  = shapes._spTree
    old_elem = old_shp._element
    old_pos  = list(spTree).index(old_elem)

    # Add new picture (appended at end of spTree)
    new_pic = shapes.add_picture(img_path, left, top, width, height)
    new_elem = new_pic._element

    # Remove old, move new to same position
    spTree.remove(old_elem)
    spTree.remove(new_elem)
    spTree.insert(old_pos, new_elem)
    print(f"  ✓ replaced picture at shape[{shape_idx}]")


def add_pain_points_overlay(slide, bg_left, bg_top, bg_w, bg_h):
    """
    Remove the pain-points image and add a dark-background text layout
    showing 6 pain points in a 3×2 grid.
    """
    pain_points = [
        ("데이터 사일로", "기기·플랫폼별 분산 데이터\n통합 인사이트 부재"),
        ("수동 캠페인 운영", "반복 마케팅 작업 수동 처리\n실시간 대응 지연"),
        ("개인화 부족", "일괄 발송 캠페인\n구독 전환율 저조"),
        ("이탈 관리 미비", "이탈 징후 실시간 감지·\n개입 체계 부재"),
        ("성과 측정 공백", "채널별 ROI 분석 도구 부재\n리포팅 지연"),
        ("비용 비효율", "벤더 의존 구조\n높은 비용·낮은 커스터마이징"),
    ]

    cols, rows = 3, 2
    pad  = Emu(180000)
    cell_w = (bg_w - pad * (cols + 1)) // cols
    cell_h = (bg_h - pad * (rows + 1)) // rows

    shapes = slide.shapes
    for i, (title, body) in enumerate(pain_points):
        col = i % cols
        row = i // cols
        x = bg_left + pad + col * (cell_w + pad)
        y = bg_top  + pad + row * (cell_h + pad)

        # dark card background
        box = shapes.add_shape(1, x, y, cell_w, cell_h)   # MSO_SHAPE_TYPE.RECTANGLE = 1
        fill = box.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(0x0a, 0x1a, 0x2e)
        box.line.color.rgb  = RGBColor(0x1e, 0x5f, 0x9a)
        box.line.width = Emu(12700)

        # title text box
        title_h = Emu(420000)
        tb_t = shapes.add_textbox(x + Emu(80000), y + Emu(80000),
                                   cell_w - Emu(160000), title_h)
        tf = tb_t.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = title
        run.font.bold  = True
        run.font.size  = Pt(13)
        run.font.color.rgb = RGBColor(0x4d, 0xa6, 0xff)

        # body text box
        tb_b = shapes.add_textbox(x + Emu(80000),
                                   y + Emu(80000) + title_h + Emu(40000),
                                   cell_w - Emu(160000),
                                   cell_h - title_h - Emu(200000))
        tf2 = tb_b.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.alignment = PP_ALIGN.LEFT
        run2 = p2.add_run()
        run2.text = body
        run2.font.size  = Pt(11)
        run2.font.color.rgb = RGBColor(0xcc, 0xdd, 0xee)

    print("  ✓ added pain-points text layout (6 cards)")


# ── main ─────────────────────────────────────────────────────────────────────

print("Loading presentation …")
prs = Presentation(PPTX_IN)

# ── download images ──────────────────────────────────────────────────────────
print("\n[Step 1] Downloading Unsplash images …")
dl(IMG_SLIDE5[1], IMG_SLIDE5[0])
dl(IMG_SLIDE6[1], IMG_SLIDE6[0])
dl(IMG_SLIDE7[1], IMG_SLIDE7[0])

# ── slide 2 (index 1): partner text ─────────────────────────────────────────
print("\n[Step 2] Slide 2 — partner text")
s2 = prs.slides[1]
replace_shape_text(s2.shapes[30],
    "Adobe 공식 파트너 · Amplitude 공식 파트너 · GA4 10년 파트너")
print("  ✓ Shape[30] updated")

# ── slide 3 (index 2): pain points image → text layout ───────────────────────
print("\n[Step 3] Slide 3 — pain points overlay")
s3 = prs.slides[2]
img_shape = s3.shapes[9]   # original pain-points image
bg_left  = img_shape.left
bg_top   = img_shape.top
bg_w     = img_shape.width
bg_h     = img_shape.height
# remove original image
spTree3 = s3.shapes._spTree
spTree3.remove(img_shape._element)
print("  ✓ removed original pain-points image")
add_pain_points_overlay(s3, bg_left, bg_top, bg_w, bg_h)

# ── slide 4 (index 3): solution text ─────────────────────────────────────────
print("\n[Step 4] Slide 4 — solution text")
s4 = prs.slides[3]
replace_shape_text(s4.shapes[7], "핵심 서비스 3가지")
print("  ✓ Shape[7] '핵심 솔루션' → '핵심 서비스'")
ok = replace_run_substring(s4.shapes[20], "CleverTap 도입/운영", "자동화 툴 도입/운영")
print(f"  {'✓' if ok else '⚠'} Shape[20] CleverTap → 자동화 툴")

# ── slide 5 (index 4): caption + image ───────────────────────────────────────
print("\n[Step 5] Slide 5 — caption + image")
s5 = prs.slides[4]
replace_shape_text(s5.shapes[13], "데이터 분석 대시보드 (실제 서비스 화면)")
print("  ✓ Shape[13] caption updated")
replace_picture(s5, 10, IMG_SLIDE5[0])

# ── slide 6 (index 5): caption + section title + image ───────────────────────
print("\n[Step 6] Slide 6 — caption + section title + image")
s6 = prs.slides[5]
replace_shape_text(s6.shapes[13], "마케팅 자동화 플랫폼 — 실제 서비스 화면")
print("  ✓ Shape[13] caption updated")
replace_shape_text(s6.shapes[16], "멀티채널 마케팅 자동화")
print("  ✓ Shape[16] section title updated")
replace_picture(s6, 10, IMG_SLIDE6[0])

# ── slide 7 (index 6): caption + image ───────────────────────────────────────
print("\n[Step 7] Slide 7 — caption + image")
s7 = prs.slides[6]
replace_shape_text(s7.shapes[13], "파트너사 레퍼런스")
print("  ✓ Shape[13] caption updated")
replace_picture(s7, 10, IMG_SLIDE7[0])

# ── save ─────────────────────────────────────────────────────────────────────
print(f"\n[Step 8] Saving → {PPTX_OUT}")
prs.save(PPTX_OUT)
print("  ✓ Saved successfully")
print("\nAll done.")
