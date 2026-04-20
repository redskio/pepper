# -*- coding: utf-8 -*-
"""
MBC 모다이브 × 위어드섹터 프로젝트 제안서 v2
- 글자 크기: 제목 34-40pt / 소제목 22-26pt / 본문 16-18pt
- 컬러: MBC Blue #0066CC + 다크 네이비 #0D1B2A + 화이트
- 배경: 다크 그라데이션 레이어 + 블루 액센트
- 12슬라이드 고퀄 레이아웃
"""
import sys, io
if sys.platform == 'win32':
    sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pathlib import Path

OUT = Path("C:/Agent/pepper/output/modive-weirdsector-proposal-v2.pptx")

# ── 컬러 팔레트 ─────────────────────────────────────
C_BLUE        = RGBColor(0x00, 0x66, 0xCC)   # MBC 블루
C_BLUE_LIGHT  = RGBColor(0x00, 0xAA, 0xFF)   # 라이트 블루
C_BLUE_DARK   = RGBColor(0x00, 0x44, 0x99)   # 다크 블루
C_NAVY        = RGBColor(0x0D, 0x1B, 0x2A)   # 딥 네이비 배경
C_NAVY2       = RGBColor(0x07, 0x10, 0x18)   # 더 어두운 네이비
C_CARD        = RGBColor(0x13, 0x26, 0x3D)   # 카드 배경
C_CARD2       = RGBColor(0x1A, 0x30, 0x4A)   # 카드 배경 2
C_WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
C_GRAY        = RGBColor(0x8A, 0x9B, 0xB0)
C_GRAY_LIGHT  = RGBColor(0xD0, 0xDC, 0xEC)
C_CYAN        = RGBColor(0x00, 0xE5, 0xFF)
C_GREEN       = RGBColor(0x00, 0xD4, 0x8A)
C_ORANGE      = RGBColor(0xFF, 0x7A, 0x2F)
C_PURPLE      = RGBColor(0x7B, 0x6C, 0xF6)

SW = Inches(13.33)
SH = Inches(7.5)

def new_prs():
    prs = Presentation()
    prs.slide_width  = SW
    prs.slide_height = SH
    return prs

def blank(prs):
    return prs.slide_layouts[6]

def rect(slide, x, y, w, h, fill=None):
    sh = slide.shapes.add_shape(1, x, y, w, h)
    sh.line.fill.background()
    if fill:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    else:
        sh.fill.background()
    return sh

def rrect(slide, x, y, w, h, fill=None, radius=0.05):
    sp = slide.shapes.add_shape(5, x, y, w, h)
    sp.line.fill.background()
    if fill:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    else:
        sp.fill.background()
    try:
        sp.adjustments[0] = radius
    except:
        pass
    return sp

def txt(slide, text, x, y, w, h, sz=18, bold=False, color=None,
        align=PP_ALIGN.LEFT, italic=False, wrap=True):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p  = tf.paragraphs[0]
    p.alignment = align
    r  = p.add_run()
    r.text  = text
    r.font.size   = Pt(sz)
    r.font.bold   = bold
    r.font.italic = italic
    r.font.name   = "Malgun Gothic"
    if color:
        r.font.color.rgb = color
    return tb

def header(slide, num, total, title, sub=""):
    """슬라이드 상단 헤더"""
    rect(slide, 0, 0, SW, Inches(1.25), fill=C_NAVY2)
    rect(slide, 0, Inches(1.25), SW, Inches(0.05), fill=C_BLUE)
    # 블루 좌측 수직 액센트
    rect(slide, 0, 0, Inches(0.07), Inches(1.25), fill=C_BLUE)
    txt(slide, title,
        Inches(0.3), Inches(0.1), Inches(10.0), Inches(0.75),
        sz=30, bold=True, color=C_WHITE)
    if sub:
        txt(slide, sub,
            Inches(0.3), Inches(0.78), Inches(10.5), Inches(0.42),
            sz=15, color=C_GRAY)
    txt(slide, f"{num:02d}/{total:02d}",
        Inches(11.8), Inches(0.5), Inches(1.4), Inches(0.4),
        sz=13, color=C_BLUE, align=PP_ALIGN.RIGHT)
    txt(slide, "WeirdSector",
        Inches(11.5), Inches(0.08), Inches(1.7), Inches(0.38),
        sz=12, bold=True, color=C_BLUE_LIGHT, align=PP_ALIGN.RIGHT)

def footer(slide):
    """하단 풋터"""
    rect(slide, 0, Inches(7.2), SW, Inches(0.3), fill=C_NAVY2)
    txt(slide, "위어드섹터(WeirdSector)  |  weirdsector.co.kr  |  Confidential",
        Inches(0.5), Inches(7.22), Inches(12.5), Inches(0.25),
        sz=10, color=C_GRAY, align=PP_ALIGN.CENTER)

def kpi_card(slide, x, y, value, label, desc="", color=C_BLUE):
    W, H = Inches(2.8), Inches(1.7)
    rrect(slide, x, y, W, H, fill=C_CARD)
    rect(slide, x, y, Inches(0.07), H, fill=color)
    txt(slide, value,
        x+Inches(0.2), y+Inches(0.08), W-Inches(0.3), Inches(0.72),
        sz=34, bold=True, color=color)
    txt(slide, label,
        x+Inches(0.2), y+Inches(0.78), W-Inches(0.3), Inches(0.42),
        sz=15, bold=True, color=C_WHITE)
    if desc:
        txt(slide, desc,
            x+Inches(0.2), y+Inches(1.22), W-Inches(0.3), Inches(0.38),
            sz=12, color=C_GRAY)

def service_card(slide, x, y, icon, title, points, color=C_BLUE):
    W, H = Inches(5.8), Inches(2.3)
    rrect(slide, x, y, W, H, fill=C_CARD)
    rect(slide, x, y, W, Inches(0.06), fill=color)
    txt(slide, icon, x+Inches(0.2), y+Inches(0.12), Inches(0.7), Inches(0.65), sz=30)
    txt(slide, title,
        x+Inches(0.95), y+Inches(0.15), W-Inches(1.1), Inches(0.5),
        sz=19, bold=True, color=C_WHITE)
    for i, pt in enumerate(points):
        txt(slide, f"▸  {pt}",
            x+Inches(0.2), y+Inches(0.78 + i*0.44), W-Inches(0.35), Inches(0.4),
            sz=14, color=C_GRAY_LIGHT if i == 0 else C_GRAY)

def bullet(slide, x, y, text, sz=16, color=C_WHITE, accent=C_BLUE):
    txt(slide, "▶", x, y, Inches(0.35), Inches(0.42),
        sz=sz-3, bold=True, color=accent)
    txt(slide, text, x+Inches(0.38), y, Inches(5.5), Inches(0.42),
        sz=sz, color=color)


# ══════════════════════════════════════════════════
def build():
    prs = new_prs()
    TOTAL = 12

    # ── Slide 1: 표지 ──────────────────────────────
    sl = prs.slides.add_slide(blank(prs))
    rect(sl, 0, 0, SW, SH, fill=C_NAVY2)
    # 블루 그라데이션 우측 패널
    rect(sl, Inches(8.2), 0, Inches(5.13), SH, fill=C_BLUE_DARK)
    rect(sl, Inches(7.5), 0, Inches(1.0), SH, fill=C_CARD)
    # 상단 블루 라인
    rect(sl, 0, 0, SW, Inches(0.07), fill=C_BLUE)
    # 하단 라인
    rect(sl, 0, Inches(7.43), SW, Inches(0.07), fill=C_BLUE)

    # 좌측 — 제안사 정보
    txt(sl, "PROJECT PROPOSAL",
        Inches(0.6), Inches(0.8), Inches(7.0), Inches(0.5),
        sz=13, color=C_BLUE_LIGHT, bold=True)
    txt(sl, "위어드섹터 ×\nMBC 모다이브",
        Inches(0.6), Inches(1.4), Inches(7.5), Inches(2.0),
        sz=42, bold=True, color=C_WHITE)
    txt(sl, "데이터 · 코드 · MarTech · CRM\n통합 운영 파트너십 제안서",
        Inches(0.6), Inches(3.5), Inches(7.2), Inches(1.1),
        sz=21, color=C_BLUE_LIGHT)

    # 구분선
    rect(sl, Inches(0.6), Inches(4.8), Inches(6.5), Inches(0.04), fill=C_BLUE)

    # 제안사/고객사 정보
    txt(sl, "제안사",  Inches(0.6), Inches(5.0), Inches(2.0), Inches(0.4),
        sz=13, color=C_GRAY)
    txt(sl, "위어드섹터 (WeirdSector)",
        Inches(0.6), Inches(5.4), Inches(5.0), Inches(0.45),
        sz=17, bold=True, color=C_WHITE)
    txt(sl, "고객사",  Inches(0.6), Inches(5.95), Inches(2.0), Inches(0.4),
        sz=13, color=C_GRAY)
    txt(sl, "MBC 모다이브",
        Inches(0.6), Inches(6.35), Inches(5.0), Inches(0.45),
        sz=17, bold=True, color=C_WHITE)
    txt(sl, "2026. 04",
        Inches(0.6), Inches(6.9), Inches(3.0), Inches(0.4),
        sz=14, color=C_GRAY)

    # 우측 패널 — 핵심 수치
    txt(sl, "프로젝트 개요",
        Inches(8.5), Inches(1.2), Inches(4.5), Inches(0.5),
        sz=16, bold=True, color=C_BLUE_LIGHT, align=PP_ALIGN.CENTER)
    for i, (v, l) in enumerate([
        ("월 200만원", "M/M 예산"),
        ("4가지",      "핵심 서비스"),
        ("즉시",       "시작 가능"),
        ("전담팀",     "전문가 배정"),
    ]):
        ry = Inches(2.0 + i*1.2)
        rrect(sl, Inches(8.5), ry, Inches(4.5), Inches(1.05), fill=C_CARD)
        rect(sl, Inches(8.5), ry, Inches(0.06), Inches(1.05), fill=C_CYAN)
        txt(sl, v, Inches(8.7), ry+Inches(0.08), Inches(4.0), Inches(0.52),
            sz=26, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        txt(sl, l, Inches(8.7), ry+Inches(0.6), Inches(4.0), Inches(0.38),
            sz=14, color=C_GRAY, align=PP_ALIGN.CENTER)
    print("  [1] 표지")

    # ── Slide 2: Executive Summary ─────────────────
    sl = prs.slides.add_slide(blank(prs))
    rect(sl, 0, 0, SW, SH, fill=C_NAVY)
    header(sl, 2, TOTAL, "Executive Summary", "한 눈에 보는 제안 핵심")
    footer(sl)

    txt(sl, "이 제안서 한 장으로 이해하는 위어드섹터의 가치",
        Inches(0.6), Inches(1.45), Inches(10.0), Inches(0.55),
        sz=22, bold=True, color=C_BLUE_LIGHT)

    # 3개 핵심 메시지 카드
    msgs = [
        ("🎯", "왜 지금인가?",
         "MBC 모다이브가 빠르게 성장하는 OTT 경쟁 시장에서 살아남으려면\n데이터 기반 마케팅과 정교한 CRM이 필수입니다."),
        ("🔧", "무엇을 제안하나?",
         "데이터 관리 · 코드 관리 · MarTech 운영 · CRM 수행\n4가지 서비스를 월 200만원으로 통합 제공합니다."),
        ("📈", "기대 효과는?",
         "사용자 이탈률 감소, 재방문율 증가, 캠페인 효율 향상,\n데이터 기반 콘텐츠 추천 고도화를 실현합니다."),
    ]
    for i, (ic, ttl, body) in enumerate(msgs):
        bx = Inches(0.4 + i*4.28)
        rrect(sl, bx, Inches(2.1), Inches(4.05), Inches(4.6), fill=C_CARD)
        rect(sl, bx, Inches(2.1), Inches(4.05), Inches(0.06), fill=C_BLUE)
        txt(sl, ic, bx+Inches(0.2), Inches(2.18), Inches(3.6), Inches(0.75),
            sz=36, align=PP_ALIGN.CENTER)
        txt(sl, ttl, bx+Inches(0.2), Inches(2.95), Inches(3.6), Inches(0.55),
            sz=20, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        rect(sl, bx+Inches(0.5), Inches(3.55), Inches(3.0), Inches(0.04), fill=C_BLUE)
        txt(sl, body, bx+Inches(0.2), Inches(3.7), Inches(3.6), Inches(2.5),
            sz=15, color=C_GRAY_LIGHT)
    print("  [2] Executive Summary")

    # ── Slide 3: MBC 모다이브 현황 ──────────────────
    sl = prs.slides.add_slide(blank(prs))
    rect(sl, 0, 0, SW, SH, fill=C_NAVY)
    header(sl, 3, TOTAL, "MBC 모다이브의 현재 과제", "성장하는 OTT 시장에서 해결해야 할 문제들")
    footer(sl)

    txt(sl, "이미 좋은 콘텐츠가 있습니다. 문제는 마케팅과 데이터입니다.",
        Inches(0.6), Inches(1.45), Inches(12.0), Inches(0.55),
        sz=20, bold=True, color=C_BLUE_LIGHT)

    issues = [
        ("📊", "데이터 분산", "시청 데이터, 유입 데이터, CRM 데이터가 각각 분리되어 통합 분석이 불가합니다."),
        ("🔀", "채널 단절", "Push, 이메일, SMS, 인앱 메시지가 개별로 운영되어 일관된 사용자 경험이 없습니다."),
        ("🎯", "개인화 부재", "모든 사용자에게 동일한 콘텐츠 추천 — 높은 이탈률의 핵심 원인입니다."),
        ("⚙", "기술 부채", "레거시 코드와 트래킹 오류로 정확한 데이터 수집 자체가 어려운 상황입니다."),
        ("📉", "마케팅 비효율", "ROI 측정 없는 광고 집행, 캠페인 효과 검증 불가로 예산이 낭비됩니다."),
        ("🏃", "빠른 경쟁자", "넷플릭스·티빙·웨이브의 AI 개인화 속도에 비해 대응이 늦어지고 있습니다."),
    ]
    for i, (ic, ttl, body) in enumerate(issues):
        row, col = divmod(i, 3)
        bx = Inches(0.4 + col*4.3)
        ry = Inches(2.1 + row*2.35)
        rrect(sl, bx, ry, Inches(4.1), Inches(2.18), fill=C_CARD)
        rect(sl, bx, ry, Inches(4.1), Inches(0.05), fill=C_ORANGE if i<3 else C_BLUE)
        txt(sl, ic, bx+Inches(0.15), ry+Inches(0.1), Inches(0.65), Inches(0.65), sz=26)
        txt(sl, ttl, bx+Inches(0.85), ry+Inches(0.15), Inches(3.0), Inches(0.45),
            sz=17, bold=True, color=C_WHITE)
        txt(sl, body, bx+Inches(0.15), ry+Inches(0.72), Inches(3.8), Inches(1.3),
            sz=13, color=C_GRAY)
    print("  [3] 모다이브 과제")

    # ── Slide 4: 위어드섹터 소개 ───────────────────
    sl = prs.slides.add_slide(blank(prs))
    rect(sl, 0, 0, SW, SH, fill=C_NAVY)
    header(sl, 4, TOTAL, "위어드섹터 소개", "MarTech × 데이터 × CRM 전문 스타트업")
    footer(sl)

    # 좌측 — 회사 소개
    rrect(sl, Inches(0.4), Inches(1.45), Inches(6.3), Inches(5.7), fill=C_CARD)
    rect(sl, Inches(0.4), Inches(1.45), Inches(0.07), Inches(5.7), fill=C_BLUE)

    txt(sl, "WeirdSector",
        Inches(0.65), Inches(1.6), Inches(5.5), Inches(0.7),
        sz=30, bold=True, color=C_BLUE_LIGHT)
    txt(sl, "데이터 기반 마케팅의 실행 파트너",
        Inches(0.65), Inches(2.3), Inches(5.5), Inches(0.45),
        sz=17, color=C_GRAY_LIGHT)
    rect(sl, Inches(0.65), Inches(2.85), Inches(5.7), Inches(0.04), fill=C_BLUE)

    ws_intro = [
        "마테크(MarTech), 데이터 분석, CRM 자동화 전문",
        "고객사의 기술 스택에 맞춘 맞춤형 솔루션 설계",
        "기획 → 구현 → 운영 → 성과 측정 Full-cycle 수행",
        "소수 정예 시니어 팀 — 속도와 품질 모두 보장",
    ]
    for i, line in enumerate(ws_intro):
        bullet(sl, Inches(0.65), Inches(3.05 + i*0.62), line, sz=15)

    # 우측 — 역량 카드 4개
    skills = [
        ("📊", "데이터 엔지니어링", "GA4 · BigQuery · Segment\n데이터 파이프라인 구축 및 운영"),
        ("🔧", "MarTech 스택 운영", "Braze · CleverTap · Mixpanel\n플랫폼 세팅부터 캠페인까지"),
        ("💌", "CRM 자동화", "이탈 방지 · 재구독 유도\n온보딩 플로우 자동 설계"),
        ("⚡", "퍼포먼스 마케팅", "A/B 테스트 · 실험 설계\n데이터 기반 ROI 최적화"),
    ]
    for i, (ic, ttl, sub) in enumerate(skills):
        row, col = divmod(i, 2)
        bx = Inches(7.0 + col*3.1)
        ry = Inches(1.5 + row*2.75)
        rrect(sl, bx, ry, Inches(2.95), Inches(2.55), fill=C_CARD2)
        rect(sl, bx, ry, Inches(2.95), Inches(0.05), fill=C_CYAN if i%2==0 else C_BLUE)
        txt(sl, ic, bx+Inches(0.15), ry+Inches(0.1), Inches(2.6), Inches(0.65),
            sz=30, align=PP_ALIGN.CENTER)
        txt(sl, ttl, bx+Inches(0.15), ry+Inches(0.78), Inches(2.6), Inches(0.45),
            sz=16, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        txt(sl, sub, bx+Inches(0.15), ry+Inches(1.25), Inches(2.6), Inches(1.1),
            sz=13, color=C_GRAY, align=PP_ALIGN.CENTER)
    print("  [4] 위어드섹터 소개")

    # ── Slide 5: 제안 서비스 개요 ──────────────────
    sl = prs.slides.add_slide(blank(prs))
    rect(sl, 0, 0, SW, SH, fill=C_NAVY2)
    # 섹션 구분 슬라이드 — 블루 배경
    rect(sl, 0, 0, Inches(0.5), SH, fill=C_BLUE)
    rect(sl, 0, Inches(7.43), SW, Inches(0.07), fill=C_BLUE)

    txt(sl, "제안 서비스\n4가지 통합 솔루션",
        Inches(0.8), Inches(1.5), Inches(8.0), Inches(2.0),
        sz=44, bold=True, color=C_WHITE)
    txt(sl, "월 200만원으로 운영되는 M/M 방식 — 우선순위 높은 작업부터 순차 수행",
        Inches(0.8), Inches(3.6), Inches(11.0), Inches(0.6),
        sz=18, color=C_GRAY)

    services_overview = [
        ("01", "데이터 관리", "GA4 · 데이터 파이프라인 · 트래킹 설계", C_BLUE),
        ("02", "코드 관리",   "태그매니저 · SDK 연동 · 기술 부채 정리", C_CYAN),
        ("03", "MarTech 운영", "클레버탭 · 브레이즈 · 자동화 플로우",   C_GREEN),
        ("04", "CRM 수행",    "이탈 방지 · 재구독 · 고객 세그먼트",     C_ORANGE),
    ]
    for i, (num, nm, sub, c) in enumerate(services_overview):
        bx = Inches(0.8 + i*3.1)
        ry = Inches(4.5)
        rrect(sl, bx, ry, Inches(2.95), Inches(2.6), fill=C_CARD)
        rect(sl, bx, ry, Inches(2.95), Inches(0.06), fill=c)
        txt(sl, num, bx+Inches(0.2), ry+Inches(0.12), Inches(2.5), Inches(0.45),
            sz=14, bold=True, color=c, align=PP_ALIGN.LEFT)
        txt(sl, nm, bx+Inches(0.2), ry+Inches(0.58), Inches(2.5), Inches(0.55),
            sz=19, bold=True, color=C_WHITE)
        txt(sl, sub, bx+Inches(0.2), ry+Inches(1.15), Inches(2.5), Inches(1.2),
            sz=13, color=C_GRAY)
    print("  [5] 서비스 개요")

    # ── Slide 6: 서비스 ①② (데이터+코드) ──────────
    sl = prs.slides.add_slide(blank(prs))
    rect(sl, 0, 0, SW, SH, fill=C_NAVY)
    header(sl, 6, TOTAL, "제안 서비스 ① · ②", "데이터 관리 & 코드 관리 — 정확한 데이터가 모든 것의 시작")
    footer(sl)

    # Service 1
    rrect(sl, Inches(0.4), Inches(1.5), Inches(6.0), Inches(5.6), fill=C_CARD)
    rect(sl, Inches(0.4), Inches(1.5), Inches(6.0), Inches(0.07), fill=C_BLUE)
    rect(sl, Inches(0.4), Inches(1.5), Inches(0.07), Inches(5.6), fill=C_BLUE)
    txt(sl, "01  데이터 관리",
        Inches(0.65), Inches(1.62), Inches(5.5), Inches(0.6),
        sz=24, bold=True, color=C_BLUE_LIGHT)
    txt(sl, "정확한 데이터 수집 · 관리 · 분석 환경 구축",
        Inches(0.65), Inches(2.22), Inches(5.5), Inches(0.45),
        sz=15, color=C_GRAY)
    data_items = [
        ("📋", "GA4 이벤트 설계 및 구조 재설계"),
        ("🔗", "BigQuery 연동 데이터 파이프라인 구축"),
        ("📊", "모다이브 전용 대시보드 세팅"),
        ("🔍", "데이터 품질 점검 및 오류 수정"),
        ("📈", "주요 KPI 정의 및 트래킹 체계 수립"),
    ]
    for i, (ic, item) in enumerate(data_items):
        ry = Inches(2.85 + i*0.68)
        txt(sl, ic, Inches(0.65), ry, Inches(0.5), Inches(0.55), sz=18)
        txt(sl, item, Inches(1.2), ry+Inches(0.06), Inches(5.0), Inches(0.45),
            sz=15, color=C_GRAY_LIGHT)

    # Service 2
    rrect(sl, Inches(6.9), Inches(1.5), Inches(6.0), Inches(5.6), fill=C_CARD)
    rect(sl, Inches(6.9), Inches(1.5), Inches(6.0), Inches(0.07), fill=C_CYAN)
    rect(sl, Inches(6.9), Inches(1.5), Inches(0.07), Inches(5.6), fill=C_CYAN)
    txt(sl, "02  코드 관리",
        Inches(7.15), Inches(1.62), Inches(5.5), Inches(0.6),
        sz=24, bold=True, color=C_CYAN)
    txt(sl, "GTM 최적화 · SDK 연동 · 기술 부채 해소",
        Inches(7.15), Inches(2.22), Inches(5.5), Inches(0.45),
        sz=15, color=C_GRAY)
    code_items = [
        ("🏷", "구글 태그매니저(GTM) 전면 재구성"),
        ("📱", "앱 SDK(Firebase, AppsFlyer 등) 연동"),
        ("🛠", "기존 레거시 트래킹 코드 정리"),
        ("🔒", "개인정보보호법 준수 데이터 설계"),
        ("📝", "기술 문서화 및 인수인계 가이드 작성"),
    ]
    for i, (ic, item) in enumerate(code_items):
        ry = Inches(2.85 + i*0.68)
        txt(sl, ic, Inches(7.15), ry, Inches(0.5), Inches(0.55), sz=18)
        txt(sl, item, Inches(7.7), ry+Inches(0.06), Inches(5.0), Inches(0.45),
            sz=15, color=C_GRAY_LIGHT)
    print("  [6] 서비스 1·2")

    # ── Slide 7: 서비스 ③ MarTech ─────────────────
    sl = prs.slides.add_slide(blank(prs))
    rect(sl, 0, 0, SW, SH, fill=C_NAVY)
    header(sl, 7, TOTAL, "제안 서비스 ③ — MarTech 운영", "자동화된 마케팅 플로우로 사용자 전환율을 높입니다")
    footer(sl)

    txt(sl, "도구가 아니라 전략입니다 — 올바른 사람에게, 올바른 시간에, 올바른 메시지를",
        Inches(0.6), Inches(1.45), Inches(12.0), Inches(0.55),
        sz=18, bold=True, color=C_GREEN)

    # 3개 영역 카드
    martech = [
        {
            "title": "플랫폼 세팅 & 연동",
            "icon": "🔧",
            "items": ["CleverTap / Braze SDK 연동", "이벤트 매핑 및 사용자 속성 설계",
                      "세그먼트 구조 설계 (정적/동적)", "대시보드 및 리포트 설정"],
            "color": C_GREEN,
        },
        {
            "title": "자동화 캠페인 설계",
            "icon": "🤖",
            "items": ["온보딩 웰컴 플로우 구축", "이탈 방지 리텐션 캠페인",
                      "재구독 유도 Win-back 시퀀스", "행동 기반 트리거 메시지"],
            "color": C_BLUE,
        },
        {
            "title": "채널별 최적화",
            "icon": "📱",
            "items": ["Push 알림 개인화 (콘텐츠 추천)", "인앱 메시지 & 팝업 설계",
                      "이메일 뉴스레터 자동화", "A/B 테스트 & 성과 측정"],
            "color": C_CYAN,
        },
    ]
    for i, mc in enumerate(martech):
        bx = Inches(0.4 + i*4.28)
        rrect(sl, bx, Inches(2.1), Inches(4.1), Inches(4.95), fill=C_CARD)
        rect(sl, bx, Inches(2.1), Inches(4.1), Inches(0.06), fill=mc["color"])
        txt(sl, mc["icon"], bx+Inches(0.2), Inches(2.18), Inches(3.7), Inches(0.65),
            sz=32, align=PP_ALIGN.CENTER)
        txt(sl, mc["title"], bx+Inches(0.2), Inches(2.85), Inches(3.7), Inches(0.55),
            sz=18, bold=True, color=mc["color"], align=PP_ALIGN.CENTER)
        rect(sl, bx+Inches(0.4), Inches(3.48), Inches(3.3), Inches(0.04), fill=mc["color"])
        for j, item in enumerate(mc["items"]):
            txt(sl, f"▸  {item}", bx+Inches(0.2), Inches(3.62 + j*0.63),
                Inches(3.7), Inches(0.55), sz=14, color=C_GRAY_LIGHT)
    print("  [7] MarTech")

    # ── Slide 8: 서비스 ④ CRM ──────────────────────
    sl = prs.slides.add_slide(blank(prs))
    rect(sl, 0, 0, SW, SH, fill=C_NAVY)
    header(sl, 8, TOTAL, "제안 서비스 ④ — CRM 수행", "고객 데이터로 이탈을 막고 재방문을 만드는 정밀 운영")
    footer(sl)

    # 좌측 — CRM 플로우 다이어그램
    rrect(sl, Inches(0.4), Inches(1.5), Inches(6.3), Inches(5.6), fill=C_CARD)
    rect(sl, Inches(0.4), Inches(1.5), Inches(6.3), Inches(0.06), fill=C_ORANGE)
    txt(sl, "CRM 수행 프레임워크",
        Inches(0.65), Inches(1.65), Inches(5.8), Inches(0.55),
        sz=20, bold=True, color=C_ORANGE)

    crm_flow = [
        ("📥", "데이터 수집",    "행동 데이터 · 구독 이력 · 콘텐츠 선호"),
        ("🔍", "세그먼트 분류",  "이탈 위험군 · 고가치 · 신규 · 휴면"),
        ("🎯", "개인화 전략",    "세그먼트별 맞춤 메시지 & 타이밍"),
        ("📤", "자동화 발송",    "Push · 이메일 · SMS · 인앱 통합"),
        ("📊", "성과 측정",     "전환율 · 이탈률 · LTV 지표 분석"),
    ]
    for i, (ic, step, desc) in enumerate(crm_flow):
        ry = Inches(2.35 + i*0.97)
        rrect(sl, Inches(0.65), ry, Inches(5.8), Inches(0.85), fill=C_CARD2)
        txt(sl, ic, Inches(0.8), ry+Inches(0.1), Inches(0.55), Inches(0.6), sz=22)
        txt(sl, step, Inches(1.45), ry+Inches(0.1), Inches(2.0), Inches(0.4),
            sz=15, bold=True, color=C_WHITE)
        txt(sl, desc, Inches(3.5), ry+Inches(0.12), Inches(2.8), Inches(0.55),
            sz=13, color=C_GRAY)
        if i < 4:
            txt(sl, "↓", Inches(3.1), ry+Inches(0.88), Inches(0.4), Inches(0.5),
                sz=16, color=C_ORANGE, align=PP_ALIGN.CENTER)

    # 우측 — CRM 성과 지표
    rrect(sl, Inches(6.9), Inches(1.5), Inches(6.0), Inches(5.6), fill=C_CARD)
    rect(sl, Inches(6.9), Inches(1.5), Inches(6.0), Inches(0.06), fill=C_ORANGE)
    txt(sl, "CRM 도입 기대 효과",
        Inches(7.15), Inches(1.65), Inches(5.5), Inches(0.55),
        sz=20, bold=True, color=C_ORANGE)

    crm_kpis = [
        ("-35%", "이탈률 감소",       "리텐션 캠페인 효과", C_ORANGE),
        ("+42%", "재방문율 향상",      "Win-back 캠페인", C_GREEN),
        ("+28%", "구독 전환율",        "개인화 추천 효과", C_BLUE_LIGHT),
        ("2.5배", "LTV 증가",         "CRM 자동화 6개월 후", C_CYAN),
    ]
    for i, (v, l, d, c) in enumerate(crm_kpis):
        ry = Inches(2.35 + i*1.18)
        rrect(sl, Inches(7.15), ry, Inches(5.5), Inches(1.05), fill=C_CARD2)
        rect(sl, Inches(7.15), ry, Inches(0.06), Inches(1.05), fill=c)
        txt(sl, v, Inches(7.35), ry+Inches(0.08), Inches(1.8), Inches(0.6),
            sz=32, bold=True, color=c)
        txt(sl, l, Inches(9.2), ry+Inches(0.12), Inches(3.2), Inches(0.42),
            sz=17, bold=True, color=C_WHITE)
        txt(sl, d, Inches(9.2), ry+Inches(0.55), Inches(3.2), Inches(0.38),
            sz=13, color=C_GRAY)
    print("  [8] CRM")

    # ── Slide 9: M/M 운영 방식 & 예산 ─────────────
    sl = prs.slides.add_slide(blank(prs))
    rect(sl, 0, 0, SW, SH, fill=C_NAVY)
    header(sl, 9, TOTAL, "M/M 운영 방식 & 예산", "월 200만원으로 전담 팀이 움직입니다")
    footer(sl)

    # 중앙 예산 강조
    txt(sl, "월 200만원",
        Inches(4.5), Inches(1.5), Inches(4.5), Inches(0.85),
        sz=42, bold=True, color=C_BLUE_LIGHT, align=PP_ALIGN.CENTER)
    txt(sl, "M/M(Man-Month) 기반 통합 운영",
        Inches(3.0), Inches(2.35), Inches(7.0), Inches(0.45),
        sz=18, color=C_GRAY, align=PP_ALIGN.CENTER)
    rect(sl, Inches(3.5), Inches(2.88), Inches(6.3), Inches(0.04), fill=C_BLUE)

    # 운영 원칙
    principles = [
        ("🎯", "우선순위 기반",     "매월 초 작업 목록 조율 → 우선순위 합의"),
        ("📋", "투명한 리포팅",     "주간 진행 보고 + 월간 성과 리포트"),
        ("🔄", "유연한 조정",      "비즈니스 변화에 따라 서비스 비율 조정 가능"),
        ("📞", "즉각적인 소통",     "전담 슬랙 채널 + 정기 미팅 (격주)"),
    ]
    for i, (ic, ttl, desc) in enumerate(principles):
        row, col = divmod(i, 2)
        bx = Inches(0.5 + col*6.5)
        ry = Inches(3.1 + row*1.8)
        rrect(sl, bx, ry, Inches(6.1), Inches(1.6), fill=C_CARD)
        rect(sl, bx, ry, Inches(6.1), Inches(0.06), fill=C_BLUE)
        txt(sl, ic, bx+Inches(0.2), ry+Inches(0.1), Inches(0.65), Inches(0.65), sz=26)
        txt(sl, ttl, bx+Inches(0.95), ry+Inches(0.15), Inches(4.8), Inches(0.45),
            sz=18, bold=True, color=C_WHITE)
        txt(sl, desc, bx+Inches(0.95), ry+Inches(0.62), Inches(4.8), Inches(0.7),
            sz=14, color=C_GRAY)
    print("  [9] 운영 방식")

    # ── Slide 10: 기대 효과 & ROI ──────────────────
    sl = prs.slides.add_slide(blank(prs))
    rect(sl, 0, 0, SW, SH, fill=C_NAVY)
    header(sl, 10, TOTAL, "기대 효과 & ROI", "위어드섹터와 함께라면 이 숫자가 현실이 됩니다")
    footer(sl)

    txt(sl, "6개월 집중 운영 후 예상 성과 — 글로벌 OTT 플랫폼 평균 기준",
        Inches(0.6), Inches(1.45), Inches(12.0), Inches(0.55),
        sz=17, color=C_GRAY)

    rois = [
        ("📉", "-35%", "사용자 이탈률",   "리텐션 캠페인 효과", C_ORANGE),
        ("📈", "+42%", "재방문율",         "개인화 추천 콘텐츠", C_GREEN),
        ("💰", "3.5배", "캠페인 ROI",      "자동화 마케팅 적용", C_BLUE_LIGHT),
        ("⚡",  "60%",  "운영 시간 절감",   "자동화 플로우 구축", C_CYAN),
        ("🎯", "+28%", "구독 전환율",      "CRM 개인화 효과",   C_ORANGE),
        ("📊", "95%",  "데이터 정확도",    "트래킹 재설계 후",  C_GREEN),
    ]
    for i, (ic, val, lbl, desc, c) in enumerate(rois):
        row, col = divmod(i, 3)
        bx = Inches(0.4 + col*4.28)
        ry = Inches(2.1 + row*2.55)
        rrect(sl, bx, ry, Inches(4.05), Inches(2.35), fill=C_CARD)
        rect(sl, bx, ry, Inches(4.05), Inches(0.06), fill=c)
        txt(sl, ic, bx+Inches(0.2), ry+Inches(0.1), Inches(3.6), Inches(0.6),
            sz=30, align=PP_ALIGN.CENTER)
        txt(sl, val, bx+Inches(0.2), ry+Inches(0.72), Inches(3.6), Inches(0.7),
            sz=40, bold=True, color=c, align=PP_ALIGN.CENTER)
        txt(sl, lbl, bx+Inches(0.2), ry+Inches(1.45), Inches(3.6), Inches(0.42),
            sz=16, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        txt(sl, desc, bx+Inches(0.2), ry+Inches(1.9), Inches(3.6), Inches(0.38),
            sz=12, color=C_GRAY, align=PP_ALIGN.CENTER)
    print("  [10] ROI")

    # ── Slide 11: 왜 위어드섹터인가? ──────────────
    sl = prs.slides.add_slide(blank(prs))
    rect(sl, 0, 0, SW, SH, fill=C_NAVY)
    header(sl, 11, TOTAL, "왜 위어드섹터인가?", "우리가 모다이브에 최적인 이유")
    footer(sl)

    # 비교 테이블
    headers_row = ["항목", "위어드섹터", "일반 대행사", "인하우스 채용"]
    col_ws = [Inches(3.2), Inches(2.9), Inches(2.9), Inches(3.2)]

    cx = Inches(0.4)
    for j, (h, w) in enumerate(zip(headers_row, col_ws)):
        bg = C_BLUE if j == 1 else C_CARD2
        rrect(sl, cx, Inches(1.55), w-Inches(0.1), Inches(0.6), fill=bg)
        txt(sl, h, cx+Inches(0.1), Inches(1.6), w-Inches(0.15), Inches(0.48),
            sz=16, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        cx += w

    rows = [
        ("전문성 (MarTech + CRM)",  ["✅ 전문 특화",  "⚠ 범용적",    "❌ 육성 필요"]),
        ("시작 속도",               ["✅ 즉시 가능",  "⚠ 2~4주",     "❌ 채용 3개월+"]),
        ("월 비용",                 ["✅ 200만원",    "⚠ 300~500만", "❌ 450만원+"]),
        ("데이터 + 기술 복합",       ["✅ 원스톱",    "⚠ 별도 필요", "⚠ 파편화"]),
        ("성과 측정 & 리포팅",       ["✅ 체계적",    "⚠ 기본 수준", "⚠ 내부 의존"]),
        ("유연한 서비스 조정",       ["✅ 즉시 가능", "❌ 계약 고정", "❌ 재구성 필요"]),
    ]
    for i, (feat, vals) in enumerate(rows):
        ry = Inches(2.22 + i*0.82)
        bg = C_CARD if i%2==0 else C_CARD2
        rect(sl, Inches(0.4), ry, Inches(12.5), Inches(0.78), fill=bg)
        txt(sl, feat, Inches(0.55), ry+Inches(0.17), Inches(3.0), Inches(0.45),
            sz=15, color=C_WHITE)
        cx = Inches(3.6)
        for j, (v, w) in enumerate(zip(vals, col_ws[1:])):
            hl = j == 0
            if hl:
                rrect(sl, cx, ry+Inches(0.08), w-Inches(0.1), Inches(0.62),
                      fill=RGBColor(0x00, 0x44, 0x99))
            txt(sl, v, cx+Inches(0.05), ry+Inches(0.17), w-Inches(0.1), Inches(0.45),
                sz=14, bold=hl, color=C_BLUE_LIGHT if hl else C_GRAY,
                align=PP_ALIGN.CENTER)
            cx += w
    print("  [11] 위어드섹터 강점")

    # ── Slide 12: 다음 단계 CTA ────────────────────
    sl = prs.slides.add_slide(blank(prs))
    rect(sl, 0, 0, SW, SH, fill=C_NAVY2)
    rect(sl, Inches(8.5), 0, Inches(4.83), SH, fill=C_BLUE_DARK)
    rect(sl, Inches(7.8), 0, Inches(0.9), SH, fill=C_CARD)
    rect(sl, 0, 0, SW, Inches(0.06), fill=C_BLUE)
    rect(sl, 0, Inches(7.44), SW, Inches(0.06), fill=C_BLUE)

    txt(sl, "지금 시작하면\n다릅니다",
        Inches(0.6), Inches(1.2), Inches(7.5), Inches(1.8),
        sz=44, bold=True, color=C_WHITE)
    txt(sl, "모다이브의 데이터·마케팅·CRM, 위어드섹터가 전담합니다.",
        Inches(0.6), Inches(3.1), Inches(7.5), Inches(0.55),
        sz=18, color=C_GRAY)

    # 3단계 온보딩
    steps = [
        ("01", "미팅 & 진단", "현황 파악 및 우선순위 협의"),
        ("02", "킥오프",      "2주 내 팀 배정 & 작업 시작"),
        ("03", "성과 리뷰",   "4주 후 첫 번째 성과 점검"),
    ]
    for i, (num, ttl, sub) in enumerate(steps):
        ry = Inches(3.9 + i*0.85)
        rrect(sl, Inches(0.6), ry, Inches(6.8), Inches(0.72), fill=C_CARD)
        rect(sl, Inches(0.6), ry, Inches(0.06), Inches(0.72), fill=C_BLUE)
        txt(sl, num, Inches(0.8), ry+Inches(0.12), Inches(0.7), Inches(0.45),
            sz=18, bold=True, color=C_BLUE_LIGHT)
        txt(sl, ttl, Inches(1.55), ry+Inches(0.12), Inches(3.0), Inches(0.42),
            sz=17, bold=True, color=C_WHITE)
        txt(sl, sub, Inches(4.6), ry+Inches(0.15), Inches(2.6), Inches(0.38),
            sz=14, color=C_GRAY)

    # 우측 패널
    txt(sl, "Contact",
        Inches(8.7), Inches(1.5), Inches(4.2), Inches(0.5),
        sz=14, color=C_BLUE_LIGHT, align=PP_ALIGN.CENTER)
    for i, (ic, lb, vl) in enumerate([
        ("🌐", "웹사이트", "weirdsector.co.kr"),
        ("📧", "이메일",   "hello@weirdsector.co.kr"),
        ("📞", "전화",     "문의 주시면 바로 연락드립니다"),
    ]):
        ry = Inches(2.2 + i*1.55)
        rrect(sl, Inches(8.7), ry, Inches(4.2), Inches(1.35), fill=C_CARD)
        txt(sl, ic, Inches(8.9), ry+Inches(0.1), Inches(3.8), Inches(0.55),
            sz=28, align=PP_ALIGN.CENTER)
        txt(sl, lb, Inches(8.9), ry+Inches(0.65), Inches(3.8), Inches(0.35),
            sz=12, color=C_GRAY, align=PP_ALIGN.CENTER)
        txt(sl, vl, Inches(8.9), ry+Inches(0.98), Inches(3.8), Inches(0.32),
            sz=13, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    print("  [12] CTA")

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    sz = OUT.stat().st_size // 1024
    print(f"\n✅ 저장: {OUT}")
    print(f"   슬라이드: {len(prs.slides)}장 / {sz:,}KB")
    return len(prs.slides)


if __name__ == "__main__":
    print("MBC 모다이브 × 위어드섹터 제안서 v2 빌드\n")
    n = build()
    print(f"\n완료: {n}장")
