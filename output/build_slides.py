"""
Pepper Potts — Fintech Business Analysis Lecture Slides
Based on Hulk lecture draft: C:\Agent\HULK\fintech_lecture\slides_draft.md
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ─── Design System ──────────────────────────────────────────────────────────
NAVY       = RGBColor(0x1E, 0x3A, 0x5F)   # Primary dark navy
BLUE       = RGBColor(0x25, 0x63, 0xEB)   # Bright blue accent
EMERALD    = RGBColor(0x05, 0x96, 0x69)   # Emerald green
AMBER      = RGBColor(0xD9, 0x77, 0x06)   # Amber highlight
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF1, 0xF5, 0xF9)
MID_GRAY   = RGBColor(0x94, 0xA3, 0xB8)
DARK_TEXT  = RGBColor(0x1F, 0x29, 0x37)
DIVIDER    = RGBColor(0xE2, 0xE8, 0xF0)
SECTION_BG = RGBColor(0x0F, 0x24, 0x44)   # Deep navy for section dividers
RED        = RGBColor(0xDC, 0x26, 0x26)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

blank_layout = prs.slide_layouts[6]  # Completely blank


# ─── Helper Functions ────────────────────────────────────────────────────────
def add_rect(slide, x, y, w, h, fill_color, line_color=None, line_width=0):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()
    return shape


def add_textbox(slide, text, x, y, w, h,
                font_size=18, bold=False, italic=False,
                color=DARK_TEXT, align=PP_ALIGN.LEFT,
                font_name="Pretendard", wrap=True):
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font_name
    return txBox


def add_textbox_multiline(slide, lines, x, y, w, h,
                          font_size=16, bold=False, color=DARK_TEXT,
                          align=PP_ALIGN.LEFT, line_spacing=1.2,
                          font_name="Pretendard"):
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for line in lines:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = font_name
    return txBox


def header_bar(slide, title, subtitle=None):
    """Top header bar with navy background"""
    add_rect(slide, 0, 0, 13.33, 1.4, NAVY)
    add_rect(slide, 0, 1.4, 13.33, 0.05, BLUE)
    add_textbox(slide, title, 0.5, 0.15, 11, 0.7,
                font_size=28, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        add_textbox(slide, subtitle, 0.5, 0.85, 11, 0.5,
                    font_size=14, color=MID_GRAY, align=PP_ALIGN.LEFT)


def slide_number(slide, num, total=45):
    add_textbox(slide, f"{num} / {total}", 11.8, 7.1, 1.4, 0.3,
                font_size=11, color=MID_GRAY, align=PP_ALIGN.RIGHT)
    # Bottom accent line
    add_rect(slide, 0, 7.35, 13.33, 0.07, BLUE)


def kpi_card(slide, x, y, w, h, value, label, color=BLUE):
    add_rect(slide, x, y, w, h, LIGHT_GRAY)
    add_rect(slide, x, y, 0.08, h, color)
    add_textbox(slide, value, x + 0.2, y + 0.1, w - 0.3, h * 0.55,
                font_size=22, bold=True, color=NAVY)
    add_textbox(slide, label, x + 0.2, y + h * 0.55, w - 0.3, h * 0.4,
                font_size=11, color=MID_GRAY)


def bullet_block(slide, items, x, y, w, h, bullet="•",
                 font_size=15, indent_color=BLUE):
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = f"{bullet}  {item}"
        run.font.size = Pt(font_size)
        run.font.color.rgb = DARK_TEXT
        run.font.name = "Pretendard"


# ─── Slide Builders ──────────────────────────────────────────────────────────

def slide_01_title():
    """표지"""
    slide = prs.slides.add_slide(blank_layout)
    # Full background
    add_rect(slide, 0, 0, 13.33, 7.5, SECTION_BG)
    # Blue accent stripe
    add_rect(slide, 0, 5.8, 13.33, 0.12, BLUE)
    # Decorative circles
    for i, (cx, cy, r) in enumerate([
        (10.5, 1.2, 2.5), (11.8, 3.0, 3.5), (9.0, 5.5, 1.8)
    ]):
        sh = slide.shapes.add_shape(9, Inches(cx - r/2), Inches(cy - r/2),
                                    Inches(r), Inches(r))
        sh.fill.solid()
        sh.fill.fore_color.rgb = RGBColor(0x25, 0x63, 0xEB)
        sh.line.fill.background()

    # Tag line
    add_rect(slide, 0.6, 1.0, 0.12, 0.6, EMERALD)
    add_textbox(slide, "FINTECH BUSINESS ANALYSIS", 0.85, 0.95, 9, 0.5,
                font_size=12, color=EMERALD, bold=True)

    # Main title
    add_textbox(slide, "핀테크 비즈니스 분석 실전", 0.6, 1.6, 10, 1.1,
                font_size=48, bold=True, color=WHITE)
    add_textbox(slide, "데이터로 읽는 핀테크\n비즈니스 모델 분석부터 대시보드 구축까지",
                0.6, 2.75, 10, 1.0, font_size=22, color=MID_GRAY)

    # Info row
    add_rect(slide, 0.6, 4.0, 3.5, 0.6, RGBColor(0x25, 0x3A, 0x5E))
    add_textbox(slide, "📅  2026-04-21", 0.75, 4.05, 3.2, 0.5,
                font_size=14, color=WHITE)
    add_rect(slide, 4.3, 4.0, 4.0, 0.6, RGBColor(0x25, 0x3A, 0x5E))
    add_textbox(slide, "⏱  3시간 강의  |  45 슬라이드", 4.45, 4.05, 3.8, 0.5,
                font_size=14, color=WHITE)

    # Target
    add_textbox(slide, "대상: 비즈니스 데이터 분석 실무 입문자  |  도구: Amplitude · Looker Studio",
                0.6, 5.0, 12, 0.4, font_size=13, color=MID_GRAY)
    # Branding
    add_textbox(slide, "HULK × PEPPER", 11.0, 7.1, 2.1, 0.3,
                font_size=11, color=EMERALD, align=PP_ALIGN.RIGHT, bold=True)
    add_rect(slide, 0, 7.35, 13.33, 0.07, BLUE)


def slide_02_agenda():
    """오늘의 여정"""
    slide = prs.slides.add_slide(blank_layout)
    header_bar(slide, "오늘의 여정", "Today's Roadmap — 3 Hours")
    add_rect(slide, 0, 0, 0.1, 7.5, BLUE)

    parts = [
        ("Part 1", "15분", "오리엔테이션", BLUE, "0:00 – 0:15"),
        ("Part 2", "45분", "비즈니스 모델 분석 방법론", EMERALD, "0:15 – 1:00"),
        ("Part 3", "60분", "케이스 스터디 (토스·카카오페이·Stripe)", AMBER, "1:00 – 2:00"),
        ("Part 4", "60분", "실습: Amplitude & Looker Studio", RGBColor(0x7C, 0x3A, 0xED), "2:00 – 3:00"),
        ("Part 5", "20분", "마무리 & Q&A", RED, "2:40 – 3:00"),
    ]

    for i, (part, dur, desc, color, time) in enumerate(parts):
        y = 1.65 + i * 1.05
        add_rect(slide, 0.3, y, 0.07, 0.75, color)
        add_rect(slide, 0.5, y, 12.5, 0.75, LIGHT_GRAY)
        add_textbox(slide, part, 0.65, y + 0.05, 1.0, 0.35,
                    font_size=12, bold=True, color=color)
        add_textbox(slide, desc, 0.65, y + 0.38, 7.5, 0.32,
                    font_size=15, bold=True, color=NAVY)
        add_textbox(slide, f"⏱ {dur}  |  {time}", 8.5, y + 0.22, 4.0, 0.3,
                    font_size=12, color=MID_GRAY, align=PP_ALIGN.RIGHT)

    add_textbox(slide, "🎯  오늘의 목표: \"핀테크 데이터를 보는 눈 장착하기\"",
                0.5, 6.95, 12.5, 0.4, font_size=13, bold=True, color=NAVY)
    slide_number(slide, 2)


def section_divider(num, part_num, title, subtitle, color=BLUE):
    """Section divider slide"""
    slide = prs.slides.add_slide(blank_layout)
    add_rect(slide, 0, 0, 13.33, 7.5, SECTION_BG)
    add_rect(slide, 0, 0, 0.5, 7.5, color)
    add_rect(slide, 0, 6.5, 13.33, 0.12, color)

    add_textbox(slide, part_num, 0.8, 1.5, 5, 0.7,
                font_size=14, bold=True, color=color)
    add_textbox(slide, title, 0.8, 2.2, 11, 1.2,
                font_size=44, bold=True, color=WHITE)
    add_textbox(slide, subtitle, 0.8, 3.5, 11, 0.6,
                font_size=18, color=MID_GRAY)
    add_textbox(slide, f"{num} / 45", 11.8, 7.1, 1.4, 0.3,
                font_size=11, color=MID_GRAY, align=PP_ALIGN.RIGHT)
    add_rect(slide, 0, 7.35, 13.33, 0.07, color)
    return slide


def content_slide(num, title, subtitle=None):
    """Standard content slide — returns slide object for further population"""
    slide = prs.slides.add_slide(blank_layout)
    header_bar(slide, title, subtitle)
    add_rect(slide, 0, 0, 0.1, 7.5, BLUE)
    slide_number(slide, num)
    return slide


# ══════════════════════════════════════════════════════════════════════════════
# BUILD ALL 45 SLIDES
# ══════════════════════════════════════════════════════════════════════════════

# 1 — Title
slide_01_title()

# 2 — Agenda
slide_02_agenda()

# ─── PART 1 SECTION ────────────────────────────────────────────────────────
section_divider(3, "PART 1  ·  0:00 – 0:15",
                "오프닝 & 오리엔테이션",
                "핀테크 데이터 분석의 세계에 오신 것을 환영합니다", BLUE)

# 3 — 왜 지금 핀테크 데이터 분석인가?
slide = content_slide(4, "왜 지금 핀테크 데이터 분석인가?",
                      "Fintech Market Explosion — Why Data Literacy Matters Now")
kpi_card(slide, 0.3, 1.6, 3.0, 1.5, "$340B → $1.1T", "글로벌 핀테크 시장\n(2024 → 2030)", BLUE)
kpi_card(slide, 3.5, 1.6, 3.0, 1.5, "CAGR 22%", "연평균 성장률", EMERALD)
kpi_card(slide, 6.7, 1.6, 3.0, 1.5, "200조원+", "한국 간편결제 연간 거래액", AMBER)
kpi_card(slide, 9.9, 1.6, 3.0, 1.5, "필수 역량", "핀테크 PM/마케터의\n데이터 리터러시", RED)
add_rect(slide, 0.3, 3.3, 12.6, 1.0, LIGHT_GRAY)
add_textbox(slide, "💡  \"핀테크 PM/마케터에게 데이터 리터러시는 선택이 아닌 필수\"",
            0.5, 3.42, 12.2, 0.5, font_size=16, bold=True, color=NAVY)
bullet_block(slide, [
    "여러분의 서비스에서 오늘 얼마나 많은 결제가 일어났을까요?",
    "그 결제 데이터는 어떤 비즈니스 인사이트를 담고 있을까요?",
    "경쟁사 대비 우리 서비스의 리텐션은 어디에 있을까요?",
], 0.3, 4.5, 12.6, 2.3, font_size=15)

# 4 — 수강생 소개 / 강의 목표
slide = content_slide(5, "강의 후 여러분이 할 수 있는 것",
                      "Learning Objectives — 오늘 장착하게 될 역량")
goals = [
    "핀테크 비즈니스 모델 유형별로 분류하고 KPI를 정의할 수 있다",
    "비즈니스 모델 캔버스(BMC)를 핀테크 서비스에 적용할 수 있다",
    "토스·카카오페이·Stripe 전략을 데이터 관점에서 해석할 수 있다",
    "Amplitude로 퍼널·리텐션·코호트 분석을 수행할 수 있다",
    "Looker Studio로 핵심 지표 대시보드를 만들 수 있다",
]
for i, goal in enumerate(goals):
    y = 1.7 + i * 1.0
    add_rect(slide, 0.3, y, 0.6, 0.7, BLUE)
    add_textbox(slide, "✅", 0.3, y + 0.1, 0.6, 0.5, font_size=18,
                color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(slide, 1.05, y, 11.8, 0.7, LIGHT_GRAY)
    add_textbox(slide, goal, 1.2, y + 0.15, 11.5, 0.45,
                font_size=15, bold=True, color=NAVY)

# ─── PART 2 SECTION ────────────────────────────────────────────────────────
section_divider(6, "PART 2  ·  0:15 – 1:00",
                "비즈니스 모델 분석 방법론",
                "핀테크 비즈니스를 구조화해서 분석하는 눈 기르기", EMERALD)

# 6 — 핀테크 5대 유형
slide = content_slide(7, "핀테크 비즈니스 모델 5대 유형",
                      "어떤 수익 구조를 가지고 있는가?")
types = [
    ("결제/송금", "토스, 카카오페이, Stripe", "거래 수수료, MDR", BLUE),
    ("대출/크레딧", "토스뱅크, 핀다, 카카오뱅크", "이자 수익, 중개 수수료", EMERALD),
    ("투자/자산관리", "토스증권, Robinhood", "거래 수수료, 구독료", AMBER),
    ("보험(인슈어테크)", "카카오손보, Lemonade", "보험료 수입, 손해율 관리", RED),
    ("B2B 인프라", "Plaid, Marqeta, AWS Pay", "API 사용료, 라이선스", RGBColor(0x7C, 0x3A, 0xED)),
]
for i, (type_name, examples, revenue, color) in enumerate(types):
    y = 1.7 + i * 0.95
    add_rect(slide, 0.3, y, 0.1, 0.75, color)
    add_rect(slide, 0.5, y, 12.5, 0.75, LIGHT_GRAY)
    add_textbox(slide, type_name, 0.65, y + 0.05, 2.5, 0.3,
                font_size=15, bold=True, color=color)
    add_textbox(slide, f"예: {examples}", 0.65, y + 0.4, 5.0, 0.3,
                font_size=12, color=MID_GRAY)
    add_textbox(slide, revenue, 7.0, y + 0.2, 5.8, 0.35,
                font_size=13, color=DARK_TEXT, align=PP_ALIGN.RIGHT)

# 7 — 수익 모델 심화
slide = content_slide(8, "수익 모델 심화 — 핀테크는 어떻게 돈을 버나?",
                      "4가지 수익 패턴")
models = [
    ("💳 수수료형", "거래당 % 또는 고정액 과금", "Stripe: 2.9% + $0.30/건", BLUE),
    ("🏦 이자 마진형", "예대마진 (대출이자 - 예금이자)", "카카오뱅크 NIM 2~4%", EMERALD),
    ("📱 구독형", "월정액 프리미엄 서비스", "Robinhood Gold $5/월", AMBER),
    ("📊 데이터/광고형", "금융 행동 데이터 활용·판매", "신용평가 데이터, 타겟 광고", RED),
]
for i, (title, desc, example, color) in enumerate(models):
    col = i % 2
    row = i // 2
    x = 0.3 + col * 6.5
    y = 1.65 + row * 2.5
    add_rect(slide, x, y, 6.2, 2.2, LIGHT_GRAY)
    add_rect(slide, x, y, 6.2, 0.55, color)
    add_textbox(slide, title, x + 0.2, y + 0.1, 5.8, 0.4,
                font_size=16, bold=True, color=WHITE)
    add_textbox(slide, desc, x + 0.2, y + 0.7, 5.8, 0.5,
                font_size=14, bold=True, color=NAVY)
    add_textbox(slide, f"예시: {example}", x + 0.2, y + 1.3, 5.8, 0.5,
                font_size=12, color=MID_GRAY)

# 8 — 비즈니스 모델 캔버스
slide = content_slide(9, "비즈니스 모델 캔버스 (BMC) — 핀테크 적용",
                      "9개 블록으로 비즈니스 전체 구조 한눈에 파악")
bmc_cells = [
    # (label, x, y, w, h, color)
    ("핵심 파트너십\nKey Partnerships", 0.3, 1.55, 2.4, 3.6, NAVY),
    ("핵심 활동\nKey Activities", 2.85, 1.55, 2.4, 1.7, BLUE),
    ("가치 제안\nValue Propositions\n\n빠르고 / 저렴하고\n쉽게 / 안전하게", 5.4, 1.55, 2.8, 3.6, RGBColor(0x1D, 0x4E, 0x89)),
    ("고객 관계\nCustomer Relationships", 8.35, 1.55, 2.4, 1.7, BLUE),
    ("고객 세그먼트\nCustomer Segments", 10.9, 1.55, 2.1, 3.6, NAVY),
    ("핵심 자원\nKey Resources", 2.85, 3.4, 2.4, 1.75, BLUE),
    ("채널\nChannels", 8.35, 3.4, 2.4, 1.75, BLUE),
    ("비용 구조\nCost Structure\n마케팅 CAC / 인프라 / 충당금 / 인건비", 0.3, 5.3, 6.15, 1.85, RGBColor(0x14, 0x53, 0x85)),
    ("수익 흐름\nRevenue Streams\n수수료 / 이자 / 구독 / 광고", 6.6, 5.3, 6.4, 1.85, RGBColor(0x05, 0x60, 0x44)),
]
for (label, x, y, w, h, color) in bmc_cells:
    add_rect(slide, x, y, w, h, color)
    add_textbox(slide, label, x + 0.1, y + 0.1, w - 0.2, h - 0.2,
                font_size=10, color=WHITE, bold=True)
    # border
    sh = slide.shapes[-1]

# 9 — 핵심 KPI 10개
slide = content_slide(10, "핀테크 핵심 KPI 10선",
                      "반드시 추적해야 할 지표들")
kpis = [
    ("GMV", "총 거래액", "우리 플랫폼을 통해 얼마나 거래되나?"),
    ("Take Rate", "GMV 대비 수익률", "거래에서 얼마를 가져가나?"),
    ("CAC", "고객 획득 비용", "신규 고객 1명 데려오는데 얼마?"),
    ("LTV", "고객 생애 가치", "고객 1명이 평생 얼마나 기여하나?"),
    ("LTV/CAC", "수익성 비율", "> 3이어야 건강한 사업"),
    ("MAU/DAU", "활성 사용자", "얼마나 자주 오나?"),
    ("Stickiness", "DAU/MAU 비율", "> 20%이면 좋은 앱"),
    ("Churn Rate", "이탈률", "얼마나 빠져나가나?"),
    ("NPS", "추천 지수", "만족하고 추천하나?"),
    ("Conversion Rate", "전환율", "행동 유도가 잘 되나?"),
]
for i, (kpi, definition, question) in enumerate(kpis):
    col = i % 2
    row = i // 2
    x = 0.3 + col * 6.6
    y = 1.65 + row * 1.05
    add_rect(slide, x, y, 6.3, 0.9, LIGHT_GRAY)
    add_rect(slide, x, y, 0.08, 0.9, BLUE if col == 0 else EMERALD)
    add_textbox(slide, kpi, x + 0.2, y + 0.05, 1.8, 0.4,
                font_size=14, bold=True, color=NAVY)
    add_textbox(slide, definition, x + 0.2, y + 0.5, 1.8, 0.35,
                font_size=11, color=MID_GRAY)
    add_textbox(slide, question, x + 2.1, y + 0.22, 4.0, 0.45,
                font_size=12, color=DARK_TEXT)

# 10 — 데이터로 비즈니스 모델 읽는 법
slide = content_slide(11, "데이터로 비즈니스 모델을 읽는 법",
                      "KPI 간의 관계로 비즈니스 건강도를 진단한다")
diagnostics = [
    ("📈 성장 진단", "MAU 증가율 + CAC 추이", "→ 효율적 성장인가?", BLUE),
    ("💰 수익성 진단", "Take Rate × GMV vs. 고정비", "→ 단위 경제성 확보 여부", EMERALD),
    ("🔄 리텐션 진단", "Churn Rate + LTV/CAC", "→ 장기 지속 가능성", AMBER),
    ("❤️ 참여도 진단", "DAU/MAU + 기능별 사용률", "→ 핵심 기능 파악", RED),
]
for i, (title, metric, insight, color) in enumerate(diagnostics):
    y = 1.7 + i * 1.2
    add_rect(slide, 0.3, y, 12.7, 1.0, LIGHT_GRAY)
    add_rect(slide, 0.3, y, 0.1, 1.0, color)
    add_textbox(slide, title, 0.55, y + 0.08, 2.8, 0.4,
                font_size=15, bold=True, color=color)
    add_textbox(slide, metric, 3.5, y + 0.08, 5.0, 0.4,
                font_size=13, color=NAVY, bold=True)
    add_textbox(slide, insight, 8.6, y + 0.08, 4.2, 0.4,
                font_size=13, color=EMERALD, align=PP_ALIGN.RIGHT)
    add_textbox(slide, "Growth Accounting: 신규 + 복귀 - 이탈 = 순증가" if i == 3 else "",
                0.55, y + 0.55, 12.0, 0.35, font_size=11, color=MID_GRAY)

add_textbox(slide, "📌  \"비즈니스 구조를 이해하고, KPI를 정의하고, 데이터로 검증하라\"",
            0.3, 6.6, 12.7, 0.4, font_size=14, bold=True, color=NAVY)

# 11 — 경쟁 우위 & 해자
slide = content_slide(12, "핀테크 경쟁 우위의 원천",
                      "네트워크 효과 · 데이터 해자 · 규제 · 전환비용")
moats = [
    ("🌐 네트워크 효과", "사용자 증가 → 서비스 가치 증가", "카카오페이: 카카오톡 5,000만 기반", BLUE),
    ("📊 데이터 해자", "거래 데이터 → 신용평가 정교화 → 리스크 감소", "토스: 송금 → 금융 데이터 → LTV 극대화", EMERALD),
    ("📋 규제 해자", "금융 라이선스 획득 비용 → 진입 장벽", "은행·보험 라이선스 획득 3~5년 소요", AMBER),
    ("🔒 전환 비용", "금융 데이터 이전 번거로움 → Lock-in", "계좌·자동이체·금융 히스토리 이전 불편", RED),
]
for i, (title, desc, example, color) in enumerate(moats):
    y = 1.65 + i * 1.15
    add_rect(slide, 0.3, y, 12.7, 1.0, LIGHT_GRAY)
    add_rect(slide, 0.3, y, 0.12, 1.0, color)
    add_textbox(slide, title, 0.55, y + 0.05, 3.0, 0.4,
                font_size=15, bold=True, color=color)
    add_textbox(slide, desc, 3.7, y + 0.05, 5.5, 0.4,
                font_size=13, color=NAVY)
    add_textbox(slide, f"예: {example}", 0.55, y + 0.55, 12.0, 0.35,
                font_size=11, color=MID_GRAY)

# 12 — Part 2 요약
slide = content_slide(13, "[중간 점검] 분석 프레임워크 한 장 요약",
                      "케이스 스터디 준비 완료")
steps = [
    ("① 비즈니스 모델 유형 분류", "결제/대출/투자/보험/B2B"),
    ("② 수익 구조 매핑", "수수료/이자/구독/데이터"),
    ("③ BMC로 전체 구조 시각화", "9개 블록 완성"),
    ("④ 핵심 KPI 정의", "GMV, Take Rate, CAC, LTV, Churn"),
    ("⑤ 경쟁 우위 분석", "네트워크 효과, 데이터 해자, 규제 해자"),
]
for i, (step, detail) in enumerate(steps):
    y = 1.7 + i * 0.95
    add_rect(slide, 0.3, y, 12.7, 0.8, LIGHT_GRAY if i % 2 == 0 else WHITE)
    add_textbox(slide, step, 0.5, y + 0.18, 6.5, 0.45,
                font_size=15, bold=True, color=NAVY)
    add_textbox(slide, detail, 7.2, y + 0.18, 5.5, 0.45,
                font_size=14, color=EMERALD, align=PP_ALIGN.RIGHT)
add_textbox(slide, "\"이 프레임워크를 들고 케이스 스터디로 갑니다 →\"",
            0.3, 6.65, 12.7, 0.4, font_size=13, bold=True, color=BLUE)

# ─── PART 3 SECTION ────────────────────────────────────────────────────────
section_divider(14, "PART 3  ·  1:00 – 2:00",
                "케이스 스터디",
                "토스 · 카카오페이 · Stripe — 데이터로 전략을 읽는다", AMBER)

# Toss overview
slide = content_slide(15, "Case 1 — 토스 (Toss) 개요",
                      "무료 송금으로 시작해 금융 슈퍼앱으로 진화")
add_rect(slide, 0.3, 1.6, 12.7, 1.3, LIGHT_GRAY)
timeline = [
    ("2014", "간편송금 출시", BLUE),
    ("2018", "카드 비교\n보험 레퍼럴", EMERALD),
    ("2021", "토스뱅크\n토스증권", AMBER),
    ("2024", "MAU 2,400만\n수익화 가속", RED),
]
for i, (year, event, color) in enumerate(timeline):
    x = 0.5 + i * 3.1
    add_rect(slide, x, 1.65, 2.8, 1.2, color)
    add_textbox(slide, year, x + 0.1, 1.7, 2.6, 0.4,
                font_size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, event, x + 0.1, 2.1, 2.6, 0.65,
                font_size=12, color=WHITE, align=PP_ALIGN.CENTER)

kpi_card(slide, 0.3, 3.1, 2.9, 1.2, "2,400만", "현재 MAU", BLUE)
kpi_card(slide, 3.35, 3.1, 2.9, 1.2, "30개+", "금융상품 수", EMERALD)
kpi_card(slide, 6.4, 3.1, 2.9, 1.2, "80조원+", "연간 결제처리금액", AMBER)
kpi_card(slide, 9.45, 3.1, 2.9, 1.2, "8~12%", "교차판매 전환율", RED)

bullet_block(slide, [
    "핵심 전략: \"금융의 모든 것을 하나의 앱에서\" — 슈퍼앱 전략",
    "수익 구조: 금융상품 중개 수수료 + 대출 이자 + 결제 수수료 + 광고",
    "Aha Moment: 첫 간편송금 완료 → 7일 리텐션 70%↑",
], 0.3, 4.45, 12.7, 2.3, font_size=14)

# Toss funnel
slide = content_slide(16, "Case 1 — 토스 퍼널 분석",
                      "가입 → 첫 거래 구간 이탈 최소화가 최우선 과제")
funnel_steps = [
    ("앱 다운로드", "100%", 12.0, BLUE),
    ("가입 완료 (본인인증)", "50%", 8.0, RGBColor(0x25, 0x63, 0xEB)),
    ("첫 송금/결제 완료  ← Aha Moment", "35%", 6.0, EMERALD),
    ("7일 내 재사용", "28%", 4.5, EMERALD),
    ("금융상품 1개+ 가입  ← 수익화 전환점", "14%", 2.5, AMBER),
]
for i, (label, pct, bar_w, color) in enumerate(funnel_steps):
    y = 1.65 + i * 0.95
    center_x = (13.33 - bar_w) / 2
    add_rect(slide, center_x, y, bar_w, 0.75, color)
    add_textbox(slide, label, center_x + 0.2, y + 0.15, bar_w - 2.5, 0.45,
                font_size=13, color=WHITE, bold=True)
    add_textbox(slide, pct, center_x + bar_w - 1.5, y + 0.15, 1.3, 0.45,
                font_size=18, bold=True, color=WHITE, align=PP_ALIGN.RIGHT)
add_textbox(slide, "💡  분석 포인트: 가입 → 첫 거래 구간 이탈 최소화가 최우선 과제",
            0.3, 6.6, 12.7, 0.4, font_size=13, bold=True, color=NAVY)

# Kakao Pay
slide = content_slide(17, "Case 2 — 카카오페이",
                      "결제에서 종합 금융 플랫폼으로 전환")
kpi_card(slide, 0.3, 1.6, 2.9, 1.2, "2,300만+", "MAU", BLUE)
kpi_card(slide, 3.35, 1.6, 2.9, 1.2, "120조원+", "연간 결제 거래액", EMERALD)
kpi_card(slide, 6.4, 1.6, 2.9, 1.2, "5,000만", "카카오톡 사용자\n기반", AMBER)
kpi_card(slide, 9.45, 1.6, 2.9, 1.2, "결제→금융", "수익 믹스 전환 중", RED)

add_textbox(slide, "플랫폼 전환 단계", 0.3, 3.0, 5.0, 0.4,
            font_size=14, bold=True, color=NAVY)
phases = [
    ("① 결제 MAU 확보", "무료/편의성"),
    ("② 결제 데이터 → 금융 프로파일링", "금융 행동 분석"),
    ("③ 신용/보험 크로스셀", "다중 상품 전환"),
    ("④ 투자 서비스 → 자산관리 플랫폼", "LTV 극대화"),
]
for i, (phase, detail) in enumerate(phases):
    y = 3.5 + i * 0.75
    add_rect(slide, 0.3, y, 12.7, 0.65, LIGHT_GRAY if i % 2 == 0 else WHITE)
    add_textbox(slide, phase, 0.5, y + 0.12, 7.0, 0.42,
                font_size=13, bold=True, color=NAVY)
    add_textbox(slide, detail, 8.0, y + 0.12, 4.8, 0.42,
                font_size=12, color=EMERALD, align=PP_ALIGN.RIGHT)

# GMV vs Take Rate Dilemma
slide = content_slide(18, "Case 2 — GMV vs. Take Rate 딜레마",
                      "성장과 수익성 사이의 핵심 긴장 관계")
add_rect(slide, 0.3, 1.6, 6.0, 4.5, LIGHT_GRAY)
add_textbox(slide, "📈 TPV 성장", 0.5, 1.75, 5.5, 0.5, font_size=16, bold=True, color=EMERALD)
add_textbox(slide, "120조원+ 거래액\nYoY 빠른 성장\n\n✅ 좋음", 0.5, 2.4, 5.5, 2.5,
            font_size=22, bold=True, color=EMERALD)
add_rect(slide, 6.5, 1.6, 6.5, 4.5, LIGHT_GRAY)
add_textbox(slide, "📉 Take Rate 압박", 6.7, 1.75, 6.0, 0.5, font_size=16, bold=True, color=RED)
add_textbox(slide, "카드 MDR 규제 인하\nTake Rate 낮음\n\n⚠️ 개선 필요", 6.7, 2.4, 6.0, 2.5,
            font_size=22, bold=True, color=RED)
add_textbox(slide, "해결책: 고마진 부문(보험·투자)으로 믹스 전환 → \"Revenue Quality\" 개선",
            0.3, 6.3, 12.7, 0.5, font_size=14, bold=True, color=NAVY)

# Stripe
slide = content_slide(19, "Case 3 — Stripe 개요",
                      "개발자 친화적 API로 글로벌 결제 인프라를 장악")
kpi_card(slide, 0.3, 1.6, 2.9, 1.2, "$65B+", "기업 가치", BLUE)
kpi_card(slide, 3.35, 1.6, 2.9, 1.2, "135개국", "서비스 국가", EMERALD)
kpi_card(slide, 6.4, 1.6, 2.9, 1.2, "100만+", "비즈니스 고객", AMBER)
kpi_card(slide, 9.45, 1.6, 2.9, 1.2, "2.9%+$0.30", "기본 수수료/건", RED)

bullet_block(slide, [
    "핵심 가치 제안: \"7줄의 코드로 결제 구현\" — 개발자 경험 최우선",
    "주요 고객: Shopify, Lyft, DoorDash, Amazon",
    "확장: Stripe Atlas (법인설립), Stripe Radar (사기탐지), Stripe Treasury (금융 서비스형)",
    "GTM 전략: Bottom-up — 개발자가 먼저 도입 → 기업 전체로 확산",
    "LTV 모델: 고객이 성장할수록 GMV 증가 → Stripe 수익 비례 증가",
], 0.3, 3.0, 12.7, 3.8, font_size=14)

# Stripe Radar (Data Flywheel)
slide = content_slide(20, "Case 3 — 데이터로 만든 경쟁 우위: Stripe Radar",
                      "거래 데이터가 AI 제품이 되는 선순환 구조")
flywheel = [
    ("수억 건\n거래 데이터", BLUE, 1.0, 3.2),
    ("ML 모델\n학습", EMERALD, 4.0, 2.0),
    ("Radar\n사기 탐지", AMBER, 7.5, 2.0),
    ("추가 수익\n창출", RED, 10.5, 3.2),
    ("데이터\n확장", RGBColor(0x7C, 0x3A, 0xED), 7.5, 4.8),
    ("더 많은\n가맹점", BLUE, 4.0, 4.8),
]
for (label, color, x, y) in flywheel:
    sh = slide.shapes.add_shape(9, Inches(x), Inches(y), Inches(1.8), Inches(1.4))
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    sh.line.fill.background()
    add_textbox(slide, label, x + 0.1, y + 0.25, 1.6, 0.9,
                font_size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_textbox(slide, "🔄 데이터 플라이휠 선순환",
            0.3, 1.65, 12.7, 0.4, font_size=14, bold=True, color=NAVY)
add_textbox(slide, "결과: 99.9%+ 결제 성공률, 사기 탐지 정확도 업계 최고\n\"데이터 자산 → AI 제품 → 추가 수익 → 데이터 확장\" 선순환",
            0.3, 6.25, 12.7, 0.8, font_size=13, color=NAVY)

# Comparison
slide = content_slide(21, "케이스 스터디 비교 분석",
                      "공통점: 데이터를 핵심 자산으로 활용 | 차이점: 고객·진입·수익화 전략")
headers = ["구분", "토스", "카카오페이", "Stripe"]
rows = [
    ("고객", "B2C 개인", "B2C 개인", "B2B 기업"),
    ("진입 전략", "무료 기능", "생태계 기반", "개발자 경험"),
    ("수익화", "크로스셀", "믹스 전환", "거래 과금"),
    ("데이터 활용", "신용평가", "금융 프로파일링", "사기 탐지 AI"),
    ("핵심 해자", "사용자 데이터", "카카오 생태계", "개발자 네트워크"),
]
col_w = [2.5, 3.2, 3.2, 3.2]
col_x = [0.3, 2.9, 6.2, 9.5]
col_colors = [NAVY, BLUE, EMERALD, AMBER]

for j, (header, cx, cw, color) in enumerate(zip(headers, col_x, col_w, col_colors)):
    add_rect(slide, cx, 1.65, cw - 0.1, 0.55, color)
    add_textbox(slide, header, cx + 0.1, 1.72, cw - 0.2, 0.42,
                font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

for i, row in enumerate(rows):
    for j, (cell, cx, cw, color) in enumerate(zip(row, col_x, col_w, col_colors)):
        y = 2.3 + i * 0.85
        bg = LIGHT_GRAY if i % 2 == 0 else WHITE
        add_rect(slide, cx, y, cw - 0.1, 0.75, bg)
        bold = (j == 0)
        clr = NAVY if j == 0 else (color if j > 0 else DARK_TEXT)
        add_textbox(slide, cell, cx + 0.1, y + 0.15, cw - 0.2, 0.45,
                    font_size=12, bold=bold, color=clr, align=PP_ALIGN.CENTER)

add_textbox(slide, "\"여러분의 서비스는 어떤 해자를 가지고 있나요?\"",
            0.3, 6.65, 12.7, 0.4, font_size=13, bold=True, color=NAVY)

# Checklist
slide = content_slide(22, "내 서비스에 적용하기 — 분석 체크리스트",
                      "오늘 배운 프레임워크를 내 업무에 즉시 적용")
checks = [
    "우리 서비스의 비즈니스 모델 유형은?",
    "주요 수익원 Top 3는?",
    "핵심 고객 세그먼트는? 각각의 LTV는?",
    "지금 추적하고 있는 KPI는? 빠진 KPI는?",
    "우리의 경쟁 우위(해자)는 무엇인가?",
    "데이터 자산을 어떻게 활용하고 있는가?",
]
for i, check in enumerate(checks):
    y = 1.7 + i * 0.85
    add_rect(slide, 0.3, y, 12.7, 0.75, LIGHT_GRAY if i % 2 == 0 else WHITE)
    add_rect(slide, 0.3, y, 0.75, 0.75, BLUE)
    add_textbox(slide, "□", 0.3, y + 0.1, 0.75, 0.55,
                font_size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, check, 1.2, y + 0.18, 11.5, 0.42,
                font_size=14, color=NAVY)

# ─── PART 4 SECTION ────────────────────────────────────────────────────────
section_divider(23, "PART 4  ·  2:00 – 3:00",
                "실습: Amplitude + Looker Studio",
                "가상 서비스 FinFlow 데이터로 직접 인사이트를 도출하라",
                RGBColor(0x7C, 0x3A, 0xED))

# FinFlow intro
slide = content_slide(24, "실습 개요 — 가상 서비스 FinFlow",
                      "\"실제 스타트업에 입사 첫 주에 이 데이터를 받았다고 상상하세요\"")
kpi_card(slide, 0.3, 1.6, 2.9, 1.2, "5만 MAU", "현재 규모", RGBColor(0x7C, 0x3A, 0xED))
kpi_card(slide, 3.35, 1.6, 2.9, 1.2, "저조", "결제 전환율", RED)
kpi_card(slide, 6.4, 1.6, 2.9, 1.2, "하락 중", "30일 리텐션", AMBER)
kpi_card(slide, 9.45, 1.6, 2.9, 1.2, "90일", "데이터 기간", BLUE)

add_textbox(slide, "🎯  미션", 0.3, 3.0, 3.0, 0.45, font_size=14, bold=True, color=NAVY)
missions = [
    "① 결제 퍼널에서 가장 큰 이탈 구간 찾기 (Amplitude)",
    "② 30일 리텐션 패턴 분석 (Amplitude)",
    "③ GMV & 결제 성공률 대시보드 구성 (Looker Studio)",
]
bullet_block(slide, missions, 0.3, 3.5, 12.7, 1.8, font_size=15)

add_textbox(slide, "📊  데이터셋", 0.3, 5.4, 3.0, 0.45, font_size=14, bold=True, color=NAVY)
datasets = [
    "user_events.csv — 사용자 이벤트 10,000건 (app_open, signup, 계좌연결, 결제 등)",
    "transactions.csv — 거래 데이터 5,000건 (금액, 상태, 결제수단, 카테고리)",
]
bullet_block(slide, datasets, 0.3, 5.9, 12.7, 1.3, font_size=13)

# Amplitude Funnel Lab
slide = content_slide(25, "실습 A-1: Amplitude 퍼널 분석",
                      "전환 퍼널에서 가장 큰 이탈 구간 찾기")
funnel_events = [
    ("Event 1", "app_open", "앱 실행", "100%"),
    ("Event 2", "signup_complete", "회원가입 완료", "62%"),
    ("Event 3", "account_link", "계좌 연결  ← 주요 이탈 구간", "34%"),
    ("Event 4", "first_payment_attempt", "첫 결제 시도", "28%"),
    ("Event 5", "first_payment_success", "첫 결제 성공  ← Aha Moment", "21%"),
]
for i, (event_id, event_name, label, pct) in enumerate(funnel_events):
    y = 1.65 + i * 0.95
    bar_w = 11.0 * (5 - i) / 5 + 1.0
    cx = (13.33 - bar_w) / 2
    color = BLUE if i < 2 else (RED if i == 2 else (EMERALD if i == 4 else AMBER))
    add_rect(slide, cx, y, bar_w, 0.75, color)
    add_textbox(slide, f"{event_id}: {label}", cx + 0.2, y + 0.15, bar_w - 2.0, 0.45,
                font_size=12, color=WHITE, bold=True)
    add_textbox(slide, pct, cx + bar_w - 1.5, y + 0.15, 1.3, 0.45,
                font_size=18, bold=True, color=WHITE, align=PP_ALIGN.RIGHT)
add_textbox(slide, "⚠️  계좌 연결 단계에서 최대 이탈 (62% → 34% = -28%p) — 가장 큰 병목 구간",
            0.3, 6.6, 12.7, 0.4, font_size=13, bold=True, color=RED)

# Retention
slide = content_slide(26, "실습 A-2: Amplitude 리텐션 분석",
                      "가입 후 사용자가 얼마나 돌아오는가?")
add_textbox(slide, "리텐션 커브 — 가입 코호트별 재방문률", 0.3, 1.65, 12.7, 0.4,
            font_size=14, bold=True, color=NAVY)
# Simplified retention curve visualization
for week, vals in enumerate([100, 52, 38, 30, 24, 20, 18]):
    x = 0.5 + week * 1.8
    # Bar for general
    add_rect(slide, x, 6.0 - vals * 0.04, 0.7, vals * 0.04, BLUE)
    add_textbox(slide, f"W{week}", x, 6.1, 0.7, 0.3, font_size=10,
                color=DARK_TEXT, align=PP_ALIGN.CENTER)
    add_textbox(slide, f"{vals}%", x, 5.85 - vals * 0.04, 0.7, 0.3,
                font_size=10, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

add_textbox(slide, "코호트 비교: iOS vs Android | 첫 결제 금액대별 리텐션 차이",
            0.3, 6.35, 12.7, 0.35, font_size=12, color=MID_GRAY)
add_textbox(slide, "💡  핵심 질문: \"어떤 세그먼트의 리텐션이 가장 높은가? 왜?\"",
            0.3, 6.75, 12.7, 0.35, font_size=13, bold=True, color=NAVY)

# Cohort
slide = content_slide(27, "실습 A-3: 코호트 분석",
                      "첫 거래 패턴으로 미래 LTV 예측하기")
cohorts = [
    ("Group A", "첫 결제 1만원 미만", "낮은 LTV", BLUE),
    ("Group B", "첫 결제 1~5만원", "중간 LTV", EMERALD),
    ("Group C", "첫 결제 5만원+", "높은 LTV", AMBER),
]
for i, (group, definition, ltv, color) in enumerate(cohorts):
    x = 0.3 + i * 4.3
    add_rect(slide, x, 1.65, 4.0, 3.5, color)
    add_textbox(slide, group, x + 0.2, 1.8, 3.6, 0.6,
                font_size=22, bold=True, color=WHITE)
    add_textbox(slide, definition, x + 0.2, 2.5, 3.6, 0.6,
                font_size=14, color=WHITE)
    add_textbox(slide, ltv, x + 0.2, 4.4, 3.6, 0.5,
                font_size=14, bold=True, color=WHITE)

add_textbox(slide, "가설: \"첫 결제 금액이 높을수록 LTV가 높다\"",
            0.3, 5.4, 12.7, 0.45, font_size=14, bold=True, color=NAVY)
add_textbox(slide, "비교 지표: 재결제율 | 평균 결제 횟수 | 30일 GMV",
            0.3, 5.95, 12.7, 0.35, font_size=13, color=MID_GRAY)
add_textbox(slide, "마케팅 적용: 고가 첫 결제 사용자를 더 많이 획득하는 채널 집중 투자",
            0.3, 6.45, 12.7, 0.35, font_size=13, color=EMERALD)

# Looker Studio
slide = content_slide(28, "실습 B: Looker Studio 대시보드",
                      "핵심 지표를 한눈에 볼 수 있는 경영진 대시보드")
add_textbox(slide, "대시보드 레이아웃 설계", 0.3, 1.65, 12.7, 0.4,
            font_size=14, bold=True, color=NAVY)
# Dashboard mockup
add_rect(slide, 0.3, 2.1, 12.7, 1.0, NAVY)
kpi_items = ["MAU: 50,124", "TPV: 48.2억", "결제성공률: 94.2%", "CAC: 3,200원", "LTV/CAC: 4.1"]
for i, kpi in enumerate(kpi_items):
    x = 0.5 + i * 2.5
    add_rect(slide, x, 2.2, 2.2, 0.8, BLUE)
    add_textbox(slide, kpi, x + 0.1, 2.3, 2.0, 0.6,
                font_size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_rect(slide, 0.3, 3.2, 7.8, 2.3, LIGHT_GRAY)
add_textbox(slide, "월별 MAU 추이 (Line Chart)", 0.5, 3.3, 7.4, 0.4, font_size=12, bold=True, color=NAVY)
add_rect(slide, 8.2, 3.2, 4.8, 2.3, LIGHT_GRAY)
add_textbox(slide, "채널별 유저 비율 (Bar)", 8.4, 3.3, 4.4, 0.4, font_size=12, bold=True, color=NAVY)
add_rect(slide, 0.3, 5.65, 12.7, 1.4, LIGHT_GRAY)
add_textbox(slide, "퍼널 전환율 (Funnel Chart)", 0.5, 5.75, 12.3, 0.4, font_size=12, bold=True, color=NAVY)

# Dashboard principles
slide = content_slide(29, "대시보드 설계 원칙",
                      "좋은 대시보드의 3가지 조건")
principles = [
    ("👁  한 눈에", "5초 안에 핵심이 보여야 한다",
     "KPI 카드를 상단에 배치, 가장 중요한 숫자를 큰 글씨로", BLUE),
    ("⚡ 행동 가능", "\"그래서 뭘 해야 하는가?\"를 알 수 있어야 한다",
     "이상 감지 알림, 임계값 표시, 드릴다운 가능한 구조", EMERALD),
    ("✅ 신뢰 가능", "데이터 정의와 집계 방식이 명확해야 한다",
     "지표 정의 툴팁, 데이터 업데이트 시간 표시, 데이터 소스 명시", AMBER),
]
for i, (title, desc, impl, color) in enumerate(principles):
    y = 1.7 + i * 1.6
    add_rect(slide, 0.3, y, 12.7, 1.4, LIGHT_GRAY)
    add_rect(slide, 0.3, y, 0.12, 1.4, color)
    add_textbox(slide, title, 0.55, y + 0.1, 3.5, 0.55,
                font_size=18, bold=True, color=color)
    add_textbox(slide, desc, 0.55, y + 0.65, 12.0, 0.45,
                font_size=14, color=NAVY, bold=True)
    add_textbox(slide, f"실천: {impl}", 0.55, y + 1.05, 12.0, 0.35,
                font_size=11, color=MID_GRAY)

# Lab results sharing
slide = content_slide(30, "실습 결과 공유 — 인사이트 발표",
                      "다양한 팀의 관점을 비교하는 것 자체가 학습")
add_textbox(slide, "팀별 발표 (1팀 2분)", 0.3, 1.65, 12.7, 0.4, font_size=14, bold=True, color=NAVY)
items = [
    ("발견한 핵심 이탈 구간", "퍼널 분석 결과", BLUE),
    ("리텐션이 낮은 세그먼트", "코호트 분석 결과", EMERALD),
    ("대시보드에서 발견한 이상 패턴", "Looker Studio 결과", AMBER),
]
for i, (title, sub, color) in enumerate(items):
    y = 2.2 + i * 1.1
    add_rect(slide, 0.3, y, 12.7, 0.9, LIGHT_GRAY)
    add_rect(slide, 0.3, y, 0.12, 0.9, color)
    add_textbox(slide, f"{i+1}. {title}", 0.55, y + 0.1, 8.0, 0.4,
                font_size=15, bold=True, color=NAVY)
    add_textbox(slide, sub, 0.55, y + 0.5, 8.0, 0.3, font_size=12, color=MID_GRAY)

add_rect(slide, 0.3, 5.7, 12.7, 1.0, RGBColor(0xEF, 0xF6, 0xFF))
add_textbox(slide, "💬  토론: \"만약 이 데이터가 실제 서비스라면 다음 액션은?\"",
            0.5, 5.85, 12.3, 0.5, font_size=14, bold=True, color=NAVY)

# ─── PART 5 SECTION ────────────────────────────────────────────────────────
section_divider(31, "PART 5  ·  2:40 – 3:00",
                "마무리 & Q&A",
                "3시간의 학습을 정리하고 실무에 적용합니다", RED)

# Wrap-up summary
slide = content_slide(32, "오늘 배운 것 정리",
                      "3시간 학습 내용을 한 장으로 압축")
parts = [
    ("Part 1", "핀테크 데이터 분석의 필요성과 오늘의 목표", BLUE),
    ("Part 2", "비즈니스 모델 5대 유형 + BMC + 10대 KPI + 경쟁 우위 분석", EMERALD),
    ("Part 3", "토스(슈퍼앱) · 카카오페이(플랫폼 전환) · Stripe(API 경제) 케이스", AMBER),
    ("Part 4", "Amplitude 퍼널/리텐션/코호트 + Looker Studio 대시보드 실습", RGBColor(0x7C, 0x3A, 0xED)),
]
for i, (part, desc, color) in enumerate(parts):
    y = 1.7 + i * 1.15
    add_rect(slide, 0.3, y, 12.7, 1.0, LIGHT_GRAY)
    add_rect(slide, 0.3, y, 0.12, 1.0, color)
    add_textbox(slide, part, 0.55, y + 0.1, 1.5, 0.4, font_size=14, bold=True, color=color)
    add_textbox(slide, desc, 2.1, y + 0.1, 10.7, 0.8, font_size=14, color=NAVY)

add_rect(slide, 0.3, 6.4, 12.7, 0.7, NAVY)
add_textbox(slide, "🎯  핵심 메시지: \"비즈니스 구조를 이해하고, KPI를 정의하고, 데이터로 검증하라\"",
            0.5, 6.55, 12.3, 0.45, font_size=14, bold=True, color=WHITE)

# Action items
slide = content_slide(33, "실무 적용 팁 — 내일 당장 할 수 있는 것",
                      "강의 내용을 월요일 업무에 바로 적용하는 구체적 액션")
action_times = [
    ("🌙  오늘 (Tonight)", ["내 서비스 BMC 빈칸 채워보기", "비즈니스 모델 유형 분류 해보기"], NAVY),
    ("📅  이번 주 (This Week)", [
        "Amplitude / Looker Studio 무료 계정 세팅",
        "서비스의 KPI 5개 정의하고 현재 수치 파악",
    ], BLUE),
    ("📆  이번 달 (This Month)", [
        "퍼널 분석 1개 완료 → 이탈 구간 가설 제시",
        "리텐션 분석 → 고LTV 세그먼트 정의",
    ], EMERALD),
]
for i, (time_label, actions, color) in enumerate(action_times):
    y = 1.7 + i * 1.65
    add_rect(slide, 0.3, y, 12.7, 1.5, LIGHT_GRAY)
    add_rect(slide, 0.3, y, 0.12, 1.5, color)
    add_textbox(slide, time_label, 0.55, y + 0.1, 4.0, 0.45,
                font_size=15, bold=True, color=color)
    for j, action in enumerate(actions):
        add_textbox(slide, f"•  {action}", 0.55, y + 0.65 + j * 0.4, 12.0, 0.35,
                    font_size=13, color=DARK_TEXT)

# Resources
slide = content_slide(34, "추천 학습 리소스",
                      "더 깊이 공부하고 싶은 분들을 위한 다음 단계")
resource_cols = [
    ("📚 책", [
        "Lean Analytics — 스타트업 지표 바이블",
        "Hacking Growth — 그로스 해킹 실전",
        "Fintech in Korea — 핀테크 산업 개요",
    ], BLUE),
    ("📰 뉴스레터/블로그", [
        "핀다 블로그 (핀테크 분석)",
        "a16z Fintech Newsletter",
        "Andreessen Horowitz Fintech",
    ], EMERALD),
    ("🛠 도구 학습", [
        "Amplitude Academy (무료 온라인)",
        "Looker Studio 공식 튜토리얼",
        "Kaggle Fintech Datasets",
    ], AMBER),
]
for i, (title, items, color) in enumerate(resource_cols):
    x = 0.3 + i * 4.3
    add_rect(slide, x, 1.65, 4.0, 4.8, LIGHT_GRAY)
    add_rect(slide, x, 1.65, 4.0, 0.55, color)
    add_textbox(slide, title, x + 0.15, 1.72, 3.7, 0.42,
                font_size=14, bold=True, color=WHITE)
    for j, item in enumerate(items):
        add_textbox(slide, f"•  {item}", x + 0.15, 2.35 + j * 0.7, 3.7, 0.6,
                    font_size=12, color=DARK_TEXT)

# Q&A
slide = content_slide(35, "Q & A",
                      "궁금한 것을 풀어가는 시간 — 15분")
faqs = [
    ("Q", "Amplitude vs. Mixpanel, 어떤 걸 써야 하나요?",
     "사용량 기준: Amplitude(이벤트 기반 코호트/리텐션 강점), Mixpanel(실시간 분석 강점). 핀테크라면 Amplitude 권장."),
    ("Q", "Take Rate는 어떻게 벤치마킹하나요?",
     "공개 IR 자료 (SEC 10-K, 분기보고서) + 업계 리포트 (CB Insights, PitchBook). 결제사 평균 0.5~2.5%."),
    ("Q", "소규모 팀에서도 이런 분석이 가능한가요?",
     "Amplitude 무료 플랜 (월 10만 이벤트) + Looker Studio (무료) 조합으로 충분히 시작 가능합니다."),
]
for i, (tag, q, a) in enumerate(faqs):
    y = 1.7 + i * 1.65
    add_rect(slide, 0.3, y, 12.7, 1.5, LIGHT_GRAY)
    add_rect(slide, 0.3, y, 0.8, 1.5, BLUE)
    add_textbox(slide, tag, 0.3, y + 0.45, 0.8, 0.6,
                font_size=24, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, q, 1.25, y + 0.1, 11.5, 0.45,
                font_size=14, bold=True, color=NAVY)
    add_textbox(slide, a, 1.25, y + 0.6, 11.5, 0.75, font_size=12, color=DARK_TEXT)

# Closing
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, 13.33, 7.5, SECTION_BG)
add_rect(slide, 0, 0, 0.12, 7.5, EMERALD)
add_rect(slide, 0, 6.5, 13.33, 0.12, EMERALD)
add_textbox(slide, "감사합니다", 0.8, 1.5, 11, 1.2,
            font_size=52, bold=True, color=WHITE)
add_textbox(slide,
            "\"오늘 배운 프레임워크로 여러분 서비스의 데이터를 다시 보세요\"\n데이터로 더 나은 결정을 — Good luck!",
            0.8, 2.9, 11, 1.2, font_size=20, color=MID_GRAY)
add_textbox(slide, "📧  강사 연락처 및 LinkedIn — 강의 자료 다운로드 QR", 0.8, 4.5, 11, 0.5,
            font_size=14, color=MID_GRAY)
add_textbox(slide, "HULK × PEPPER  |  Jarvis System", 0.8, 6.6, 11, 0.4,
            font_size=12, bold=True, color=EMERALD)
add_rect(slide, 0, 7.35, 13.33, 0.07, EMERALD)

# ─── APPENDIX SECTION ───────────────────────────────────────────────────────
section_divider(37, "APPENDIX",
                "부록 — 심화 & 참고용",
                "규제 환경 · 사기 탐지 · Unit Economics · 데이터 파이프라인 · 커리어", MID_GRAY)

# A1 — Regulation
slide = content_slide(38, "[부록] 핀테크 규제 환경과 데이터 분석",
                      "규제가 데이터 활용과 비즈니스 모델에 미치는 영향")
regs = [
    ("전자금융거래법", "전자결제 서비스 허가 및 운영 규제", "PG사 허가 필수"),
    ("개인정보보호법", "금융 데이터 수집·처리·분석 제한", "익명화·가명화 필수"),
    ("금융소비자보호법", "금융상품 추천/판매 프로세스 규제", "적합성 원칙"),
    ("마이데이터(본인신용정보관리업)", "금융 데이터 통합 → 분석 기회", "오픈 API 기반"),
    ("오픈뱅킹", "API 기반 계좌 접근 서비스", "서비스 확장 기반"),
]
for i, (reg, desc, note) in enumerate(regs):
    y = 1.7 + i * 0.95
    add_rect(slide, 0.3, y, 12.7, 0.8, LIGHT_GRAY if i % 2 == 0 else WHITE)
    add_textbox(slide, reg, 0.5, y + 0.18, 3.0, 0.45, font_size=13, bold=True, color=NAVY)
    add_textbox(slide, desc, 3.7, y + 0.18, 6.0, 0.45, font_size=13, color=DARK_TEXT)
    add_textbox(slide, note, 9.8, y + 0.18, 3.0, 0.45, font_size=12, color=EMERALD, align=PP_ALIGN.RIGHT)
add_textbox(slide, "💡  \"규제는 제약이자 새로운 시장의 출발점\"",
            0.3, 6.6, 12.7, 0.4, font_size=13, bold=True, color=NAVY)

# A2 — Fraud
slide = content_slide(39, "[부록] 사기 탐지 (Fraud Detection) 기초",
                      "핀테크 데이터 분석의 중요한 응용 — 이상 탐지")
fraud_types = [
    ("계정 탈취\n(ATO)", "비정상 로그인", RED),
    ("신용카드 사기", "이상 결제 패턴", AMBER),
    ("보이스피싱", "비정상 송금", BLUE),
    ("자금세탁\n(AML)", "구조화 거래", RGBColor(0x7C, 0x3A, 0xED)),
]
for i, (fraud, signal, color) in enumerate(fraud_types):
    x = 0.3 + i * 3.1
    add_rect(slide, x, 1.65, 2.9, 1.5, color)
    add_textbox(slide, fraud, x + 0.15, 1.8, 2.6, 0.65,
                font_size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, f"신호: {signal}", x + 0.15, 2.55, 2.6, 0.45,
                font_size=11, color=WHITE, align=PP_ALIGN.CENTER)

bullet_block(slide, [
    "이상 탐지 지표: 비정상적 거래 시간대, 위치 불일치, 금액 패턴 이상",
    "Rule-based → ML-based 단계적 고도화 (Stripe Radar 사례)",
    "Amplitude 활용: 비정상 행동 코호트 정의 → 세그먼트 분리 분석",
    "머신러닝: 이진 분류 모델 (정상/사기) — 수억 건 학습 → 정확도 향상",
], 0.3, 3.35, 12.7, 2.8, font_size=14)

# A3 — Unit Economics Advanced
slide = content_slide(40, "[부록] Unit Economics 심화",
                      "핀테크 투자/평가에서 핵심 지표인 단위 경제성 계산")
# LTV/CAC calculation example
add_rect(slide, 0.3, 1.65, 5.8, 2.8, LIGHT_GRAY)
add_textbox(slide, "LTV 고급 모델", 0.5, 1.75, 5.4, 0.4, font_size=14, bold=True, color=NAVY)
add_textbox(slide,
            "LTV = (ARPU × Gross Margin) ÷ Churn Rate\n\n예: ARPU 5,000원 × 70% ÷ 4%\n    = 87,500원",
            0.5, 2.25, 5.4, 2.0, font_size=14, color=DARK_TEXT)
add_rect(slide, 6.4, 1.65, 5.8, 2.8, LIGHT_GRAY)
add_textbox(slide, "CAC 분해", 6.6, 1.75, 5.4, 0.4, font_size=14, bold=True, color=NAVY)
add_textbox(slide,
            "CAC = (마케팅비 + 영업비) ÷ 신규 유료 고객 수\n\n예: 1억 5천만원 ÷ 50,000명\n    = 3,000원/명",
            6.6, 2.25, 5.4, 2.0, font_size=14, color=DARK_TEXT)

benchmark = [
    ("LTV/CAC < 1", "위험 — 돈 잃으며 성장 중", RED),
    ("LTV/CAC 1~3", "주의 — 효율 개선 필요", AMBER),
    ("LTV/CAC > 3", "건강 — 지속 투자 가능", EMERALD),
    ("LTV/CAC > 5", "확장 — 마케팅 투자 여력 충분", BLUE),
]
for i, (level, desc, color) in enumerate(benchmark):
    y = 4.6 + i * 0.55
    add_rect(slide, 0.3, y, 12.7, 0.5, LIGHT_GRAY if i % 2 == 0 else WHITE)
    add_rect(slide, 0.3, y, 3.5, 0.5, color)
    add_textbox(slide, level, 0.5, y + 0.08, 3.1, 0.35, font_size=13, bold=True, color=WHITE)
    add_textbox(slide, desc, 4.0, y + 0.08, 9.0, 0.35, font_size=13, color=DARK_TEXT)

# A4 — Data Pipeline
slide = content_slide(41, "[부록] 핀테크 데이터 파이프라인 구조",
                      "\"데이터 파이프라인이 없으면 분석도 없다\"")
pipeline = [
    ("📱 데이터 소스", "앱 이벤트 / 결제 시스템 / 고객 DB / 외부 신용 데이터", BLUE),
    ("📡 수집", "Amplitude SDK · 서버 이벤트 → Event Tracking", EMERALD),
    ("🗄 저장", "Data Warehouse: BigQuery · Snowflake · Redshift", AMBER),
    ("🔍 분석", "Amplitude (행동 분석) + Looker Studio (BI) + SQL", RGBColor(0x7C, 0x3A, 0xED)),
    ("💡 활용", "마케팅 세그먼트 · 제품 개선 · 리스크 관리 · 경영 보고", RED),
]
for i, (stage, desc, color) in enumerate(pipeline):
    y = 1.65 + i * 1.05
    if i < len(pipeline) - 1:
        add_textbox(slide, "↓", 6.4, y + 0.85, 0.5, 0.3, font_size=16, color=MID_GRAY, align=PP_ALIGN.CENTER)
    add_rect(slide, 0.3, y, 12.7, 0.85, LIGHT_GRAY)
    add_rect(slide, 0.3, y, 0.12, 0.85, color)
    add_textbox(slide, stage, 0.55, y + 0.1, 3.0, 0.45, font_size=14, bold=True, color=color)
    add_textbox(slide, desc, 3.7, y + 0.1, 9.0, 0.65, font_size=13, color=DARK_TEXT)

# A5 — Career Path
slide = content_slide(42, "[부록] 핀테크 분석가 커리어 패스",
                      "오늘 배운 역량이 어떤 직무로 연결되는가")
careers = [
    ("그로스 분석가", "Growth Analyst", "MAU/CAC/LTV 분석, 그로스 실험 설계", BLUE),
    ("프로덕트 분석가", "Product Analyst", "퍼널·리텐션·코호트, UX 최적화", EMERALD),
    ("BI 분석가", "Business Intelligence", "경영진 대시보드, 데이터 시각화", AMBER),
    ("데이터 사이언티스트", "Data Scientist (Fintech)", "신용평가 모델, 사기 탐지 AI", RED),
    ("전략 기획", "Strategy Planning", "비즈니스 모델 분석, M&A 실사", RGBColor(0x7C, 0x3A, 0xED)),
]
for i, (kr_title, en_title, desc, color) in enumerate(careers):
    col = i % 2
    row = i // 2 if i < 4 else 2
    if i == 4:
        x, y = 3.5, 1.65 + row * 1.2
    else:
        x = 0.3 + col * 6.5
        y = 1.65 + row * 1.2
    add_rect(slide, x, y, 6.2, 1.0, LIGHT_GRAY)
    add_rect(slide, x, y, 0.1, 1.0, color)
    add_textbox(slide, kr_title, x + 0.25, y + 0.08, 3.5, 0.4, font_size=14, bold=True, color=color)
    add_textbox(slide, en_title, x + 0.25, y + 0.5, 3.0, 0.35, font_size=11, color=MID_GRAY)
    add_textbox(slide, desc, x + 3.5, y + 0.25, 2.5, 0.5, font_size=11, color=DARK_TEXT)

add_textbox(slide, "필수 역량: SQL · Python(pandas) · BI 도구 · 금융 도메인 이해",
            0.3, 6.3, 12.7, 0.35, font_size=12, color=NAVY, bold=True)
add_textbox(slide, "포트폴리오: Kaggle · Dacon · 공공데이터포털 오픈 데이터 활용",
            0.3, 6.7, 12.7, 0.35, font_size=12, color=MID_GRAY)

# Final Q&A back cover
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, 13.33, 7.5, NAVY)
add_rect(slide, 0, 3.5, 13.33, 0.08, BLUE)
add_textbox(slide, "Q & A", 0.8, 0.8, 11, 2.0,
            font_size=80, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_textbox(slide, "질문은 언제든지 환영합니다",
            0.8, 3.0, 11, 0.6, font_size=20, color=MID_GRAY, align=PP_ALIGN.CENTER)
add_textbox(slide, "강의 자료 및 샘플 데이터 공유: 구글 드라이브 링크 (강사가 현장 공유)",
            0.8, 4.2, 11, 0.5, font_size=14, color=MID_GRAY, align=PP_ALIGN.CENTER)
add_textbox(slide, "연락처: 슬랙 채널 / 이메일 (강사 입력)",
            0.8, 4.85, 11, 0.45, font_size=14, color=MID_GRAY, align=PP_ALIGN.CENTER)
add_rect(slide, 0, 7.35, 13.33, 0.07, BLUE)
add_textbox(slide, "HULK × PEPPER  ·  Jarvis System  ·  2026", 0.8, 7.1, 11, 0.3,
            font_size=11, color=MID_GRAY, align=PP_ALIGN.CENTER)

# ─── Save ───────────────────────────────────────────────────────────────────
output_path = r"C:\Agent\pepper\output\fintech-bizanalysis-lecture-slides.pptx"
prs.save(output_path)
print(f"[OK] Saved: {output_path}")
print(f"     Slides: {len(prs.slides)}")
