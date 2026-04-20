# -*- coding: utf-8 -*-
"""
CleverTap 한국어 제안서 v5 — 슬라이드별 컨텍스트 최적화 이미지 교체판
- 이전 v4 슬라이드 구조 동일
- 모든 이미지를 슬라이드 주제에 맞는 UI 목업으로 교체
  ctx_analytics.png  → 실시간 분석 슬라이드
  ctx_segmentation   → 세그멘테이션 슬라이드
  ctx_engagement     → 멀티채널 인게이지먼트 슬라이드
  ctx_ai             → AI/ML 슬라이드
  ctx_ab             → A/B 테스트 슬라이드
  ctx_architecture   → 제품 아키텍처 슬라이드
  ctx_customers      → 고객사 슬라이드
  ctx_fintech        → 핀테크 슬라이드
  ctx_intro          → CleverTap 소개 슬라이드
  screenshot_main    → 표지 (homepage 유지)
"""
import sys, io
if sys.platform == 'win32':
    sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pathlib import Path
from PIL import Image
import os

ASSETS = Path("C:/Agent/pepper/output/clevertap_assets")
OUTPUT = Path("C:/Agent/pepper/output/clevertap_proposal_kr_v6.pptx")

C_PURPLE      = RGBColor(0xFF, 0x6B, 0x35)
C_PURPLE_DARK = RGBColor(0xCC, 0x44, 0x11)
C_PURPLE_LIGHT= RGBColor(0xFF, 0xA0, 0x70)
C_DARK        = RGBColor(0x1A, 0x1A, 0x2E)
C_DARK2       = RGBColor(0x16, 0x21, 0x3E)
C_MID         = RGBColor(0x2D, 0x2D, 0x4E)
C_WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
C_GRAY        = RGBColor(0x6C, 0x75, 0x8D)
C_LIGHT_GRAY  = RGBColor(0xF4, 0xF4, 0xF8)
C_ACCENT      = RGBColor(0xFF, 0x6B, 0x35)
C_GREEN       = RGBColor(0x00, 0xC9, 0x85)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)
IMG_X     = Inches(7.1)
IMG_MAX_W = Inches(5.7)
IMG_Y     = Inches(1.7)
IMG_MAX_H = Inches(5.0)


def new_prs():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs

def blank(prs):
    return prs.slide_layouts[6]

def rect(slide, x, y, w, h, fill=None, line=None):
    s = slide.shapes.add_shape(1, x, y, w, h)
    if fill:
        s.fill.solid(); s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    s.line.fill.background() if not line else setattr(s.line.color, 'rgb', line)
    return s

def txt(slide, text, x, y, w, h, sz=18, bold=False, color=None,
        align=PP_ALIGN.LEFT, italic=False, wrap=True):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = wrap
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(sz); r.font.bold = bold; r.font.italic = italic
    r.font.name = "Malgun Gothic"
    if color: r.font.color.rgb = color
    return tb

def add_img(slide, img_path, x, y, max_w, max_h):
    try:
        img_path = str(img_path)
        if not os.path.exists(img_path):
            print(f"  [skip] {img_path} not found")
            return None
        img = Image.open(img_path)
        iw, ih = img.size
        scale_w = max_w / (iw * 914400 / 96)
        scale_h = max_h / (ih * 914400 / 96)
        scale = min(scale_w, scale_h, 1.0)
        fw = int(iw * 914400 / 96 * scale)
        fh = int(ih * 914400 / 96 * scale)
        return slide.shapes.add_picture(img_path, x, y, fw, fh)
    except Exception as e:
        print(f"  [img error] {img_path}: {e}")
        return None

def cap(slide, label, x, y, w):
    txt(slide, label, x, y, w, Inches(0.3), sz=9, color=C_GRAY, align=PP_ALIGN.CENTER, italic=True)

def header(slide, num, title, subtitle=None):
    rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill=C_DARK)
    rect(slide, 0, 0, Inches(0.3), SLIDE_H, fill=C_PURPLE)
    rect(slide, 0, 0, SLIDE_W, Inches(1.5), fill=C_DARK2)
    txt(slide, f"{num:02d}  {title}", Inches(0.6), Inches(0.3),
        Inches(9.0), Inches(0.8), sz=30, bold=True, color=C_WHITE)
    if subtitle:
        txt(slide, subtitle, Inches(0.6), Inches(1.05),
            Inches(8.0), Inches(0.4), sz=13, color=C_PURPLE_LIGHT)


def build():
    prs = new_prs()

    # ── Slide 1: 표지 ─────────────────────────────────
    sl = prs.slides.add_slide(blank(prs))
    rect(sl, 0, 0, SLIDE_W, SLIDE_H, fill=C_DARK)
    rect(sl, 0, 0, Inches(0.3), SLIDE_H, fill=C_PURPLE)
    iw, ih = Inches(5.6), Inches(4.2)
    ix, iy = Inches(7.3), Inches(1.6)
    add_img(sl, ASSETS/"screenshot_main.png", ix, iy, iw, ih)
    cap(sl, "clevertap.com — 글로벌 모바일 마케팅 플랫폼", ix, iy+ih+Inches(0.05), iw)
    rect(sl, Inches(0.5), Inches(0.4), Inches(2.0), Inches(0.45), fill=C_PURPLE)
    txt(sl, "CleverTap", Inches(0.55), Inches(0.43), Inches(1.9), Inches(0.4), sz=17, bold=True, color=C_WHITE)
    txt(sl, "CleverTap", Inches(0.6), Inches(1.5), Inches(6.5), Inches(0.9), sz=50, bold=True, color=C_WHITE)
    txt(sl, "솔루션 제안서", Inches(0.6), Inches(2.38), Inches(6.5), Inches(0.85), sz=44, bold=True, color=C_PURPLE_LIGHT)
    rect(sl, Inches(0.6), Inches(3.4), Inches(3.2), Inches(0.04), fill=C_PURPLE)
    txt(sl, "모바일 마케팅 자동화 & 실시간 고객 인게이지먼트 플랫폼",
        Inches(0.6), Inches(3.6), Inches(6.3), Inches(0.55), sz=14, color=C_LIGHT_GRAY)
    for i, (v, l) in enumerate([("10,000+","글로벌 앱"),("12조+","월간 이벤트"),("100+","서비스 국가")]):
        bx = Inches(0.6 + i*2.05)
        rect(sl, bx, Inches(5.0), Inches(1.8), Inches(1.0), fill=C_MID)
        txt(sl, v, bx+Inches(0.1), Inches(5.1), Inches(1.6), Inches(0.45), sz=20, bold=True, color=C_PURPLE_LIGHT, align=PP_ALIGN.CENTER)
        txt(sl, l, bx+Inches(0.1), Inches(5.55), Inches(1.6), Inches(0.35), sz=11, color=C_GRAY, align=PP_ALIGN.CENTER)
    txt(sl, "2026년 4월", Inches(0.6), Inches(6.9), Inches(3.0), Inches(0.4), sz=13, color=C_GRAY)
    print("  [1] 표지")

    # ── Slide 2: 목차 ─────────────────────────────────
    sl = prs.slides.add_slide(blank(prs))
    rect(sl, 0, 0, SLIDE_W, SLIDE_H, fill=C_DARK)
    rect(sl, 0, 0, Inches(0.3), SLIDE_H, fill=C_PURPLE)
    rect(sl, 0, 0, SLIDE_W, Inches(1.5), fill=C_DARK2)
    txt(sl, "목 차", Inches(0.6), Inches(0.3), Inches(6.0), Inches(0.7), sz=34, bold=True, color=C_WHITE)
    txt(sl, "CONTENTS", Inches(0.6), Inches(1.0), Inches(4.0), Inches(0.35), sz=12, color=C_PURPLE_LIGHT)
    for i, (n, t) in enumerate([
        ("01","CleverTap 소개"),("02","주요 글로벌 고객사"),("03","핵심 가치 제안"),
        ("04","제품 아키텍처 개요"),("05","실시간 분석 (Analytics)"),
        ("06","사용자 세그멘테이션 & 개인화"),("07","멀티채널 인게이지먼트"),
        ("08","AI/ML 기반 개인화"),("09","A/B 테스트 & 최적화"),
        ("10","핀테크/금융 활용 사례"),("11","도입 효과 & ROI"),
        ("12","도입 프로세스 & 타임라인"),("13","가격 정책 & 다음 단계"),
    ]):
        col=i%2; row=i//2
        x=Inches(0.7+col*6.3); y=Inches(1.7+row*0.72)
        rect(sl, x, y+Inches(0.05), Inches(0.48), Inches(0.48), fill=C_PURPLE)
        txt(sl, n, x+Inches(0.02), y+Inches(0.08), Inches(0.44), Inches(0.36), sz=11, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        txt(sl, t, x+Inches(0.58), y, Inches(5.5), Inches(0.52), sz=16, color=C_WHITE)
    print("  [2] 목차")

    # ── Slide 3: 소개 — ctx_intro (플랫폼 오버뷰 다이어그램)
    sl = prs.slides.add_slide(blank(prs))
    header(sl, 1, "CleverTap 소개", "모바일 마케팅 자동화의 글로벌 리더")
    txt(sl, "CleverTap은 2013년 설립된 글로벌 모바일 마케팅 자동화 & 분석 플랫폼으로,\n"
        "앱 비즈니스가 사용자를 유지하고 성장시킬 수 있도록 돕는 올인원 솔루션입니다.\n\n"
        "실시간 행동 데이터를 기반으로 개인화된 메시지를 적시에 전달하여\n고객 참여도와 매출을 극대화합니다.",
        Inches(0.6), Inches(1.7), Inches(5.8), Inches(2.0), sz=14, color=C_LIGHT_GRAY, wrap=True)
    for i, (v, l) in enumerate([("2013년","설립"),("10,000+","글로벌 고객사"),("100+","국가 서비스"),
                                  ("12조+","월간 이벤트"),("600+","임직원"),("$105M+","누적 투자")]):
        col=i%3; row=i//3
        bx=Inches(0.6+col*2.1); by=Inches(3.9+row*1.3)
        rect(sl, bx, by, Inches(1.9), Inches(1.1), fill=C_MID)
        rect(sl, bx, by, Inches(1.9), Inches(0.05), fill=C_PURPLE)
        txt(sl, v, bx+Inches(0.1), by+Inches(0.12), Inches(1.7), Inches(0.5), sz=20, bold=True, color=C_PURPLE_LIGHT, align=PP_ALIGN.CENTER)
        txt(sl, l, bx+Inches(0.1), by+Inches(0.62), Inches(1.7), Inches(0.38), sz=10, color=C_GRAY, align=PP_ALIGN.CENTER)
    add_img(sl, ASSETS/"ctx_intro.png", IMG_X, IMG_Y, IMG_MAX_W, IMG_MAX_H)
    cap(sl, "CleverTap 통합 플랫폼 — Analytics · AI · Campaign · Journeys", IMG_X, IMG_Y+IMG_MAX_H+Inches(0.05), IMG_MAX_W)
    print("  [3] 소개")

    # ── Slide 4: 고객사 — ctx_customers
    sl = prs.slides.add_slide(blank(prs))
    header(sl, 2, "주요 글로벌 고객사", "전 세계 10,000+ 앱이 CleverTap을 신뢰합니다")
    for i, (n, d, c) in enumerate([
        ("Jio","인도 최대 통신사","통신"),("Swiggy","푸드 딜리버리 1위","이커머스"),
        ("PhonePe","UPI 결제 플랫폼","핀테크"),("Axis Bank","인도 3대 민간은행","금융"),
        ("Disney+ Hotstar","OTT 스트리밍","미디어"),("Vodafone","글로벌 통신사","통신"),
    ]):
        col=i%2; row=i//2
        x=Inches(0.6+col*3.3); y=Inches(1.8+row*1.75)
        rect(sl, x, y, Inches(3.0), Inches(1.55), fill=C_MID)
        rect(sl, x, y, Inches(0.1), Inches(1.55), fill=C_PURPLE)
        rect(sl, x+Inches(0.2), y+Inches(0.12), Inches(0.75), Inches(0.25), fill=C_PURPLE)
        txt(sl, c, x+Inches(0.22), y+Inches(0.13), Inches(0.71), Inches(0.22), sz=8, bold=True, color=C_WHITE)
        txt(sl, n, x+Inches(0.2), y+Inches(0.45), Inches(2.7), Inches(0.42), sz=18, bold=True, color=C_WHITE)
        txt(sl, d, x+Inches(0.2), y+Inches(0.9), Inches(2.7), Inches(0.38), sz=11, color=C_GRAY)
    add_img(sl, ASSETS/"ctx_customers.png", Inches(7.1), Inches(1.7), Inches(5.7), Inches(5.2))
    cap(sl, "10,000+ 글로벌 브랜드 — 전 세계 100개국 서비스", Inches(7.1), Inches(6.95), Inches(5.7))
    print("  [4] 고객사")

    # ── Slide 5: 핵심 가치 (이미지 없음)
    sl = prs.slides.add_slide(blank(prs))
    header(sl, 3, "핵심 가치 제안", "왜 CleverTap인가?")
    for i, (t, d) in enumerate([
        ("올인원 플랫폼","분석→세그멘테이션→캠페인→최적화\n하나의 플랫폼에서 완결"),
        ("실시간 처리","초당 수백만 이벤트 실시간 처리\n즉각적인 의사결정 지원"),
        ("AI 기반 개인화","ML 알고리즘으로 최적 메시지 자동 전달"),
        ("멀티채널 통합","푸시/이메일/SMS/인앱/WhatsApp 단일 관리"),
        ("강력한 분석","퍼널·코호트·플로우 분석 깊이 있는 인사이트"),
        ("엔터프라이즈 보안","SOC2·GDPR·ISO27001 금융 수준 보안"),
    ]):
        col=i%2; row=i//2
        x=Inches(0.6+col*6.5); y=Inches(1.8+row*1.8)
        rect(sl, x, y, Inches(6.1), Inches(1.6), fill=C_MID)
        rect(sl, x, y, Inches(6.1), Inches(0.05), fill=C_PURPLE)
        txt(sl, t, x+Inches(0.2), y+Inches(0.15), Inches(5.7), Inches(0.4), sz=15, bold=True, color=C_WHITE)
        txt(sl, d, x+Inches(0.2), y+Inches(0.6), Inches(5.7), Inches(0.88), sz=12, color=C_LIGHT_GRAY, wrap=True)
    print("  [5] 가치 제안")

    # ── Slide 6: 아키텍처 — ctx_architecture (Journey 오케스트레이션)
    sl = prs.slides.add_slide(blank(prs))
    header(sl, 4, "제품 아키텍처 개요", "데이터 수집부터 실행까지 완전 통합 플랫폼")
    for i, (t, d, c) in enumerate([
        ("데이터 수집 레이어","SDK (iOS/Android/Web) · REST API · S3 커넥터 · 서드파티 통합", C_PURPLE),
        ("데이터 처리 레이어","실시간 스트림 처리 · 사용자 프로필 병합 · 이벤트 저장소", C_PURPLE_DARK),
        ("분석 & AI 레이어","행동 분석 · ML 모델 · 예측 세그멘테이션 · 최적 전송 시간", RGBColor(0x2B,0x4E,0xCC)),
        ("실행 레이어","캠페인 관리 · 멀티채널 전송 · A/B 테스트 · 자동화 워크플로우", C_MID),
    ]):
        y=Inches(1.9+i*1.1)
        rect(sl, Inches(0.6), y, Inches(6.0), Inches(0.95), fill=c)
        rect(sl, Inches(0.6), y, Inches(0.38), Inches(0.95), fill=C_PURPLE_DARK)
        txt(sl, str(i+1), Inches(0.62), y+Inches(0.2), Inches(0.34), Inches(0.5), sz=13, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        txt(sl, t, Inches(1.1), y+Inches(0.05), Inches(5.3), Inches(0.38), sz=13, bold=True, color=C_WHITE)
        txt(sl, d, Inches(1.1), y+Inches(0.45), Inches(5.3), Inches(0.42), sz=10, color=C_LIGHT_GRAY)
    add_img(sl, ASSETS/"ctx_architecture.png", IMG_X, IMG_Y, IMG_MAX_W, IMG_MAX_H)
    cap(sl, "Journey Orchestration — 실시간 자동화 워크플로우 빌더", IMG_X, IMG_Y+IMG_MAX_H+Inches(0.05), IMG_MAX_W)
    print("  [6] 아키텍처")

    # ── Slide 7: Analytics — ctx_analytics (대시보드)
    sl = prs.slides.add_slide(blank(prs))
    header(sl, 5, "실시간 분석 (Analytics)", "모든 데이터를 실시간으로 — 즉각적인 비즈니스 인사이트")
    for i, (t, d) in enumerate([
        ("퍼널 분석","전환율 추적, 이탈 지점 파악, 최적화"),
        ("코호트 분석","사용자 그룹별 리텐션 트렌드 분석"),
        ("플로우 분석","사용자 탐색 패턴 시각화"),
        ("이벤트 분석","커스텀 이벤트 실시간 트래킹"),
        ("피벗 분석","다차원 데이터 분석 및 리포팅"),
    ]):
        y=Inches(1.85+i*0.92)
        rect(sl, Inches(0.6), y, Inches(5.7), Inches(0.78), fill=C_MID)
        rect(sl, Inches(0.6), y, Inches(0.05), Inches(0.78), fill=C_GREEN)
        txt(sl, t, Inches(0.8), y+Inches(0.06), Inches(2.5), Inches(0.32), sz=13, bold=True, color=C_WHITE)
        txt(sl, d, Inches(0.8), y+Inches(0.4), Inches(5.2), Inches(0.32), sz=11, color=C_GRAY)
    add_img(sl, ASSETS/"ctx_analytics.png", IMG_X, IMG_Y, IMG_MAX_W, IMG_MAX_H)
    cap(sl, "Analytics Dashboard — 퍼널 분석 · 리텐션 트렌드 실시간 조회", IMG_X, IMG_Y+IMG_MAX_H+Inches(0.05), IMG_MAX_W)
    print("  [7] 분석")

    # ── Slide 8: 세그멘테이션 — ctx_segmentation (세그먼트 빌더)
    sl = prs.slides.add_slide(blank(prs))
    header(sl, 6, "사용자 세그멘테이션 & 개인화", "정밀한 타겟팅으로 관련성 높은 경험 제공")
    for i, (t, d) in enumerate([
        ("정적 세그먼트","고정 기준으로 사용자 그룹화\n연령대·지역·가입일 등"),
        ("동적 세그먼트","실시간 행동에 따라 자동 업데이트\n최근 7일 미접속자 등"),
        ("RFM 세그먼트","Recency·Frequency·Monetary 기반\n고가치 사용자 자동 식별"),
        ("예측 세그먼트","ML로 이탈 가능성·구매 의향 예측\n선제적 대응 자동화"),
    ]):
        col=i%2; row=i//2
        x=Inches(0.6+col*3.2); y=Inches(1.9+row*2.1)
        rect(sl, x, y, Inches(2.9), Inches(1.9), fill=C_MID)
        rect(sl, x, y, Inches(2.9), Inches(0.05), fill=C_PURPLE)
        txt(sl, t, x+Inches(0.15), y+Inches(0.15), Inches(2.6), Inches(0.4), sz=14, bold=True, color=C_WHITE)
        txt(sl, d, x+Inches(0.15), y+Inches(0.62), Inches(2.6), Inches(1.1), sz=11, color=C_LIGHT_GRAY, wrap=True)
    add_img(sl, ASSETS/"ctx_segmentation.png", IMG_X, IMG_Y, IMG_MAX_W, IMG_MAX_H)
    cap(sl, "Segment Builder — Dynamic/RFM/Predictive 세그먼트 실시간 구성", IMG_X, IMG_Y+IMG_MAX_H+Inches(0.05), IMG_MAX_W)
    print("  [8] 세그멘테이션")

    # ── Slide 9: 인게이지먼트 — ctx_engagement (캠페인 매니저)
    sl = prs.slides.add_slide(blank(prs))
    header(sl, 7, "멀티채널 인게이지먼트", "고객이 있는 모든 채널에서 일관된 경험 제공")
    for i, (n, d) in enumerate([
        ("푸시 알림","iOS/Android 리치 푸시·딥링크"),
        ("이메일","드래그&드롭 에디터·AMP 지원"),
        ("SMS/MMS","글로벌 통신사 연동·개인화"),
        ("인앱 메시지","풀스크린/배너/팝업·인터랙티브"),
        ("웹 푸시","브라우저 알림·세션 외 재참여"),
        ("WhatsApp","공식 비즈니스 API·챗봇 연동"),
    ]):
        col=i%2; row=i//2
        x=Inches(0.6+col*3.3); y=Inches(1.8+row*1.72)
        rect(sl, x, y, Inches(3.0), Inches(1.55), fill=C_MID)
        rect(sl, x, y, Inches(3.0), Inches(0.05), fill=C_PURPLE)
        txt(sl, n, x+Inches(0.18), y+Inches(0.15), Inches(2.7), Inches(0.38), sz=14, bold=True, color=C_WHITE)
        txt(sl, d, x+Inches(0.18), y+Inches(0.6), Inches(2.7), Inches(0.78), sz=11, color=C_GRAY, wrap=True)
    add_img(sl, ASSETS/"ctx_engagement.png", IMG_X, IMG_Y, IMG_MAX_W, IMG_MAX_H)
    cap(sl, "Campaign Manager — 옴니채널 캠페인 성과 실시간 모니터링", IMG_X, IMG_Y+IMG_MAX_H+Inches(0.05), IMG_MAX_W)
    print("  [9] 인게이지먼트")

    # ── Slide 10: AI/ML — ctx_ai (AI 대시보드)
    sl = prs.slides.add_slide(blank(prs))
    header(sl, 8, "AI/ML 기반 개인화", "머신러닝으로 구현하는 1:1 초개인화 마케팅")
    for i, (t, d) in enumerate([
        ("Clever.AI","CleverTap 독자 AI 엔진 — 고객 행동 학습·예측"),
        ("최적 전송 시간 (STO)","사용자별 최적 수신 시간 자동 선택 → 오픈율 47% 향상"),
        ("예측 세그먼트","이탈·구매·구독 해지 사전 예측 → 선제 대응"),
        ("동적 콘텐츠","프로필 기반 실시간 콘텐츠 변환 · 제품 추천 자동화"),
        ("이탈 방지","이탈 위험 감지 → 재참여 캠페인 자동 트리거"),
    ]):
        y=Inches(1.85+i*0.92)
        rect(sl, Inches(0.6), y, Inches(5.7), Inches(0.78), fill=C_MID)
        rect(sl, Inches(0.6), y, Inches(0.05), Inches(0.78), fill=C_PURPLE)
        txt(sl, t, Inches(0.8), y+Inches(0.06), Inches(2.4), Inches(0.32), sz=13, bold=True, color=C_WHITE)
        txt(sl, d, Inches(0.8), y+Inches(0.4), Inches(5.2), Inches(0.32), sz=11, color=C_GRAY)
    add_img(sl, ASSETS/"ctx_ai.png", IMG_X, IMG_Y, IMG_MAX_W, IMG_MAX_H)
    cap(sl, "CleverAI™ — 예측 세그먼트 · STO · 동적 개인화 엔진", IMG_X, IMG_Y+IMG_MAX_H+Inches(0.05), IMG_MAX_W)
    print("  [10] AI/ML")

    # ── Slide 11: A/B 테스트 — ctx_ab (실험 대시보드)
    sl = prs.slides.add_slide(blank(prs))
    header(sl, 9, "A/B 테스트 & 최적화", "데이터 기반 의사결정으로 마케팅 성과 극대화")
    for i, (t, d) in enumerate([
        ("메시지 A/B 테스트","제목·본문·CTA·이미지·이모지 테스트"),
        ("전송 시간 테스트","요일·시간대별 성과 비교 및 최적화"),
        ("채널 A/B 테스트","푸시 vs 이메일 vs SMS 효과 측정"),
        ("자동 승자 선택","통계적 유의성 기반 자동 선택 · 트래픽 손실 최소화"),
    ]):
        x=Inches(0.6+(i%2)*3.2); y=Inches(1.9+(i//2)*1.95)
        rect(sl, x, y, Inches(2.9), Inches(1.75), fill=C_MID)
        rect(sl, x, y, Inches(2.9), Inches(0.05), fill=C_ACCENT)
        txt(sl, t, x+Inches(0.15), y+Inches(0.15), Inches(2.6), Inches(0.4), sz=13, bold=True, color=C_WHITE)
        txt(sl, d, x+Inches(0.15), y+Inches(0.62), Inches(2.6), Inches(0.95), sz=11, color=C_LIGHT_GRAY, wrap=True)
    for i, (v, l) in enumerate([("47%","오픈율 향상"),("32%","전환율 개선"),("2.3배","ROI 증가")]):
        bx=Inches(0.6+i*2.15); by=Inches(5.9)
        rect(sl, bx, by, Inches(1.9), Inches(1.15), fill=C_PURPLE_DARK)
        txt(sl, v, bx+Inches(0.1), by+Inches(0.1), Inches(1.7), Inches(0.55), sz=26, bold=True, color=C_GREEN, align=PP_ALIGN.CENTER)
        txt(sl, l, bx+Inches(0.1), by+Inches(0.7), Inches(1.7), Inches(0.38), sz=11, color=C_GRAY, align=PP_ALIGN.CENTER)
    add_img(sl, ASSETS/"ctx_ab.png", IMG_X, IMG_Y, IMG_MAX_W, IMG_MAX_H)
    cap(sl, "Experiment Dashboard — 실험 CTR 비교 · 통계 유의성 자동 검증", IMG_X, IMG_Y+IMG_MAX_H+Inches(0.05), IMG_MAX_W)
    print("  [11] A/B 테스트")

    # ── Slide 12: 핀테크 — ctx_fintech (온보딩 퍼널)
    sl = prs.slides.add_slide(blank(prs))
    header(sl, 10, "핀테크/금융 업계 활용 사례", "금융 서비스에 최적화된 마케팅 자동화")
    for i, (t, d) in enumerate([
        ("신규 가입 온보딩","7일 온보딩 시나리오 · 첫 거래 유도"),
        ("KYC 완료 독려","미완료 사용자 자동 리마인드"),
        ("투자 상품 추천","거래 이력 기반 AI 상품 추천"),
        ("이탈 방지 캠페인","30일 미거래 사용자 자동 감지"),
    ]):
        y=Inches(1.85+i*1.02)
        rect(sl, Inches(0.6), y, Inches(5.8), Inches(0.88), fill=C_MID)
        rect(sl, Inches(0.6), y, Inches(0.05), Inches(0.88), fill=C_ACCENT)
        txt(sl, t, Inches(0.8), y+Inches(0.07), Inches(2.8), Inches(0.34), sz=13, bold=True, color=C_WHITE)
        txt(sl, d, Inches(0.8), y+Inches(0.45), Inches(5.3), Inches(0.35), sz=11, color=C_GRAY)
    for i, (v, l) in enumerate([("40%","신규 가입 전환율"),("28%","KYC 완료율"),("3.2배","온보딩 ROI"),("65%","이탈 방지율")]):
        col=i%2; row=i//2
        x=Inches(0.6+col*3.0); y=Inches(5.25+row*1.05)
        rect(sl, x, y, Inches(2.7), Inches(0.9), fill=C_PURPLE_DARK)
        txt(sl, f"{v}  {l}", x+Inches(0.1), y+Inches(0.15), Inches(2.5), Inches(0.6), sz=14, bold=True, color=C_GREEN)
    add_img(sl, ASSETS/"ctx_fintech.png", IMG_X, IMG_Y, IMG_MAX_W, IMG_MAX_H)
    cap(sl, "FinTech Journey — 온보딩 전환 퍼널 · 라이프사이클 캠페인 성과", IMG_X, IMG_Y+IMG_MAX_H+Inches(0.05), IMG_MAX_W)
    print("  [12] 핀테크")

    # ── Slide 13: ROI (이미지 없음)
    sl = prs.slides.add_slide(blank(prs))
    header(sl, 11, "도입 효과 & ROI", "측정된 성과 — 실제 고객사 기준 평균 수치")
    for i, (v, t, d) in enumerate([
        ("43%","리텐션율 향상","이탈 방지 캠페인 적용 후"),
        ("2.8배","마케팅 ROI","전통 채널 대비"),
        ("62%","인게이지먼트 증가","개인화 메시지 적용 후"),
        ("30%","비용 절감","자동화로 운영 효율 개선"),
        ("5배","캠페인 실행 속도","수동 대비 자동화 효과"),
        ("93%","고객 만족도","NPS 기준 도입사 평균"),
    ]):
        col=i%3; row=i//3
        x=Inches(0.6+col*4.15); y=Inches(1.9+row*2.4)
        rect(sl, x, y, Inches(3.8), Inches(2.15), fill=C_MID)
        rect(sl, x, y, Inches(3.8), Inches(0.05), fill=C_GREEN)
        txt(sl, v, x+Inches(0.1), y+Inches(0.2), Inches(3.6), Inches(0.85), sz=44, bold=True, color=C_GREEN, align=PP_ALIGN.CENTER)
        txt(sl, t, x+Inches(0.1), y+Inches(1.05), Inches(3.6), Inches(0.42), sz=14, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        txt(sl, d, x+Inches(0.1), y+Inches(1.5), Inches(3.6), Inches(0.42), sz=11, color=C_GRAY, align=PP_ALIGN.CENTER)
    print("  [13] ROI")

    # ── Slide 14: 타임라인 (이미지 없음)
    sl = prs.slides.add_slide(blank(prs))
    header(sl, 12, "도입 프로세스 & 타임라인", "빠르고 안전한 구축 — 평균 4주 내 라이브")
    rect(sl, Inches(1.2), Inches(3.5), Inches(11.0), Inches(0.06), fill=C_PURPLE)
    for i, (wk, t, d) in enumerate([
        ("Week 1","킥오프 & 설계","요구사항 정의\nSDK 통합 계획\n이벤트 텍소노미 설계"),
        ("Week 2","SDK 통합","iOS/Android SDK 설치\n이벤트 태깅\n데이터 검증"),
        ("Week 3","캠페인 구성","세그멘테이션 설정\n캠페인 템플릿\nA/B 테스트 설계"),
        ("Week 4","라이브 & 최적화","첫 캠페인 발송\n성과 모니터링\n피드백 반영"),
    ]):
        x=Inches(0.8+i*3.1)
        circ=sl.shapes.add_shape(9, x+Inches(1.1), Inches(3.32), Inches(0.45), Inches(0.45))
        circ.fill.solid(); circ.fill.fore_color.rgb=C_PURPLE; circ.line.fill.background()
        if i%2==0:
            txt(sl, wk, x+Inches(0.1), Inches(1.9), Inches(2.8), Inches(0.32), sz=11, color=C_PURPLE_LIGHT, bold=True, align=PP_ALIGN.CENTER)
            txt(sl, t, x+Inches(0.1), Inches(2.28), Inches(2.8), Inches(0.38), sz=13, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        else:
            txt(sl, wk, x+Inches(0.1), Inches(4.05), Inches(2.8), Inches(0.32), sz=11, color=C_PURPLE_LIGHT, bold=True, align=PP_ALIGN.CENTER)
            txt(sl, t, x+Inches(0.1), Inches(4.43), Inches(2.8), Inches(0.38), sz=13, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        rect(sl, x+Inches(0.1), Inches(5.2), Inches(2.8), Inches(1.9), fill=C_MID)
        txt(sl, d, x+Inches(0.2), Inches(5.32), Inches(2.6), Inches(1.65), sz=11, color=C_LIGHT_GRAY, wrap=True)
    rect(sl, Inches(0.6), Inches(7.1), Inches(12.0), Inches(0.3), fill=C_PURPLE_DARK)
    txt(sl, "전담 CSM 배정  |  한국어 기술 지원  |  24/7 모니터링  |  온보딩 플레이북 제공",
        Inches(0.8), Inches(7.15), Inches(11.8), Inches(0.25), sz=11, color=C_WHITE, align=PP_ALIGN.CENTER)
    print("  [14] 타임라인")

    # ── Slide 15: 가격 정책 (이미지 없음)
    sl = prs.slides.add_slide(blank(prs))
    header(sl, 13, "가격 정책 & 다음 단계", "비즈니스 규모에 맞는 유연한 플랜")
    for i, (nm, ds, feats, rec) in enumerate([
        ("Essentials","스타트업·MAU 50만",["핵심 분석 기능","푸시 + 이메일","기본 세그멘테이션","이메일 지원"],False),
        ("Advanced","Mid-Market·MAU 500만",["Essentials 전체","AI 예측 세그먼트","멀티채널 전체","전담 CSM"],True),
        ("Enterprise","대기업·MAU 무제한",["Advanced 전체","커스텀 SLA","온프레미스 옵션","전용 인프라"],False),
    ]):
        x=Inches(0.6+i*4.15)
        rect(sl, x, Inches(1.9), Inches(3.8), Inches(4.9), fill=C_PURPLE if rec else C_MID)
        if rec:
            rect(sl, x, Inches(1.9), Inches(3.8), Inches(0.36), fill=C_PURPLE_DARK)
            txt(sl, "추천", x+Inches(0.1), Inches(1.92), Inches(3.6), Inches(0.3), sz=11, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        txt(sl, nm, x+Inches(0.2), Inches(2.35), Inches(3.4), Inches(0.5), sz=24, bold=True, color=C_WHITE)
        txt(sl, ds, x+Inches(0.2), Inches(2.88), Inches(3.4), Inches(0.42), sz=11, color=C_LIGHT_GRAY)
        for j, f in enumerate(feats):
            txt(sl, f"v  {f}", x+Inches(0.2), Inches(3.45+j*0.52), Inches(3.4), Inches(0.46),
                sz=12, color=C_WHITE if rec else C_LIGHT_GRAY)
    rect(sl, Inches(0.6), Inches(7.05), Inches(12.0), Inches(0.35), fill=C_MID)
    txt(sl, "모든 플랜: 무료 파일럿 2주 제공 · 연간 계약 할인 · 한국 전담 지원팀",
        Inches(0.8), Inches(7.1), Inches(11.8), Inches(0.28), sz=12, color=C_PURPLE_LIGHT, align=PP_ALIGN.CENTER)
    print("  [15] 가격")

    # ── Slide 16: CTA (이미지 없음)
    sl = prs.slides.add_slide(blank(prs))
    rect(sl, 0, 0, SLIDE_W, SLIDE_H, fill=C_DARK)
    rect(sl, 0, 0, Inches(0.3), SLIDE_H, fill=C_PURPLE)
    rect(sl, 0, 0, SLIDE_W, Inches(1.5), fill=C_DARK2)
    rect(sl, Inches(1.5), Inches(1.5), Inches(10.0), Inches(4.8), fill=C_MID)
    rect(sl, Inches(1.5), Inches(1.5), Inches(10.0), Inches(0.07), fill=C_PURPLE)
    txt(sl, "지금 바로 시작하세요", Inches(2.0), Inches(2.0), Inches(9.0), Inches(0.9), sz=40, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    txt(sl, "2주 무료 파일럿 프로그램 신청", Inches(2.0), Inches(2.95), Inches(9.0), Inches(0.6), sz=24, color=C_PURPLE_LIGHT, align=PP_ALIGN.CENTER)
    txt(sl, "실제 데이터로 CleverTap의 가치를 직접 확인하세요", Inches(2.0), Inches(3.6), Inches(9.0), Inches(0.4), sz=16, color=C_LIGHT_GRAY, align=PP_ALIGN.CENTER)
    rect(sl, Inches(4.5), Inches(4.4), Inches(4.0), Inches(0.65), fill=C_PURPLE)
    txt(sl, "데모 신청하기  ->", Inches(4.5), Inches(4.44), Inches(4.0), Inches(0.57), sz=17, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    for i, (lb, vl) in enumerate([("웹사이트","clevertap.com"),("이메일","korea@clevertap.com"),("전화","+82-2-XXXX-XXXX")]):
        bx=Inches(2.0+i*3.5)
        txt(sl, lb, bx, Inches(6.5), Inches(3.0), Inches(0.3), sz=11, color=C_GRAY, align=PP_ALIGN.CENTER)
        txt(sl, vl, bx, Inches(6.82), Inches(3.0), Inches(0.38), sz=13, bold=True, color=C_PURPLE_LIGHT, align=PP_ALIGN.CENTER)
    txt(sl, "CleverTap", Inches(5.4), Inches(0.3), Inches(2.7), Inches(0.65), sz=28, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    print("  [16] CTA")

    prs.save(str(OUTPUT))
    size = OUTPUT.stat().st_size // 1024
    print(f"\n저장 완료: {OUTPUT}")
    print(f"  슬라이드: {len(prs.slides)}장 / 파일 크기: {size}KB")
    return len(prs.slides)


if __name__ == "__main__":
    print("CleverTap 제안서 v5 빌드 시작\n")
    n = build()
    print(f"\n완료: {n}장")
