# -*- coding: utf-8 -*-
"""전북은행 Amplitude 제안서 — 검증 결과 반영 수정"""
import sys, io
if sys.platform == 'win32':
    sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from copy import deepcopy
import re

SRC = r'C:\Agent\pepper\outputs\전북은행_Amplitude_도입제안서.pptx'
DST = r'C:\Agent\pepper\outputs\전북은행_Amplitude_도입제안서_v2.pptx'

prs = Presentation(SRC)

# ── Helper: find and replace text in a shape, preserving formatting ──
def replace_in_shape(shape, old, new):
    """Replace text in a shape while preserving run formatting."""
    if not shape.has_text_frame:
        return False
    found = False
    for para in shape.text_frame.paragraphs:
        full_text = para.text
        if old in full_text:
            found = True
            # If single run, simple replace
            if len(para.runs) == 1:
                para.runs[0].text = para.runs[0].text.replace(old, new)
            elif len(para.runs) > 1:
                # Concatenate all runs, replace, put into first run, clear rest
                combined = ''.join(r.text for r in para.runs)
                new_text = combined.replace(old, new)
                para.runs[0].text = new_text
                for r in para.runs[1:]:
                    r.text = ''
            else:
                # No runs but has text? Unlikely but handle
                pass
    return found

def shape_contains(shape, text):
    if not shape.has_text_frame:
        return False
    return text in shape.text_frame.text

def set_shape_text(shape, new_text):
    """Replace entire text content of a shape, keeping first run's formatting."""
    if not shape.has_text_frame:
        return
    for para in shape.text_frame.paragraphs:
        if para.runs:
            para.runs[0].text = new_text
            for r in para.runs[1:]:
                r.text = ''
            return
    # fallback
    shape.text_frame.paragraphs[0].text = new_text

changes = []

# ================================================================
# SLIDE 3 — Background: MAU 71만 + 고령층 34.8%
# ================================================================
slide3 = prs.slides[2]
for shape in slide3.shapes:
    # MAU big number: '71만명' → 모바일 확대
    if shape_contains(shape, '71만명'):
        set_shape_text(shape, '모바일 확대')
        changes.append('S3: 71만명 -> 모바일 확대')
    if shape_contains(shape, 'MAU(2024 Q1 기준)'):
        set_shape_text(shape, '쏙뱅크 이용 증가')
        changes.append('S3: MAU subtitle changed')
    if shape_contains(shape, '전년 대비 12.7% 증가'):
        set_shape_text(shape, '앱 행동 데이터 분석 필요성 증대')
        changes.append('S3: 12.7% -> 행동 데이터 분석 필요성')

    # 고령층 34.8%
    if shape_contains(shape, '34.8%') and not shape_contains(shape, '고령층'):
        # The big number shape
        full = shape.text_frame.text
        if full.strip() == '34.8%':
            set_shape_text(shape, '고령 고객')
            changes.append('S3: 34.8% -> 고령 고객')
    if shape_contains(shape, '60세 이상고객 비중'):
        set_shape_text(shape, '디지털 이용 편의성')
        changes.append('S3: 60세 이상 subtitle changed')
    if shape_contains(shape, '고령층 맞춤 전략 필요'):
        set_shape_text(shape, '이탈 구간 파악이 중요')
        changes.append('S3: 고령층 맞춤 -> 이탈 구간')

    # Bullet: 고령층(34.8%) reference
    if shape_contains(shape, '고령층(34.8%)'):
        replace_in_shape(shape,
            '고령층(34.8%)과 MZ세대의 상이한 행동 패턴을 반영한 맞춤형 전략 수립 필요',
            '지역 기반 은행 특성상 고령 고객의 디지털 이용 편의성과 이탈 구간 파악이 중요')
        changes.append('S3: bullet 고령층 34.8% -> 지역 기반 은행 특성')

# ================================================================
# SLIDE 4 — 쏙뱅크 앱 MAU 71만명
# ================================================================
slide4 = prs.slides[3]
for shape in slide4.shapes:
    if shape_contains(shape, '쏙뱅크 앱 MAU 71만명'):
        replace_in_shape(shape,
            '쏙뱅크 앱 MAU 71만명 — 행동 분석 시 즉각 효과 확인 가능',
            '쏙뱅크 중심의 모바일 금융 이용 확대에 따라 앱 내 행동 데이터 분석 중요성 증가')
        changes.append('S4: MAU 71만명 -> 모바일 금융 이용 확대')

# ================================================================
# SLIDE 5 — 페이지뷰 기반 단정 표현 완화
# ================================================================
slide5 = prs.slides[4]
for shape in slide5.shapes:
    if shape_contains(shape, '"왜 이탈했는가"에 대한 답을 줄 수 없음'):
        replace_in_shape(shape,
            '"왜 이탈했는가"에 대한 답을 줄 수 없음',
            '"왜 이탈했는가"에 대한 답을 얻기 어려움')
        changes.append('S5: 줄 수 없음 -> 얻기 어려움 (단정 완화)')

    if shape_contains(shape, '사용자의 의도와 맥락을 파악하기 어려움'):
        # This is fine - already soft
        pass

# ================================================================
# SLIDE 9 — 시나리오 1: 전환율 수치 완화
# ================================================================
slide9 = prs.slides[8]
for shape in slide9.shapes:
    if shape_contains(shape, '전환율 10~15%p 향상 목표'):
        replace_in_shape(shape,
            'Pathfinder로 이탈 원인 분석 후 UX 개선 시 전환율 10~15%p 향상 목표',
            'Pathfinder로 이탈 원인 분석 후 UX 개선을 통한 전환율 개선 기대')
        changes.append('S9: 10~15%p -> 전환율 개선 기대')

    # 'Amplitude 적용 시 기대 효과' title stays, but mark as scenario
    if shape_contains(shape, 'Amplitude 적용 시 기대 효과'):
        replace_in_shape(shape, 'Amplitude 적용 시 기대 효과', 'Amplitude 적용 시 기대 효과 (적용 가능 시나리오)')
        changes.append('S9: 기대 효과 -> 적용 가능 시나리오 명시')

# ================================================================
# SLIDE 10 — 세그먼트별: 50대 이상 (34.8%) reference
# ================================================================
slide10 = prs.slides[9]
for shape in slide10.shapes:
    if shape_contains(shape, '50대 이상 (34.8%)'):
        replace_in_shape(shape, '50대 이상 (34.8%)', '50대 이상')
        changes.append('S10: 34.8% 비율 제거')

# ================================================================
# SLIDE 11 — 시나리오 3: 대출 이탈 분석 수치 완화
# ================================================================
slide11 = prs.slides[10]
for shape in slide11.shapes:
    # 예상 개선 효과 수치 완화
    if shape_contains(shape, '대출 신청 완료율 38%'):
        replace_in_shape(shape,
            '대출 신청 완료율 38% → 48% (+10%p)',
            '대출 신청 완료율: 퍼널 병목 구간 개선을 통한 단계별 전환율 향상 목표')
        changes.append('S11: 대출 완료율 수치 -> 퍼널 병목 개선')
    if shape_contains(shape, '대출 실행률 32%'):
        replace_in_shape(shape,
            '대출 실행률 32% → 42% (+10%p)',
            '대출 실행률: 서류 제출 단계 간소화 및 안내 강화로 이탈률 감소 기대')
        changes.append('S11: 대출 실행률 수치 -> 이탈률 감소 기대')

    # '적용 가능 시나리오' 명시 - 제목이 있으면
    if shape_contains(shape, '시나리오 3: 이탈 구간 분석 및 전환율 최적화'):
        replace_in_shape(shape,
            '시나리오 3: 이탈 구간 분석 및 전환율 최적화',
            '시나리오 3: 이탈 구간 분석 및 전환율 최적화 (적용 가능 시나리오)')
        changes.append('S11: 시나리오 제목에 적용 가능 시나리오 표기')

# ================================================================
# SLIDE 12 — ROI 분석: 전환율 +10%p, 분석시간 -70%, ROI 10~15배
# ================================================================
slide12 = prs.slides[11]
for shape in slide12.shapes:
    # 전환율 +10~15%p big number
    if shape_contains(shape, '+10~15%p'):
        set_shape_text(shape, '퍼널 개선')
        changes.append('S12: +10~15%p -> 퍼널 개선')
    if shape_contains(shape, '핵심 퍼널(가입, 대출, 이체)') and shape_contains(shape, '전환율 개선 효과'):
        replace_in_shape(shape,
            '핵심 퍼널(가입, 대출, 이체)전환율 개선 효과',
            '대출·상품가입 퍼널의 병목 구간 식별 및단계별 전환율 개선 목표 수립')
        changes.append('S12: 전환율 subtitle -> 퍼널 병목 구간 식별')

    # 분석 시간 -70%
    if shape_contains(shape, '-70%') and not shape_contains(shape, '수일'):
        full = shape.text_frame.text.strip()
        if full == '-70%':
            set_shape_text(shape, '셀프서비스')
            changes.append('S12: -70% -> 셀프서비스')
    if shape_contains(shape, '데이터 요청→분석 프로세스') and shape_contains(shape, '수일 → 수분 단축'):
        replace_in_shape(shape,
            '데이터 요청→분석 프로세스수일 → 수분 단축',
            'SQL/수작업 리포팅 의존도 감소 및현업 셀프서비스 분석 체계 구축')
        changes.append('S12: 분석시간 subtitle -> 셀프서비스')

    # ROI 시뮬레이션 section
    if shape_contains(shape, '연 $120,000~$200,000'):
        replace_in_shape(shape,
            'Amplitude 도입 비용: 연 $120,000~$200,000 (약 1.5~2.5억원)',
            'Amplitude 도입 비용: Enterprise Custom Pricing (MTU·이벤트 볼륨·기능 범위에 따라 산정)')
        changes.append('S12: 도입 비용 수치 -> Custom Pricing')

    if shape_contains(shape, '투자 대비 10~15배'):
        replace_in_shape(shape,
            '예상 ROI: 투자 대비 10~15배 수익 창출',
            '예상 ROI: Amplitude Forrester TEI 기준 217% ROI 및 6개월 내 투자 회수 사례 제시')
        changes.append('S12: 10~15배 ROI -> Forrester TEI 217%')

    if shape_contains(shape, '업계 평균 Amplitude 도입 후 12개월'):
        replace_in_shape(shape,
            '(참고: 업계 평균 Amplitude 도입 후 12개월 내 투자비 회수)',
            '(출처: Forrester Total Economic Impact of Amplitude, 2023)')
        changes.append('S12: 업계 평균 참고 -> Forrester 출처')

    # 전환율 10%p 개선 시 추가 대출 시뮬레이션 완화
    if shape_contains(shape, '전환율 10%p 개선 시 추가 대출 실행 건수'):
        replace_in_shape(shape,
            '전환율 10%p 개선 시 추가 대출 실행 건수: 월 약 500건 → 연간 6,000건',
            '퍼널 전환율 개선 시 추가 대출 실행 건수 증가 기대 (구체적 수치는 PoC 후 산정)')
        changes.append('S12: 대출 시뮬레이션 수치 완화')
    if shape_contains(shape, '건당 대출 이자 수익 평균 50만원'):
        replace_in_shape(shape,
            '건당 대출 이자 수익 평균 50만원 기준 → 연간 추가 수익 약 30억원',
            '이벤트 기반 분석 도입으로 데이터 기반 의사결정 체계 확립')
        changes.append('S12: 30억원 시뮬레이션 -> 데이터 기반 의사결정')
    if shape_contains(shape, '마케팅 예산 효율화: 연간 약 3~5억원'):
        replace_in_shape(shape,
            '마케팅 예산 효율화: 연간 약 3~5억원 절감 예상',
            '마케팅 캠페인 ROI 측정 및 예산 배분 최적화 가능')
        changes.append('S12: 3~5억 절감 -> 예산 최적화 가능')

# ================================================================
# SLIDE 13 — Roadmap: '검증된 사실' → '제안 로드맵'
# ================================================================
slide13 = prs.slides[12]
for shape in slide13.shapes:
    if shape_contains(shape, '도입 로드맵 (Phase 1~3)'):
        replace_in_shape(shape, '도입 로드맵 (Phase 1~3)', '제안 로드맵 (Phase 1~3)')
        changes.append('S13: 도입 로드맵 -> 제안 로드맵')

# ================================================================
# SLIDE 14 — Pricing: Enterprise $100K~$250K
# ================================================================
slide14 = prs.slides[13]
for shape in slide14.shapes:
    if shape_contains(shape, '$100,000~$250,000'):
        replace_in_shape(shape,
            '* Enterprise: 500K~1M MTU 기준 연 $100,000~$250,000 (Vendr 데이터 기반). 다년 계약 시 20~35% 할인 가능. 정확한 견적은 별도 협의.',
            '* Enterprise는 Custom Pricing / Contact Sales 기준이며, 실제 견적은 MTU·이벤트 볼륨·기능 범위에 따라 산정됩니다. 정확한 견적은 별도 협의.')
        changes.append('S14: Enterprise 가격 -> Custom Pricing')

# ================================================================
# SLIDE 15 — References: 하나증권, 키움증권, Revolut/Chime 삭제
# ================================================================
slide15 = prs.slides[14]
for shape in slide15.shapes:
    # 하나증권 title: 국내 최초 → 국내 금융권 활용 사례
    if shape_contains(shape, '하나증권 (국내)'):
        replace_in_shape(shape, '하나증권 (국내)', '하나증권·키움증권 (국내)')
        changes.append('S15: 하나증권 -> 하나증권·키움증권')

    if shape_contains(shape, '국내 금융사 최초 Amplitude 도입 사례'):
        replace_in_shape(shape,
            '국내 금융사 최초 Amplitude 도입 사례',
            '하나증권, 키움증권 등 국내 금융권 Amplitude 활용 사례')
        changes.append('S15: 국내 최초 -> 국내 금융권 활용 사례')

    # Add 키움증권 info where relevant
    if shape_contains(shape, '작은 범위 PoC에서 시작'):
        replace_in_shape(shape,
            '작은 범위 PoC에서 시작 → 전사 확산 성공',
            '2026년 1월 키움증권 MTS 영웅문S#에 Amplitude 공급, 데이터 기반 고객 분석 환경 구축')
        changes.append('S15: PoC 전사확산 -> 키움증권 영웅문S#')

    # 글로벌 금융사 — Revolut, Chime 삭제, PayPal 수정
    if shape_contains(shape, 'PayPal') and shape_contains(shape, '결제 전환율 12%'):
        replace_in_shape(shape,
            'PayPal — 사용자 행동 분석으로 결제 전환율 12% 개선',
            'PayPal, Square, Capital One, Intuit 등 글로벌 금융권 고객')
        changes.append('S15: PayPal -> 글로벌 금융권 고객 목록')

    if shape_contains(shape, 'Revolut'):
        replace_in_shape(shape,
            'Revolut — 코호트 분석 기반 리텐션 전략 수립',
            'Forrester TEI 분석 기준 217% ROI 및 6개월 내 투자 회수 달성')
        changes.append('S15: Revolut -> Forrester ROI')

    if shape_contains(shape, 'Chime'):
        replace_in_shape(shape,
            'Chime — 퍼널 분석으로 계좌 개설 전환율 20% 향상',
            '금융사 특화 보안·컴플라이언스 요건 충족 (SOC 2, GDPR)')
        changes.append('S15: Chime -> 보안 컴플라이언스')

    # 하나증권 도입 이후 다수 금융사 → 하나증권·키움증권
    if shape_contains(shape, '하나증권 도입 이후 다수 금융사가'):
        replace_in_shape(shape,
            '하나증권 도입 이후 다수 금융사가',
            '하나증권·키움증권 도입 이후 다수 금융사가')
        changes.append('S15: 키움증권 추가')

    # 인사이트 인용문도 수정
    if shape_contains(shape, '"작은 범위에서 실제 인사이트를 보여주자'):
        replace_in_shape(shape,
            '"작은 범위에서 실제 인사이트를 보여주자',
            '"핵심 퍼널에서 빠르게 인사이트를 확인하자')
        changes.append('S15: 인용문 수정')
    if shape_contains(shape, '다른 팀에서도 관심을 갖게 되었습니다"'):
        replace_in_shape(shape,
            '다른 팀에서도 관심을 갖게 되었습니다"',
            '조직 전체로 데이터 문화가 확산되었습니다"')
        changes.append('S15: 인용문 후반 수정')

# ================================================================
# Save
# ================================================================
prs.save(DST)
print(f'Saved: {DST}')
print(f'\nTotal changes: {len(changes)}')
for c in changes:
    print(f'  - {c}')
