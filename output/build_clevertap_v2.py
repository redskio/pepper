# -*- coding: utf-8 -*-
"""
CleverTap 한국어 제안서 v2 — 고퀄 디자인 빌드
- 글자 크기: 제목 36-44pt / 소제목 24-28pt / 본문 16-18pt
- 컬러: #FF6B35 오렌지, #1A1F36 딥네이비, #FFFFFF
- 배경: 그라데이션 + 카드형 레이아웃
- 이미지: 실제 CleverTap 대시보드 목업 사용
"""
import sys, io
if sys.platform == 'win32':
    sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pathlib import Path
from PIL import Image
import os, copy
from lxml import etree

ASSETS = Path("C:/Agent/pepper/output/clevertap_assets")
OUT    = Path("C:/Agent/pepper/output/clevertap-proposal-kr-v2.pptx")

# ── 컬러 팔레트 ──────────────────────────────────────────
C_ORANGE      = RGBColor(0xFF, 0x6B, 0x35)   # CleverTap 오렌지
C_ORANGE_DARK = RGBColor(0xE5, 0x55, 0x20)
C_NAVY        = RGBColor(0x1A, 0x1F, 0x36)   # 딥 네이비
C_NAVY2       = RGBColor(0x0D, 0x11, 0x1F)   # 더 어두운 네이비
C_NAVY_MID    = RGBColor(0x25, 0x2D, 0x4E)   # 미드 네이비
C_NAVY_CARD   = RGBColor(0x2A, 0x32, 0x54)   # 카드 배경
C_WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
C_LIGHT       = RGBColor(0xF5, 0xF7, 0xFF)   # 밝은 배경
C_GRAY        = RGBColor(0x8C, 0x95, 0xB0)
C_GRAY_LIGHT  = RGBColor(0xE8, 0xEC, 0xF5)
C_GREEN       = RGBColor(0x00, 0xC9, 0x85)
C_PURPLE      = RGBColor(0x6C, 0x63, 0xFF)

SW = Inches(13.33)
SH = Inches(7.5)


def new_prs():
    prs = Presentation()
    prs.slide_width  = SW
    prs.slide_height = SH
    return prs

def blank(prs):
    return prs.slide_layouts[6]

def rect(slide, x, y, w, h, fill=None, alpha=None):
    """사각형 추가. fill=RGBColor or None(투명)"""
    sh = slide.shapes.add_shape(1, x, y, w, h)
    sh.line.fill.background()
    if fill:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    else:
        sh.fill.background()
    return sh

def rrect(slide, x, y, w, h, fill=None, radius_emu=Emu(91440)):
    """둥근 모서리 사각형"""
    sp = slide.shapes.add_shape(5, x, y, w, h)  # 5 = rounded rect
    sp.line.fill.background()
    if fill:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    else:
        sp.fill.background()
    # 모서리 radius
    try:
        sp.adjustments[0] = 0.05
    except:
        pass
    return sp

def txt(slide, text, x, y, w, h, sz=18, bold=False, color=None,
        align=PP_ALIGN.LEFT, italic=False, wrap=True):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p  = tf.paragraphs[0]
    p.alignment = align
    r  = p.add_run()
    r.text = text
    r.font.size  = Pt(sz)
    r.font.bold  = bold
    r.font.italic = italic
    r.font.name  = "Malgun Gothic"
    if color:
        r.font.color.rgb = color
    return tb

def multi_line(slide, lines, x, y, w, h, sz=18, bold=False, color=None,
               align=PP_ALIGN.LEFT, line_spacing=1.3):
    """여러 줄 텍스트 (각 줄은 별도 paragraph)"""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    from pptx.util import Pt as _Pt
    from pptx.oxml.ns import qn as _qn
    import copy
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = line
        r.font.size  = _Pt(sz)
        r.font.bold  = bold
        r.font.name  = "Malgun Gothic"
        if color:
            r.font.color.rgb = color
    return tb

def img(slide, path, x, y, w, h=None):
    """이미지 삽입, h=None이면 비율 유지"""
    p = Path(path)
    if not p.exists():
        return None
    try:
        im = Image.open(p)
        iw, ih = im.size
        if h is None:
            h = Emu(int(w.emu * ih / iw))
        slide.shapes.add_picture(str(p), x, y, w, h)
    except Exception as e:
        print(f"  [WARN] img failed: {path} — {e}")
    return None

def header_bar(slide, slide_num, total=15, title="", subtitle=""):
    """각 슬라이드 상단 헤더 바"""
    # 상단 네이비 헤더
    rect(slide, 0, 0, SW, Inches(1.3), fill=C_NAVY)
    # 오렌지 액센트 라인
    rect(slide, 0, Inches(1.3), SW, Inches(0.06), fill=C_ORANGE)
    # 제목
    txt(slide, title,
        Inches(0.6), Inches(0.15), Inches(10.0), Inches(0.75),
        sz=32, bold=True, color=C_WHITE)
    # 소제목
    if subtitle:
        txt(slide, subtitle,
            Inches(0.6), Inches(0.82), Inches(9.5), Inches(0.45),
            sz=16, color=C_GRAY)
    # 슬라이드 번호
    txt(slide, f"{slide_num:02d} / {total:02d}",
        Inches(12.0), Inches(0.5), Inches(1.2), Inches(0.4),
        sz=13, color=C_GRAY, align=PP_ALIGN.RIGHT)
    # CleverTap 브랜드 작은 텍스트
    txt(slide, "CleverTap",
        Inches(12.0), Inches(0.05), Inches(1.2), Inches(0.4),
        sz=12, bold=True, color=C_ORANGE, align=PP_ALIGN.RIGHT)

def footer_bar(slide):
    """하단 풋터 바"""
    rect(slide, 0, Inches(7.2), SW, Inches(0.3), fill=C_NAVY2)
    txt(slide, "clevertap.com  |  korea@clevertap.com  |  © 2026 CleverTap Inc.",
        Inches(0.5), Inches(7.22), Inches(12.5), Inches(0.26),
        sz=10, color=C_GRAY, align=PP_ALIGN.CENTER)

def stat_card(slide, x, y, value, label, desc="", accent=C_ORANGE):
    """KPI 통계 카드"""
    W, H = Inches(2.8), Inches(1.8)
    # 카드 배경
    sp = rrect(slide, x, y, W, H, fill=C_NAVY_CARD)
    # 액센트 라인
    rect(slide, x, y, Inches(0.08), H, fill=accent)
    # 값
    txt(slide, value, x+Inches(0.2), y+Inches(0.1), W-Inches(0.3), Inches(0.75),
        sz=36, bold=True, color=accent)
    # 레이블
    txt(slide, label, x+Inches(0.2), y+Inches(0.78), W-Inches(0.3), Inches(0.42),
        sz=15, bold=True, color=C_WHITE)
    if desc:
        txt(slide, desc, x+Inches(0.2), y+Inches(1.22), W-Inches(0.3), Inches(0.45),
            sz=12, color=C_GRAY)

def feature_card(slide, x, y, icon, title, desc, accent=C_ORANGE):
    """기능 소개 카드"""
    W, H = Inches(2.4), Inches(2.4)
    sp = rrect(slide, x, y, W, H, fill=C_NAVY_CARD)
    rect(slide, x, y, W, Inches(0.05), fill=accent)
    txt(slide, icon, x+Inches(0.2), y+Inches(0.15), W-Inches(0.3), Inches(0.65),
        sz=36, color=accent)
    txt(slide, title, x+Inches(0.2), y+Inches(0.8), W-Inches(0.3), Inches(0.5),
        sz=16, bold=True, color=C_WHITE)
    txt(slide, desc, x+Inches(0.2), y+Inches(1.35), W-Inches(0.3), Inches(0.95),
        sz=13, color=C_GRAY)

def check_item(slide, x, y, text, sz=16):
    """체크 아이템 (▶ bullet)"""
    txt(slide, "▶", x, y, Inches(0.35), Inches(0.4), sz=sz-2, color=C_ORANGE, bold=True)
    txt(slide, text, x+Inches(0.35), y, Inches(4.5), Inches(0.4), sz=sz, color=C_WHITE)


# ══════════════════════════════════════════════════
# 슬라이드 빌드 시작
# ══════════════════════════════════════════════════
def build():
    prs = new_prs()
    TOTAL = 15

    # ─────────────────────────────────────────────
    # Slide 1: 표지 (Cover)
    # ─────────────────────────────────────────────
    sl = prs.slides.add_slide(blank(prs))
    # 전체 배경 — 진한 네이비
    rect(sl, 0, 0, SW, SH, fill=C_NAVY2)
    # 우측 오렌지 배경 패널
    rect(sl, Inches(8.5), 0, Inches(4.83), SH, fill=C_ORANGE_DARK)
    # 오버레이 반투명 네이비 그라데이션 효과 (레이어)
    rect(sl, Inches(7.0), 0, Inches(2.0), SH, fill=C_NAVY)
    # 실제 CleverTap 홈페이지 스크린샷 (우측)
    img(sl, ASSETS/"ctx_cover.png", Inches(8.6), Inches(0.5), Inches(4.5), Inches(6.5))

    # 좌측 텍스트
    # 로고 영역
    rect(sl, Inches(0.5), Inches(0.5), Inches(0.08), Inches(0.9), fill=C_ORANGE)
    txt(sl, "CleverTap",
        Inches(0.7), Inches(0.5), Inches(4.0), Inches(0.9),
        sz=28, bold=True, color=C_ORANGE)

    txt(sl, "고객 평생 가치를\n극대화하는",
        Inches(0.5), Inches(1.7), Inches(7.5), Inches(1.6),
        sz=40, bold=False, color=C_WHITE)
    txt(sl, "All-in-One\n모바일 마케팅 플랫폼",
        Inches(0.5), Inches(3.1), Inches(7.5), Inches(1.8),
        sz=44, bold=True, color=C_ORANGE)

    txt(sl, "한국 기업 CRM/마케팅 담당자를 위한\nCleverTap 서비스 제안서 2026",
        Inches(0.5), Inches(5.05), Inches(7.5), Inches(1.0),
        sz=17, color=C_GRAY)

    # 하단 스탯 바
    rect(sl, 0, Inches(6.3), SW, Inches(0.9), fill=C_NAVY_CARD)
    rect(sl, 0, Inches(6.3), SW, Inches(0.04), fill=C_ORANGE)
    for i, (v, l) in enumerate([("10,000+","글로벌 고객사"),("1조+","이벤트/월 처리"),("50+","국가 지원"),("2014","창립")]):
        bx = Inches(1.5 + i*2.8)
        txt(sl, v, bx, Inches(6.38), Inches(2.5), Inches(0.45),
            sz=22, bold=True, color=C_ORANGE, align=PP_ALIGN.CENTER)
        txt(sl, l, bx, Inches(6.83), Inches(2.5), Inches(0.38),
            sz=13, color=C_GRAY, align=PP_ALIGN.CENTER)
    print("  [1] 표지")

    # ─────────────────────────────────────────────
    # Slide 2: 목차
    # ─────────────────────────────────────────────
    sl = prs.slides.add_slide(blank(prs))
    rect(sl, 0, 0, SW, SH, fill=C_LIGHT)
    header_bar(sl, 2, TOTAL, "목차", "CleverTap 제안서 구성")
    footer_bar(sl)

    toc = [
        ("01", "CleverTap 회사 소개", "글로벌 현황 및 핵심 수치"),
        ("02", "한국 시장에서의 필요성", "모바일 마케팅 트렌드 & 과제"),
        ("03", "CleverTap 핵심 기능", "5가지 핵심 기능 개요"),
        ("04", "실시간 세그멘테이션", "정밀 타겟팅으로 전환율 향상"),
        ("05", "옴니채널 캠페인", "Push · Email · SMS · In-App 통합"),
        ("06", "AI 기반 개인화", "Clever.AI — 예측 분석 & STO"),
        ("07", "CleverTap 대시보드", "통합 데이터 분석 & 인사이트"),
        ("08", "A/B 테스트 & 분석", "데이터 기반 의사결정"),
        ("09", "고객 성공 사례", "글로벌 & 아시아 레퍼런스"),
        ("10", "도입 효과 & ROI", "수치로 증명된 임팩트"),
        ("11", "요금제 & 도입 절차", "유연한 플랜 & 온보딩"),
        ("12", "Why CleverTap?", "경쟁사 대비 차별화 강점"),
        ("13", "CTA", "지금 바로 시작하세요"),
    ]
    cols = [toc[:7], toc[7:]]
    for col_idx, col in enumerate(cols):
        cx = Inches(0.7 + col_idx * 6.4)
        for row_idx, (num, title, sub) in enumerate(col):
            ry = Inches(1.55 + row_idx * 0.8)
            # 번호 박스
            rrect(sl, cx, ry+Inches(0.04), Inches(0.5), Inches(0.5), fill=C_ORANGE)
            txt(sl, num, cx, ry+Inches(0.04), Inches(0.5), Inches(0.5),
                sz=14, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
            txt(sl, title, cx+Inches(0.6), ry, Inches(5.5), Inches(0.38),
                sz=17, bold=True, color=C_NAVY)
            txt(sl, sub, cx+Inches(0.6), ry+Inches(0.38), Inches(5.5), Inches(0.38),
                sz=13, color=C_GRAY)
    print("  [2] 목차")

    # ─────────────────────────────────────────────
    # Slide 3: CleverTap 회사 소개
    # ─────────────────────────────────────────────
    sl = prs.slides.add_slide(blank(prs))
    rect(sl, 0, 0, SW, SH, fill=C_NAVY)
    header_bar(sl, 3, TOTAL, "CleverTap 회사 소개", "고객 평생 가치를 높이는 글로벌 1위 올인원 플랫폼")
    footer_bar(sl)

    # 배경 장식
    rect(sl, Inches(9.5), Inches(1.5), Inches(4.0), Inches(5.5), fill=C_NAVY_MID)

    txt(sl, "CleverTap은 무엇인가?",
        Inches(0.6), Inches(1.5), Inches(7.0), Inches(0.7),
        sz=26, bold=True, color=C_ORANGE)

    desc_lines = [
        "CleverTap은 2014년 설립된 글로벌 All-in-One 고객 인게이지먼트 플랫폼으로,",
        "기업이 고객 데이터를 수집 · 분석 · 활용하여 개인화된 경험을 제공하고",
        "고객 유지율(Retention)과 평생 가치(LTV)를 극대화할 수 있도록 지원합니다.",
    ]
    for i, line in enumerate(desc_lines):
        txt(sl, line,
            Inches(0.6), Inches(2.28 + i*0.42), Inches(8.5), Inches(0.4),
            sz=16, color=C_WHITE if i == 0 else C_GRAY)

    # 수치 카드 4개
    stats = [
        ("10,000+", "글로벌 고객사", "50개국 이상"),
        ("1조+",    "월 처리 이벤트", "초당 수백만 건"),
        ("95%+",    "고객 유지율", "업계 최고 수준"),
        ("2026",    "Gartner 리더", "개인화 엔진 부문"),
    ]
    accents = [C_ORANGE, C_GREEN, C_PURPLE, C_ORANGE]
    for i, (v, l, d) in enumerate(stats):
        stat_card(sl, Inches(0.5 + i*3.1), Inches(4.0), v, l, d, accents[i])

    # 우측 이미지
    img(sl, ASSETS/"ctx_intro.png", Inches(9.6), Inches(1.6), Inches(3.5), Inches(5.3))
    print("  [3] 회사 소개")

    # ─────────────────────────────────────────────
    # Slide 4: 한국 시장 필요성
    # ─────────────────────────────────────────────
    sl = prs.slides.add_slide(blank(prs))
    rect(sl, 0, 0, SW, SH, fill=C_LIGHT)
    header_bar(sl, 4, TOTAL, "한국 시장에서의 필요성", "모바일 퍼스트 시대, 개인화 마케팅이 답입니다")
    footer_bar(sl)

    # 왼쪽 — 문제점
    rect(sl, Inches(0.5), Inches(1.5), Inches(6.0), Inches(5.5), fill=C_NAVY)
    rect(sl, Inches(0.5), Inches(1.5), Inches(6.0), Inches(0.05), fill=C_ORANGE)
    txt(sl, "한국 마케터가 직면한 과제",
        Inches(0.7), Inches(1.6), Inches(5.5), Inches(0.6),
        sz=22, bold=True, color=C_ORANGE)

    problems = [
        ("📉", "앱 이탈률 급증", "신규 설치 후 30일 내 75% 이탈"),
        ("🔀", "채널 분산 관리", "Push·Email·SMS 각각 개별 운영"),
        ("🔒", "데이터 사일로",  "고객 데이터 통합 분석 불가"),
        ("⏱",  "분석 지연",     "실시간 인사이트 부재로 의사결정 지연"),
        ("🎯", "개인화 부족",    "일괄 발송으로 고객 이탈 가속화"),
    ]
    for i, (ic, ttl, sub) in enumerate(problems):
        ry = Inches(2.35 + i*0.85)
        txt(sl, ic, Inches(0.75), ry, Inches(0.5), Inches(0.55), sz=22)
        txt(sl, ttl, Inches(1.3), ry+Inches(0.02), Inches(4.5), Inches(0.35),
            sz=16, bold=True, color=C_WHITE)
        txt(sl, sub, Inches(1.3), ry+Inches(0.38), Inches(4.5), Inches(0.38),
            sz=13, color=C_GRAY)

    # 오른쪽 — 시장 데이터
    rect(sl, Inches(6.8), Inches(1.5), Inches(6.0), Inches(5.5), fill=C_WHITE)
    rect(sl, Inches(6.8), Inches(1.5), Inches(6.0), Inches(0.05), fill=C_GREEN)
    txt(sl, "한국 모바일 시장 현황 (2025)",
        Inches(7.0), Inches(1.6), Inches(5.5), Inches(0.6),
        sz=22, bold=True, color=C_NAVY)

    market_data = [
        ("5,200만+", "스마트폰 사용자", C_ORANGE),
        ("87%",       "앱으로 쇼핑 경험", C_GREEN),
        ("3배",       "개인화 마케팅 ROI", C_PURPLE),
        ("68%",       "Push 수신 후 재방문율", C_ORANGE),
    ]
    for i, (v, l, c) in enumerate(market_data):
        ry = Inches(2.35 + i*1.1)
        rect(sl, Inches(7.0), ry, Inches(5.5), Inches(0.9), fill=C_LIGHT)
        rect(sl, Inches(7.0), ry, Inches(0.07), Inches(0.9), fill=c)
        txt(sl, v, Inches(7.2), ry+Inches(0.05), Inches(2.0), Inches(0.5),
            sz=28, bold=True, color=c)
        txt(sl, l, Inches(9.3), ry+Inches(0.18), Inches(3.0), Inches(0.5),
            sz=16, color=C_NAVY)
    print("  [4] 한국 시장")

    # ─────────────────────────────────────────────
    # Slide 5: 핵심 기능 개요 (섹션 슬라이드)
    # ─────────────────────────────────────────────
    sl = prs.slides.add_slide(blank(prs))
    rect(sl, 0, 0, SW, SH, fill=C_NAVY2)
    # 오렌지 사선 장식
    rect(sl, 0, Inches(6.5), SW, Inches(1.0), fill=C_ORANGE)
    rect(sl, Inches(0), Inches(0), Inches(0.5), SH, fill=C_ORANGE)

    txt(sl, "CleverTap\n핵심 기능 5가지",
        Inches(1.0), Inches(1.5), Inches(7.0), Inches(2.0),
        sz=44, bold=True, color=C_WHITE)
    txt(sl, "하나의 플랫폼에서 데이터 수집 · 분석 · 세그멘테이션 · 캠페인 · AI 자동화",
        Inches(1.0), Inches(3.6), Inches(10.0), Inches(0.6),
        sz=18, color=C_GRAY)

    # 5개 기능 아이콘 배지
    features = [
        ("📊", "실시간 분석"),
        ("🎯", "세그멘테이션"),
        ("📱", "옴니채널"),
        ("🤖", "AI 개인화"),
        ("🧪", "A/B 테스트"),
    ]
    for i, (ic, nm) in enumerate(features):
        bx = Inches(1.0 + i*2.4)
        rrect(sl, bx, Inches(4.5), Inches(2.1), Inches(1.7), fill=C_NAVY_CARD)
        rect(sl, bx, Inches(4.5), Inches(2.1), Inches(0.05), fill=C_ORANGE)
        txt(sl, ic, bx+Inches(0.1), Inches(4.55), Inches(1.9), Inches(0.7), sz=32, align=PP_ALIGN.CENTER)
        txt(sl, nm, bx+Inches(0.05), Inches(5.2), Inches(2.0), Inches(0.45),
            sz=14, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    print("  [5] 핵심 기능 개요")

    # ─────────────────────────────────────────────
    # Slide 6: 실시간 세그멘테이션
    # ─────────────────────────────────────────────
    sl = prs.slides.add_slide(blank(prs))
    rect(sl, 0, 0, SW, SH, fill=C_NAVY)
    header_bar(sl, 6, TOTAL, "실시간 세그멘테이션", "정밀 타겟팅으로 전환율과 리텐션을 동시에 향상")
    footer_bar(sl)

    # 왼쪽 설명
    txt(sl, "누구에게, 언제, 무엇을?",
        Inches(0.6), Inches(1.5), Inches(5.8), Inches(0.6),
        sz=24, bold=True, color=C_ORANGE)

    points = [
        ("Dynamic Segment", "실시간 행동 기반 세그먼트 자동 업데이트"),
        ("RFM 분석",         "Recency·Frequency·Monetary 기반 고가치 고객 분류"),
        ("Predictive",      "이탈 위험 / 구매 의도 AI 예측 세그먼트"),
        ("South Korea 필터", "국가·언어·플랜 등 속성 기반 정밀 조합"),
    ]
    for i, (title, desc) in enumerate(points):
        ry = Inches(2.25 + i*1.1)
        rect(sl, Inches(0.6), ry+Inches(0.05), Inches(0.06), Inches(0.7), fill=C_ORANGE)
        txt(sl, title, Inches(0.85), ry+Inches(0.02), Inches(5.2), Inches(0.4),
            sz=17, bold=True, color=C_WHITE)
        txt(sl, desc, Inches(0.85), ry+Inches(0.42), Inches(5.2), Inches(0.45),
            sz=14, color=C_GRAY)

    # 오른쪽 대시보드 이미지
    rect(sl, Inches(6.7), Inches(1.45), Inches(6.45), Inches(5.6), fill=C_NAVY_CARD)
    img(sl, ASSETS/"ctx_segmentation.png", Inches(6.8), Inches(1.55), Inches(6.25), Inches(5.35))
    print("  [6] 세그멘테이션")

    # ─────────────────────────────────────────────
    # Slide 7: 옴니채널 캠페인
    # ─────────────────────────────────────────────
    sl = prs.slides.add_slide(blank(prs))
    rect(sl, 0, 0, SW, SH, fill=C_LIGHT)
    header_bar(sl, 7, TOTAL, "옴니채널 캠페인 자동화", "하나의 대시보드로 모든 채널을 통합 운영하세요")
    footer_bar(sl)

    # 상단 설명
    txt(sl, "채널별 분산 관리는 이제 그만 — CleverTap 하나로 통합",
        Inches(0.6), Inches(1.5), Inches(10.0), Inches(0.55),
        sz=20, bold=True, color=C_NAVY)

    # 채널 배지 5개
    channels = [
        ("📲", "Push\nNotification", "38.4% 오픈율", C_PURPLE),
        ("📧", "Email", "41.2% 오픈율", C_ORANGE),
        ("💬", "In-App\nMessage", "52% 도달율", C_GREEN),
        ("📩", "SMS", "높은 즉시성", C_NAVY_MID),
        ("💚", "WhatsApp", "비즈니스 API", C_GREEN),
    ]
    for i, (ic, nm, stat, c) in enumerate(channels):
        bx = Inches(0.4 + i*2.56)
        rrect(sl, bx, Inches(2.1), Inches(2.35), Inches(1.9), fill=C_WHITE)
        rect(sl, bx, Inches(2.1), Inches(2.35), Inches(0.07), fill=c)
        txt(sl, ic, bx+Inches(0.1), Inches(2.18), Inches(2.1), Inches(0.65), sz=34, align=PP_ALIGN.CENTER)
        txt(sl, nm, bx+Inches(0.05), Inches(2.82), Inches(2.2), Inches(0.6),
            sz=14, bold=True, color=C_NAVY, align=PP_ALIGN.CENTER)
        txt(sl, stat, bx+Inches(0.05), Inches(3.45), Inches(2.2), Inches(0.45),
            sz=12, color=c, align=PP_ALIGN.CENTER)

    # 하단 — 대시보드 이미지
    rect(sl, Inches(0.5), Inches(4.1), Inches(12.3), Inches(2.9), fill=C_NAVY)
    img(sl, ASSETS/"ctx_engagement.png", Inches(0.6), Inches(4.15), Inches(12.1), Inches(2.75))
    print("  [7] 옴니채널")

    # ─────────────────────────────────────────────
    # Slide 8: AI 개인화 (Clever.AI)
    # ─────────────────────────────────────────────
    sl = prs.slides.add_slide(blank(prs))
    rect(sl, 0, 0, SW, SH, fill=C_NAVY)
    header_bar(sl, 8, TOTAL, "Clever.AI™ — AI 기반 개인화 엔진", "머신러닝으로 모든 고객에게 최적의 경험을 자동 제공")
    footer_bar(sl)

    # 왼쪽 — AI 기능 카드들
    ai_features = [
        ("🔮", "이탈 예측 (Churn Prediction)", "7일 이내 이탈 위험 고객을 91.3% 정확도로 예측"),
        ("⏰", "최적 발송 시간 (STO)",         "개인별 가장 높은 오픈율 시간대 자동 계산"),
        ("🛍",  "구매 의도 예측",               "다음 구매 가능성 87.6% 정확도 예측"),
        ("🎁", "동적 콘텐츠 개인화",            "각 유저에게 맞는 제품·메시지 자동 생성"),
    ]
    for i, (ic, title, desc) in enumerate(ai_features):
        ry = Inches(1.55 + i*1.35)
        rrect(sl, Inches(0.5), ry, Inches(6.0), Inches(1.22), fill=C_NAVY_CARD)
        rect(sl, Inches(0.5), ry, Inches(0.07), Inches(1.22), fill=C_PURPLE)
        txt(sl, ic, Inches(0.7), ry+Inches(0.1), Inches(0.7), Inches(0.7), sz=26)
        txt(sl, title, Inches(1.45), ry+Inches(0.1), Inches(4.8), Inches(0.45),
            sz=16, bold=True, color=C_WHITE)
        txt(sl, desc, Inches(1.45), ry+Inches(0.55), Inches(4.8), Inches(0.55),
            sz=13, color=C_GRAY)

    # 오른쪽 — AI 대시보드 이미지
    rect(sl, Inches(6.8), Inches(1.45), Inches(6.35), Inches(5.6), fill=C_NAVY_CARD)
    img(sl, ASSETS/"ctx_ai.png", Inches(6.9), Inches(1.55), Inches(6.15), Inches(5.35))
    print("  [8] AI 개인화")

    # ─────────────────────────────────────────────
    # Slide 9: CleverTap 대시보드
    # ─────────────────────────────────────────────
    sl = prs.slides.add_slide(blank(prs))
    rect(sl, 0, 0, SW, SH, fill=C_NAVY2)
    header_bar(sl, 9, TOTAL, "통합 Analytics 대시보드", "모든 데이터를 하나의 화면에서 — 실시간 인사이트")
    footer_bar(sl)

    # 풀사이즈 대시보드 이미지
    rect(sl, Inches(0.4), Inches(1.45), Inches(12.5), Inches(5.55), fill=C_NAVY_CARD)
    img(sl, ASSETS/"ctx_analytics.png", Inches(0.5), Inches(1.5), Inches(12.3), Inches(5.4))

    # 오버레이 배지
    for i, (v, l, c) in enumerate([("2.4M","총 사용자",C_GREEN),("8.7M","월간 세션",C_ORANGE),("42.1%","D7 리텐션",C_PURPLE)]):
        bx = Inches(0.8 + i*3.5)
        rrect(sl, bx, Inches(6.55), Inches(3.0), Inches(0.5), fill=C_NAVY)
        rect(sl, bx, Inches(6.55), Inches(3.0), Inches(0.04), fill=c)
        txt(sl, f"{v}  {l}", bx+Inches(0.15), Inches(6.58), Inches(2.7), Inches(0.4),
            sz=16, bold=True, color=C_WHITE)
    print("  [9] 대시보드")

    # ─────────────────────────────────────────────
    # Slide 10: A/B 테스트
    # ─────────────────────────────────────────────
    sl = prs.slides.add_slide(blank(prs))
    rect(sl, 0, 0, SW, SH, fill=C_LIGHT)
    header_bar(sl, 10, TOTAL, "A/B 테스트 & 실험", "데이터로 증명하고, 최선의 메시지만 발송하세요")
    footer_bar(sl)

    txt(sl, "무엇이 더 효과적인지 데이터로 알아보세요",
        Inches(0.6), Inches(1.5), Inches(8.0), Inches(0.55),
        sz=20, bold=True, color=C_NAVY)

    # 왼쪽 — A/B 결과 카드
    ab_items = [
        ("Control (A)", "\"지금 업그레이드!\"", "18.3% CTR", C_GRAY),
        ("Variant B ★", "\"오늘 밤 만료됩니다\"", "24.7% CTR", C_PURPLE),
        ("Variant C",   "\"30% 할인 받기\"",    "21.1% CTR", C_GREEN),
    ]
    for i, (lbl, copy_txt, ctr, c) in enumerate(ab_items):
        ry = Inches(2.2 + i*1.45)
        rrect(sl, Inches(0.5), ry, Inches(5.5), Inches(1.28), fill=C_WHITE)
        rect(sl, Inches(0.5), ry, Inches(0.07), Inches(1.28), fill=c)
        txt(sl, lbl, Inches(0.72), ry+Inches(0.1), Inches(3.5), Inches(0.4),
            sz=15, bold=True, color=C_NAVY)
        txt(sl, copy_txt, Inches(0.72), ry+Inches(0.5), Inches(3.5), Inches(0.4),
            sz=14, color=C_GRAY, italic=True)
        txt(sl, ctr, Inches(4.2), ry+Inches(0.3), Inches(1.6), Inches(0.5),
            sz=24, bold=True, color=c, align=PP_ALIGN.RIGHT)

    # 오른쪽 — A/B 이미지
    rect(sl, Inches(6.3), Inches(1.45), Inches(6.8), Inches(5.6), fill=C_NAVY)
    img(sl, ASSETS/"ctx_ab.png", Inches(6.4), Inches(1.55), Inches(6.6), Inches(5.35))
    print("  [10] A/B 테스트")

    # ─────────────────────────────────────────────
    # Slide 11: 고객 성공 사례
    # ─────────────────────────────────────────────
    sl = prs.slides.add_slide(blank(prs))
    rect(sl, 0, 0, SW, SH, fill=C_NAVY)
    header_bar(sl, 11, TOTAL, "글로벌 고객 성공 사례", "전세계 10,000+ 기업이 CleverTap을 선택한 이유")
    footer_bar(sl)

    cases = [
        {
            "company": "Gojek",
            "region":  "동남아 최대 슈퍼앱",
            "results": ["리텐션 +40%", "캠페인 ROI 3배", "이탈률 25% 감소"],
            "color":   C_GREEN,
        },
        {
            "company": "Jio",
            "region":  "인도 1위 통신사",
            "results": ["MAU +35%", "Push CTR 2.5배", "운영비 60% 절감"],
            "color":   C_ORANGE,
        },
        {
            "company": "Hotstar",
            "region":  "아시아 스트리밍 플랫폼",
            "results": ["재구독률 +28%", "알림 오픈율 45%", "세그먼트 300+ 운영"],
            "color":   C_PURPLE,
        },
        {
            "company": "PhonePe",
            "region":  "인도 핀테크 유니콘",
            "results": ["전환율 +52%", "LTV 2.1배 증가", "AI 예측 정확도 91%"],
            "color":   C_GREEN,
        },
    ]
    for i, c in enumerate(cases):
        bx = Inches(0.4 + i*3.22)
        rrect(sl, bx, Inches(1.55), Inches(3.0), Inches(5.4), fill=C_NAVY_CARD)
        rect(sl, bx, Inches(1.55), Inches(3.0), Inches(0.06), fill=c["color"])
        txt(sl, c["company"],
            bx+Inches(0.15), Inches(1.68), Inches(2.7), Inches(0.6),
            sz=26, bold=True, color=C_WHITE)
        txt(sl, c["region"],
            bx+Inches(0.15), Inches(2.3), Inches(2.7), Inches(0.45),
            sz=13, color=C_GRAY)
        rect(sl, bx+Inches(0.15), Inches(2.8), Inches(2.6), Inches(0.04), fill=c["color"])
        for j, r in enumerate(c["results"]):
            txt(sl, f"✓  {r}",
                bx+Inches(0.15), Inches(3.0 + j*0.78), Inches(2.7), Inches(0.65),
                sz=15, color=C_WHITE if j == 0 else C_GRAY)
    print("  [11] 성공 사례")

    # ─────────────────────────────────────────────
    # Slide 12: 도입 효과 & ROI
    # ─────────────────────────────────────────────
    sl = prs.slides.add_slide(blank(prs))
    rect(sl, 0, 0, SW, SH, fill=C_LIGHT)
    header_bar(sl, 12, TOTAL, "도입 효과 & ROI", "CleverTap 도입 후 평균 성과 — 수치로 증명된 임팩트")
    footer_bar(sl)

    roi_stats = [
        ("3.5배",   "마케팅 ROI",           "도입 6개월 내 달성"),
        ("+42%",    "고객 리텐션",           "이탈률 평균 42% 감소"),
        ("60%↓",    "캠페인 운영 비용",      "자동화로 인건비 절감"),
        ("+35%",    "전환율",               "개인화 메시징 효과"),
        ("2.5배",   "고객 평생 가치 (LTV)", "데이터 기반 업셀"),
        ("91%",     "AI 예측 정확도",        "이탈 예측 모델 기준"),
    ]
    for i, (v, l, d) in enumerate(roi_stats):
        row, col = divmod(i, 3)
        bx = Inches(0.5 + col*4.28)
        ry = Inches(1.55 + row*2.7)
        rrect(sl, bx, ry, Inches(4.0), Inches(2.4), fill=C_NAVY)
        rect(sl, bx, ry, Inches(4.0), Inches(0.07), fill=C_ORANGE if i%2==0 else C_GREEN)
        txt(sl, v, bx+Inches(0.25), ry+Inches(0.15), Inches(3.5), Inches(0.9),
            sz=48, bold=True, color=C_ORANGE if i%2==0 else C_GREEN)
        txt(sl, l, bx+Inches(0.25), ry+Inches(1.08), Inches(3.5), Inches(0.45),
            sz=17, bold=True, color=C_WHITE)
        txt(sl, d, bx+Inches(0.25), ry+Inches(1.55), Inches(3.5), Inches(0.42),
            sz=13, color=C_GRAY)
    print("  [12] ROI")

    # ─────────────────────────────────────────────
    # Slide 13: 요금제 & 도입 절차
    # ─────────────────────────────────────────────
    sl = prs.slides.add_slide(blank(prs))
    rect(sl, 0, 0, SW, SH, fill=C_NAVY)
    header_bar(sl, 13, TOTAL, "요금제 & 도입 절차", "비즈니스 규모에 맞는 유연한 플랜으로 시작하세요")
    footer_bar(sl)

    # 상단 — 3개 요금제
    plans = [
        ("Essentials", "스타트업·MAU 50만", ["핵심 분석 기능", "Push + Email 발송", "기본 세그멘테이션", "이메일 지원"], False),
        ("Advanced",   "Mid-Market·MAU 500만", ["Essentials 전체", "AI 예측 세그먼트", "멀티채널 전체", "전담 CSM 배정"], True),
        ("Enterprise", "대기업·무제한 MAU", ["Advanced 전체", "커스텀 SLA", "온프레미스 옵션", "24/7 전용 지원"], False),
    ]
    for i, (nm, sub, feats, rec) in enumerate(plans):
        bx = Inches(0.4 + i*4.3)
        h = Inches(3.5)
        rrect(sl, bx, Inches(1.55), Inches(4.05), h,
              fill=C_ORANGE if rec else C_NAVY_CARD)
        if rec:
            txt(sl, "★ 추천", bx+Inches(0.15), Inches(1.6), Inches(3.7), Inches(0.35),
                sz=13, bold=True, color=C_NAVY, align=PP_ALIGN.CENTER)
        txt(sl, nm, bx+Inches(0.2), Inches(1.97 if rec else 1.72),
            Inches(3.6), Inches(0.6), sz=24, bold=True,
            color=C_NAVY if rec else C_WHITE)
        txt(sl, sub, bx+Inches(0.2), Inches(2.6 if rec else 2.35),
            Inches(3.6), Inches(0.42), sz=13,
            color=C_NAVY2 if rec else C_GRAY)
        for j, f in enumerate(feats):
            txt(sl, f"✓  {f}", bx+Inches(0.2),
                Inches(3.1 + j*0.55 if rec else 2.85 + j*0.55),
                Inches(3.6), Inches(0.48), sz=14,
                color=C_NAVY2 if rec else C_WHITE)

    # 하단 — 4단계 온보딩
    rect(sl, 0, Inches(5.3), SW, Inches(0.04), fill=C_ORANGE)
    txt(sl, "온보딩 4단계 프로세스",
        Inches(0.5), Inches(5.4), Inches(6.0), Inches(0.45),
        sz=18, bold=True, color=C_ORANGE)

    steps = [("킥오프 미팅","목표 & KPI 설정"),("SDK 연동","기술 팀 지원 포함"),("세그먼트 설정","초기 캠페인 구성"),("성과 리뷰","2주 후 파일럿 결과")]
    for i, (s, d) in enumerate(steps):
        bx = Inches(0.5 + i*3.25)
        rrect(sl, bx, Inches(5.95), Inches(3.0), Inches(1.1), fill=C_NAVY_CARD)
        txt(sl, f"0{i+1}", bx+Inches(0.15), Inches(6.0), Inches(0.5), Inches(0.45),
            sz=22, bold=True, color=C_ORANGE)
        txt(sl, s, bx+Inches(0.65), Inches(6.02), Inches(2.2), Inches(0.42),
            sz=15, bold=True, color=C_WHITE)
        txt(sl, d, bx+Inches(0.65), Inches(6.45), Inches(2.2), Inches(0.38),
            sz=12, color=C_GRAY)
    print("  [13] 요금제")

    # ─────────────────────────────────────────────
    # Slide 14: Why CleverTap? (경쟁사 비교)
    # ─────────────────────────────────────────────
    sl = prs.slides.add_slide(blank(prs))
    rect(sl, 0, 0, SW, SH, fill=C_LIGHT)
    header_bar(sl, 14, TOTAL, "Why CleverTap?", "경쟁사 대비 CleverTap만의 차별화된 강점")
    footer_bar(sl)

    # 헤더 행
    headers = ["기능 / 항목", "CleverTap", "Braze", "MoEngage", "Amplitude"]
    col_w = [Inches(3.5), Inches(2.2), Inches(2.2), Inches(2.2), Inches(2.2)]
    cx = Inches(0.4)
    for j, (h, w) in enumerate(zip(headers, col_w)):
        bg = C_ORANGE if j==1 else C_NAVY
        rrect(sl, cx, Inches(1.55), w-Inches(0.1), Inches(0.55), fill=bg)
        txt(sl, h, cx+Inches(0.1), Inches(1.58), w-Inches(0.15), Inches(0.45),
            sz=15, bold=True, color=C_WHITE if j!=0 else C_NAVY,
            align=PP_ALIGN.CENTER)
        cx += w

    rows = [
        ("통합 플랫폼 (올인원)",    ["✅", "⚠ 부분", "⚠ 부분", "❌"]),
        ("실시간 세그멘테이션",      ["✅", "✅",     "✅",      "⚠"]),
        ("AI/ML 예측 분석",         ["✅", "⚠",      "⚠",      "❌"]),
        ("한국어 & 현지 지원",       ["✅", "❌",     "❌",      "❌"]),
        ("가격 경쟁력",             ["✅ 높음", "낮음", "보통", "높음"]),
        ("데이터 거버넌스",          ["✅", "⚠",      "⚠",      "✅"]),
    ]
    for i, (feature, vals) in enumerate(rows):
        ry = Inches(2.18 + i*0.82)
        bg = C_WHITE if i%2==0 else C_GRAY_LIGHT
        rect(sl, Inches(0.4), ry, Inches(12.5), Inches(0.78), fill=bg)
        txt(sl, feature, Inches(0.55), ry+Inches(0.15), Inches(3.2), Inches(0.45),
            sz=15, color=C_NAVY)
        cx = Inches(3.9)
        for j, v in enumerate(vals):
            bg2 = RGBColor(0xFF, 0xF0, 0xE8) if j==0 else None
            if bg2:
                rect(sl, cx, ry, col_w[j+1]-Inches(0.1), Inches(0.78), fill=bg2)
            txt(sl, v, cx+Inches(0.1), ry+Inches(0.15), col_w[j+1]-Inches(0.2), Inches(0.45),
                sz=14, bold=(j==0), color=C_ORANGE if j==0 else C_NAVY, align=PP_ALIGN.CENTER)
            cx += col_w[j+1]
    print("  [14] 경쟁사 비교")

    # ─────────────────────────────────────────────
    # Slide 15: CTA
    # ─────────────────────────────────────────────
    sl = prs.slides.add_slide(blank(prs))
    rect(sl, 0, 0, SW, SH, fill=C_NAVY2)
    # 오렌지 우측 패널
    rect(sl, Inches(9.0), 0, Inches(4.33), SH, fill=C_ORANGE)
    # 장식 라인
    rect(sl, 0, Inches(6.8), SW, Inches(0.7), fill=C_NAVY_CARD)

    # 좌측 텍스트
    txt(sl, "지금 바로 시작하세요",
        Inches(0.7), Inches(1.5), Inches(8.0), Inches(1.0),
        sz=42, bold=True, color=C_WHITE)
    txt(sl, "2주 무료 파일럿으로 CleverTap의 가치를 직접 확인하세요",
        Inches(0.7), Inches(2.6), Inches(8.0), Inches(0.65),
        sz=20, color=C_GRAY)

    # CTA 버튼
    rrect(sl, Inches(0.7), Inches(3.5), Inches(3.5), Inches(0.75), fill=C_ORANGE)
    txt(sl, "무료 데모 신청하기  →",
        Inches(0.7), Inches(3.54), Inches(3.5), Inches(0.65),
        sz=18, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    # 연락처
    contacts = [
        ("🌐", "웹사이트",  "clevertap.com"),
        ("📧", "이메일",    "korea@clevertap.com"),
        ("📞", "전화",      "+82-2-XXXX-XXXX"),
    ]
    for i, (ic, lb, vl) in enumerate(contacts):
        ry = Inches(4.6 + i*0.75)
        txt(sl, f"{ic}  {lb}:", Inches(0.8), ry, Inches(2.0), Inches(0.5),
            sz=15, color=C_GRAY)
        txt(sl, vl, Inches(2.8), ry, Inches(5.5), Inches(0.5),
            sz=15, bold=True, color=C_WHITE)

    # 우측 패널 문구
    txt(sl, "무료\n파일럿\n2주",
        Inches(9.3), Inches(1.5), Inches(3.7), Inches(3.0),
        sz=52, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    txt(sl, "실제 데이터로 효과 검증",
        Inches(9.0), Inches(4.6), Inches(4.2), Inches(0.6),
        sz=18, color=C_NAVY2, align=PP_ALIGN.CENTER)
    txt(sl, "계약 없이 시작 · 언제든 해지 가능",
        Inches(9.0), Inches(5.2), Inches(4.2), Inches(0.55),
        sz=14, color=C_NAVY2, align=PP_ALIGN.CENTER)

    # 하단 풋터
    txt(sl, "CleverTap Inc.  |  clevertap.com  |  Powered by Clever.AI™",
        Inches(0.5), Inches(6.85), Inches(12.0), Inches(0.35),
        sz=11, color=C_GRAY, align=PP_ALIGN.CENTER)
    print("  [15] CTA")

    # ── 저장
    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    sz = OUT.stat().st_size // 1024
    print(f"\n✅ 저장 완료: {OUT}")
    print(f"   슬라이드: {len(prs.slides)}장 / 파일 크기: {sz:,}KB")
    return len(prs.slides)


if __name__ == "__main__":
    print("CleverTap 한국어 제안서 v2 빌드\n")
    n = build()
    print(f"\n완료: {n}장")
