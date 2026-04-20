"""
Pepper Potts — Hulk 핀테크 강의 슬라이드 생성기
핀테크 비즈니스 분석 실전 — 45장 슬라이드
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import pptx.util as util
from copy import deepcopy
import sys

# ─────────────────────────────────────────────
# 디자인 팔레트 (핀테크 다크 블루 테마)
# ─────────────────────────────────────────────
C_DARK_BG     = RGBColor(0x0D, 0x1B, 0x2A)   # 매우 진한 네이비
C_NAVY        = RGBColor(0x1A, 0x2E, 0x4A)   # 진한 네이비 (카드 배경)
C_ACCENT_BLUE = RGBColor(0x00, 0x8B, 0xFF)   # 밝은 파란색 강조
C_ACCENT_TEAL = RGBColor(0x00, 0xD4, 0xAA)   # 청록색 강조
C_ACCENT_ORG  = RGBColor(0xFF, 0x8C, 0x00)   # 오렌지 강조 (소수 활용)
C_WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
C_LIGHT_GRAY  = RGBColor(0xB0, 0xC4, 0xDE)   # 연한 강철 파랑 텍스트
C_MID_GRAY    = RGBColor(0x7A, 0x8D, 0xA0)
C_PART_HEADER = RGBColor(0x00, 0x4E, 0x92)   # 파트 헤더 박스 색
C_GREEN_BG    = RGBColor(0x00, 0x3D, 0x2E)   # 체크리스트 배경

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


def new_prs():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def add_rect(slide, l, t, w, h, fill_rgb=None, line_rgb=None, line_width=None):
    shape = slide.shapes.add_shape(
        pptx.enum.shapes.MSO_SHAPE_TYPE.RECTANGLE if False else 1,  # MSO_CONNECTOR_TYPE
        Inches(l), Inches(t), Inches(w), Inches(h)
    )
    if fill_rgb:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_rgb
    else:
        shape.fill.background()
    if line_rgb:
        shape.line.color.rgb = line_rgb
        if line_width:
            shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()
    return shape


def add_text_box(slide, text, l, t, w, h,
                  font_size=18, bold=False, italic=False,
                  color=C_WHITE, align=PP_ALIGN.LEFT,
                  word_wrap=True, font_name="맑은 고딕"):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = word_wrap
    tf  = txb.text_frame
    tf.word_wrap = word_wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(font_size)
    run.font.bold  = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name  = font_name
    return txb


def set_bg(slide, color: RGBColor):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_bullet_box(slide, lines, l, t, w, h,
                   font_size=15, color=C_WHITE, bullet_char="•",
                   line_spacing=1.2):
    """여러 줄 bullet point 텍스트 박스"""
    from pptx.util import Pt
    from pptx.oxml.ns import qn
    from lxml import etree

    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True

    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(4)
        run = p.add_run()
        if line.startswith("**") and line.endswith("**"):
            run.text = line[2:-2]
            run.font.bold = True
            run.font.size = Pt(font_size + 1)
            run.font.color.rgb = C_ACCENT_TEAL
        elif line.startswith("- ") or line.startswith("• "):
            run.text = f"  {bullet_char} {line[2:]}"
            run.font.size = Pt(font_size)
            run.font.color.rgb = color
        elif line.startswith("✅"):
            run.text = line
            run.font.size = Pt(font_size)
            run.font.color.rgb = C_ACCENT_TEAL
        elif line.startswith("□"):
            run.text = line
            run.font.size = Pt(font_size)
            run.font.color.rgb = C_LIGHT_GRAY
        else:
            run.text = line
            run.font.size = Pt(font_size)
            run.font.color.rgb = color
        run.font.name = "맑은 고딕"
    return txb


def slide_cover(prs):
    """슬라이드 1: 표지"""
    sl = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    set_bg(sl, C_DARK_BG)

    # 상단 그라디언트 느낌 바
    add_rect(sl, 0, 0, 13.33, 0.08, C_ACCENT_BLUE)
    add_rect(sl, 0, 0.08, 13.33, 0.04, C_ACCENT_TEAL)

    # 좌측 강조 바
    add_rect(sl, 0, 0, 0.12, 7.5, C_ACCENT_BLUE)

    # 로고/아이콘 대신 텍스트 레이블
    add_rect(sl, 0.5, 0.5, 2.8, 0.55, C_ACCENT_BLUE)
    add_text_box(sl, "FINTECH MASTERCLASS", 0.5, 0.52, 2.8, 0.5,
                 font_size=10, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    # 메인 제목
    add_text_box(sl, "핀테크 비즈니스 분석 실전", 0.5, 1.5, 12, 1.2,
                 font_size=48, bold=True, color=C_WHITE)

    # 부제목
    add_text_box(sl, "데이터로 읽는 핀테크 — 비즈니스 모델 분석부터 대시보드 구축까지",
                 0.5, 2.9, 12, 0.7, font_size=22, color=C_ACCENT_TEAL)

    # 구분선
    add_rect(sl, 0.5, 3.75, 4, 0.04, C_ACCENT_BLUE)

    # 강의 정보
    info_lines = ["총 강의 시간: 3시간  |  총 슬라이드: 45장  |  실습 포함",
                  "2026.04.20  |  강사: 이재우"]
    add_bullet_box(sl, info_lines, 0.5, 3.9, 12, 0.9,
                   font_size=14, color=C_LIGHT_GRAY)

    # 하단 파트 미리보기 박스 5개
    parts = ["Part 1\n오리엔테이션", "Part 2\n비즈니스 모델", "Part 3\n케이스 스터디",
             "Part 4\n실습", "Part 5\nQ&A"]
    colors_list = [C_ACCENT_BLUE, C_ACCENT_TEAL, C_ACCENT_BLUE, C_ACCENT_TEAL, C_ACCENT_BLUE]
    for i, (part, col) in enumerate(zip(parts, colors_list)):
        x = 0.4 + i * 2.55
        add_rect(sl, x, 5.7, 2.35, 1.3, C_NAVY)
        add_rect(sl, x, 5.7, 2.35, 0.08, col)
        add_text_box(sl, part, x, 5.78, 2.35, 1.2,
                     font_size=12, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    # 하단 바
    add_rect(sl, 0, 7.3, 13.33, 0.2, C_NAVY)
    add_text_box(sl, "© 2026 Fintech Masterclass  |  HULK × PEPPER", 0, 7.32, 13.33, 0.18,
                 font_size=9, color=C_MID_GRAY, align=PP_ALIGN.CENTER)
    return sl


def make_section_divider(prs, part_num, title, subtitle="", duration=""):
    """파트 구분 슬라이드"""
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sl, C_DARK_BG)

    # 좌측 강조 바 (굵게)
    add_rect(sl, 0, 0, 0.5, 7.5, C_ACCENT_BLUE)

    # 파트 번호
    add_rect(sl, 1, 1.5, 3, 0.7, C_ACCENT_BLUE)
    add_text_box(sl, part_num, 1, 1.52, 3, 0.66,
                 font_size=24, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    # 제목
    add_text_box(sl, title, 1, 2.4, 11.3, 1.2,
                 font_size=40, bold=True, color=C_WHITE)

    # 부제목
    if subtitle:
        add_text_box(sl, subtitle, 1, 3.7, 11.3, 0.7,
                     font_size=20, color=C_ACCENT_TEAL)

    # 시간 배지
    if duration:
        add_rect(sl, 1, 4.6, 2, 0.55, C_NAVY)
        add_text_box(sl, f"⏱ {duration}", 1, 4.62, 2, 0.5,
                     font_size=14, color=C_ACCENT_TEAL, align=PP_ALIGN.CENTER)

    add_rect(sl, 0, 7.3, 13.33, 0.2, C_NAVY)
    return sl


def make_content_slide(prs, slide_num, title, bullets, note="", accent_color=None):
    """일반 내용 슬라이드"""
    if accent_color is None:
        accent_color = C_ACCENT_BLUE
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sl, C_DARK_BG)

    # 상단 헤더
    add_rect(sl, 0, 0, 13.33, 1.1, C_NAVY)
    add_rect(sl, 0, 0, 0.12, 1.1, accent_color)
    add_text_box(sl, title, 0.25, 0.05, 11.5, 0.95,
                 font_size=26, bold=True, color=C_WHITE)
    # 슬라이드 번호
    add_text_box(sl, f"{slide_num:02d}", 12.5, 0.1, 0.7, 0.5,
                 font_size=18, bold=True, color=accent_color, align=PP_ALIGN.RIGHT)

    # 콘텐츠 카드
    add_rect(sl, 0.3, 1.25, 12.73, 5.8, C_NAVY)
    add_rect(sl, 0.3, 1.25, 0.08, 5.8, accent_color)

    add_bullet_box(sl, bullets, 0.55, 1.4, 12.4, 5.5,
                   font_size=16, color=C_WHITE)

    # 하단
    add_rect(sl, 0, 7.3, 13.33, 0.2, C_NAVY)
    add_text_box(sl, "핀테크 비즈니스 분석 실전", 0, 7.32, 13.33, 0.18,
                 font_size=8, color=C_MID_GRAY, align=PP_ALIGN.CENTER)
    return sl


def make_two_col_slide(prs, slide_num, title, left_title, left_bullets,
                        right_title, right_bullets, accent_color=None):
    """두 컬럼 슬라이드"""
    if accent_color is None:
        accent_color = C_ACCENT_TEAL
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sl, C_DARK_BG)

    # 상단 헤더
    add_rect(sl, 0, 0, 13.33, 1.1, C_NAVY)
    add_rect(sl, 0, 0, 0.12, 1.1, accent_color)
    add_text_box(sl, title, 0.25, 0.05, 11.5, 0.95,
                 font_size=26, bold=True, color=C_WHITE)
    add_text_box(sl, f"{slide_num:02d}", 12.5, 0.1, 0.7, 0.5,
                 font_size=18, bold=True, color=accent_color, align=PP_ALIGN.RIGHT)

    # 왼쪽 카드
    add_rect(sl, 0.3, 1.25, 6.2, 5.8, C_NAVY)
    add_rect(sl, 0.3, 1.25, 6.2, 0.08, accent_color)
    add_text_box(sl, left_title, 0.45, 1.35, 5.9, 0.5,
                 font_size=16, bold=True, color=accent_color)
    add_bullet_box(sl, left_bullets, 0.45, 1.9, 5.9, 5.0, font_size=14)

    # 오른쪽 카드
    add_rect(sl, 6.8, 1.25, 6.2, 5.8, C_NAVY)
    add_rect(sl, 6.8, 1.25, 6.2, 0.08, C_ACCENT_BLUE)
    add_text_box(sl, right_title, 6.95, 1.35, 5.9, 0.5,
                 font_size=16, bold=True, color=C_ACCENT_BLUE)
    add_bullet_box(sl, right_bullets, 6.95, 1.9, 5.9, 5.0, font_size=14)

    # 하단
    add_rect(sl, 0, 7.3, 13.33, 0.2, C_NAVY)
    add_text_box(sl, "핀테크 비즈니스 분석 실전", 0, 7.32, 13.33, 0.18,
                 font_size=8, color=C_MID_GRAY, align=PP_ALIGN.CENTER)
    return sl


def make_table_slide(prs, slide_num, title, headers, rows, accent_color=None):
    """테이블 슬라이드"""
    if accent_color is None:
        accent_color = C_ACCENT_BLUE
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sl, C_DARK_BG)

    add_rect(sl, 0, 0, 13.33, 1.1, C_NAVY)
    add_rect(sl, 0, 0, 0.12, 1.1, accent_color)
    add_text_box(sl, title, 0.25, 0.05, 11.5, 0.95,
                 font_size=26, bold=True, color=C_WHITE)
    add_text_box(sl, f"{slide_num:02d}", 12.5, 0.1, 0.7, 0.5,
                 font_size=18, bold=True, color=accent_color, align=PP_ALIGN.RIGHT)

    n_cols = len(headers)
    n_rows = len(rows)
    col_w_inch = 12.5 / n_cols
    row_h_inch = min(5.5 / (n_rows + 1), 0.75)

    # 헤더 행
    for ci, h in enumerate(headers):
        x = 0.4 + ci * col_w_inch
        add_rect(sl, x, 1.3, col_w_inch - 0.05, row_h_inch - 0.02, accent_color)
        add_text_box(sl, h, x + 0.05, 1.32, col_w_inch - 0.1, row_h_inch - 0.05,
                     font_size=13, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    # 데이터 행
    for ri, row in enumerate(rows):
        row_bg = C_NAVY if ri % 2 == 0 else RGBColor(0x15, 0x24, 0x3A)
        for ci, cell in enumerate(row):
            x = 0.4 + ci * col_w_inch
            y = 1.3 + (ri + 1) * row_h_inch
            add_rect(sl, x, y, col_w_inch - 0.05, row_h_inch - 0.02, row_bg)
            font_col = C_ACCENT_TEAL if ci == 0 else C_LIGHT_GRAY
            add_text_box(sl, cell, x + 0.05, y + 0.02, col_w_inch - 0.1,
                         row_h_inch - 0.05, font_size=12, color=font_col)

    add_rect(sl, 0, 7.3, 13.33, 0.2, C_NAVY)
    add_text_box(sl, "핀테크 비즈니스 분석 실전", 0, 7.32, 13.33, 0.18,
                 font_size=8, color=C_MID_GRAY, align=PP_ALIGN.CENTER)
    return sl


def make_kpi_slide(prs, slide_num, title, kpis, accent_color=None):
    """KPI 카드 그리드 슬라이드"""
    if accent_color is None:
        accent_color = C_ACCENT_TEAL
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sl, C_DARK_BG)

    add_rect(sl, 0, 0, 13.33, 1.1, C_NAVY)
    add_rect(sl, 0, 0, 0.12, 1.1, accent_color)
    add_text_box(sl, title, 0.25, 0.05, 11.5, 0.95,
                 font_size=26, bold=True, color=C_WHITE)
    add_text_box(sl, f"{slide_num:02d}", 12.5, 0.1, 0.7, 0.5,
                 font_size=18, bold=True, color=accent_color, align=PP_ALIGN.RIGHT)

    # kpis: list of (name, value, desc)
    per_row = 5
    rows_n = (len(kpis) + per_row - 1) // per_row
    card_w = 12.5 / per_row
    card_h = 5.8 / rows_n

    for i, (kname, kval, kdesc) in enumerate(kpis):
        r, c = divmod(i, per_row)
        x = 0.4 + c * card_w
        y = 1.25 + r * card_h
        add_rect(sl, x, y, card_w - 0.1, card_h - 0.1, C_NAVY)
        add_rect(sl, x, y, card_w - 0.1, 0.06,
                 C_ACCENT_TEAL if i % 2 == 0 else C_ACCENT_BLUE)
        add_text_box(sl, kname, x + 0.05, y + 0.1, card_w - 0.2, 0.4,
                     font_size=13, bold=True, color=C_ACCENT_TEAL if i % 2 == 0 else C_ACCENT_BLUE)
        add_text_box(sl, kval, x + 0.05, y + 0.55, card_w - 0.2, 0.5,
                     font_size=11, color=C_WHITE)
        add_text_box(sl, kdesc, x + 0.05, y + 1.1, card_w - 0.2, card_h - 1.2,
                     font_size=10, color=C_LIGHT_GRAY)

    add_rect(sl, 0, 7.3, 13.33, 0.2, C_NAVY)
    add_text_box(sl, "핀테크 비즈니스 분석 실전", 0, 7.32, 13.33, 0.18,
                 font_size=8, color=C_MID_GRAY, align=PP_ALIGN.CENTER)
    return sl


def make_quote_slide(prs, slide_num, title, quote, source="", accent_color=None):
    """명언/핵심 메시지 슬라이드"""
    if accent_color is None:
        accent_color = C_ACCENT_ORG
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sl, C_DARK_BG)

    add_rect(sl, 0, 0, 13.33, 1.1, C_NAVY)
    add_rect(sl, 0, 0, 0.12, 1.1, accent_color)
    add_text_box(sl, title, 0.25, 0.05, 11.5, 0.95,
                 font_size=26, bold=True, color=C_WHITE)
    add_text_box(sl, f"{slide_num:02d}", 12.5, 0.1, 0.7, 0.5,
                 font_size=18, bold=True, color=accent_color, align=PP_ALIGN.RIGHT)

    # 큰따옴표 장식
    add_text_box(sl, '"', 0.5, 1.3, 2, 2, font_size=120, color=accent_color)

    add_text_box(sl, quote, 1.5, 2.2, 10.5, 2.5,
                 font_size=26, italic=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    if source:
        add_text_box(sl, f"— {source}", 0, 5.2, 13.33, 0.5,
                     font_size=14, color=C_LIGHT_GRAY, align=PP_ALIGN.CENTER)

    add_rect(sl, 0, 7.3, 13.33, 0.2, C_NAVY)
    add_text_box(sl, "핀테크 비즈니스 분석 실전", 0, 7.32, 13.33, 0.18,
                 font_size=8, color=C_MID_GRAY, align=PP_ALIGN.CENTER)
    return sl


# ─────────────────────────────────────────────
# 슬라이드 데이터 정의
# ─────────────────────────────────────────────

def build_all_slides(prs):
    # Slide 1: 표지
    slide_cover(prs)
    print("  [1/45] 표지 완료")

    # ─── PART 1 (슬라이드 2~6) ───────────────────
    make_section_divider(prs, "PART 1", "오프닝 & 오리엔테이션",
                         "왜 지금 핀테크 데이터 분석인가?", "15분")
    print("  [Part 1 구분] 완료")

    # Slide 2: 오늘의 여정
    make_content_slide(prs, 2, "오늘의 여정 — 3시간 로드맵", [
        "**오늘 우리가 함께 걷는 길:**",
        "- Part 1 (15분): 오리엔테이션 — 왜 핀테크 데이터 분석인가?",
        "- Part 2 (45분): 비즈니스 모델 분석 방법론 심화",
        "- Part 3 (60분): 케이스 스터디 — 토스 · 카카오페이 · Stripe",
        "- Part 4 (60분): 실습 — Amplitude & Looker Studio 핸즈온",
        "- Part 5 (20분): 마무리 & Q&A",
        "",
        "**오늘의 목표:**",
        '  "핀테크 데이터를 보는 눈 장착하기"',
    ])
    print("  [2/45] 완료")

    # Slide 3: 왜 지금 핀테크 데이터 분석인가?
    make_content_slide(prs, 3, "왜 지금 핀테크 데이터 분석인가?", [
        "**글로벌 핀테크 시장 폭발적 성장:**",
        "  • 2024년 $340B → 2030년 $1.1T (연평균 성장률 22%)",
        "  • 한국 간편결제 연간 거래액: 200조원 돌파",
        "",
        "**현실:**",
        '  "핀테크 PM/마케터에게 데이터 리터러시는 선택이 아닌 필수"',
        "",
        "**질문:**",
        '  "여러분의 서비스에서 오늘 얼마나 많은 결제가 일어났을까요?"',
    ], accent_color=C_ACCENT_TEAL)
    print("  [3/45] 완료")

    # Slide 4: 수강생 여러분은 누구인가요?
    make_content_slide(prs, 4, "수강생 여러분은 누구인가요?", [
        "**현재 직무는?**",
        "  • PM / 마케터 / 기획자 / 전략 / 개발 / 기타",
        "",
        "**현재 근무 업종은?**",
        "  • 핀테크 / 이커머스 / 일반 IT / 금융 / 기타",
        "",
        "**오늘 가장 얻어가고 싶은 것은?**",
        "",
        "  → Slido 또는 구두 질문으로 진행",
    ], accent_color=C_ACCENT_ORG)
    print("  [4/45] 완료")

    # Slide 5: 강의 후 여러분이 할 수 있는 것
    make_content_slide(prs, 5, "강의 후 여러분이 할 수 있는 것", [
        "✅ 핀테크 비즈니스 모델 유형별로 분류하고 KPI를 정의할 수 있다",
        "✅ 비즈니스 모델 캔버스를 핀테크 서비스에 적용할 수 있다",
        "✅ 토스·카카오페이·Stripe 전략을 데이터 관점에서 해석할 수 있다",
        "✅ Amplitude로 퍼널·리텐션·코호트 분석을 수행할 수 있다",
        "✅ Looker Studio로 핵심 지표 대시보드를 만들 수 있다",
    ], accent_color=C_ACCENT_TEAL)
    print("  [5/45] 완료")

    # ─── PART 2 (슬라이드 6~14) ──────────────────
    make_section_divider(prs, "PART 2", "비즈니스 모델 분석 방법론 심화",
                         "데이터를 보기 전에, 비즈니스 구조를 먼저 이해하라", "45분")
    print("  [Part 2 구분] 완료")

    # Slide 6
    make_content_slide(prs, 6, "Part 2 — 비즈니스 모델을 해부하는 법", [
        '"데이터를 보기 전에, 비즈니스 구조를 먼저 이해해야 한다"',
        "",
        "**비즈니스 모델 = 누가 / 무엇을 / 어떻게 / 얼마에 / 왜 구매하는가**",
        "",
        "- 핀테크는 전통 금융과 달리 데이터 기반 수익화가 핵심",
        "",
        "**이번 파트 목표:**",
        "  핀테크 비즈니스를 구조화해서 분석하는 눈 기르기",
        "",
        '  → "토스는 돈을 어떻게 버나요?" — 한번 생각해보세요',
    ])
    print("  [6/45] 완료")

    # Slide 7: 핀테크 비즈니스 모델 5대 유형 (테이블)
    make_table_slide(prs, 7, "핀테크 비즈니스 모델 5대 유형",
                     ["유형", "대표 서비스", "수익 모델"],
                     [
                         ["결제/송금", "토스, 카카오페이, Stripe", "거래 수수료, MDR"],
                         ["대출/크레딧", "토스뱅크, 핀다, 카카오뱅크", "이자 수익, 중개 수수료"],
                         ["투자/자산관리", "토스증권, Robinhood", "거래 수수료, 구독료"],
                         ["보험 (인슈어테크)", "카카오손보, Lemonade", "보험료 수입"],
                         ["B2B 인프라", "Plaid, Marqeta, AWS Pay", "API 사용료, 라이선스"],
                     ])
    print("  [7/45] 완료")

    # Slide 8: 수익 모델 심화
    make_content_slide(prs, 8, "수익 모델 심화 — 핀테크는 어떻게 돈을 버나?", [
        "**수수료형**: 거래당 % 또는 고정액 과금",
        "  • 예: Stripe — 2.9% + $0.30 per transaction",
        "",
        "**이자 마진형**: 예대마진 (대출 이자 - 예금 이자)",
        "  • 예: 토스뱅크, 카카오뱅크",
        "",
        "**구독형**: 월정액 프리미엄 서비스",
        "  • 예: Robinhood Gold $5/월",
        "",
        "**데이터/광고형**: 금융 행동 데이터 활용",
        "  • 예: 신용평가 데이터 판매, 맞춤형 금융광고",
        "",
        "**핵심 질문:** 이 서비스의 주 수익원은 무엇인가? 부 수익원은?",
    ], accent_color=C_ACCENT_TEAL)
    print("  [8/45] 완료")

    # Slide 9: BMC
    make_content_slide(prs, 9, "비즈니스 모델 캔버스 (BMC) — 핀테크 적용", [
        "**BMC 9개 블록:**",
        "  1. 고객 세그먼트 (Customer Segments)",
        "  2. 가치 제안 (Value Propositions)",
        "  3. 채널 (Channels)",
        "  4. 고객 관계 (Customer Relationships)",
        "  5. 수익 흐름 (Revenue Streams)",
        "  6. 핵심 자원 (Key Resources)",
        "  7. 핵심 활동 (Key Activities)",
        "  8. 핵심 파트너십 (Key Partnerships)",
        "  9. 비용 구조 (Cost Structure)",
        "",
        "**핀테크 특수 요소:** 라이선스/규제, 데이터 자산, 네트워크 효과",
    ])
    print("  [9/45] 완료")

    # Slide 10: 고객 세그먼트와 LTV
    make_content_slide(prs, 10, "고객 세그먼트와 LTV — 누가 진짜 핵심 고객인가?", [
        "**모든 고객이 동등하지 않다 — LTV로 핵심 세그먼트 찾기**",
        "",
        "**세그먼트 분류 기준:**",
        "  • 인구통계 / 행동 / 거래 패턴 / 금융 프로파일",
        "",
        "**LTV 계산 공식:**",
        "  LTV = ARPU × (1 / Churn Rate)",
        "",
        "**핀테크 세그먼트 예시:**",
        "  • 소액 다빈도 사용자 vs. 고액 저빈도 사용자",
        "  • 단일 기능 사용자 vs. 다기능 크로스셀 사용자",
        "",
        '**인사이트:** "어떤 세그먼트의 Whale 고객을 획득하고 유지해야 하는가?"',
    ], accent_color=C_ACCENT_TEAL)
    print("  [10/45] 완료")

    # Slide 11: 핵심 KPI 총정리 (테이블)
    make_table_slide(prs, 11, "핀테크 핵심 KPI 총정리",
                     ["KPI", "정의", "핵심 질문"],
                     [
                         ["GMV", "총 거래액", "우리 플랫폼을 통해 얼마나 거래되나?"],
                         ["Take Rate", "GMV 대비 수익률", "거래에서 얼마를 가져가나?"],
                         ["CAC", "고객 획득 비용", "신규 고객 1명 데려오는데 얼마?"],
                         ["LTV", "고객 생애 가치", "고객 1명이 평생 얼마나 기여하나?"],
                         ["LTV/CAC", "수익성 비율", ">3이어야 건강한 사업"],
                         ["DAU/MAU", "점착성 (Stickiness)", ">20%이면 좋은 앱"],
                         ["Churn Rate", "이탈률", "얼마나 빠져나가나?"],
                         ["NPS", "추천 지수", "만족하고 추천하나?"],
                     ], accent_color=C_ACCENT_TEAL)
    print("  [11/45] 완료")

    # Slide 12: 데이터로 비즈니스 모델을 읽는 법
    make_two_col_slide(prs, 12, "데이터로 비즈니스 모델을 읽는 법",
                        "4가지 진단 축",
                        [
                            "성장 진단",
                            "  MAU 증가율 + CAC 추이",
                            "  → 효율적 성장인가?",
                            "",
                            "수익성 진단",
                            "  Take Rate × GMV vs. 고정비",
                            "  → 단위 경제성 확보 여부",
                        ],
                        "핵심 프레임워크",
                        [
                            "리텐션 진단",
                            "  Churn Rate + LTV/CAC",
                            "  → 장기 지속 가능성",
                            "",
                            "참여도 진단",
                            "  DAU/MAU + 기능별 사용률",
                            "  → 핵심 기능 파악",
                            "",
                            "Growth Accounting",
                            "  신규 + 복귀 - 이탈 = 순증가",
                        ])
    print("  [12/45] 완료")

    # Slide 13: 네트워크 효과와 데이터 해자
    make_content_slide(prs, 13, "핀테크 경쟁 우위 — 네트워크 효과와 데이터 해자", [
        "**네트워크 효과:**",
        "  사용자 증가 → 서비스 가치 증가",
        "  예: 카카오페이 — 카카오톡 5천만 사용자 기반",
        "",
        "**데이터 해자:**",
        "  거래 데이터 축적 → 신용평가 정교화 → 리스크 감소 선순환",
        "",
        "**규제 해자:**",
        "  금융 라이선스 획득 비용 → 진입 장벽",
        "",
        "**전환 비용:**",
        "  금융 데이터 이전 번거로움 → Lock-in 효과",
        "",
        '**분석 질문:** "이 회사의 해자(Moat)는 무엇인가?"',
    ], accent_color=C_ACCENT_ORG)
    print("  [13/45] 완료")

    # Slide 14: 중간 점검
    make_content_slide(prs, 14, "[중간 점검] 분석 프레임워크 한 장 요약", [
        "**핀테크 비즈니스 분석 5단계:**",
        "",
        "  1️⃣  비즈니스 모델 유형 분류 (결제/대출/투자/보험/B2B)",
        "  2️⃣  수익 구조 매핑 (수수료/이자/구독/데이터)",
        "  3️⃣  BMC로 전체 구조 시각화",
        "  4️⃣  핵심 KPI 정의 (GMV, Take Rate, CAC, LTV, Churn)",
        "  5️⃣  경쟁 우위 분석 (네트워크 효과, 데이터 해자, 규제 해자)",
        "",
        '  → "이 프레임워크를 들고 케이스 스터디로 갑니다"',
    ], accent_color=C_ACCENT_TEAL)
    print("  [14/45] 완료")

    # ─── PART 3 (슬라이드 15~27) ─────────────────
    make_section_divider(prs, "PART 3", "케이스 스터디",
                         "토스 · 카카오페이 · Stripe — 실전 분석 워크숍", "60분")
    print("  [Part 3 구분] 완료")

    # Slide 15
    make_content_slide(prs, 15, "Part 3 — 케이스 스터디 진행 방식", [
        "**진행 방식:**",
        "  강의 설명 (10분) → 팀별 BMC 작성 (5분) → 발표 및 토론 (5분)",
        "",
        "**케이스 1:** 토스 (Toss) — 슈퍼앱 전략",
        "**케이스 2:** 카카오페이 — 플랫폼 전환",
        "**케이스 3:** Stripe — B2B API 경제",
        "",
        "**준비물:** BMC 양식 핸드아웃, 포스트잇, 마커",
        "",
        "  → 팀 구성: 4~5인 1팀",
    ])
    print("  [15/45] 완료")

    # Slide 16: 토스 개요
    make_content_slide(prs, 16, "케이스 1 — 토스 (Toss) 개요", [
        "**설립:** 2014년  |  운영사: 비바리퍼블리카",
        "",
        "**시작:** 무료 간편송금 (은행 앱의 복잡함을 해결)",
        "",
        "**현재:**",
        "  • MAU 2,400만  |  금융상품 30개+",
        "  • 토스뱅크 + 토스증권 + 토스페이먼츠",
        "",
        '**핵심 전략:** "금융의 모든 것을 하나의 앱에서"',
        "",
        "**수익 구조:** 금융상품 중개 수수료 + 대출 이자 + 결제 수수료 + 광고",
    ], accent_color=C_ACCENT_BLUE)
    print("  [16/45] 완료")

    # Slide 17: 토스 BMC
    make_two_col_slide(prs, 17, "케이스 1 — 토스 비즈니스 모델 캔버스",
                        "고객/가치/채널/수익",
                        [
                            "고객 세그먼트",
                            "  20~40대 모바일 금융 사용자",
                            "  중소 사업자 (토스페이)",
                            "",
                            "가치 제안",
                            "  무료 송금, 원스톱 금융 관리",
                            "  쉬운 UX",
                            "",
                            "수익 흐름",
                            "  금융상품 수수료, 대출 이자",
                            "  가맹점 MDR, 광고",
                        ],
                        "자원/활동/채널",
                        [
                            "핵심 자원",
                            "  금융 라이선스",
                            "  2,400만 사용자 데이터",
                            "  기술 인재",
                            "",
                            "핵심 활동",
                            "  상품 중개, 신용평가",
                            "  앱 개발, 규제 대응",
                            "",
                            "채널",
                            "  앱스토어, 입소문, 카드 혜택",
                        ])
    print("  [17/45] 완료")

    # Slide 18: 토스 핵심 지표
    make_content_slide(prs, 18, "케이스 1 — 토스 핵심 지표 분석", [
        "**MAU 성장 추이:**",
        "  2018년 500만 → 2023년 2,400만 (5배 성장)",
        "",
        "**핵심 교차판매 전환율:**",
        "  송금 사용자 → 금융상품 가입 전환율 %",
        "",
        "**GMV:**",
        "  토스페이먼츠 연간 결제처리금액 80조원+",
        "",
        "**Take Rate 계산 예시:**",
        "  수수료 수익 / GMV",
        "",
        '**인사이트:** "사용자를 잠그고 (Lock-in), 교차판매로 LTV를 극대화"',
    ], accent_color=C_ACCENT_TEAL)
    print("  [18/45] 완료")

    # Slide 19: 토스 토론
    make_quote_slide(prs, 19, "케이스 1 — 토론: 토스의 다음 성장 동력은?",
                     "토스가 아직 공략하지 못한 고객 세그먼트는?\nTake Rate를 높이려면 어떤 전략이 필요한가?",
                     "팀별 2분 논의 → 한 문장 결론 공유")
    print("  [19/45] 완료")

    # Slide 20: 카카오페이 개요
    make_content_slide(prs, 20, "케이스 2 — 카카오페이 개요", [
        "**설립:** 2017년 분사  |  모기업: 카카오",
        "",
        "**시작:** 카카오톡 간편결제",
        "",
        "**현재:** 결제 + 보험 + 투자 + 대출 + 인증 서비스",
        "",
        "**MAU:** 2,300만+  |  연간 결제 거래액 120조원+",
        "",
        "**핵심 강점:** 카카오톡 5,000만 사용자 기반 즉시 침투",
        "",
        "**수익 구조:** 결제 수수료 + 보험료 + 금융상품 수수료",
    ], accent_color=C_ACCENT_TEAL)
    print("  [20/45] 완료")

    # Slide 21: 카카오페이 전환 전략
    make_content_slide(prs, 21, "케이스 2 — 결제에서 플랫폼으로의 전환 전략", [
        "**전환 단계:**",
        "  1. 결제 MAU 확보 (무료/편의성)",
        "  2. 결제 데이터 → 금융 프로파일링",
        "  3. 신용/보험 상품 크로스셀",
        "  4. 투자 서비스 추가 → 자산 관리 플랫폼화",
        "",
        "**데이터 관점:**",
        "  결제 빈도/금액대 → 보험 추천 정확도 → 전환율 향상",
        "",
        "**핵심 지표 변화:**",
        "  결제 MAU → 다기능 사용자 비율 (Multi-product Ratio)",
        "",
        '**왜 결제부터 시작했을까?** "가장 빈번한 접점이 가장 많은 데이터를 생성"',
    ])
    print("  [21/45] 완료")

    # Slide 22: GMV vs Take Rate 딜레마
    make_content_slide(prs, 22, "케이스 2 — GMV vs. Take Rate 딜레마", [
        "**결제 GMV:** 빠르게 성장 (120조원+)",
        "",
        "**결제 Take Rate:** 규제로 인해 압박 (카드 MDR 인하 정책)",
        "",
        "**해결책:** 고마진 부문으로 믹스 전환",
        "  보험 Take Rate > 결제 Take Rate",
        "",
        "**분석 포인트:**",
        '  "성장하는 GMV에서 어떻게 마진을 지키는가?"',
        "",
        "  → 이 딜레마는 모든 결제 플랫폼의 공통 과제",
    ], accent_color=C_ACCENT_ORG)
    print("  [22/45] 완료")

    # Slide 23: Stripe 개요
    make_content_slide(prs, 23, "케이스 3 — Stripe 개요", [
        "**설립:** 2010년  |  창업자: Patrick & John Collison (아일랜드)",
        "",
        '**핵심 가치 제안:** "7줄의 코드로 결제 구현"',
        "",
        "**현재:** 기업 가치 $65B+  |  135개국 서비스  |  100만+ 비즈니스 고객",
        "",
        "**수익 모델:** 거래당 2.9% + $0.30 (카드) + 볼륨 디스카운트",
        "",
        "**주요 고객:** Shopify, Lyft, DoorDash, Amazon",
        "",
        "**확장:** Stripe Atlas, Stripe Radar, Stripe Treasury",
    ], accent_color=C_ACCENT_BLUE)
    print("  [23/45] 완료")

    # Slide 24: Stripe API 경제
    make_content_slide(prs, 24, "케이스 3 — API 경제와 개발자 GTM 전략", [
        "**Bottom-up Sales:**",
        "  개발자가 먼저 도입 → 회사 전체로 확산",
        "",
        "**Developer Experience First:**",
        "  문서화, SDK, Sandbox 환경 최우선",
        "",
        "**Network Effect:**",
        "  더 많은 가맹점 → Stripe Connect 생태계 강화",
        "",
        "**데이터 자산화:**",
        "  수백만 거래 데이터 → Radar AI 사기 탐지 정확도 향상",
        "",
        "**LTV 모델:**",
        "  기업이 성장할수록 GMV 증가 → Stripe 수익 비례 증가",
    ], accent_color=C_ACCENT_TEAL)
    print("  [24/45] 완료")

    # Slide 25: Stripe Radar
    make_content_slide(prs, 25, "케이스 3 — 데이터로 만든 경쟁 우위: Stripe Radar", [
        "**문제:** 온라인 결제 사기(Fraud) → 가맹점 손실, 고객 불만",
        "",
        "**솔루션:** 수억 건의 거래 데이터로 ML 모델 학습 → Radar",
        "",
        "**결과:**",
        "  99.9%+ 결제 성공률  |  사기 탐지 정확도 업계 최고",
        "",
        "**비즈니스 효과:** Radar 과금 → 추가 수익 창출",
        "",
        "**핵심 인사이트:**",
        '  "데이터 자산 → AI 제품 → 추가 수익 → 데이터 확장" 선순환',
        "",
        '  → "여러분 회사의 데이터 자산으로 만들 수 있는 제품은?"',
    ], accent_color=C_ACCENT_ORG)
    print("  [25/45] 완료")

    # Slide 26: 비교 분석 (테이블)
    make_table_slide(prs, 26, "케이스 스터디 비교 분석",
                     ["항목", "토스", "카카오페이", "Stripe"],
                     [
                         ["고객", "B2C", "B2C", "B2B"],
                         ["진입 전략", "무료 기능", "생태계 기반", "개발자 경험"],
                         ["수익화", "크로스셀", "믹스 전환", "거래 과금"],
                         ["데이터 활용", "신용평가", "금융 프로파일링", "사기 탐지 AI"],
                         ["핵심 해자", "사용자 데이터", "카카오 생태계", "개발자 네트워크"],
                     ], accent_color=C_ACCENT_TEAL)
    print("  [26/45] 완료")

    # Slide 27: 내 서비스 체크리스트
    make_content_slide(prs, 27, "내 서비스에 적용하기 — 분석 체크리스트", [
        "□ 우리 서비스의 비즈니스 모델 유형은?",
        "□ 주요 수익원 Top 3는?",
        "□ 핵심 고객 세그먼트는? 각각의 LTV는?",
        "□ 지금 추적하고 있는 KPI는? 빠진 KPI는?",
        "□ 우리의 경쟁 우위(해자)는 무엇인가?",
        "□ 데이터 자산을 어떻게 활용하고 있는가?",
        "",
        "  → 2분 개인 메모 시간  |  강의 후 액션 아이템으로!",
    ], accent_color=C_ACCENT_TEAL)
    print("  [27/45] 완료")

    # ─── PART 4 (슬라이드 28~35) ─────────────────
    make_section_divider(prs, "PART 4", "실습 — Amplitude & Looker Studio",
                         "FinFlow 가상 핀테크 앱 데이터로 직접 분석", "60분")
    print("  [Part 4 구분] 완료")

    # Slide 28: 실습 안내
    make_content_slide(prs, 28, "Part 4 — 실습 안내", [
        "**가상 서비스:** FinFlow — 간편결제 핀테크 앱",
        "",
        "**데이터셋:**",
        "  • 사용자 이벤트 10,000건",
        "  • 트랜잭션 5,000건",
        "",
        "**미션:**",
        "  1️⃣  결제 퍼널에서 가장 큰 이탈 구간 찾기 (Amplitude)",
        "  2️⃣  30일 리텐션 패턴 분석 (Amplitude)",
        "  3️⃣  GMV & 결제 성공률 대시보드 구성 (Looker Studio)",
        "",
        "  → 팀별 협력 실습  |  15분마다 인사이트 공유",
    ])
    print("  [28/45] 완료")

    # Slide 29: FinFlow 소개
    make_content_slide(prs, 29, "FinFlow 앱 소개 — 가상 데이터 컨텍스트", [
        "**FinFlow:** 20~35세 타겟 간편결제 + 가계부 기능 앱",
        "",
        "**주요 기능:**",
        "  간편결제 / 소비 분석 / 예산 설정 / 친구에게 송금",
        "",
        "**현재 상황:**",
        "  • MAU 5만",
        "  • 결제 전환율 저조",
        "  • 30일 리텐션 하락 중",
        "",
        "**데이터 기간:** 최근 90일",
        "",
        '**분석 목표:** "어디서 왜 이탈하는가? 어떻게 개선할 것인가?"',
        "",
        '  → "실제 스타트업 입사 첫 주에 이 데이터를 받았다고 상상하세요"',
    ], accent_color=C_ACCENT_TEAL)
    print("  [29/45] 완료")

    # Slide 30: 데이터셋 구조
    make_two_col_slide(prs, 30, "데이터셋 구조 소개",
                        "user_events.csv",
                        [
                            "user_id: 사용자 고유 ID",
                            "event_type: 앱 실행/회원가입/",
                            "  계좌연결/첫결제/재결제/이탈",
                            "timestamp: 이벤트 발생 시각",
                            "amount: 결제 금액 (원)",
                            "category: 식비/교통/쇼핑/구독",
                            "device_type: iOS / Android",
                            "session_id: 세션 고유 ID",
                        ],
                        "transactions.csv",
                        [
                            "transaction_id: 거래 고유 ID",
                            "user_id: 사용자 ID",
                            "merchant_id: 가맹점 ID",
                            "amount: 거래 금액",
                            "status: 성공/실패/취소",
                            "payment_method: 결제 수단",
                            "timestamp: 거래 시각",
                        ])
    print("  [30/45] 완료")

    # Slide 31: Amplitude 실습 1
    make_content_slide(prs, 31, "Amplitude 실습 1 — 퍼널 분석", [
        "**분석 목표:** 앱 설치 → 첫 결제까지의 전환율 최적화",
        "",
        "**퍼널 단계:**",
        "  1️⃣  앱 실행 (App Open)",
        "  2️⃣  회원가입 완료 (Sign Up)",
        "  3️⃣  계좌 연결 (Account Link)",
        "  4️⃣  첫 결제 시도 (First Payment Attempt)",
        "  5️⃣  첫 결제 성공 (First Payment Success)",
        "",
        "**실습 질문:** 가장 큰 이탈이 일어나는 단계는? 왜 그럴까?",
        "",
        "**Amplitude 설정:** Funnel Chart → 5단계 이벤트 순서 정의",
    ], accent_color=C_ACCENT_BLUE)
    print("  [31/45] 완료")

    # Slide 32: Amplitude 실습 2
    make_content_slide(prs, 32, "Amplitude 실습 2 — 리텐션 분석", [
        "**분석 목표:** 가입 후 사용자가 얼마나 돌아오는가",
        "",
        "**리텐션 정의:** 가입 후 N일 내 재방문 사용자 비율",
        "",
        "**Amplitude 설정:**",
        "  Retention Chart",
        "  Starting Event: Sign Up",
        "  Return Event: App Open",
        "",
        "**코호트 비교:**",
        "  iOS vs. Android",
        "  첫 결제 금액대별 리텐션 차이",
        "",
        "**인사이트 질문:** 어떤 세그먼트의 리텐션이 가장 높은가? 왜?",
    ], accent_color=C_ACCENT_TEAL)
    print("  [32/45] 완료")

    # Slide 33: Amplitude 실습 3
    make_content_slide(prs, 33, "Amplitude 실습 3 — 코호트 분석", [
        "**분석 목표:** 첫 거래 패턴으로 미래 LTV 예측하기",
        "",
        "**코호트 정의:**",
        "  Group A: 첫 결제 1만원 미만",
        "  Group B: 첫 결제 1~5만원",
        "  Group C: 첫 결제 5만원 이상",
        "",
        "**비교 지표:**",
        "  재결제율, 평균 결제 횟수, 30일 GMV",
        "",
        '**가설:** "첫 결제 금액이 높을수록 LTV가 높다"',
        "",
        "**Amplitude 설정:** Behavioral Cohorts → Funnel/Retention 결합",
    ], accent_color=C_ACCENT_BLUE)
    print("  [33/45] 완료")

    # Slide 34: Looker Studio
    make_content_slide(prs, 34, "Looker Studio 실습 — 핀테크 대시보드 구성", [
        "**데이터 소스:** Google Sheets (transactions.csv 업로드)",
        "",
        "**대시보드 구성 요소:**",
        "  1️⃣  GMV 일별 추이 (Line Chart)",
        "  2️⃣  결제 성공률 / 실패율 (Donut Chart)",
        "  3️⃣  결제 수단별 비중 (Bar Chart)",
        "  4️⃣  가맹점 카테고리별 GMV (Bar Chart)",
        "  5️⃣  기간 필터 + 결제수단 필터",
        "",
        '**목표:** "5분 안에 핵심 상황 파악 가능한 대시보드"',
    ], accent_color=C_ACCENT_TEAL)
    print("  [34/45] 완료")

    # Slide 35: 실습 결과 공유
    make_content_slide(prs, 35, "실습 결과 공유 — 인사이트 발표", [
        "**팀별 발표 (1팀 2분):**",
        "  • 발견한 핵심 이탈 구간",
        "  • 리텐션이 낮은 세그먼트",
        "  • 대시보드에서 발견한 이상 패턴",
        "",
        '**토론:** "만약 이 데이터가 실제 서비스라면 다음 액션은?"',
        "",
        "**강사 피드백:**",
        "  분석 접근법 및 인사이트 품질 평가",
        "",
        "  → 다양한 팀의 관점 비교 자체가 학습",
        "  → 정답이 아닌 사고 과정에 집중",
    ], accent_color=C_ACCENT_ORG)
    print("  [35/45] 완료")

    # ─── PART 5 (슬라이드 36~40) ─────────────────
    make_section_divider(prs, "PART 5", "마무리 & Q&A",
                         "오늘 배운 프레임워크로 내 서비스를 다시 보자", "20분")
    print("  [Part 5 구분] 완료")

    # Slide 36: 오늘 배운 것 정리
    make_content_slide(prs, 36, "Part 5 — 오늘 배운 것 정리", [
        "**Part 1:** 핀테크 데이터 분석의 필요성과 오늘의 목표",
        "",
        "**Part 2:** 비즈니스 모델 5대 유형 + BMC + 10대 KPI + 경쟁 우위 분석",
        "",
        "**Part 3:** 토스(슈퍼앱) · 카카오페이(플랫폼 전환) · Stripe(API 경제) 케이스",
        "",
        "**Part 4:** Amplitude 퍼널/리텐션/코호트 + Looker Studio 대시보드 실습",
        "",
        "**핵심 메시지:**",
        '  "비즈니스 구조를 이해하고, KPI를 정의하고, 데이터로 검증하라"',
    ], accent_color=C_ACCENT_TEAL)
    print("  [36/45] 완료")

    # Slide 37: 실무 적용 팁
    make_content_slide(prs, 37, "실무 적용 팁 — 내일 당장 할 수 있는 것", [
        "**오늘 (Tonight):**",
        "  내 서비스 BMC 빈칸 채워보기",
        "",
        "**이번 주 (This Week):**",
        "  • Amplitude / Looker Studio 무료 계정 세팅",
        "  • 서비스의 KPI 5개 정의하고 현재 수치 파악",
        "",
        "**이번 달 (This Month):**",
        "  • 퍼널 분석 1개 완료 → 이탈 구간 가설 제시",
        "  • 리텐션 분석 → 고LTV 세그먼트 정의",
        "",
        '  → "막연한 다짐"이 아닌 "실행 계획"으로',
    ], accent_color=C_ACCENT_BLUE)
    print("  [37/45] 완료")

    # Slide 38: 추천 학습 리소스
    make_two_col_slide(prs, 38, "추천 학습 리소스",
                        "📚 책 & 뉴스레터",
                        [
                            "책",
                            "  Fintech in Korea",
                            "  Lean Analytics (스타트업 지표 바이블)",
                            "  Hacking Growth (그로스 해킹 실전)",
                            "",
                            "뉴스레터/블로그",
                            "  핀다 블로그 (핀테크 분석)",
                            "  a16z Fintech Newsletter",
                            "  Andreessen Horowitz Fintech",
                        ],
                        "🛠️ 도구 & 데이터셋",
                        [
                            "도구 공식 학습",
                            "  Amplitude Academy (무료 온라인)",
                            "  Looker Studio 공식 튜토리얼",
                            "",
                            "데이터셋",
                            "  Kaggle Fintech Datasets",
                            "  FinFlow 실습 데이터 (강의자료)",
                            "",
                            "→ QR코드로 링크 제공 예정",
                        ])
    print("  [38/45] 완료")

    # Slide 39: Q&A
    make_content_slide(prs, 39, "Q&A", [
        "**Q&A 시간:** 15분",
        "",
        "**질문 방식:** 구두 또는 Slido 익명 질문",
        "",
        "**자주 묻는 질문 (FAQ):**",
        "  Q. Amplitude vs. Mixpanel, 어떤 걸 써야 하나요?",
        "  A. 스타트업은 Amplitude, 엔터프라이즈는 Mixpanel — 무료 플랜 먼저 시작",
        "",
        "  Q. Take Rate는 어떻게 벤치마킹하나요?",
        "  A. 동종 상장사 IR 자료 + CB Insights 리포트 활용",
        "",
        "  Q. 소규모 팀에서도 이런 분석이 가능한가요?",
        "  A. 가능합니다 — Amplitude 무료 플랜 + Looker Studio (무료)로 시작",
    ], accent_color=C_ACCENT_TEAL)
    print("  [39/45] 완료")

    # Slide 40: 마무리
    make_quote_slide(prs, 40, "마무리 — 감사합니다",
                     "오늘 배운 프레임워크로\n여러분 서비스의 데이터를 다시 보세요.",
                     "데이터로 더 나은 결정을 — Good luck!")
    print("  [40/45] 완료")

    # ─── 부록 (슬라이드 41~45) ───────────────────
    make_section_divider(prs, "부록", "심화 자료 & 참고",
                         "규제 환경, 사기 탐지, Unit Economics, 데이터 파이프라인, 커리어")
    print("  [부록 구분] 완료")

    # Slide 41: 규제 환경
    make_content_slide(prs, 41, "[부록] 핀테크 규제 환경과 데이터 분석", [
        "**주요 규제:**",
        "  전자금융거래법, 개인정보보호법, 금융소비자보호법",
        "",
        "**마이데이터:** 금융 데이터 통합 → 분석 기회",
        "",
        "**오픈뱅킹:** API 기반 계좌 접근 → 서비스 확장 기반",
        "",
        "**분석 시 주의:**",
        "  개인정보 익명화 / 동의 기반 데이터 활용",
        "",
        '**인사이트:** "규제는 제약이자 새로운 시장의 출발점"',
    ], accent_color=C_ACCENT_ORG)
    print("  [41/45] 완료")

    # Slide 42: 사기 탐지
    make_content_slide(prs, 42, "[부록] 사기 탐지 (Fraud Detection) 기초", [
        "**핀테크 사기 유형:**",
        "  계정 탈취 / 신용카드 사기 / 보이스피싱 / 자금세탁",
        "",
        "**이상 탐지 지표:**",
        "  비정상적 거래 시간대, 위치 불일치, 금액 패턴",
        "",
        "**Amplitude 활용:**",
        "  비정상 행동 코호트 정의 → 세그먼트 분리",
        "",
        "**머신러닝 적용:**",
        "  이진 분류 모델 (정상/사기) — Stripe Radar 사례",
        "",
        "**실무 팁:** Rule-based → ML-based 단계적 고도화",
    ], accent_color=C_ACCENT_ORG)
    print("  [42/45] 완료")

    # Slide 43: Unit Economics
    make_content_slide(prs, 43, "[부록] Unit Economics 심화 — LTV/CAC 모델링", [
        "**LTV 고급 모델:**",
        "  LTV = (ARPU × Gross Margin) / Churn Rate",
        "",
        "**CAC 분해:**",
        "  (마케팅비 + 영업비) / 신규 유료 고객 수",
        "",
        "**LTV/CAC 벤치마크:**",
        "  < 1: 위험 (돈 잃으며 성장)",
        "  1 ~ 3: 주의 (개선 필요)",
        "  > 3: 건강 (지속 투자 가능)",
        "",
        "**Payback Period:**",
        "  CAC / (ARPU × Gross Margin) → 몇 달이면 회수?",
    ], accent_color=C_ACCENT_TEAL)
    print("  [43/45] 완료")

    # Slide 44: 데이터 파이프라인
    make_content_slide(prs, 44, "[부록] 핀테크 데이터 파이프라인 구조", [
        "**데이터 소스:**",
        "  앱 이벤트 / 결제 시스템 / 고객 DB / 외부 신용 데이터",
        "",
        "**수집:** Amplitude SDK, 서버 이벤트 → Event Tracking",
        "",
        "**저장:** Data Warehouse (BigQuery, Snowflake, Redshift)",
        "",
        "**분석:**",
        "  Amplitude (행동 분석)",
        "  Looker Studio (BI 대시보드)",
        "  SQL (커스텀 쿼리)",
        "",
        "**활용:** 마케팅 세그먼트, 제품 개선, 리스크 관리, 경영 보고",
        "",
        '**인사이트:** "데이터 파이프라인이 없으면 분석도 없다"',
    ])
    print("  [44/45] 완료")

    # Slide 45: 커리어 패스
    make_content_slide(prs, 45, "[부록] 핀테크 분석가 커리어 패스", [
        "**관련 직무:**",
        "  • 그로스 분석가 (Growth Analyst)",
        "  • 프로덕트 분석가 (Product Analyst)",
        "  • 비즈니스 인텔리전스 (BI Analyst)",
        "  • 데이터 사이언티스트 (핀테크 특화)",
        "  • 전략 기획 (핀테크 도메인)",
        "",
        "**필요 역량:** SQL, Python(pandas), BI 도구, 금융 도메인 이해",
        "",
        "**취업 팁:**",
        "  오픈 데이터로 포트폴리오 프로젝트 만들기",
        "  추천 플랫폼: Kaggle, Dacon, 공공데이터포털",
    ], accent_color=C_ACCENT_TEAL)
    print("  [45/45] 완료")

    print(f"\n  ✅ 총 {len(prs.slides)}장 슬라이드 생성 완료")
    return prs


if __name__ == "__main__":
    print("🎨 Pepper Potts — 핀테크 강의 슬라이드 생성 시작")
    print("=" * 60)

    prs = new_prs()
    build_all_slides(prs)

    out_path = r"C:\Agent\pepper\output\hulk_fintech_lecture_slides.pptx"
    prs.save(out_path)
    print(f"\n✅ 저장 완료: {out_path}")
    print(f"   파일 크기: {__import__('os').path.getsize(out_path) / 1024:.1f} KB")
    print("=" * 60)
    print("💄 Pepper Potts 슬라이드 제작 완료!")
