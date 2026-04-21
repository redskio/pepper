# -*- coding: utf-8 -*-
"""
위어드섹터 회사소개서 v4
- 실제 브랜드 이미지 반영 (Weirdsector 웹 스크랩 + Labbit 스크랩)
- 인포그래픽 8종 (기존 4 + 신규 4)
- 이모지 완전 제거 / 가격 삭제
- 15슬라이드 풀블리드 야경
"""
import sys, io
if sys.platform == 'win32':
    sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree
from pathlib import Path
from PIL import Image
import os

ASSETS  = Path("C:/Agent/pepper/output/wirdsector_assets")
SCRAPE  = Path("C:/Agent/pepper/output/brand_scrape")
OUT     = Path("C:/Agent/pepper/output/wirdsector_company_profile_v4.pptx")

C_NAVY    = RGBColor(0x0D, 0x1B, 0x2A)
C_NAVY2   = RGBColor(0x06, 0x0E, 0x18)
C_TEAL    = RGBColor(0x00, 0xD4, 0xAA)
C_WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
C_GRAY    = RGBColor(0xB0, 0xBC, 0xCC)
C_GOLD    = RGBColor(0xFF, 0xC8, 0x4B)
C_PURPLE  = RGBColor(0x9B, 0x7A, 0xFF)
C_BLUE    = RGBColor(0x38, 0xBD, 0xF8)
C_RED     = RGBColor(0xFF, 0x6B, 0x6B)
C_GREEN   = RGBColor(0x34, 0xD3, 0x99)

SW = Inches(13.33)
SH = Inches(7.5)
TOTAL = 15

def new_prs():
    prs = Presentation()
    prs.slide_width  = SW
    prs.slide_height = SH
    return prs

def blank(prs):  return prs.slide_layouts[6]

def bg_img(slide, img_path):
    p = Path(img_path)
    if not p.exists(): return
    try: slide.shapes.add_picture(str(p), Inches(0), Inches(0), SW, SH)
    except Exception as e: print(f"  [WARN] bg: {e}")

def overlay(slide, x, y, w, h, r=0, g=0, b=0, alpha=65):
    sp = slide.shapes.add_shape(1, x, y, w, h)
    sp.line.fill.background()
    sp.fill.solid()
    sp.fill.fore_color.rgb = RGBColor(r, g, b)
    spPr = sp.element.spPr if hasattr(sp.element, 'spPr') else sp.element.find(qn('p:spPr'))
    sf = spPr.find('.//' + qn('a:solidFill'))
    if sf is not None:
        clr = sf.find(qn('a:srgbClr'))
        if clr is not None:
            a = etree.SubElement(clr, qn('a:alpha'))
            a.set('val', str(int(alpha * 1000)))
    return sp

def rrect(slide, x, y, w, h, fill=None, alpha=100):
    sp = slide.shapes.add_shape(5, x, y, w, h)
    sp.line.fill.background()
    if fill:
        sp.fill.solid(); sp.fill.fore_color.rgb = fill
    else:
        sp.fill.background()
    if alpha < 100 and fill:
        spPr = sp.element.spPr if hasattr(sp.element, 'spPr') else sp.element.find(qn('p:spPr'))
        sf = spPr.find('.//' + qn('a:solidFill'))
        if sf is not None:
            clr = sf.find(qn('a:srgbClr'))
            if clr is not None:
                a = etree.SubElement(clr, qn('a:alpha'))
                a.set('val', str(int(alpha * 1000)))
    return sp

def rect(slide, x, y, w, h, fill=None):
    sp = slide.shapes.add_shape(1, x, y, w, h)
    sp.line.fill.background()
    if fill:
        sp.fill.solid(); sp.fill.fore_color.rgb = fill
    else:
        sp.fill.background()
    return sp

def txt(slide, text, x, y, w, h, sz=20, bold=False, color=None,
        align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    p  = tf.paragraphs[0]; p.alignment = align
    r  = p.add_run()
    r.text = text; r.font.size = Pt(sz); r.font.bold = bold
    r.font.italic = italic; r.font.name = "Malgun Gothic"
    if color: r.font.color.rgb = color
    return tb

def pic(slide, path, x, y, w, h):
    p = Path(path)
    if not p.exists(): return
    try: slide.shapes.add_picture(str(p), x, y, w, h)
    except Exception as e: print(f"  [WARN] pic: {e}")

def infographic_full(slide, img_path, top=Inches(2.2), avail_h=Inches(5.0)):
    """인포그래픽 PNG를 슬라이드 중앙에 최대 크기로 삽입"""
    p = Path(img_path)
    if not p.exists(): return
    try:
        im = Image.open(str(p))
        iw, ih = im.size
        aspect = iw / ih
        avail_w = Inches(12.4)
        if aspect > avail_w / avail_h:
            w = avail_w; h = Emu(int(avail_w / aspect))
        else:
            h = avail_h; w = Emu(int(avail_h * aspect))
        left = Emu(int((SW - w) / 2))
        t    = top + Emu(int((avail_h - h) / 2))
        slide.shapes.add_picture(str(p), left, t, w, h)
    except Exception as e:
        print(f"  [WARN] infographic: {e}")

def slide_num(slide, n):
    overlay(slide, Inches(12.5), Inches(7.1), Inches(0.83), Inches(0.35), 0, 0, 0, 40)
    txt(slide, f"{n:02d} / {TOTAL:02d}", Inches(12.5), Inches(7.12),
        Inches(0.8), Inches(0.3), sz=11, color=C_GRAY, align=PP_ALIGN.CENTER)

def logo_badge(slide):
    rrect(slide, Inches(0.4), Inches(0.28), Inches(2.5), Inches(0.52), fill=C_TEAL, alpha=90)
    txt(slide, "Weirdsector", Inches(0.45), Inches(0.3), Inches(2.4), Inches(0.48),
        sz=15, bold=True, color=C_NAVY2, align=PP_ALIGN.CENTER)

def teal_line(slide, y=Inches(1.0)):
    rect(slide, Inches(0.5), y, Inches(1.5), Inches(0.05), fill=C_TEAL)

def sec_hdr(slide, en_label, ko_title, subtitle=None):
    txt(slide, en_label, Inches(0.5), Inches(0.85), Inches(6.0), Inches(0.4),
        sz=13, color=C_TEAL, bold=True)
    txt(slide, ko_title, Inches(0.5), Inches(1.25), Inches(9.5), Inches(0.85),
        sz=42, bold=True, color=C_WHITE)
    teal_line(slide, Inches(2.2))
    if subtitle:
        txt(slide, subtitle, Inches(0.5), Inches(2.38), Inches(9.0), Inches(0.5),
            sz=19, color=C_TEAL)

def num_badge(slide, num_str, x, y, c=None):
    c = c or C_TEAL
    rrect(slide, x, y, Inches(0.58), Inches(0.58), fill=c, alpha=90)
    txt(slide, num_str, x, y+Inches(0.03), Inches(0.58), Inches(0.52),
        sz=15, bold=True, color=C_NAVY2, align=PP_ALIGN.CENTER)

def infographic_page(prs, img_path, en_label, ko_title, bg_file, n,
                     bg_top=Inches(2.18), bg_avail=Inches(5.05)):
    """인포그래픽 전용 슬라이드"""
    sl = prs.slides.add_slide(blank(prs))
    bg_img(sl, ASSETS / bg_file)
    overlay(sl, 0, 0, SW, SH, 4, 10, 20, 82)
    logo_badge(sl)
    slide_num(sl, n)
    txt(sl, en_label, Inches(0.5), Inches(0.85), Inches(6.0), Inches(0.4),
        sz=13, color=C_TEAL, bold=True)
    txt(sl, ko_title, Inches(0.5), Inches(1.25), Inches(9.5), Inches(0.65),
        sz=36, bold=True, color=C_WHITE)
    teal_line(sl, Inches(2.0))
    infographic_full(sl, img_path, top=bg_top, avail_h=bg_avail)
    return sl


# ══════════════════════════════════════════════════
def build():
    prs = new_prs()

    # ── Slide 1: 커버 ─────────────────────────────
    sl = prs.slides.add_slide(blank(prs))
    bg_img(sl, ASSETS/"city_city_aerial.jpg")
    overlay(sl, 0, 0, SW, SH, 6, 14, 26, 70)
    overlay(sl, 0, Inches(4.2), SW, Inches(3.3), 0, 0, 0, 55)

    rect(sl, Inches(0.5), Inches(1.4), Inches(0.06), Inches(4.2), fill=C_TEAL)

    # Weirdsector 로고 이미지 삽입
    ws_logo = SCRAPE / "ws_logo.jpg"
    if ws_logo.exists():
        try:
            overlay(sl, Inches(0.7), Inches(1.38), Inches(3.0), Inches(0.55), 0, 0, 0, 70)
            sl.shapes.add_picture(str(ws_logo), Inches(0.75), Inches(1.4),
                                  Inches(2.5), Inches(0.42))
        except:
            pass

    txt(sl, "Weirdsector", Inches(0.8), Inches(1.5), Inches(9.5), Inches(1.3),
        sz=60, bold=True, color=C_WHITE)
    txt(sl, "위어드섹터", Inches(0.8), Inches(2.82), Inches(9.0), Inches(0.75),
        sz=34, color=C_TEAL)
    txt(sl, "데이터로 성장을 만드는 마케팅 테크 파트너",
        Inches(0.8), Inches(3.65), Inches(9.5), Inches(0.7),
        sz=24, color=C_GRAY)
    txt(sl, '"Be yourself, Embrace your weirdness!"',
        Inches(0.8), Inches(4.55), Inches(9.0), Inches(0.5),
        sz=16, italic=True, color=C_TEAL)
    rect(sl, Inches(0.8), Inches(4.45), Inches(5.0), Inches(0.03), fill=C_TEAL)

    for i, tag in enumerate(["Data & Code", "MarTech", "CRM", "Labbit", "DataNugget"]):
        rrect(sl, Inches(0.8 + i*2.35), Inches(6.38), Inches(2.2), Inches(0.45),
              fill=C_NAVY2, alpha=80)
        txt(sl, tag, Inches(0.8 + i*2.35), Inches(6.4), Inches(2.2), Inches(0.42),
            sz=13, color=C_TEAL, align=PP_ALIGN.CENTER)

    txt(sl, "2026", Inches(12.0), Inches(7.0), Inches(1.2), Inches(0.38),
        sz=13, color=C_GRAY, align=PP_ALIGN.RIGHT)
    slide_num(sl, 1)
    print("  [1] 커버")

    # ── Slide 2: 회사 소개 + 실제 웹 이미지 ──────────
    sl = prs.slides.add_slide(blank(prs))
    bg_img(sl, ASSETS/"city_city_night.jpg")
    overlay(sl, 0, 0, SW, SH, 6, 14, 26, 78)
    logo_badge(sl); slide_num(sl, 2)

    sec_hdr(sl, "Company Overview", "위어드섹터 소개",
            "데이터 기반 마케팅의 실행 파트너")

    # 실제 Weirdsector 웹 이미지 — 우측 패널
    ws_img = SCRAPE / "ws_index1.jpg"
    if ws_img.exists():
        # 우측에 실제 이미지 배치 (어두운 오버레이 포함)
        pic(sl, ws_img, Inches(9.0), Inches(1.0), Inches(4.0), Inches(6.2))
        overlay(sl, Inches(9.0), Inches(1.0), Inches(4.0), Inches(6.2), 6, 14, 26, 50)

    txt(sl, "개발력, 마케팅 감각, 데이터 분석 — 세 가지를 한 팀에서 제공합니다.\n단순한 대행사가 아닌, 함께 성장하는 테크 파트너를 경험하세요.",
        Inches(0.5), Inches(2.95), Inches(8.2), Inches(1.1),
        sz=17, color=C_GRAY)

    facts = [
        ("01", "서울 & 창원", "노원 본사 · 도봉 연구소\n창원 지사 3거점 운영"),
        ("02", "하이브리드 팀", "개발 + 마케팅 + 데이터\n원스톱 제공"),
        ("03", "자체 서비스", "Labbit · DataNugget\n자체 SaaS 보유"),
    ]
    for i, (num, ttl, body) in enumerate(facts):
        by = Inches(4.3 + i * 0.95)
        rrect(sl, Inches(0.5), by, Inches(8.2), Inches(0.85), fill=C_NAVY2, alpha=80)
        rect(sl, Inches(0.5), by, Inches(0.05), Inches(0.85), fill=C_TEAL)
        num_badge(sl, num, Inches(0.6), by+Inches(0.12), C_TEAL)
        txt(sl, ttl, Inches(1.3), by+Inches(0.08), Inches(2.8), Inches(0.38),
            sz=16, bold=True, color=C_WHITE)
        txt(sl, body, Inches(4.3), by+Inches(0.08), Inches(4.2), Inches(0.7),
            sz=13, color=C_GRAY)
    print("  [2] 회사 소개")

    # ── Slide 3: 문제 정의 + Before/After 인포그래픽 ─
    infographic_page(prs, ASSETS/"infographic_before_after.png",
                     "Pain Points", "우리가 해결하는 문제",
                     "city_data_viz.jpg", 3,
                     bg_top=Inches(2.1), bg_avail=Inches(5.1))
    print("  [3] 문제 정의 (Before/After)")

    # ── Slide 4: 서비스 개요 ───────────────────────
    sl = prs.slides.add_slide(blank(prs))
    bg_img(sl, ASSETS/"bg_technology.jpg")
    overlay(sl, 0, 0, SW, SH, 4, 10, 20, 82)
    logo_badge(sl); slide_num(sl, 4)

    sec_hdr(sl, "Our Solutions", "위어드섹터의 5가지 솔루션")

    solutions = [
        ("01", "데이터 &\n코드 관리", "GA4 이벤트 설계\n로그정의서 관리\nGTM 최적화", C_TEAL),
        ("02", "마테크\n(MarTech)", "CleverTap 운영\n자동화 캠페인\nA/B 테스트", C_GOLD),
        ("03", "CRM\n운영", "세그멘테이션\n리텐션 캠페인\nLTV 분석", C_TEAL),
        ("04", "Labbit", "그로스해킹\nA/B 테스트\n퍼포먼스", C_PURPLE),
        ("05", "DataNugget", "데이터 인사이트\n브랜드 리포트\n대시보드", C_BLUE),
    ]
    for i, (num, ttl, body, c) in enumerate(solutions):
        bx = Inches(0.35 + i*2.6)
        rrect(sl, bx, Inches(2.5), Inches(2.45), Inches(4.65), fill=C_NAVY2, alpha=80)
        rect(sl, bx, Inches(2.5), Inches(2.45), Inches(0.06), fill=c)
        num_badge(sl, num, bx+Inches(0.88), Inches(2.6), c)
        txt(sl, ttl, bx+Inches(0.1), Inches(3.35), Inches(2.25), Inches(0.75),
            sz=15, bold=True, color=WHITE_COL(c), align=PP_ALIGN.CENTER)
        txt(sl, body, bx+Inches(0.1), Inches(4.15), Inches(2.25), Inches(2.0),
            sz=12, color=C_GRAY, align=PP_ALIGN.CENTER)
    print("  [4] 서비스 개요")

    # ── Slide 5: Labbit — 실제 서비스 이미지 ────────
    sl = prs.slides.add_slide(blank(prs))
    bg_img(sl, ASSETS/"city_neon_tech.jpg")
    overlay(sl, 0, 0, SW, SH, 4, 8, 18, 78)
    logo_badge(sl); slide_num(sl, 5)

    txt(sl, "Own Service 01", Inches(0.5), Inches(0.85), Inches(6.0), Inches(0.38),
        sz=13, color=C_TEAL, bold=True)
    txt(sl, "Labbit", Inches(0.5), Inches(1.25), Inches(5.5), Inches(0.9),
        sz=50, bold=True, color=C_WHITE)
    txt(sl, "래빗  —  그로스해킹 파트너",
        Inches(0.5), Inches(2.2), Inches(7.0), Inches(0.5),
        sz=22, color=C_TEAL)
    teal_line(sl, Inches(2.82))

    # 실제 Labbit 소개 이미지
    labbit_img = SCRAPE / "labbit_intro.jpg"
    if labbit_img.exists():
        pic(sl, labbit_img, Inches(7.8), Inches(1.0), Inches(5.1), Inches(6.2))
        overlay(sl, Inches(7.8), Inches(1.0), Inches(5.1), Inches(6.2), 4, 8, 18, 50)

    txt(sl, "고객 행동 데이터를 기반으로 전환율과\n리텐션을 개선하는 그로스해킹 서비스",
        Inches(0.5), Inches(3.0), Inches(7.0), Inches(0.9),
        sz=16, color=C_GRAY)

    labbit_svcs = [
        ("01", "구매 단계 트래킹",    "전환 퍼널 분석"),
        ("02", "광고 성과 추적",       "채널별 ROI 측정"),
        ("03", "커스텀 타겟 발굴",     "행동 기반 세그먼트"),
        ("04", "사이트 UX 개선",       "A/B 테스트"),
        ("05", "데이터 시각화 리포트", "대시보드 설계"),
        ("06", "커스텀 솔루션",        "전담 그로스해커"),
    ]
    for i, (num, ttl, sub) in enumerate(labbit_svcs):
        row, col = divmod(i, 2)
        bx = Inches(0.5 + col*3.55)
        ry = Inches(4.05 + row*1.12)
        rrect(sl, bx, ry, Inches(3.35), Inches(1.0), fill=C_NAVY2, alpha=80)
        rect(sl, bx, ry, Inches(0.05), Inches(1.0), fill=C_TEAL)
        num_badge(sl, num, bx+Inches(0.1), ry+Inches(0.1), C_TEAL)
        txt(sl, ttl, bx+Inches(0.82), ry+Inches(0.1), Inches(2.3), Inches(0.38),
            sz=14, bold=True, color=C_WHITE)
        txt(sl, sub, bx+Inches(0.82), ry+Inches(0.52), Inches(2.3), Inches(0.35),
            sz=12, color=C_GRAY)

    rrect(sl, Inches(0.5), Inches(7.1), Inches(2.8), Inches(0.33), fill=C_TEAL, alpha=88)
    txt(sl, "labbit.kr", Inches(0.5), Inches(7.12), Inches(2.8), Inches(0.3),
        sz=15, bold=True, color=C_NAVY2, align=PP_ALIGN.CENTER)
    print("  [5] Labbit")

    # ── Slide 6: DataNugget — 실제 서비스 이미지 ─────
    sl = prs.slides.add_slide(blank(prs))
    bg_img(sl, ASSETS/"city_data_viz.jpg")
    overlay(sl, 0, 0, SW, SH, 2, 8, 18, 78)
    logo_badge(sl); slide_num(sl, 6)

    txt(sl, "Own Service 02", Inches(0.5), Inches(0.85), Inches(6.0), Inches(0.38),
        sz=13, color=C_GOLD, bold=True)
    txt(sl, "DataNugget", Inches(0.5), Inches(1.25), Inches(7.5), Inches(0.9),
        sz=50, bold=True, color=C_WHITE)
    txt(sl, "데이터너겟  —  브랜드 성장박스",
        Inches(0.5), Inches(2.2), Inches(7.5), Inches(0.5),
        sz=22, color=C_GOLD)
    rect(sl, Inches(0.5), Inches(2.82), Inches(1.5), Inches(0.05), fill=C_GOLD)

    # 실제 DataNugget 스크린샷
    dn_img = SCRAPE / "datanugget_main.png"
    if dn_img.exists():
        pic(sl, dn_img, Inches(7.8), Inches(1.0), Inches(5.1), Inches(6.2))
        overlay(sl, Inches(7.8), Inches(1.0), Inches(5.1), Inches(6.2), 2, 8, 18, 45)

    txt(sl, "광고 성과 조회부터 홈페이지 리뉴얼까지\n데이터에 대한 모든 질문을 해결합니다.",
        Inches(0.5), Inches(3.05), Inches(7.0), Inches(0.9),
        sz=16, color=C_GRAY)

    nugget_features = [
        ("01", "광고 성과 통합 분석", "모든 매체를 하나의 대시보드로"),
        ("02", "홈페이지 데이터 진단", "방문자 행동 패턴 및 UX 개선"),
        ("03", "데이터 시각화 자동화", "커스텀 리포트 자동 생성"),
        ("04", "인사이트 리포팅",      "실행 가능한 인사이트 추출"),
    ]
    for i, (num, ttl, sub) in enumerate(nugget_features):
        by = Inches(4.05 + i*0.82)
        rrect(sl, Inches(0.5), by, Inches(7.0), Inches(0.72), fill=C_NAVY2, alpha=80)
        rect(sl, Inches(0.5), by, Inches(0.05), Inches(0.72), fill=C_GOLD)
        num_badge(sl, num, Inches(0.6), by+Inches(0.06), C_GOLD)
        txt(sl, ttl, Inches(1.32), by+Inches(0.06), Inches(3.0), Inches(0.35),
            sz=15, bold=True, color=C_WHITE)
        txt(sl, sub, Inches(4.5), by+Inches(0.06), Inches(2.8), Inches(0.35),
            sz=13, color=C_GRAY, align=PP_ALIGN.RIGHT)

    rrect(sl, Inches(0.5), Inches(7.1), Inches(2.8), Inches(0.33), fill=C_GOLD, alpha=88)
    txt(sl, "datanugget.io", Inches(0.5), Inches(7.12), Inches(2.8), Inches(0.3),
        sz=15, bold=True, color=C_NAVY2, align=PP_ALIGN.CENTER)
    print("  [6] DataNugget")

    # ── Slide 7: 서비스 프로세스 인포그래픽 ─────────
    infographic_page(prs, ASSETS/"infographic_process.png",
                     "Service Process", "위어드섹터 서비스 프로세스",
                     "bg_technology.jpg", 7)
    print("  [7] 서비스 프로세스")

    # ── Slide 8: 핵심 수치 인포그래픽 ───────────────
    infographic_page(prs, ASSETS/"infographic_stats.png",
                     "Key Metrics", "주요 성과 지표",
                     "bg_growth.jpg", 8)
    print("  [8] 핵심 수치")

    # ── Slide 9: 기술 역량 매트릭스 인포그래픽 ──────
    infographic_page(prs, ASSETS/"infographic_team.png",
                     "Team Capabilities", "팀 기술 역량 매트릭스",
                     "bg_technology.jpg", 9)
    print("  [9] 팀 역량 매트릭스")

    # ── Slide 10: 레퍼런스 — 실제 이미지 활용 ────────
    sl = prs.slides.add_slide(blank(prs))
    bg_img(sl, ASSETS/"bg_growth.jpg")
    overlay(sl, 0, 0, SW, SH, 4, 10, 20, 82)
    logo_badge(sl); slide_num(sl, 10)

    sec_hdr(sl, "Case Studies", "실적 & 케이스 스터디")

    # 실제 Weirdsector 웹 이미지를 배경 패널로 활용
    ws2 = SCRAPE / "ws_index2.jpg"
    ws3 = SCRAPE / "ws_index3.jpg"
    if ws2.exists():
        pic(sl, ws2, Inches(9.5), Inches(2.3), Inches(3.4), Inches(4.8))
        overlay(sl, Inches(9.5), Inches(2.3), Inches(3.4), Inches(4.8), 4, 10, 20, 60)
    if ws3.exists():
        pic(sl, ws3, Inches(10.3), Inches(2.3), Inches(2.6), Inches(4.8))
        overlay(sl, Inches(10.3), Inches(2.3), Inches(2.6), Inches(4.8), 4, 10, 20, 65)

    cases = [
        {
            "industry": "OTT · 미디어 플랫폼",
            "challenge": "신규 가입자 이탈률 60%, CRM 자동화 부재",
            "solution": "온보딩 재설계 + CleverTap 세그먼트 + 이탈방지 자동화",
            "results": [("-38%", "이탈률"), ("+52%", "D30 리텐션"), ("2.4x", "ROI")],
            "color": C_TEAL,
        },
        {
            "industry": "이커머스 · 쇼핑몰",
            "challenge": "GA4 이전 후 데이터 신뢰도 0%, 광고 최적화 불가",
            "solution": "GA4 전면 재설계 + GTM 정리 + DataNugget 대시보드",
            "results": [("100%", "데이터 정확도"), ("-45%", "광고 낭비"), ("+28%", "전환율")],
            "color": C_GOLD,
        },
    ]
    for i, c in enumerate(cases):
        bx = Inches(0.4 + i*4.55)
        rrect(sl, bx, Inches(2.5), Inches(4.4), Inches(4.65), fill=C_NAVY2, alpha=82)
        rect(sl, bx, Inches(2.5), Inches(4.4), Inches(0.07), fill=c["color"])
        txt(sl, c["industry"], bx+Inches(0.2), Inches(2.62), Inches(4.0), Inches(0.45),
            sz=17, bold=True, color=c["color"])
        txt(sl, "과제: " + c["challenge"],
            bx+Inches(0.2), Inches(3.14), Inches(4.0), Inches(0.65),
            sz=13, color=C_GRAY)
        txt(sl, "솔루션: " + c["solution"],
            bx+Inches(0.2), Inches(3.85), Inches(4.0), Inches(0.75),
            sz=13, color=C_WHITE)
        rect(sl, bx+Inches(0.2), Inches(4.72), Inches(3.85), Inches(0.04), fill=c["color"])
        for j, (v, l) in enumerate(c["results"]):
            kx = bx + Inches(0.2 + j*1.4)
            txt(sl, v, kx, Inches(4.88), Inches(1.3), Inches(0.62),
                sz=30, bold=True, color=c["color"])
            txt(sl, l, kx, Inches(5.52), Inches(1.3), Inches(0.35),
                sz=12, color=C_GRAY)
    print("  [10] 케이스 스터디")

    # ── Slide 11: 파트너 인포그래픽 ─────────────────
    infographic_page(prs, ASSETS/"infographic_references.png",
                     "Our References", "주요 파트너 & 레퍼런스",
                     "city_skyline_blue.jpg", 11)
    print("  [11] 파트너 그리드")

    # ── Slide 12: 팀 소개 인포그래픽 ────────────────
    infographic_page(prs, ASSETS/"infographic_capabilities.png",
                     "Core Capabilities", "핵심 역량 그리드",
                     "city_office_night.jpg", 12)
    print("  [12] 핵심 역량 그리드")

    # ── Slide 13: 비즈니스 모델 다이어그램 ──────────
    infographic_page(prs, ASSETS/"infographic_bizmodel.png",
                     "Business Model", "비즈니스 모델",
                     "bg_technology.jpg", 13,
                     bg_top=Inches(2.15), bg_avail=Inches(5.0))
    print("  [13] 비즈니스 모델")

    # ── Slide 14: 비전 & 로드맵 타임라인 ─────────────
    infographic_page(prs, ASSETS/"infographic_roadmap.png",
                     "Vision & Roadmap", "비전 & 성장 로드맵",
                     "city_city_night.jpg", 14,
                     bg_top=Inches(2.15), bg_avail=Inches(5.0))
    print("  [14] 비전 & 로드맵")

    # ── Slide 15: CTA ─────────────────────────────
    sl = prs.slides.add_slide(blank(prs))
    bg_img(sl, ASSETS/"city_city_dawn.jpg")
    overlay(sl, 0, 0, SW, SH, 4, 10, 22, 72)
    overlay(sl, 0, Inches(3.3), SW, Inches(4.2), 2, 6, 14, 65)
    slide_num(sl, 15)

    txt(sl, "지금 시작하세요",
        Inches(1.5), Inches(1.1), Inches(10.5), Inches(1.2),
        sz=52, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    txt(sl, "위어드섹터와 함께 데이터로 성장하는 마케팅을 경험하세요",
        Inches(1.5), Inches(2.4), Inches(10.5), Inches(0.6),
        sz=21, color=C_TEAL, align=PP_ALIGN.CENTER)

    rrect(sl, Inches(4.7), Inches(3.25), Inches(4.0), Inches(0.78), fill=C_TEAL, alpha=95)
    txt(sl, "미팅 신청하기  ->",
        Inches(4.7), Inches(3.29), Inches(4.0), Inches(0.7),
        sz=19, bold=True, color=C_NAVY2, align=PP_ALIGN.CENTER)

    contacts = [
        ("W", "웹사이트",   "weirdsector.co.kr"),
        ("E", "이메일",     "info@weirdsector.co.kr"),
        ("T", "전화",       "02-458-0601"),
    ]
    for i, (badge, lb, vl) in enumerate(contacts):
        bx = Inches(1.0 + i*3.8)
        rrect(sl, bx, Inches(4.5), Inches(3.5), Inches(2.25), fill=C_NAVY2, alpha=80)
        rect(sl, bx, Inches(4.5), Inches(3.5), Inches(0.05), fill=C_TEAL)
        rrect(sl, bx+Inches(1.35), Inches(4.62), Inches(0.78), Inches(0.6),
              fill=C_TEAL, alpha=90)
        txt(sl, badge, bx+Inches(1.35), Inches(4.64), Inches(0.78), Inches(0.56),
            sz=20, bold=True, color=C_NAVY2, align=PP_ALIGN.CENTER)
        txt(sl, lb,   bx+Inches(0.2), Inches(5.28), Inches(3.1), Inches(0.36),
            sz=12, color=C_GRAY, align=PP_ALIGN.CENTER)
        txt(sl, vl,   bx+Inches(0.2), Inches(5.65), Inches(3.1), Inches(0.4),
            sz=14, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    txt(sl, '"Be yourself, Embrace your weirdness!"',
        Inches(1.5), Inches(6.95), Inches(10.5), Inches(0.36),
        sz=14, italic=True, color=C_TEAL, align=PP_ALIGN.CENTER)
    print("  [15] CTA")

    # ── 저장
    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    sz = OUT.stat().st_size // 1024
    print(f"\nSaved: {OUT}")
    print(f"Slides: {len(prs.slides)} / {sz:,}KB")
    return len(prs.slides)


def WHITE_COL(c):
    return C_WHITE


if __name__ == "__main__":
    print("위어드섹터 회사소개서 v4 빌드 시작\n")
    n = build()
    print(f"\nDone: {n} slides")
