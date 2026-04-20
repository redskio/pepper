# -*- coding: utf-8 -*-
"""
CleverTap 한국어 제안서 생성 스크립트
python-pptx 기반 / CleverTap 브랜드 컬러 적용
"""
import sys, io
if sys.platform == 'win32':
    sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')
    sys.stderr = io.TextIOWrapper(sys.stderr.buffer, encoding='utf-8', errors='replace')

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml
import copy
from lxml import etree
from pathlib import Path
from PIL import Image
import os

# ── 경로 설정 ──────────────────────────────────────────────
ASSETS = Path("C:/Agent/pepper/output/clevertap_assets")
OUTPUT = Path("C:/Agent/pepper/output/clevertap_proposal_kr.pptx")

# ── 브랜드 컬러 ────────────────────────────────────────────
C_PURPLE      = RGBColor(0x6B, 0x4E, 0xFF)   # CleverTap 메인 퍼플
C_PURPLE_DARK = RGBColor(0x3D, 0x1F, 0xCC)   # 다크 퍼플
C_PURPLE_LIGHT= RGBColor(0x9B, 0x8A, 0xFF)   # 라이트 퍼플
C_DARK        = RGBColor(0x1A, 0x1A, 0x2E)   # 거의 검정
C_DARK2       = RGBColor(0x16, 0x21, 0x3E)   # 네이비 다크
C_MID         = RGBColor(0x2D, 0x2D, 0x4E)   # 미드 다크
C_WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
C_GRAY        = RGBColor(0x6C, 0x75, 0x8D)
C_LIGHT_GRAY  = RGBColor(0xF4, 0xF4, 0xF8)
C_ACCENT      = RGBColor(0xFF, 0x6B, 0x35)   # 오렌지 액센트
C_GREEN       = RGBColor(0x00, 0xC9, 0x85)   # 성공/성장 그린

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

def new_prs():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs

def blank_layout(prs):
    return prs.slide_layouts[6]  # blank

def add_rect(slide, x, y, w, h, fill_color=None, line_color=None, line_width=None, alpha=None):
    shape = slide.shapes.add_shape(1, x, y, w, h)
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        if line_width:
            shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, text, x, y, w, h,
             font_size=24, bold=False, color=None,
             align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    if color:
        run.font.color.rgb = color
    run.font.name = "Malgun Gothic"
    return txBox

def add_image(slide, img_path, x, y, w, h):
    try:
        img_path = str(img_path)
        if not os.path.exists(img_path):
            return None
        # Check it's a valid image
        img = Image.open(img_path)
        if img.mode != 'RGB':
            img = img.convert('RGB')
            tmp_path = img_path + "_tmp.jpg"
            img.save(tmp_path, 'JPEG', quality=85)
            img_path = tmp_path
        pic = slide.shapes.add_picture(img_path, x, y, w, h)
        return pic
    except Exception as e:
        print(f"  [이미지 오류] {img_path}: {e}")
        return None

def add_gradient_bg(slide, color1, color2):
    """좌→우 그라데이션 배경 (두 개 rect 겹침으로 시뮬레이션)"""
    # 왼쪽 진한 색
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill_color=color1)
    # 오른쪽 밝은 색 (반투명 시뮬레이션 — 실제 그라데이션 대신 오버레이)
    overlay = slide.shapes.add_shape(1, SLIDE_W//2, 0, SLIDE_W//2, SLIDE_H)
    overlay.fill.solid()
    overlay.fill.fore_color.rgb = color2
    overlay.line.fill.background()

def build_slides():
    prs = new_prs()

    # ══════════════════════════════════════════════════════
    # SLIDE 1: 표지
    # ══════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout(prs))

    # 배경 - 다크 그라데이션
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill_color=C_DARK)

    # 좌측 퍼플 사이드바
    add_rect(slide, 0, 0, Inches(0.3), SLIDE_H, fill_color=C_PURPLE)

    # 우측에 메인 페이지 스크린샷 (반투명 효과)
    ss = ASSETS / "screenshot_main.png"
    if ss.exists():
        pic = add_image(slide, ss, Inches(6.8), Inches(0.3), Inches(6.2), Inches(6.9))
        if pic:
            # 투명도 조절 — 슬라이드 오버레이
            add_rect(slide, Inches(6.8), 0, Inches(6.53), SLIDE_H, fill_color=C_DARK)
            # re-add image on top with overlay
            # Instead just add a semi-transparent overlay
            overlay2 = add_rect(slide, Inches(6.5), 0, Inches(6.83), SLIDE_H, fill_color=C_DARK2)
            # Add actual screenshot behind
    # 실제 메인 스크린샷 배치 (오른쪽)
    if ss.exists():
        add_image(slide, ss, Inches(7), Inches(0.5), Inches(5.8), Inches(6.5))
        # 오른쪽 다크 오버레이
        add_rect(slide, Inches(7), 0, Inches(6.33), SLIDE_H, fill_color=RGBColor(0x1A, 0x1A, 0x2E))

    # 상단 로고 영역
    add_rect(slide, Inches(0.3), Inches(0.4), Inches(2.2), Inches(0.5), fill_color=C_PURPLE)
    add_text(slide, "CleverTap", Inches(0.4), Inches(0.42), Inches(2.0), Inches(0.45),
             font_size=18, bold=True, color=C_WHITE, align=PP_ALIGN.LEFT)

    # 메인 타이틀
    add_text(slide, "CleverTap", Inches(0.5), Inches(1.8), Inches(6.0), Inches(0.9),
             font_size=52, bold=True, color=C_WHITE, align=PP_ALIGN.LEFT)
    add_text(slide, "솔루션 제안서", Inches(0.5), Inches(2.6), Inches(6.0), Inches(0.9),
             font_size=48, bold=True, color=C_PURPLE_LIGHT, align=PP_ALIGN.LEFT)

    # 구분선
    add_rect(slide, Inches(0.5), Inches(3.7), Inches(3.5), Inches(0.04), fill_color=C_PURPLE)

    # 서브타이틀
    add_text(slide, "모바일 마케팅 자동화 & 실시간 고객 인게이지먼트 플랫폼",
             Inches(0.5), Inches(3.9), Inches(6.0), Inches(0.6),
             font_size=16, color=C_LIGHT_GRAY, align=PP_ALIGN.LEFT)

    # 날짜
    add_text(slide, "2026년 4월", Inches(0.5), Inches(6.5), Inches(3.0), Inches(0.5),
             font_size=14, color=C_GRAY, align=PP_ALIGN.LEFT)

    # 통계 하이라이트 (하단)
    stats = [("10,000+", "글로벌 앱"), ("12조+", "월간 이벤트"), ("25개국+", "서비스 국가")]
    for i, (num, label) in enumerate(stats):
        x = Inches(0.5 + i * 2.0)
        add_rect(slide, x, Inches(5.2), Inches(1.7), Inches(0.9),
                 fill_color=C_MID)
        add_text(slide, num, x + Inches(0.1), Inches(5.3), Inches(1.5), Inches(0.4),
                 font_size=20, bold=True, color=C_PURPLE_LIGHT, align=PP_ALIGN.CENTER)
        add_text(slide, label, x + Inches(0.1), Inches(5.65), Inches(1.5), Inches(0.35),
                 font_size=11, color=C_GRAY, align=PP_ALIGN.CENTER)

    print("  [1/15] 표지 완료")

    # ══════════════════════════════════════════════════════
    # SLIDE 2: 목차
    # ══════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout(prs))
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill_color=C_DARK)
    add_rect(slide, 0, 0, Inches(0.3), SLIDE_H, fill_color=C_PURPLE)
    add_rect(slide, 0, 0, SLIDE_W, Inches(1.5), fill_color=C_DARK2)

    add_text(slide, "목 차", Inches(0.6), Inches(0.35), Inches(6.0), Inches(0.7),
             font_size=36, bold=True, color=C_WHITE)
    add_text(slide, "CONTENTS", Inches(0.6), Inches(1.0), Inches(4.0), Inches(0.4),
             font_size=13, color=C_PURPLE_LIGHT, bold=False)

    toc_items = [
        ("01", "CleverTap 소개"),
        ("02", "주요 글로벌 고객사"),
        ("03", "핵심 가치 제안"),
        ("04", "제품 아키텍처"),
        ("05", "실시간 분석 (Analytics)"),
        ("06", "사용자 세그멘테이션"),
        ("07", "멀티채널 인게이지먼트"),
        ("08", "AI/ML 기반 개인화"),
        ("09", "A/B 테스트 & 최적화"),
        ("10", "핀테크 활용 사례"),
        ("11", "도입 효과 & ROI"),
        ("12", "도입 프로세스 & 타임라인"),
        ("13", "가격 정책 & 다음 단계"),
    ]

    cols = 2
    col_w = Inches(5.8)
    for i, (num, title) in enumerate(toc_items):
        col = i % cols
        row = i // cols
        x = Inches(0.8 + col * 6.3)
        y = Inches(1.7 + row * 0.72)
        # 번호 배지
        add_rect(slide, x, y + Inches(0.05), Inches(0.5), Inches(0.5), fill_color=C_PURPLE)
        add_text(slide, num, x + Inches(0.02), y + Inches(0.08), Inches(0.46), Inches(0.38),
                 font_size=12, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        add_text(slide, title, x + Inches(0.6), y, col_w - Inches(0.7), Inches(0.55),
                 font_size=17, color=C_WHITE, bold=False)

    # 스크린샷 배치 (오른쪽 하단)
    add_image(slide, ASSETS / "viewport_product.png",
              Inches(9.5), Inches(5.5), Inches(3.5), Inches(1.8))

    print("  [2/15] 목차 완료")

    # ══════════════════════════════════════════════════════
    # SLIDE 3: CleverTap 소개
    # ══════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout(prs))
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill_color=C_DARK)
    add_rect(slide, 0, 0, SLIDE_W, Inches(1.6), fill_color=C_PURPLE_DARK)
    add_rect(slide, 0, 0, Inches(0.3), SLIDE_H, fill_color=C_PURPLE)

    add_text(slide, "01  CleverTap 소개", Inches(0.6), Inches(0.35), Inches(10.0), Inches(0.8),
             font_size=32, bold=True, color=C_WHITE)
    add_text(slide, "모바일 마케팅 자동화의 글로벌 리더", Inches(0.6), Inches(1.05), Inches(8.0), Inches(0.45),
             font_size=14, color=C_PURPLE_LIGHT)

    # 회사 소개 텍스트
    desc = (
        "CleverTap은 2013년 설립된 글로벌 모바일 마케팅 자동화 & 분석 플랫폼으로,\n"
        "앱 비즈니스가 사용자를 유지하고 성장시킬 수 있도록 돕는 올인원 솔루션입니다.\n\n"
        "실시간 행동 데이터를 기반으로 개인화된 메시지를 적시에 전달하여\n"
        "고객 참여도와 매출을 극대화합니다."
    )
    add_text(slide, desc, Inches(0.6), Inches(1.8), Inches(5.8), Inches(2.2),
             font_size=15, color=C_LIGHT_GRAY, wrap=True)

    # 핵심 수치 카드
    key_stats = [
        ("2013년", "설립"),
        ("10,000+", "글로벌 고객사"),
        ("50+", "국가 서비스"),
        ("12조+", "월간 이벤트 처리"),
        ("600+", "임직원"),
        ("$105M+", "누적 투자 유치"),
    ]
    for i, (val, label) in enumerate(key_stats):
        col = i % 3
        row = i // 3
        x = Inches(0.6 + col * 2.1)
        y = Inches(4.0 + row * 1.3)
        add_rect(slide, x, y, Inches(1.9), Inches(1.1), fill_color=C_MID)
        add_rect(slide, x, y, Inches(1.9), Inches(0.06), fill_color=C_PURPLE)
        add_text(slide, val, x + Inches(0.1), y + Inches(0.1), Inches(1.7), Inches(0.55),
                 font_size=22, bold=True, color=C_PURPLE_LIGHT, align=PP_ALIGN.CENTER)
        add_text(slide, label, x + Inches(0.1), y + Inches(0.62), Inches(1.7), Inches(0.4),
                 font_size=11, color=C_GRAY, align=PP_ALIGN.CENTER)

    # 오른쪽 스크린샷
    add_image(slide, ASSETS / "screenshot_main.png",
              Inches(6.8), Inches(1.7), Inches(6.2), Inches(4.8))
    # 반투명 오버레이
    add_rect(slide, Inches(6.8), Inches(1.7), Inches(6.2), Inches(4.8),
             fill_color=RGBColor(0x1A, 0x1A, 0x2E))
    add_image(slide, ASSETS / "viewport_main.png",
              Inches(7.0), Inches(2.0), Inches(5.8), Inches(4.0))

    print("  [3/15] CleverTap 소개 완료")

    # ══════════════════════════════════════════════════════
    # SLIDE 4: 주요 고객사
    # ══════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout(prs))
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill_color=C_DARK)
    add_rect(slide, 0, 0, SLIDE_W, Inches(1.6), fill_color=C_PURPLE_DARK)
    add_rect(slide, 0, 0, Inches(0.3), SLIDE_H, fill_color=C_PURPLE)

    add_text(slide, "02  주요 글로벌 고객사", Inches(0.6), Inches(0.35), Inches(10.0), Inches(0.8),
             font_size=32, bold=True, color=C_WHITE)
    add_text(slide, "전 세계 10,000+ 앱이 CleverTap을 신뢰합니다", Inches(0.6), Inches(1.05), Inches(8.0), Inches(0.45),
             font_size=14, color=C_PURPLE_LIGHT)

    # 고객사 카드
    customers = [
        ("Jio", "인도 최대 통신사", "통신"),
        ("Swiggy", "푸드 딜리버리 1위", "이커머스"),
        ("PhonePe", "UPI 결제 플랫폼", "핀테크"),
        ("Axis Bank", "인도 3대 민간은행", "금융"),
        ("Disney+ Hotstar", "OTT 스트리밍", "미디어"),
        ("Vodafone", "글로벌 통신사", "통신"),
        ("Puma", "스포츠 브랜드", "리테일"),
        ("OLX", "중고거래 플랫폼", "마켓플레이스"),
        ("Carousell", "동남아 커머스", "이커머스"),
    ]

    cols = 3
    for i, (name, desc, category) in enumerate(customers):
        col = i % cols
        row = i // cols
        x = Inches(0.6 + col * 4.1)
        y = Inches(1.9 + row * 1.7)
        add_rect(slide, x, y, Inches(3.7), Inches(1.5), fill_color=C_MID)
        add_rect(slide, x, y, Inches(0.12), Inches(1.5), fill_color=C_PURPLE)
        # 카테고리 뱃지
        add_rect(slide, x + Inches(0.25), y + Inches(0.15), Inches(0.85), Inches(0.28),
                 fill_color=C_PURPLE)
        add_text(slide, category, x + Inches(0.27), y + Inches(0.16), Inches(0.81), Inches(0.25),
                 font_size=9, color=C_WHITE, bold=True, align=PP_ALIGN.LEFT)
        add_text(slide, name, x + Inches(0.25), y + Inches(0.5), Inches(3.3), Inches(0.45),
                 font_size=20, bold=True, color=C_WHITE)
        add_text(slide, desc, x + Inches(0.25), y + Inches(0.95), Inches(3.3), Inches(0.4),
                 font_size=12, color=C_GRAY)

    # 고객사 스크린샷
    add_image(slide, ASSETS / "screenshot_customers.png",
              Inches(0.6), Inches(6.5), Inches(12.0), Inches(0.8))

    print("  [4/15] 고객사 완료")

    # ══════════════════════════════════════════════════════
    # SLIDE 5: 핵심 가치 제안
    # ══════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout(prs))
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill_color=C_DARK)
    add_rect(slide, 0, 0, SLIDE_W, Inches(1.6), fill_color=C_PURPLE_DARK)
    add_rect(slide, 0, 0, Inches(0.3), SLIDE_H, fill_color=C_PURPLE)

    add_text(slide, "03  핵심 가치 제안", Inches(0.6), Inches(0.35), Inches(10.0), Inches(0.8),
             font_size=32, bold=True, color=C_WHITE)
    add_text(slide, "왜 CleverTap인가?", Inches(0.6), Inches(1.05), Inches(8.0), Inches(0.45),
             font_size=14, color=C_PURPLE_LIGHT)

    values = [
        ("🔄", "올인원 플랫폼", "분석 → 세그멘테이션 → 캠페인 → 최적화까지\n하나의 플랫폼에서 완결"),
        ("⚡", "실시간 처리", "초당 수백만 이벤트를 실시간 처리\n즉각적인 의사결정 지원"),
        ("🎯", "AI 기반 개인화", "ML 알고리즘으로 각 사용자에게\n최적의 메시지를 자동 전달"),
        ("📱", "멀티채널 통합", "푸시, 이메일, SMS, 인앱, WhatsApp\n모든 채널을 단일 플랫폼에서 관리"),
        ("📊", "강력한 분석", "퍼널, 코호트, 플로우 분석\n깊이 있는 사용자 행동 인사이트"),
        ("🔒", "엔터프라이즈 보안", "SOC2, GDPR, ISO 27001 인증\n금융/헬스케어 수준 보안"),
    ]

    for i, (icon, title, desc) in enumerate(values):
        col = i % 2
        row = i // 2
        x = Inches(0.6 + col * 6.2)
        y = Inches(1.8 + row * 1.8)
        add_rect(slide, x, y, Inches(5.8), Inches(1.6), fill_color=C_MID)
        add_rect(slide, x, y, Inches(5.8), Inches(0.05), fill_color=C_PURPLE)

        # 아이콘 박스
        add_rect(slide, x + Inches(0.2), y + Inches(0.25), Inches(0.7), Inches(0.7),
                 fill_color=C_PURPLE_DARK)
        add_text(slide, icon, x + Inches(0.2), y + Inches(0.25), Inches(0.7), Inches(0.7),
                 font_size=22, align=PP_ALIGN.CENTER)
        add_text(slide, title, x + Inches(1.05), y + Inches(0.2), Inches(4.5), Inches(0.4),
                 font_size=16, bold=True, color=C_WHITE)
        add_text(slide, desc, x + Inches(1.05), y + Inches(0.62), Inches(4.5), Inches(0.85),
                 font_size=12, color=C_LIGHT_GRAY, wrap=True)

    print("  [5/15] 핵심 가치 제안 완료")

    # ══════════════════════════════════════════════════════
    # SLIDE 6: 제품 아키텍처
    # ══════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout(prs))
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill_color=C_DARK)
    add_rect(slide, 0, 0, SLIDE_W, Inches(1.6), fill_color=C_PURPLE_DARK)
    add_rect(slide, 0, 0, Inches(0.3), SLIDE_H, fill_color=C_PURPLE)

    add_text(slide, "04  제품 아키텍처 개요", Inches(0.6), Inches(0.35), Inches(10.0), Inches(0.8),
             font_size=32, bold=True, color=C_WHITE)
    add_text(slide, "데이터 수집부터 실행까지 — 완전 통합 플랫폼", Inches(0.6), Inches(1.05), Inches(8.0), Inches(0.45),
             font_size=14, color=C_PURPLE_LIGHT)

    # 아키텍처 레이어 다이어그램
    layers = [
        ("데이터 수집 레이어", "SDK (iOS/Android/Web) · REST API · S3 커넥터 · 서드파티 통합", C_PURPLE),
        ("데이터 처리 레이어", "실시간 스트림 처리 · 사용자 프로필 병합 · 이벤트 저장소", C_PURPLE_DARK),
        ("분석 & AI 레이어", "행동 분석 · ML 모델 · 예측 세그멘테이션 · 최적 전송 시간", RGBColor(0x3D, 0x5A, 0xCC)),
        ("실행 레이어", "캠페인 관리 · 멀티채널 전송 · A/B 테스트 · 자동화 워크플로우", C_MID),
    ]

    for i, (title, desc, color) in enumerate(layers):
        y = Inches(1.9 + i * 1.1)
        add_rect(slide, Inches(0.6), y, Inches(6.2), Inches(0.95), fill_color=color)
        # 번호
        add_rect(slide, Inches(0.6), y, Inches(0.4), Inches(0.95), fill_color=C_PURPLE_DARK)
        add_text(slide, str(i+1), Inches(0.62), y + Inches(0.2), Inches(0.36), Inches(0.5),
                 font_size=14, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        add_text(slide, title, Inches(1.15), y + Inches(0.05), Inches(5.5), Inches(0.38),
                 font_size=14, bold=True, color=C_WHITE)
        add_text(slide, desc, Inches(1.15), y + Inches(0.45), Inches(5.5), Inches(0.42),
                 font_size=11, color=C_LIGHT_GRAY)

    # 오른쪽 제품 스크린샷
    add_image(slide, ASSETS / "screenshot_product.png",
              Inches(7.1), Inches(1.7), Inches(5.8), Inches(4.5))
    add_text(slide, "CleverTap 제품 대시보드",
             Inches(7.1), Inches(6.2), Inches(5.8), Inches(0.35),
             font_size=11, color=C_GRAY, align=PP_ALIGN.CENTER, italic=True)

    print("  [6/15] 제품 아키텍처 완료")

    # ══════════════════════════════════════════════════════
    # SLIDE 7: 실시간 분석 (Analytics)
    # ══════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout(prs))
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill_color=C_DARK)
    add_rect(slide, 0, 0, SLIDE_W, Inches(1.6), fill_color=C_PURPLE_DARK)
    add_rect(slide, 0, 0, Inches(0.3), SLIDE_H, fill_color=C_PURPLE)

    add_text(slide, "05  실시간 분석 (Analytics)", Inches(0.6), Inches(0.35), Inches(10.0), Inches(0.8),
             font_size=32, bold=True, color=C_WHITE)
    add_text(slide, "모든 데이터를 실시간으로 — 즉각적인 비즈니스 인사이트", Inches(0.6), Inches(1.05), Inches(8.0), Inches(0.45),
             font_size=14, color=C_PURPLE_LIGHT)

    # 분석 기능 리스트
    analytics_features = [
        ("퍼널 분석", "전환율 추적, 이탈 지점 파악, 퍼널 최적화"),
        ("코호트 분석", "사용자 그룹별 리텐션 트렌드 분석"),
        ("플로우 분석", "사용자 탐색 패턴 시각화"),
        ("이벤트 분석", "커스텀 이벤트 실시간 트래킹"),
        ("피벗 분석", "다차원 데이터 분석 및 리포팅"),
    ]

    for i, (title, desc) in enumerate(analytics_features):
        y = Inches(1.9 + i * 0.9)
        add_rect(slide, Inches(0.6), y, Inches(5.5), Inches(0.78), fill_color=C_MID)
        add_rect(slide, Inches(0.6), y, Inches(0.05), Inches(0.78), fill_color=C_GREEN)
        add_text(slide, title, Inches(0.8), y + Inches(0.05), Inches(2.5), Inches(0.35),
                 font_size=13, bold=True, color=C_WHITE)
        add_text(slide, desc, Inches(0.8), y + Inches(0.4), Inches(5.0), Inches(0.32),
                 font_size=11, color=C_GRAY)

    # 오른쪽 Analytics 스크린샷
    add_image(slide, ASSETS / "screenshot_analytics.png",
              Inches(6.5), Inches(1.7), Inches(6.5), Inches(5.0))
    add_text(slide, "CleverTap Analytics 대시보드 화면",
             Inches(6.5), Inches(6.7), Inches(6.5), Inches(0.35),
             font_size=10, color=C_GRAY, align=PP_ALIGN.CENTER, italic=True)

    print("  [7/15] 분석 완료")

    # ══════════════════════════════════════════════════════
    # SLIDE 8: 사용자 세그멘테이션
    # ══════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout(prs))
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill_color=C_DARK)
    add_rect(slide, 0, 0, SLIDE_W, Inches(1.6), fill_color=C_PURPLE_DARK)
    add_rect(slide, 0, 0, Inches(0.3), SLIDE_H, fill_color=C_PURPLE)

    add_text(slide, "06  사용자 세그멘테이션", Inches(0.6), Inches(0.35), Inches(10.0), Inches(0.8),
             font_size=32, bold=True, color=C_WHITE)
    add_text(slide, "정밀한 타겟팅으로 관련성 높은 경험 제공", Inches(0.6), Inches(1.05), Inches(8.0), Inches(0.45),
             font_size=14, color=C_PURPLE_LIGHT)

    # 세그멘테이션 타입
    seg_types = [
        ("정적 세그먼트", "고정된 기준으로 사용자 그룹화\n예: 연령대별, 지역별, 가입일별"),
        ("동적 세그먼트", "실시간 행동에 따라 자동 업데이트\n예: 최근 7일 미접속자"),
        ("RFM 세그먼트", "Recency·Frequency·Monetary 기반\n고가치 사용자 자동 식별"),
        ("예측 세그먼트", "ML로 이탈 가능성, 구매 의향\n사전 예측 및 선제적 대응"),
    ]

    for i, (title, desc) in enumerate(seg_types):
        col = i % 2
        row = i // 2
        x = Inches(0.6 + col * 3.2)
        y = Inches(1.9 + row * 2.0)
        add_rect(slide, x, y, Inches(2.8), Inches(1.8), fill_color=C_MID)
        add_rect(slide, x, y, Inches(2.8), Inches(0.06), fill_color=C_PURPLE)
        add_text(slide, title, x + Inches(0.15), y + Inches(0.15), Inches(2.5), Inches(0.4),
                 font_size=14, bold=True, color=C_WHITE)
        add_text(slide, desc, x + Inches(0.15), y + Inches(0.6), Inches(2.5), Inches(1.1),
                 font_size=12, color=C_LIGHT_GRAY, wrap=True)

    # 세그멘테이션 스크린샷
    add_image(slide, ASSETS / "screenshot_segmentation.png",
              Inches(6.8), Inches(1.7), Inches(6.2), Inches(5.1))
    add_text(slide, "CleverTap 세그멘테이션 화면",
             Inches(6.8), Inches(6.8), Inches(6.2), Inches(0.35),
             font_size=10, color=C_GRAY, align=PP_ALIGN.CENTER, italic=True)

    print("  [8/15] 세그멘테이션 완료")

    # ══════════════════════════════════════════════════════
    # SLIDE 9: 멀티채널 인게이지먼트
    # ══════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout(prs))
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill_color=C_DARK)
    add_rect(slide, 0, 0, SLIDE_W, Inches(1.6), fill_color=C_PURPLE_DARK)
    add_rect(slide, 0, 0, Inches(0.3), SLIDE_H, fill_color=C_PURPLE)

    add_text(slide, "07  멀티채널 인게이지먼트", Inches(0.6), Inches(0.35), Inches(10.0), Inches(0.8),
             font_size=32, bold=True, color=C_WHITE)
    add_text(slide, "고객이 있는 모든 채널에서 일관된 경험 제공", Inches(0.6), Inches(1.05), Inches(8.0), Inches(0.45),
             font_size=14, color=C_PURPLE_LIGHT)

    channels = [
        ("📲", "푸시 알림", "iOS/Android\n리치 푸시, 딥링크\n세그먼트 타겟팅"),
        ("📧", "이메일", "드래그&드롭 에디터\nAMP 이메일 지원\n자동화 트리거"),
        ("💬", "SMS/MMS", "글로벌 통신사 연동\n개인화 메시지\n전환 추적"),
        ("📱", "인앱 메시지", "풀스크린/배너/팝업\n실시간 노출 조건\n인터랙티브 CTA"),
        ("🌐", "웹 푸시", "브라우저 알림\n세션 외 재참여\n퍼미션 관리"),
        ("💚", "WhatsApp", "공식 비즈니스 API\n템플릿 메시지\n챗봇 연동"),
    ]

    for i, (icon, name, desc) in enumerate(channels):
        col = i % 3
        row = i // 3
        x = Inches(0.6 + col * 2.75)
        y = Inches(1.9 + row * 2.5)
        add_rect(slide, x, y, Inches(2.5), Inches(2.2), fill_color=C_MID)
        add_rect(slide, x, y, Inches(2.5), Inches(0.06), fill_color=C_PURPLE)
        # 아이콘
        add_text(slide, icon, x + Inches(0.15), y + Inches(0.15), Inches(0.6), Inches(0.6),
                 font_size=28, align=PP_ALIGN.CENTER)
        add_text(slide, name, x + Inches(0.15), y + Inches(0.8), Inches(2.2), Inches(0.42),
                 font_size=14, bold=True, color=C_WHITE)
        add_text(slide, desc, x + Inches(0.15), y + Inches(1.25), Inches(2.2), Inches(0.85),
                 font_size=11, color=C_GRAY, wrap=True)

    # 인게이지먼트 스크린샷
    add_image(slide, ASSETS / "screenshot_engagement.png",
              Inches(9.0), Inches(1.8), Inches(4.0), Inches(5.5))

    print("  [9/15] 멀티채널 인게이지먼트 완료")

    # ══════════════════════════════════════════════════════
    # SLIDE 10: AI/ML 기반 개인화
    # ══════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout(prs))
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill_color=C_DARK)
    add_rect(slide, 0, 0, SLIDE_W, Inches(1.6), fill_color=C_PURPLE_DARK)
    add_rect(slide, 0, 0, Inches(0.3), SLIDE_H, fill_color=C_PURPLE)

    add_text(slide, "08  AI/ML 기반 개인화", Inches(0.6), Inches(0.35), Inches(10.0), Inches(0.8),
             font_size=32, bold=True, color=C_WHITE)
    add_text(slide, "머신러닝으로 구현하는 1:1 초개인화 마케팅", Inches(0.6), Inches(1.05), Inches(8.0), Inches(0.45),
             font_size=14, color=C_PURPLE_LIGHT)

    ai_features = [
        ("🤖 Clever.AI", "CleverTap 독자 AI 엔진\n고객 행동을 학습하고 예측"),
        ("⏰ 최적 전송 시간 (STO)", "각 사용자별 최적 수신 시간 자동 선택\n오픈율 평균 47% 향상"),
        ("🎯 예측 세그먼트", "이탈 예측, 구매 예측, 구독 해지 예측\nLTV 기반 자동 분류"),
        ("💡 동적 콘텐츠", "개인 프로필 기반 실시간 콘텐츠 변환\n제품 추천 자동화"),
        ("📈 자동 최적화", "A/B 테스트 자동 승자 선택\n트래픽 동적 배분"),
        ("🔮 이탈 방지", "이탈 위험 사용자 선제 감지\n재참여 캠페인 자동 트리거"),
    ]

    for i, (title, desc) in enumerate(ai_features):
        col = i % 2
        row = i // 2
        x = Inches(0.6 + col * 6.0)
        y = Inches(1.9 + row * 1.6)
        add_rect(slide, x, y, Inches(5.5), Inches(1.4), fill_color=C_MID)
        add_rect(slide, x, y, Inches(0.06), Inches(1.4), fill_color=C_PURPLE)
        add_text(slide, title, x + Inches(0.2), y + Inches(0.1), Inches(5.0), Inches(0.45),
                 font_size=14, bold=True, color=C_WHITE)
        add_text(slide, desc, x + Inches(0.2), y + Inches(0.6), Inches(5.0), Inches(0.7),
                 font_size=12, color=C_LIGHT_GRAY, wrap=True)

    # AI 스크린샷
    add_image(slide, ASSETS / "screenshot_ai_ml.png",
              Inches(7.3), Inches(5.0), Inches(5.7), Inches(2.2))

    print("  [10/15] AI/ML 개인화 완료")

    # ══════════════════════════════════════════════════════
    # SLIDE 11: A/B 테스트
    # ══════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout(prs))
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill_color=C_DARK)
    add_rect(slide, 0, 0, SLIDE_W, Inches(1.6), fill_color=C_PURPLE_DARK)
    add_rect(slide, 0, 0, Inches(0.3), SLIDE_H, fill_color=C_PURPLE)

    add_text(slide, "09  A/B 테스트 & 최적화", Inches(0.6), Inches(0.35), Inches(10.0), Inches(0.8),
             font_size=32, bold=True, color=C_WHITE)
    add_text(slide, "데이터 기반 의사결정으로 마케팅 성과 극대화", Inches(0.6), Inches(1.05), Inches(8.0), Inches(0.45),
             font_size=14, color=C_PURPLE_LIGHT)

    ab_items = [
        ("메시지 A/B 테스트", "제목, 본문, CTA 버튼\n이미지, 이모지 테스트"),
        ("전송 시간 테스트", "발송 시간대별 성과 비교\n요일/시간 최적화"),
        ("채널 A/B 테스트", "푸시 vs 이메일 vs SMS\n채널별 효과 측정"),
        ("자동 승자 선택", "통계적 유의성 기반 자동 선택\n트래픽 손실 최소화"),
    ]

    for i, (title, desc) in enumerate(ab_items):
        x = Inches(0.6 + (i % 2) * 3.3)
        y = Inches(1.9 + (i // 2) * 1.9)
        add_rect(slide, x, y, Inches(3.0), Inches(1.7), fill_color=C_MID)
        add_rect(slide, x, y, Inches(3.0), Inches(0.06), fill_color=C_ACCENT)
        add_text(slide, title, x + Inches(0.15), y + Inches(0.15), Inches(2.7), Inches(0.4),
                 font_size=14, bold=True, color=C_WHITE)
        add_text(slide, desc, x + Inches(0.15), y + Inches(0.6), Inches(2.7), Inches(0.9),
                 font_size=12, color=C_LIGHT_GRAY, wrap=True)

    # A/B 스크린샷
    add_image(slide, ASSETS / "screenshot_ab_testing.png",
              Inches(7.0), Inches(1.7), Inches(6.0), Inches(5.5))
    add_text(slide, "CleverTap A/B 테스트 화면",
             Inches(7.0), Inches(7.1), Inches(6.0), Inches(0.3),
             font_size=10, color=C_GRAY, align=PP_ALIGN.CENTER, italic=True)

    # 성과 수치
    ab_results = [("47%", "오픈율 향상"), ("32%", "전환율 개선"), ("2.3배", "ROI 증가")]
    for i, (val, label) in enumerate(ab_results):
        x = Inches(0.6 + i * 2.15)
        y = Inches(5.8)
        add_rect(slide, x, y, Inches(1.9), Inches(1.2), fill_color=C_PURPLE_DARK)
        add_text(slide, val, x + Inches(0.1), y + Inches(0.1), Inches(1.7), Inches(0.55),
                 font_size=26, bold=True, color=C_GREEN, align=PP_ALIGN.CENTER)
        add_text(slide, label, x + Inches(0.1), y + Inches(0.7), Inches(1.7), Inches(0.4),
                 font_size=11, color=C_GRAY, align=PP_ALIGN.CENTER)

    print("  [11/15] A/B 테스트 완료")

    # ══════════════════════════════════════════════════════
    # SLIDE 12: 핀테크 활용 사례
    # ══════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout(prs))
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill_color=C_DARK)
    add_rect(slide, 0, 0, SLIDE_W, Inches(1.6), fill_color=C_PURPLE_DARK)
    add_rect(slide, 0, 0, Inches(0.3), SLIDE_H, fill_color=C_PURPLE)

    add_text(slide, "10  핀테크/금융 업계 활용 사례", Inches(0.6), Inches(0.35), Inches(10.0), Inches(0.8),
             font_size=32, bold=True, color=C_WHITE)
    add_text(slide, "금융 서비스에 최적화된 마케팅 자동화", Inches(0.6), Inches(1.05), Inches(8.0), Inches(0.45),
             font_size=14, color=C_PURPLE_LIGHT)

    # 왼쪽: 사용 시나리오
    use_cases = [
        ("신규 가입 온보딩", "가입 후 7일 온보딩 시나리오\n단계별 안내로 첫 거래 유도"),
        ("KYC 완료 독려", "미완료 사용자 자동 리마인드\n채널별 순차 발송"),
        ("투자 상품 추천", "거래 이력 기반 AI 상품 추천\n맞춤형 리스크 등급별 제안"),
        ("이탈 방지 캠페인", "30일 미거래 사용자 자동 감지\n특별 혜택 오퍼 발송"),
        ("거래 알림 & 리텐션", "실시간 거래 알림 + 크로스셀\n수수료 면제 프로모션"),
    ]

    for i, (title, desc) in enumerate(use_cases):
        y = Inches(1.9 + i * 0.9)
        add_rect(slide, Inches(0.6), y, Inches(5.8), Inches(0.78), fill_color=C_MID)
        add_rect(slide, Inches(0.6), y, Inches(0.05), Inches(0.78), fill_color=C_ACCENT)
        add_text(slide, title, Inches(0.8), y + Inches(0.05), Inches(2.8), Inches(0.35),
                 font_size=13, bold=True, color=C_WHITE)
        add_text(slide, desc, Inches(0.8), y + Inches(0.4), Inches(5.3), Inches(0.32),
                 font_size=11, color=C_GRAY)

    # 고객사 로고 영역
    add_rect(slide, Inches(0.6), Inches(6.8), Inches(5.8), Inches(0.5), fill_color=C_MID)
    add_text(slide, "PhonePe · Axis Bank · BankBazaar · Groww · 카카오페이(도입 검토중)",
             Inches(0.7), Inches(6.9), Inches(5.6), Inches(0.35),
             font_size=11, color=C_PURPLE_LIGHT, align=PP_ALIGN.CENTER)

    # 오른쪽: 성과 수치
    results = [
        ("40%", "신규 가입 전환율 향상"),
        ("28%", "KYC 완료율 증가"),
        ("3.2배", "온보딩 캠페인 ROI"),
        ("65%", "이탈 방지 성공률"),
    ]
    for i, (val, label) in enumerate(results):
        x = Inches(7.0 + (i % 2) * 2.9)
        y = Inches(1.9 + (i // 2) * 2.0)
        add_rect(slide, x, y, Inches(2.5), Inches(1.7), fill_color=C_PURPLE_DARK)
        add_rect(slide, x, y, Inches(2.5), Inches(0.06), fill_color=C_GREEN)
        add_text(slide, val, x + Inches(0.1), y + Inches(0.3), Inches(2.3), Inches(0.7),
                 font_size=36, bold=True, color=C_GREEN, align=PP_ALIGN.CENTER)
        add_text(slide, label, x + Inches(0.1), y + Inches(1.05), Inches(2.3), Inches(0.5),
                 font_size=12, color=C_LIGHT_GRAY, align=PP_ALIGN.CENTER, wrap=True)

    # 스크린샷
    add_image(slide, ASSETS / "screenshot_customers.png",
              Inches(7.0), Inches(5.6), Inches(5.8), Inches(1.6))

    print("  [12/15] 핀테크 활용 사례 완료")

    # ══════════════════════════════════════════════════════
    # SLIDE 13: 도입 효과 & ROI
    # ══════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout(prs))
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill_color=C_DARK)
    add_rect(slide, 0, 0, SLIDE_W, Inches(1.6), fill_color=C_PURPLE_DARK)
    add_rect(slide, 0, 0, Inches(0.3), SLIDE_H, fill_color=C_PURPLE)

    add_text(slide, "11  도입 효과 & ROI", Inches(0.6), Inches(0.35), Inches(10.0), Inches(0.8),
             font_size=32, bold=True, color=C_WHITE)
    add_text(slide, "측정된 성과 — 실제 고객사 기준 평균 수치", Inches(0.6), Inches(1.05), Inches(8.0), Inches(0.45),
             font_size=14, color=C_PURPLE_LIGHT)

    roi_metrics = [
        ("43%", "리텐션율 향상", "이탈 방지 캠페인 적용 후"),
        ("2.8배", "마케팅 ROI", "전통 채널 대비"),
        ("62%", "인게이지먼트 증가", "개인화 메시지 적용 후"),
        ("30%", "마케팅 비용 절감", "자동화로 운영 효율 개선"),
        ("5배", "캠페인 실행 속도", "수동 대비 자동화 효과"),
        ("93%", "고객 만족도", "NPS 기준 CleverTap 도입사"),
    ]

    for i, (val, title, desc) in enumerate(roi_metrics):
        col = i % 3
        row = i // 2
        x = Inches(0.6 + col * 4.15)
        y = Inches(1.9 + (i // 3) * 2.3)
        add_rect(slide, x, y, Inches(3.8), Inches(2.0), fill_color=C_MID)
        add_rect(slide, x, y, Inches(3.8), Inches(0.06), fill_color=C_GREEN)
        add_text(slide, val, x + Inches(0.15), y + Inches(0.2), Inches(3.5), Inches(0.8),
                 font_size=42, bold=True, color=C_GREEN, align=PP_ALIGN.CENTER)
        add_text(slide, title, x + Inches(0.15), y + Inches(1.0), Inches(3.5), Inches(0.4),
                 font_size=14, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        add_text(slide, desc, x + Inches(0.15), y + Inches(1.42), Inches(3.5), Inches(0.45),
                 font_size=11, color=C_GRAY, align=PP_ALIGN.CENTER)

    print("  [13/15] ROI 완료")

    # ══════════════════════════════════════════════════════
    # SLIDE 14: 도입 프로세스 & 타임라인
    # ══════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout(prs))
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill_color=C_DARK)
    add_rect(slide, 0, 0, SLIDE_W, Inches(1.6), fill_color=C_PURPLE_DARK)
    add_rect(slide, 0, 0, Inches(0.3), SLIDE_H, fill_color=C_PURPLE)

    add_text(slide, "12  도입 프로세스 & 타임라인", Inches(0.6), Inches(0.35), Inches(10.0), Inches(0.8),
             font_size=32, bold=True, color=C_WHITE)
    add_text(slide, "빠르고 안전한 구축 — 평균 4주 내 라이브", Inches(0.6), Inches(1.05), Inches(8.0), Inches(0.45),
             font_size=14, color=C_PURPLE_LIGHT)

    # 타임라인
    phases = [
        ("Week 1", "킥오프 & 설계", "요구사항 정의\nSDK 통합 계획\n이벤트 텍소노미 설계"),
        ("Week 2", "SDK 통합", "iOS/Android SDK 설치\n이벤트 태깅\n데이터 검증"),
        ("Week 3", "캠페인 구성", "세그멘테이션 설정\n캠페인 템플릿 구성\nA/B 테스트 설계"),
        ("Week 4", "라이브 & 최적화", "첫 캠페인 발송\n성과 모니터링\n피드백 반영"),
    ]

    # 타임라인 라인
    add_rect(slide, Inches(1.2), Inches(3.2), Inches(11.0), Inches(0.06), fill_color=C_PURPLE)

    for i, (week, title, desc) in enumerate(phases):
        x = Inches(0.8 + i * 2.9)
        # 원
        circle = slide.shapes.add_shape(9, x + Inches(0.9), Inches(3.0), Inches(0.45), Inches(0.45))
        circle.fill.solid()
        circle.fill.fore_color.rgb = C_PURPLE
        circle.line.fill.background()

        # 위쪽 텍스트 (홀수)
        if i % 2 == 0:
            add_text(slide, week, x + Inches(0.15), Inches(1.8), Inches(2.5), Inches(0.35),
                     font_size=12, color=C_PURPLE_LIGHT, bold=True, align=PP_ALIGN.CENTER)
            add_text(slide, title, x + Inches(0.15), Inches(2.2), Inches(2.5), Inches(0.4),
                     font_size=14, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
            add_rect(slide, x + Inches(0.15), Inches(2.7), Inches(2.5), Inches(0.35),
                     fill_color=C_MID)
        else:
            add_text(slide, week, x + Inches(0.15), Inches(3.8), Inches(2.5), Inches(0.35),
                     font_size=12, color=C_PURPLE_LIGHT, bold=True, align=PP_ALIGN.CENTER)
            add_text(slide, title, x + Inches(0.15), Inches(4.2), Inches(2.5), Inches(0.4),
                     font_size=14, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

        # 설명 (아래)
        desc_y = Inches(5.1) if i % 2 == 0 else Inches(5.1)
        add_rect(slide, x + Inches(0.1), Inches(5.1), Inches(2.6), Inches(1.8), fill_color=C_MID)
        add_text(slide, desc, x + Inches(0.2), Inches(5.2), Inches(2.4), Inches(1.6),
                 font_size=12, color=C_LIGHT_GRAY, wrap=True)

    # 추가 지원 사항
    add_rect(slide, Inches(0.6), Inches(7.0), Inches(12.0), Inches(0.3), fill_color=C_PURPLE_DARK)
    add_text(slide, "✅  전담 CSM 배정  |  한국어 기술 지원  |  24/7 모니터링  |  전용 온보딩 플레이북 제공",
             Inches(0.8), Inches(7.05), Inches(11.8), Inches(0.25),
             font_size=11, color=C_WHITE, align=PP_ALIGN.CENTER)

    print("  [14/15] 도입 프로세스 완료")

    # ══════════════════════════════════════════════════════
    # SLIDE 15: 가격 정책 & 다음 단계
    # ══════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout(prs))
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill_color=C_DARK)
    add_rect(slide, 0, 0, SLIDE_W, Inches(1.6), fill_color=C_PURPLE_DARK)
    add_rect(slide, 0, 0, Inches(0.3), SLIDE_H, fill_color=C_PURPLE)

    add_text(slide, "13  가격 정책 & 다음 단계", Inches(0.6), Inches(0.35), Inches(10.0), Inches(0.8),
             font_size=32, bold=True, color=C_WHITE)
    add_text(slide, "비즈니스 규모에 맞는 유연한 플랜", Inches(0.6), Inches(1.05), Inches(8.0), Inches(0.45),
             font_size=14, color=C_PURPLE_LIGHT)

    # 가격 플랜 카드
    plans = [
        ("Essentials", "스타트업 & 성장기\n월간 MAU 50만 이하", [
            "핵심 분석 기능",
            "푸시 + 이메일",
            "기본 세그멘테이션",
            "이메일 지원",
        ], False),
        ("Advanced", "Mid-Market\n월간 MAU 500만 이하", [
            "Essentials 전체 포함",
            "AI 예측 세그먼트",
            "멀티채널 전체",
            "전담 CSM",
        ], True),
        ("Enterprise", "대기업 & 금융기관\nMAU 무제한", [
            "Advanced 전체 포함",
            "커스텀 SLA",
            "온프레미스 옵션",
            "전용 인프라",
        ], False),
    ]

    plan_x_positions = [Inches(0.6), Inches(4.5), Inches(8.4)]
    for i, (plan_name, desc, features, is_recommended) in enumerate(plans):
        x = plan_x_positions[i]
        bg_color = C_PURPLE if is_recommended else C_MID
        add_rect(slide, x, Inches(1.9), Inches(3.6), Inches(4.8),
                 fill_color=bg_color)
        if is_recommended:
            add_rect(slide, x, Inches(1.9), Inches(3.6), Inches(0.38),
                     fill_color=C_PURPLE_DARK)
            add_text(slide, "추천", x + Inches(0.1), Inches(1.92), Inches(3.4), Inches(0.3),
                     font_size=12, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        add_text(slide, plan_name, x + Inches(0.2), Inches(2.35), Inches(3.2), Inches(0.5),
                 font_size=22, bold=True, color=C_WHITE)
        add_text(slide, desc, x + Inches(0.2), Inches(2.85), Inches(3.2), Inches(0.55),
                 font_size=11, color=C_LIGHT_GRAY, wrap=True)
        add_rect(slide, x + Inches(0.1), Inches(3.45), Inches(3.4), Inches(0.03),
                 fill_color=C_PURPLE_LIGHT)
        for j, feature in enumerate(features):
            fy = Inches(3.6 + j * 0.55)
            add_text(slide, "✓  " + feature, x + Inches(0.2), fy, Inches(3.2), Inches(0.45),
                     font_size=12, color=C_WHITE)

    # 다음 단계
    add_text(slide, "다음 단계 (Next Steps)", Inches(0.6), Inches(6.9), Inches(5.0), Inches(0.4),
             font_size=16, bold=True, color=C_PURPLE_LIGHT)
    next_steps = "① POC 범위 확정  →  ② 무료 파일럿 세션 예약 (2주)  →  ③ 맞춤 견적 제안  →  ④ 계약 및 온보딩"
    add_text(slide, next_steps, Inches(0.6), Inches(7.2), Inches(12.0), Inches(0.4),
             font_size=13, color=C_WHITE)

    print("  [15/15] 가격 정책 완료")

    # ══════════════════════════════════════════════════════
    # SLIDE 16 (보너스): 문의 & CTA
    # ══════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout(prs))
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill_color=C_DARK)
    add_rect(slide, 0, 0, Inches(0.3), SLIDE_H, fill_color=C_PURPLE)

    # 중앙 배치
    add_rect(slide, Inches(1.5), Inches(1.5), Inches(10.0), Inches(4.5), fill_color=C_MID)
    add_rect(slide, Inches(1.5), Inches(1.5), Inches(10.0), Inches(0.08), fill_color=C_PURPLE)

    add_text(slide, "지금 바로 시작하세요", Inches(2.0), Inches(2.0), Inches(9.0), Inches(0.9),
             font_size=40, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, "2주 무료 파일럿 프로그램 신청",
             Inches(2.0), Inches(2.9), Inches(9.0), Inches(0.6),
             font_size=24, color=C_PURPLE_LIGHT, align=PP_ALIGN.CENTER)
    add_text(slide, "실제 데이터로 CleverTap의 가치를 직접 확인하세요",
             Inches(2.0), Inches(3.5), Inches(9.0), Inches(0.4),
             font_size=16, color=C_LIGHT_GRAY, align=PP_ALIGN.CENTER)

    add_rect(slide, Inches(4.5), Inches(4.3), Inches(4.0), Inches(0.7), fill_color=C_PURPLE)
    add_text(slide, "데모 신청하기  →", Inches(4.5), Inches(4.35), Inches(4.0), Inches(0.6),
             font_size=18, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    # 연락처 정보
    contact_info = [
        ("웹사이트", "clevertap.com"),
        ("이메일", "korea@clevertap.com"),
        ("전화", "+82-2-XXXX-XXXX"),
    ]
    for i, (label, value) in enumerate(contact_info):
        x = Inches(1.8 + i * 3.5)
        add_text(slide, label, x, Inches(6.3), Inches(3.0), Inches(0.3),
                 font_size=12, color=C_GRAY, align=PP_ALIGN.CENTER)
        add_text(slide, value, x, Inches(6.6), Inches(3.0), Inches(0.4),
                 font_size=14, bold=True, color=C_PURPLE_LIGHT, align=PP_ALIGN.CENTER)

    # CleverTap 로고 텍스트
    add_text(slide, "CleverTap", Inches(5.5), Inches(0.2), Inches(2.5), Inches(0.6),
             font_size=28, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    # 스크린샷 배경 (희미하게)
    add_image(slide, ASSETS / "viewport_main.png",
              Inches(1.5), Inches(1.5), Inches(10.0), Inches(4.5))
    add_rect(slide, Inches(1.5), Inches(1.5), Inches(10.0), Inches(4.5),
             fill_color=C_MID)
    # 다시 텍스트 위에 올리기
    add_text(slide, "지금 바로 시작하세요", Inches(2.0), Inches(2.0), Inches(9.0), Inches(0.9),
             font_size=40, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, "2주 무료 파일럿 프로그램 신청",
             Inches(2.0), Inches(2.9), Inches(9.0), Inches(0.6),
             font_size=24, color=C_PURPLE_LIGHT, align=PP_ALIGN.CENTER)
    add_text(slide, "실제 데이터로 CleverTap의 가치를 직접 확인하세요",
             Inches(2.0), Inches(3.5), Inches(9.0), Inches(0.4),
             font_size=16, color=C_LIGHT_GRAY, align=PP_ALIGN.CENTER)
    add_rect(slide, Inches(4.5), Inches(4.3), Inches(4.0), Inches(0.7), fill_color=C_PURPLE)
    add_text(slide, "데모 신청하기  ->", Inches(4.5), Inches(4.35), Inches(4.0), Inches(0.6),
             font_size=18, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    print("  [보너스] CTA 슬라이드 완료")

    # ── 저장 ──────────────────────────────────────────────
    prs.save(str(OUTPUT))
    print(f"\n✅ 저장 완료: {OUTPUT}")
    print(f"   총 슬라이드: {len(prs.slides)}장")
    print(f"   파일 크기: {OUTPUT.stat().st_size // 1024}KB")
    return len(prs.slides)

if __name__ == "__main__":
    print("CleverTap 한국어 제안서 PPTX 생성 시작")
    n = build_slides()
    print(f"\n완료! {n}장 슬라이드 생성")
