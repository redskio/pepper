"""
MBC 모다이브 × 위어드섹터 서비스 제안서 PPTX 빌더
11슬라이드 | 16:9 | Malgun Gothic
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml
import copy
from lxml import etree
import os
from io import BytesIO

# ── 색상 팔레트 ──────────────────────────────────────────
C_WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
C_DEEP_PURPLE= RGBColor(0x00, 0x66, 0xCC)
C_BLUE       = RGBColor(0x1D, 0x4E, 0xD8)
C_DARK       = RGBColor(0x0F, 0x17, 0x2A)
C_LIGHT_PURP = RGBColor(0xE3, 0xF2, 0xFD)
C_LIGHT_BLUE = RGBColor(0xEF, 0xF6, 0xFF)
C_GOLD       = RGBColor(0xD9, 0x77, 0x06)
C_GREEN      = RGBColor(0x05, 0x96, 0x69)
C_TABLE_EVEN = RGBColor(0xFA, 0xFA, 0xFA)
C_SUB_TEXT   = RGBColor(0x64, 0x74, 0x8B)
C_RED        = RGBColor(0xDC, 0x26, 0x26)
C_PURPLE_LIGHT2 = RGBColor(0x90, 0xCA, 0xFF)
C_MED_PURPLE = RGBColor(0x00, 0x52, 0xAA)
C_BORDER     = RGBColor(0xE2, 0xE8, 0xF0)
C_DARK_GREY  = RGBColor(0x1E, 0x29, 0x3B)
C_GRAD_MID   = RGBColor(0x00, 0x3D, 0x99)  # gradient midpoint purple

# Slide dimensions (16:9 widescreen)
SLIDE_W = Cm(33.87)
SLIDE_H = Cm(19.05)

IMG_DIR = os.path.join(os.path.dirname(__file__), "images")
OUT_PATH = os.path.join(os.path.dirname(__file__), "mbc_modive_proposal_v2.pptx")

# ── 헬퍼 함수들 ──────────────────────────────────────────

def add_rect(slide, left, top, width, height, fill_color=None, line_color=None, line_width=Pt(1)):
    from pptx.util import Pt
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape


def add_textbox(slide, left, top, width, height, text, font_name="Malgun Gothic",
                font_size=14, bold=False, italic=False, color=None, align=PP_ALIGN.LEFT,
                word_wrap=True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    if color:
        run.font.color.rgb = color
    return txBox


def add_text_in_shape(shape, text, font_name="Malgun Gothic", font_size=14,
                       bold=False, italic=False, color=None, align=PP_ALIGN.LEFT):
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    if color:
        run.font.color.rgb = color


def add_paragraph(tf, text, font_name="Malgun Gothic", font_size=14,
                  bold=False, italic=False, color=None, align=PP_ALIGN.LEFT,
                  space_before=0):
    from pptx.util import Pt
    p = tf.add_paragraph()
    p.alignment = align
    if space_before:
        p.space_before = Pt(space_before)
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    if color:
        run.font.color.rgb = color
    return p


def set_gradient_fill(shape, color1, color2, angle=0):
    """Left-to-right gradient via XML manipulation"""
    sp = shape.element
    spPr = sp.find(qn('p:spPr'))
    if spPr is None:
        spPr = etree.SubElement(sp, qn('p:spPr'))

    # Remove existing fill
    for old in spPr.findall(qn('a:solidFill')):
        spPr.remove(old)
    for old in spPr.findall(qn('a:gradFill')):
        spPr.remove(old)
    for old in spPr.findall(qn('a:noFill')):
        spPr.remove(old)

    def hex_color(rgb):
        return str(rgb)  # RGBColor.__str__ returns hex like "6B21A8"

    grad_xml = f'''<a:gradFill xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" rotWithShape="1">
  <a:gsLst>
    <a:gs pos="0">
      <a:srgbClr val="{hex_color(color1)}"/>
    </a:gs>
    <a:gs pos="100000">
      <a:srgbClr val="{hex_color(color2)}"/>
    </a:gs>
  </a:gsLst>
  <a:lin ang="{angle}" scaled="0"/>
</a:gradFill>'''
    grad_elem = parse_xml(grad_xml)
    spPr.append(grad_elem)


def add_header_bar(slide, title_text):
    """공통 상단 헤더 바 (딥 퍼플 배경)"""
    bar = add_rect(slide, 0, 0, SLIDE_W, Cm(1.6), fill_color=C_DEEP_PURPLE)
    add_text_in_shape(bar, title_text, font_size=20, bold=True,
                      color=C_WHITE, align=PP_ALIGN.CENTER)
    bar.text_frame.paragraphs[0].runs[0].font.size = Pt(18)


def add_image(slide, img_path, left, top, width, height):
    try:
        if os.path.exists(img_path):
            slide.shapes.add_picture(img_path, left, top, width, height)
            return True
    except Exception as e:
        print(f"Image load error {img_path}: {e}")
    return False


# ════════════════════════════════════════════════════════
# SLIDE BUILDERS
# ════════════════════════════════════════════════════════

def build_slide1(slide):
    """표지"""
    # 좌측 그라데이션 배경 (전체)
    bg = add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill_color=C_DEEP_PURPLE)
    set_gradient_fill(bg, C_DEEP_PURPLE, C_BLUE, angle=5400000)  # left→right ~90 deg

    # 우측 이미지 오버레이
    img_path = os.path.join(IMG_DIR, "slide1_tech.jpg")
    img_added = add_image(slide, img_path, Cm(19), 0, Cm(14.87), SLIDE_H)
    if img_added:
        # 이미지 위에 반투명 오버레이
        overlay = add_rect(slide, Cm(19), 0, Cm(14.87), SLIDE_H,
                           fill_color=RGBColor(0x00, 0x66, 0xCC))
        sp = overlay.element
        spPr = sp.find(qn('p:spPr'))
        # Set transparency
        solidFill = spPr.find('.//' + qn('a:solidFill'))
        if solidFill is not None:
            srgbClr = solidFill.find(qn('a:srgbClr'))
            if srgbClr is not None:
                alpha = etree.SubElement(srgbClr, qn('a:alpha'))
                alpha.set('val', '50000')  # 50% transparent

    # 상단 회사 크레딧
    add_textbox(slide, Cm(2), Cm(2.5), Cm(16), Cm(0.8),
                "위어드섹터 × MBC 모다이브",
                font_size=16, color=C_WHITE)

    # 메인 타이틀
    tb = slide.shapes.add_textbox(Cm(2), Cm(4.5), Cm(16), Cm(3))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = "데이터·코드·마테크·CRM"
    run.font.name = "Malgun Gothic"
    run.font.size = Pt(32)
    run.font.bold = True
    run.font.color.rgb = C_WHITE

    add_paragraph(tf, "서비스 제안서", font_size=22, bold=True, color=C_WHITE)

    # 구분선
    line = add_rect(slide, Cm(2), Cm(8.8), Cm(15), Cm(0.05), fill_color=C_PURPLE_LIGHT2)

    # 태그라인
    add_textbox(slide, Cm(2), Cm(9.2), Cm(16), Cm(1.2),
                "팬 데이터를 수익으로 — 모잇의 성장을 데이터로 가속합니다",
                font_size=14, italic=True, color=C_PURPLE_LIGHT2)

    # 하단 연락처
    add_textbox(slide, Cm(2), Cm(17.2), Cm(14), Cm(0.8),
                "위어드섹터  |  info@weirdsector.co.kr",
                font_size=11, color=C_WHITE)
    add_textbox(slide, Cm(28), Cm(17.2), Cm(5), Cm(0.8),
                "2026년 4월", font_size=11, color=C_WHITE, align=PP_ALIGN.RIGHT)


def build_slide2(slide):
    """Executive Summary"""
    add_header_bar(slide, "Executive Summary — 한 눈에 보는 제안 핵심")

    # 3줄 핵심 요약 박스
    summary_bg = add_rect(slide, Cm(1.2), Cm(2), Cm(31.5), Cm(4.8),
                           fill_color=C_WHITE, line_color=C_DEEP_PURPLE, line_width=Pt(1.5))

    badges = [
        ("①", "문제", "모잇은 급성장 중이지만, 팬 행동 데이터를 체계적으로 수집·분석·활용하는 인프라가 부족합니다."),
        ("②", "솔루션", "위어드섹터가 데이터 수집 → 코드 관리 → 마테크 자동화 → CRM까지 M/M 방식으로 순차 구축합니다."),
        ("③", "기대효과", "팬 이탈 30% 감소, 굿즈 전환율 2배, 운영팀 리포팅 시간 80% 절감 (업계 평균 기준 추정)"),
    ]

    for i, (num, label, desc) in enumerate(badges):
        y = Cm(2.3) + i * Cm(1.4)
        badge = add_rect(slide, Cm(1.5), y, Cm(0.7), Cm(0.7), fill_color=C_DEEP_PURPLE)
        add_text_in_shape(badge, num, font_size=10, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        add_textbox(slide, Cm(2.5), y, Cm(2.5), Cm(0.8),
                    label, font_size=12, bold=True, color=C_DEEP_PURPLE)
        add_textbox(slide, Cm(5.2), y, Cm(27), Cm(0.8),
                    desc, font_size=12, color=C_DARK)

    # 서비스 테이블
    table_y = Cm(7.2)
    headers = ["서비스 영역", "핵심 내용"]
    col_widths = [Cm(8), Cm(22.5)]
    rows = [
        ("📊  데이터 관리", "사용자 행동 수집 · 정제 · 대시보드 구축"),
        ("🏷️  코드 관리", "트래킹 코드 · 태그 매니저 · 측정 인프라"),
        ("🤖  마테크 (MarTech)", "마케팅 자동화 · 개인화 캠페인 도구"),
        ("👥  CRM", "팬 세그먼테이션 · 리텐션 · 생애주기 관리"),
    ]

    # Table header
    x = Cm(1.2)
    for j, (h, w) in enumerate(zip(headers, col_widths)):
        cell = add_rect(slide, x, table_y, w, Cm(0.8), fill_color=C_DEEP_PURPLE)
        add_text_in_shape(cell, h, font_size=13, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        x += w

    for i, (svc, desc) in enumerate(rows):
        row_y = table_y + Cm(0.8) + i * Cm(1.1)
        bg_color = C_LIGHT_PURP if i % 2 == 0 else C_TABLE_EVEN
        x = Cm(1.2)
        cell1 = add_rect(slide, x, row_y, col_widths[0], Cm(1.0), fill_color=bg_color,
                          line_color=C_BORDER, line_width=Pt(0.5))
        add_text_in_shape(cell1, svc, font_size=13, bold=True, color=C_DEEP_PURPLE)
        x += col_widths[0]
        cell2 = add_rect(slide, x, row_y, col_widths[1], Cm(1.0), fill_color=bg_color,
                          line_color=C_BORDER, line_width=Pt(0.5))
        add_text_in_shape(cell2, desc, font_size=12, color=C_DARK)

    # 하단 골드 강조 바
    gold_bar = add_rect(slide, 0, Cm(17.5), SLIDE_W, Cm(1.2), fill_color=C_GOLD)
    add_text_in_shape(gold_bar,
                      "운영 방식: 月 200만원 / M/M 기반 · 매월 우선순위 선정 · 실행 · 리뷰",
                      font_size=14, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)


def build_slide3(slide):
    """MBC 모다이브의 과제"""
    add_header_bar(slide, "MBC 모다이브의 과제")

    # 상단 좌측: 모잇 현황 박스
    info_box = add_rect(slide, Cm(1.2), Cm(1.9), Cm(12), Cm(5.5), fill_color=C_LIGHT_PURP)
    tf = info_box.text_frame
    tf.word_wrap = True
    add_text_in_shape(info_box, "모잇 플랫폼 현황", font_size=14, bold=True, color=C_DEEP_PURPLE)
    add_paragraph(tf, "• AI 캐릭터 채팅 · 굿즈 커머스 · 팬 커뮤니티", font_size=12, color=C_DARK)
    add_paragraph(tf, "• MBC 드라마 IP 30개+ 캐릭터 운영 중", font_size=12, color=C_DARK)
    add_paragraph(tf, "  (커피프린스, 내이름은 김삼순 등)", font_size=11, color=C_SUB_TEXT)
    add_paragraph(tf, "• 2024년 10월 정식 출시 → 2026년 글로벌 예정", font_size=12, color=C_DEEP_PURPLE, bold=True)

    # 상단 우측: 이미지
    img_path = os.path.join(IMG_DIR, "slide3_dashboard.jpg")
    if not add_image(slide, img_path, Cm(14), Cm(1.9), Cm(18.7), Cm(5.5)):
        placeholder = add_rect(slide, Cm(14), Cm(1.9), Cm(18.7), Cm(5.5),
                                fill_color=C_LIGHT_BLUE)
        add_text_in_shape(placeholder, "데이터 대시보드 이미지", font_size=14,
                          color=C_BLUE, align=PP_ALIGN.CENTER)

    # 3개 과제 카드
    challenges = [
        ("과제 1", "팬 데이터가 파편화",
         "• 채팅·굿즈·커뮤니티 데이터가 통합 분석 없이 분리\n• '가장 가치있는 팬'을 알 수 없음\n• 신규 IP 추가 기준이 데이터 없이 감(感)으로 결정"),
        ("과제 2", "마케팅 개인화 부재",
         "• 헤비유저와 이탈 위험 유저에게 동일한 푸시 발송\n• 굿즈 구매 고관여 팬 vs 채팅 팬 구분 없음\n• 개인화 캠페인 부재 → 광고 효율 낮음"),
        ("과제 3", "리포팅이 수작업",
         "• KPI 대시보드 없거나 엑셀 기반 운영\n• 경영진 보고 자료에 매주 5~10시간 소요\n• 글로벌 확장 시 더 복잡해질 전망"),
    ]

    card_w = Cm(10.2)
    for i, (badge_text, title, desc) in enumerate(challenges):
        card_x = Cm(1.2) + i * Cm(10.8)
        card_y = Cm(8)
        card = add_rect(slide, card_x, card_y, card_w, Cm(9.8), fill_color=C_LIGHT_PURP,
                         line_color=C_BORDER, line_width=Pt(0.5))
        # Red badge
        badge = add_rect(slide, card_x + Cm(0.3), card_y + Cm(0.3), Cm(2), Cm(0.7),
                          fill_color=C_RED)
        add_text_in_shape(badge, badge_text, font_size=11, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        # Title
        add_textbox(slide, card_x + Cm(0.3), card_y + Cm(1.2), card_w - Cm(0.6), Cm(0.8),
                    title, font_size=14, bold=True, color=C_DEEP_PURPLE)
        # Description
        add_textbox(slide, card_x + Cm(0.3), card_y + Cm(2.2), card_w - Cm(0.6), Cm(7),
                    desc, font_size=12, color=C_DARK)


def build_slide4(slide):
    """위어드섹터 소개"""
    add_header_bar(slide, "위어드섹터 소개")

    # 좌측: 프로파일 테이블
    profile_data = [
        ("설립", "2020년 (벤처기업 인증)"),
        ("대표", "최재우"),
        ("핵심 역량", "데이터 분석 · 마테크 · 생성형 AI"),
        ("고객사", "1,500개+ 기업고객 운영 경험"),
        ("주요 고객", "LG U+, 카카오페이, 무신사, NHN 한게임"),
    ]
    table_y = Cm(2)
    for i, (k, v) in enumerate(profile_data):
        row_y = table_y + i * Cm(1.15)
        bg = C_LIGHT_PURP if i % 2 == 0 else C_TABLE_EVEN
        key_cell = add_rect(slide, Cm(1.2), row_y, Cm(4.5), Cm(1.1), fill_color=C_DEEP_PURPLE)
        add_text_in_shape(key_cell, k, font_size=12, bold=True, color=C_WHITE)
        val_cell = add_rect(slide, Cm(5.7), row_y, Cm(7.5), Cm(1.1), fill_color=bg,
                             line_color=C_BORDER, line_width=Pt(0.5))
        add_text_in_shape(val_cell, v, font_size=12, color=C_DARK)

    # 좌측 아래: 팀 이미지
    img_path = os.path.join(IMG_DIR, "slide4_team.jpg")
    if not add_image(slide, img_path, Cm(1.2), Cm(7.8), Cm(12), Cm(5.5)):
        ph = add_rect(slide, Cm(1.2), Cm(7.8), Cm(12), Cm(5.5), fill_color=C_LIGHT_BLUE)
        add_text_in_shape(ph, "팀 이미지", font_size=14, color=C_BLUE, align=PP_ALIGN.CENTER)

    # 우측: 3가지 핵심 역량 카드
    competencies = [
        ("①", "데이터 인프라 구축",
         "GA4 · Amplitude · Firebase · BigQuery 전 스택 운영\n퍼널 · 리텐션 · LTV · 세그먼트 분석 실무 경험\nData Nugget SaaS로 실시간 대시보드 운영 중"),
        ("②", "AI 기반 자동화",
         "생성형 AI로 리포트 자동 생성 → 운영 비용 절감\n사용자 행동 패턴 예측 모델 구축 경험\n추천 알고리즘 · 개인화 마케팅 인프라 설계"),
        ("③", "운영 파트너십",
         "단순 분석 대행이 아닌 도구화·자동화·제품화 접근\n고객사 내부 데이터 역량 강화 교육 병행\nM/M 방식으로 우선순위 중심 실행, 매월 결과 공유"),
    ]

    for i, (num, title, desc) in enumerate(competencies):
        card_y = Cm(1.9) + i * Cm(4.5)
        # 좌측 보라 세로선
        add_rect(slide, Cm(14.5), card_y, Cm(0.25), Cm(4.2), fill_color=C_DEEP_PURPLE)
        card = add_rect(slide, Cm(14.75), card_y, Cm(17.9), Cm(4.2),
                         fill_color=C_LIGHT_PURP, line_color=C_BORDER, line_width=Pt(0.5))
        # 번호 배지
        badge = add_rect(slide, Cm(15), card_y + Cm(0.3), Cm(0.85), Cm(0.85),
                          fill_color=C_DEEP_PURPLE)
        add_text_in_shape(badge, num, font_size=12, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        add_textbox(slide, Cm(16.1), card_y + Cm(0.3), Cm(16), Cm(0.8),
                    title, font_size=16, bold=True, color=C_DEEP_PURPLE)
        add_textbox(slide, Cm(15.1), card_y + Cm(1.3), Cm(17.3), Cm(2.7),
                    desc, font_size=12, color=C_DARK)


def build_slide5(slide):
    """서비스 ① — 데이터 관리 & 코드 관리"""
    add_header_bar(slide, "제안 서비스 ① — 데이터 관리 & 코드 관리")

    half_w = Cm(15.5)

    # 좌측: 데이터 관리
    title_l = add_rect(slide, Cm(0.5), Cm(2), Cm(15.5), Cm(1.0), fill_color=C_BLUE)
    add_text_in_shape(title_l, "📊  데이터 관리", font_size=16, bold=True,
                      color=C_WHITE, align=PP_ALIGN.CENTER)

    content_l = [
        ("수집 인프라 구축", [
            "• 모잇 앱/웹 사용자 이벤트 택소노미(Event Taxonomy) 설계",
            "• Firebase / Amplitude 이벤트 수집 체계 표준화",
            "• BigQuery 연동 데이터 레이크 파이프라인 구축",
        ]),
        ("분석 체계 수립", [
            "• 핵심 KPI 정의: DAU · MAU · 채팅 횟수 · 굿즈 전환율 · ARPU",
            "• 팬 세그먼트 분류 (VIP / 일반 / 잠재 이탈)",
            "• 주간·월간 자동화 리포트 운영",
        ]),
        ("대시보드 구축", [
            "• 경영진용: 매출·사용자·캐릭터별 성과 요약",
            "• 마케팅팀용: 캠페인 효과, 팬 유입 경로, 전환 퍼널",
            "• 개발팀용: 기능별 사용률, 오류 발생 패턴",
        ]),
    ]

    y_pos = Cm(3.2)
    for section_title, bullets in content_l:
        add_textbox(slide, Cm(0.7), y_pos, Cm(15), Cm(0.7),
                    section_title, font_size=13, bold=True, color=C_BLUE)
        y_pos += Cm(0.7)
        for bullet in bullets:
            add_textbox(slide, Cm(0.7), y_pos, Cm(15), Cm(0.6),
                        bullet, font_size=11, color=C_DARK)
            y_pos += Cm(0.58)
        y_pos += Cm(0.2)

    # 중앙 골드 구분선
    add_rect(slide, Cm(16.5), Cm(2), Cm(0.08), Cm(14.8), fill_color=C_GOLD)

    # 우측: 코드 관리
    title_r = add_rect(slide, Cm(16.8), Cm(2), Cm(15.8), Cm(1.0), fill_color=C_DEEP_PURPLE)
    add_text_in_shape(title_r, "🏷️  코드 관리", font_size=16, bold=True,
                      color=C_WHITE, align=PP_ALIGN.CENTER)

    content_r = [
        ("태그 매니저(GTM) 구축", [
            "• Google Tag Manager 세팅 및 최적화",
            "• 이벤트 트래킹 코드 구조화 및 문서화",
            "• 마케팅 픽셀 통합 관리 (메타, 구글, 카카오)",
        ]),
        ("측정 인프라 유지보수", [
            "• 신규 기능 배포 시 트래킹 코드 검수 프로세스 수립",
            "• 데이터 품질 모니터링: 누락 이벤트, 중복 수집 감지",
            "• SDK 버전 관리 및 업데이트 대응",
        ]),
    ]

    y_pos = Cm(3.2)
    for section_title, bullets in content_r:
        add_textbox(slide, Cm(17), y_pos, Cm(15.5), Cm(0.7),
                    section_title, font_size=13, bold=True, color=C_DEEP_PURPLE)
        y_pos += Cm(0.7)
        for bullet in bullets:
            add_textbox(slide, Cm(17), y_pos, Cm(15.5), Cm(0.6),
                        bullet, font_size=11, color=C_DARK)
            y_pos += Cm(0.58)
        y_pos += Cm(0.2)

    # 하단 골드 산출물 바
    gold_bar = add_rect(slide, 0, Cm(17.2), SLIDE_W, Cm(1.5), fill_color=C_GOLD)
    add_text_in_shape(gold_bar,
                      "예상 산출물: 이벤트 택소노미 문서 + 실시간 대시보드 3종 (경영진/마케팅/개발) + 태그 구조도",
                      font_size=13, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)


def build_slide6(slide):
    """서비스 ② — 마테크 & CRM"""
    add_header_bar(slide, "제안 서비스 ② — 마테크 & CRM")

    # 좌측: 마테크
    title_l = add_rect(slide, Cm(0.5), Cm(2), Cm(15.5), Cm(1.0), fill_color=C_BLUE)
    add_text_in_shape(title_l, "🤖  마테크 (MarTech)", font_size=16, bold=True,
                      color=C_WHITE, align=PP_ALIGN.CENTER)

    content_martech = [
        ("마케팅 자동화 도구 도입", [
            "• CRM 연동: 팬 행동 트리거 기반 자동 메시지 발송",
            '  예: "신규 IP 추가 시 관련 팬에게 자동 알림"',
            '  예: "3일 미접속 팬에게 재방문 유도 푸시"',
            "• A/B 테스트 인프라: 신기능 효과 검증",
            "• 광고 캠페인 자동 최적화: 전환율 기반 예산 배분",
        ]),
        ("개인화 마케팅 구현", [
            "• 팬 프로파일 기반 맞춤 콘텐츠 추천",
            "• 굿즈 구매 이력 → 관심 IP 기반 크로스셀링",
            "• 팬 생일/가입 기념일 자동화 메시지",
        ]),
    ]

    y_pos = Cm(3.2)
    for section_title, bullets in content_martech:
        add_textbox(slide, Cm(0.7), y_pos, Cm(15), Cm(0.7),
                    section_title, font_size=13, bold=True, color=C_BLUE)
        y_pos += Cm(0.7)
        for bullet in bullets:
            add_textbox(slide, Cm(0.7), y_pos, Cm(15), Cm(0.6),
                        bullet, font_size=11, color=C_DARK)
            y_pos += Cm(0.58)
        y_pos += Cm(0.2)

    # 중앙 구분선
    add_rect(slide, Cm(16.5), Cm(2), Cm(0.08), Cm(14.8), fill_color=C_GOLD)

    # 우측: CRM
    title_r = add_rect(slide, Cm(16.8), Cm(2), Cm(15.8), Cm(1.0), fill_color=C_DEEP_PURPLE)
    add_text_in_shape(title_r, "👥  CRM", font_size=16, bold=True,
                      color=C_WHITE, align=PP_ALIGN.CENTER)

    content_crm = [
        ("팬 세그먼테이션", [
            "• RFM 분석: 최근성(R) · 빈도(F) · 구매금액(M)",
            "• VIP 팬 식별: 상위 10% 팬 특성 분석 및 우대 프로그램",
            "• 이탈 위험 탐지: 미접속 기간 · 채팅 빈도 감소 패턴",
        ]),
        ("리텐션 캠페인", [
            "• 이탈 위험 팬 자동 감지 → 맞춤 리텐션 오퍼 발송",
            "• 복귀 팬 온보딩 시퀀스 자동화",
            "• 계절별·이벤트별 팬 리인게이지먼트 캠페인 운영",
        ]),
        ("LTV 극대화", [
            "• 팬 생애주기(Fan Lifecycle) 단계별 전략 수립",
            "• 무료 → 유료 전환 최적 타이밍 분석",
            "• 굿즈 구매 예측 모델 구축",
        ]),
    ]

    y_pos = Cm(3.2)
    for section_title, bullets in content_crm:
        add_textbox(slide, Cm(17), y_pos, Cm(15.5), Cm(0.7),
                    section_title, font_size=13, bold=True, color=C_DEEP_PURPLE)
        y_pos += Cm(0.7)
        for bullet in bullets:
            add_textbox(slide, Cm(17), y_pos, Cm(15.5), Cm(0.6),
                        bullet, font_size=11, color=C_DARK)
            y_pos += Cm(0.58)
        y_pos += Cm(0.2)

    # 하단 핵심 활용 예시
    ex_y = Cm(16.2)
    add_textbox(slide, Cm(1), ex_y - Cm(0.5), Cm(31), Cm(0.6),
                "핵심 활용 예시", font_size=12, bold=True, color=C_DARK)
    ex1 = add_rect(slide, Cm(1), ex_y, Cm(15), Cm(1.5),
                   fill_color=C_LIGHT_BLUE, line_color=C_BLUE, line_width=Pt(1))
    add_text_in_shape(ex1, "재방문 유도: 3일 미접속 팬 → 자동 맞춤 푸시 발송",
                      font_size=12, color=C_BLUE)
    ex2 = add_rect(slide, Cm(17), ex_y, Cm(15.8), Cm(1.5),
                   fill_color=C_LIGHT_PURP, line_color=C_DEEP_PURPLE, line_width=Pt(1))
    add_text_in_shape(ex2, "VIP 팬 굿즈 추천: 구매 이력 기반 크로스셀링 자동화",
                      font_size=12, color=C_DEEP_PURPLE)


def build_slide7(slide):
    """M/M 운영 방식"""
    add_header_bar(slide, "M/M 운영 방식 — 月 200만원으로 어떻게 운영되나요?")

    # 4단계 사이클 다이어그램
    steps = [
        ("①", "우선순위\n협의"),
        ("②", "실행"),
        ("③", "중간\n체크"),
        ("④", "월말\n리뷰"),
    ]
    circle_y = Cm(2.5)
    circle_r = Cm(1.5)
    step_x_starts = [Cm(2), Cm(9.5), Cm(17), Cm(24.5)]

    for i, (num, label) in enumerate(steps):
        cx = step_x_starts[i]
        circle = add_rect(slide, cx, circle_y, Cm(3), Cm(2.5), fill_color=C_DEEP_PURPLE)
        sp = circle.element
        # round corners
        add_text_in_shape(circle, f"{num}\n{label}", font_size=13, bold=True,
                          color=C_WHITE, align=PP_ALIGN.CENTER)
        # 화살표
        if i < 3:
            arrow_x = cx + Cm(3.2)
            add_rect(slide, arrow_x, circle_y + Cm(1.0), Cm(2.8), Cm(0.1),
                     fill_color=C_GOLD)
            # arrowhead
            add_textbox(slide, arrow_x + Cm(2.3), circle_y + Cm(0.6), Cm(0.8), Cm(0.8),
                        "▶", font_size=14, color=C_GOLD)

    # 서브타이틀
    add_textbox(slide, Cm(1), Cm(5.4), Cm(31), Cm(0.6),
                "매월 반복 사이클 | Boss(모다이브)와 우선순위 협의 → 3주 실행 → 중간 체크 → 월말 리뷰",
                font_size=11, color=C_SUB_TEXT, align=PP_ALIGN.CENTER)

    # 월별 로드맵 테이블
    add_textbox(slide, Cm(1.2), Cm(6.2), Cm(20), Cm(0.7),
                "초기 4개월 우선순위 로드맵", font_size=15, bold=True, color=C_DARK)

    headers = ["월", "우선순위 과제", "예상 산출물"]
    col_widths = [Cm(2.5), Cm(12), Cm(16.5)]
    rows_data = [
        ("1월", "이벤트 택소노미 설계 + GTM 세팅", "측정 체계 문서 + 태그 구조도"),
        ("2월", "KPI 대시보드 구축 (Looker Studio)", "실시간 대시보드 3종 (경영진/마케팅/개발)"),
        ("3월", "팬 세그먼테이션 + RFM 분석", "팬 등급 체계 + 세그먼트 프로파일"),
        ("4월", "리텐션 자동화 첫 캠페인", "이탈 위험 감지 + 자동 메시지 발송 시스템"),
    ]
    table_y = Cm(7.1)
    x_start = Cm(1.2)
    # header row
    x = x_start
    for h, w in zip(headers, col_widths):
        cell = add_rect(slide, x, table_y, w, Cm(0.8), fill_color=C_DEEP_PURPLE)
        add_text_in_shape(cell, h, font_size=13, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        x += w
    for i, (m, task, output) in enumerate(rows_data):
        row_y = table_y + Cm(0.8) + i * Cm(1.0)
        bg = C_LIGHT_PURP if i % 2 == 0 else C_TABLE_EVEN
        x = x_start
        for val, w in zip([m, task, output], col_widths):
            cell = add_rect(slide, x, row_y, w, Cm(1.0), fill_color=bg,
                             line_color=C_BORDER, line_width=Pt(0.5))
            is_center = (w == col_widths[0])
            add_text_in_shape(cell, val, font_size=12, color=C_DARK,
                              align=PP_ALIGN.CENTER if is_center else PP_ALIGN.LEFT)
            x += w

    # 운영 구조 사이드박스
    side = add_rect(slide, Cm(1.2), Cm(12.3), Cm(31.5), Cm(5.5),
                    fill_color=C_LIGHT_BLUE, line_color=C_BLUE, line_width=Pt(1))
    tf = side.text_frame
    tf.word_wrap = True
    add_text_in_shape(side, "운영 구조", font_size=14, bold=True, color=C_BLUE)
    add_paragraph(tf, "위어드섹터 담당팀: PM(일정·산출물 관리) + 데이터 엔지니어(수집·파이프라인) + 마케팅 애널리스트(분석·인사이트)",
                  font_size=12, color=C_DARK)
    add_paragraph(tf, "계약 조건: 月 200만원 (VAT 별도) · 3개월 최소 약정 · 긴급 대응 月 +30만원",
                  font_size=12, color=C_DARK)
    add_paragraph(tf, "소통 채널: Slack 실시간 연동 + 월 2회 화상 미팅 + 월말 성과 리포트",
                  font_size=12, color=C_DARK)


def build_slide8(slide):
    """기대 효과 & ROI"""
    add_header_bar(slide, "기대 효과 & ROI")

    # 상단 좌측: 이미지
    img_path = os.path.join(IMG_DIR, "slide8_growth.jpg")
    if not add_image(slide, img_path, Cm(1.2), Cm(2), Cm(11), Cm(6)):
        ph = add_rect(slide, Cm(1.2), Cm(2), Cm(11), Cm(6), fill_color=C_LIGHT_BLUE)
        add_text_in_shape(ph, "성장 차트 이미지", font_size=14, color=C_BLUE, align=PP_ALIGN.CENTER)

    # 상단 우측: KPI 기대효과 테이블
    add_textbox(slide, Cm(13), Cm(1.8), Cm(20), Cm(0.7),
                "KPI별 기대 효과 (업계 평균 기준 추정)", font_size=13, bold=True, color=C_DARK)
    headers = ["KPI", "3개월 후", "6개월 후"]
    col_widths = [Cm(7), Cm(5.5), Cm(5.5)]
    kpi_rows = [
        ("팬 이탈률", "-15% 감소", "-30% 감소"),
        ("굿즈 전환율", "+30% 향상", "+100% 향상"),
        ("리포팅 시간", "-50% 절감", "-80% (자동화)"),
        ("VIP 팬 식별", "상위 10% 가능", "예측 포함 고도화"),
        ("캠페인 CTR", "+20% 향상", "+50% (개인화)"),
    ]
    table_y = Cm(2.6)
    x_start = Cm(13)
    x = x_start
    for h, w in zip(headers, col_widths):
        cell = add_rect(slide, x, table_y, w, Cm(0.75), fill_color=C_DEEP_PURPLE)
        add_text_in_shape(cell, h, font_size=12, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        x += w
    for i, row in enumerate(kpi_rows):
        row_y = table_y + Cm(0.75) + i * Cm(0.95)
        x = x_start
        for j, (val, w) in enumerate(zip(row, col_widths)):
            if j == 1:
                bg = C_LIGHT_BLUE
            elif j == 2:
                bg = C_LIGHT_PURP
            else:
                bg = C_TABLE_EVEN if i % 2 == 0 else C_WHITE
            cell = add_rect(slide, x, row_y, w, Cm(0.9), fill_color=bg,
                             line_color=C_BORDER, line_width=Pt(0.5))
            is_green = j > 0
            add_text_in_shape(cell, val, font_size=12,
                              bold=(j > 0), color=C_GREEN if is_green else C_DARK,
                              align=PP_ALIGN.CENTER)
            x += w

    # 하단 ROI 계산 박스
    roi_y = Cm(8.5)
    add_textbox(slide, Cm(1.2), roi_y, Cm(20), Cm(0.7),
                "ROI 계산", font_size=15, bold=True, color=C_DARK)

    # 비용 박스
    cost_box = add_rect(slide, Cm(1.2), roi_y + Cm(0.9), Cm(8), Cm(6),
                         fill_color=C_DEEP_PURPLE)
    tf = cost_box.text_frame
    tf.word_wrap = True
    add_text_in_shape(cost_box, "연간 투자 비용", font_size=13, bold=True, color=C_PURPLE_LIGHT2)
    add_paragraph(tf, "연 2,400만원", font_size=32, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_paragraph(tf, "(月 200만원 × 12개월)", font_size=11, color=C_PURPLE_LIGHT2)

    # 기대 효과 박스
    effect_box = add_rect(slide, Cm(10), roi_y + Cm(0.9), Cm(22.8), Cm(6),
                           fill_color=C_LIGHT_PURP, line_color=C_DEEP_PURPLE, line_width=Pt(1))
    tf = effect_box.text_frame
    tf.word_wrap = True
    add_text_in_shape(effect_box, "절감·증가 효과 (연간 추정)", font_size=13, bold=True,
                      color=C_DEEP_PURPLE)
    add_paragraph(tf, "• 마케팅팀 리포팅 시간 절감 (주 5시간 × 인건비): 연 720만원+",
                  font_size=12, color=C_DARK)
    add_paragraph(tf, "• 이탈 팬 리텐션으로 유료 구독 유지: 연 600만원+",
                  font_size=12, color=C_DARK)
    add_paragraph(tf, "• 굿즈 전환율 향상 (+30% × 기존 매출): 플랫폼 규모에 비례",
                  font_size=12, color=C_DARK)
    add_paragraph(tf, "", font_size=8, color=C_WHITE)
    add_paragraph(tf, "투자금의 5.5배 이상 (보수적 추정)",
                  font_size=20, bold=True, color=C_GOLD)


def build_slide9(slide):
    """왜 위어드섹터인가?"""
    add_header_bar(slide, "왜 위어드섹터인가?")

    cards = [
        ("🎬", "미디어·엔터 업계 경험", C_DEEP_PURPLE,
         "MBC 모다이브 기존 협력 이력 보유 → 맥락 이해 빠름\nNHN 한게임·스마일게이트 등 게임/엔터 플랫폼 경험\n팬덤 플랫폼의 특수성(감정 연관도, 이벤트 급등) 이해"),
        ("🔧", "실행 파트너십", C_BLUE,
         "분석 리포트만 납품하는 것이 아닌 도구를 직접 구축·운영\nData Nugget SaaS로 안정적 실시간 대시보드 제공\nAI 기반 자동 리포트로 운영 비용 지속 절감"),
        ("🏆", "1,500개 기업 운영 노하우", C_GOLD,
         "B2B·B2C·앱·웹 다양한 플랫폼 분석 경험 보유\nGA4·Amplitude·Firebase·BigQuery 全 스택 실무 운영\n이탈 예측·추천 시스템·LTV 모델 실제 구축 사례 다수"),
        ("✅", "기술 최신성 & 정부 공인", C_GREEN,
         "AI AGENT 융합·확산 지원사업 (NIA 2026) 수행사 선정\n벤처기업 인증 · 기업부설연구소 보유\n생성형 AI 기반 자동화 리포팅·추천 시스템 자체 개발"),
    ]

    positions = [
        (Cm(1.2), Cm(2.0)),
        (Cm(17.5), Cm(2.0)),
        (Cm(1.2), Cm(10.5)),
        (Cm(17.5), Cm(10.5)),
    ]

    for (icon, title, accent, desc), (cx, cy) in zip(cards, positions):
        card = add_rect(slide, cx, cy, Cm(15), Cm(7.5),
                         fill_color=C_WHITE, line_color=accent, line_width=Pt(1.5))
        # 아이콘
        add_textbox(slide, cx + Cm(0.4), cy + Cm(0.3), Cm(1.5), Cm(1.5),
                    icon, font_size=26, color=accent)
        # 제목
        add_textbox(slide, cx + Cm(2), cy + Cm(0.4), Cm(12.5), Cm(0.9),
                    title, font_size=15, bold=True, color=accent)
        # 설명
        add_textbox(slide, cx + Cm(0.4), cy + Cm(1.5), Cm(14.2), Cm(5.5),
                    desc, font_size=12, color=C_DARK)

    # 하단 주요 고객사
    client_bar = add_rect(slide, 0, Cm(18.0), SLIDE_W, Cm(0.75), fill_color=C_DARK)
    add_text_in_shape(client_bar,
                      "LG U+  ·  카카오페이  ·  무신사  ·  NHN 한게임  ·  MBC 모다이브",
                      font_size=13, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)


def build_slide10(slide):
    """계약 조건 & 다음 단계"""
    add_header_bar(slide, "계약 조건 & 다음 단계")

    # 좌측: 패키지 구성 테이블
    add_textbox(slide, Cm(1.2), Cm(1.9), Cm(15), Cm(0.7),
                "패키지 구성", font_size=15, bold=True, color=C_DARK)

    pkg_rows = [
        ("월 계약금", "200만원 (VAT 별도)"),
        ("운영 방식", "M/M 기반 — 우선순위 중심 실행"),
        ("최소 약정", "3개월 (이후 월 단위 자동 갱신)"),
        ("담당팀", "PM 1명 + 데이터 엔지니어 + 애널리스트"),
        ("소통 채널", "Slack + 월 2회 화상 미팅"),
        ("산출물", "월말 리포트 + 대시보드 + 태스크 완료 보고"),
    ]

    table_y = Cm(2.7)
    for i, (k, v) in enumerate(pkg_rows):
        row_y = table_y + i * Cm(1.2)
        is_gold = (k == "월 계약금")
        key_cell = add_rect(slide, Cm(1.2), row_y, Cm(5), Cm(1.1),
                             fill_color=C_GOLD if is_gold else C_DEEP_PURPLE)
        add_text_in_shape(key_cell, k, font_size=12, bold=True, color=C_WHITE)
        val_cell = add_rect(slide, Cm(6.2), row_y, Cm(9.5), Cm(1.1),
                             fill_color=RGBColor(0xFF, 0xF7, 0xE0) if is_gold else C_LIGHT_PURP,
                             line_color=C_BORDER, line_width=Pt(0.5))
        add_text_in_shape(val_cell, v, font_size=12,
                          bold=is_gold, color=C_GOLD if is_gold else C_DARK)

    # 우측: 계약 프로세스
    add_textbox(slide, Cm(17.5), Cm(1.9), Cm(15), Cm(0.7),
                "계약 프로세스", font_size=15, bold=True, color=C_DARK)

    steps_proc = [
        ("STEP 1", "킥오프 미팅 (1~2시간)",
         "현황 파악 · 우선순위 협의 · 접근 권한 공유"),
        ("STEP 2", "계약서 서명 (이후 1주일)",
         "서비스 범위 · 산출물 · 기간 확정"),
        ("STEP 3", "온보딩 (1주일)",
         "Slack 채널 세팅 · 기존 데이터 현황 감사(Audit)"),
        ("STEP 4", "1차 스프린트 시작",
         "우선순위 1위 과제부터 즉시 실행"),
    ]

    proc_y = Cm(2.7)
    for i, (step, title, desc) in enumerate(steps_proc):
        y = proc_y + i * Cm(3.5)
        badge = add_rect(slide, Cm(17.5), y, Cm(2), Cm(0.8), fill_color=C_DEEP_PURPLE)
        add_text_in_shape(badge, step, font_size=11, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        add_textbox(slide, Cm(20), y, Cm(12), Cm(0.8),
                    title, font_size=13, bold=True, color=C_DEEP_PURPLE)
        add_textbox(slide, Cm(17.5), y + Cm(0.9), Cm(14.5), Cm(1.5),
                    desc, font_size=12, color=C_DARK)
        if i < 3:
            add_rect(slide, Cm(18.3), y + Cm(2.5), Cm(0.1), Cm(0.9), fill_color=C_DEEP_PURPLE)

    # 하단 CTA 박스
    cta = add_rect(slide, 0, Cm(16.8), SLIDE_W, Cm(2.0), fill_color=C_DEEP_PURPLE)
    tf = cta.text_frame
    tf.word_wrap = True
    add_text_in_shape(cta,
                      "지금 바로 킥오프 미팅을 잡으세요.",
                      font_size=16, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_paragraph(tf, "미팅 1회로 현재 데이터 현황을 무료로 점검해 드립니다.",
                  font_size=13, color=C_PURPLE_LIGHT2, align=PP_ALIGN.CENTER)


def build_slide11(slide):
    """Q&A / 마무리"""
    # 전체 배경 그라데이션
    bg = add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill_color=C_DEEP_PURPLE)
    set_gradient_fill(bg, C_DEEP_PURPLE, C_BLUE, angle=5400000)

    # 감사합니다
    add_textbox(slide, Cm(2), Cm(1.5), Cm(29.8), Cm(2),
                "감사합니다",
                font_size=40, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    # 인용 박스
    quote_box = add_rect(slide, Cm(3), Cm(4), Cm(27.8), Cm(5.5),
                          fill_color=RGBColor(0x00, 0x2E, 0x80))
    tf = quote_box.text_frame
    tf.word_wrap = True
    add_text_in_shape(quote_box,
                      "모잇은 팬들에게 AI로 새로운 경험을 주고 있습니다.",
                      font_size=16, italic=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_paragraph(tf,
                  "위어드섹터는 그 경험 뒤에서 데이터가 말하는 것을 듣고, 실행하는 팀이 되겠습니다.",
                  font_size=16, italic=True, color=C_PURPLE_LIGHT2, align=PP_ALIGN.CENTER)
    add_paragraph(tf,
                  "팬 한 명 한 명을 더 잘 이해하고, 더 오래 함께할 수 있도록.",
                  font_size=14, italic=True, color=C_PURPLE_LIGHT2, align=PP_ALIGN.CENTER)

    # 연락처 테이블
    contact_y = Cm(10.2)
    add_textbox(slide, Cm(3), contact_y - Cm(0.5), Cm(27), Cm(0.6),
                "연락처", font_size=13, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    contacts = [
        ("이메일", "info@weirdsector.co.kr"),
        ("전화", "010-6527-1218"),
        ("웹사이트", "weirdsector.co.kr"),
        ("슬랙", "요청 시 공유 워크스페이스 개설"),
    ]
    for i, (label, val) in enumerate(contacts):
        row_y = contact_y + i * Cm(1.2)
        label_cell = add_rect(slide, Cm(5), row_y, Cm(5), Cm(1.0),
                               fill_color=RGBColor(0x00, 0x2E, 0x80),
                               line_color=C_PURPLE_LIGHT2, line_width=Pt(0.75))
        add_text_in_shape(label_cell, label, font_size=13, bold=True,
                          color=C_PURPLE_LIGHT2, align=PP_ALIGN.CENTER)
        val_cell = add_rect(slide, Cm(10), row_y, Cm(16), Cm(1.0),
                             fill_color=RGBColor(0x00, 0x2E, 0x80),
                             line_color=C_PURPLE_LIGHT2, line_width=Pt(0.75))
        add_text_in_shape(val_cell, val, font_size=13, color=C_WHITE)

    # 하단 기밀 표시
    add_textbox(slide, Cm(1), Cm(17.8), Cm(31), Cm(0.6),
                "이 제안서는 기밀 자료입니다. 무단 배포를 금합니다.  |  Weird Sector © 2026",
                font_size=10, italic=True, color=C_PURPLE_LIGHT2, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════
# MAIN BUILD
# ════════════════════════════════════════════════════════

def build_presentation():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    blank_layout = prs.slide_layouts[6]  # Blank layout

    builders = [
        build_slide1,
        build_slide2,
        build_slide3,
        build_slide4,
        build_slide5,
        build_slide6,
        build_slide7,
        build_slide8,
        build_slide9,
        build_slide10,
        build_slide11,
    ]

    for i, builder in enumerate(builders, 1):
        slide = prs.slides.add_slide(blank_layout)
        try:
            builder(slide)
            print(f"  [OK] Slide {i} built: {builder.__name__}")
        except Exception as e:
            print(f"  [ERR] Slide {i} error ({builder.__name__}): {e}")
            import traceback
            traceback.print_exc()

    prs.save(OUT_PATH)
    print(f"\n[DONE] Saved: {OUT_PATH}")
    return OUT_PATH


if __name__ == "__main__":
    print("Building MBC Modive x WeirsSector PPTX...")
    build_presentation()
