# -*- coding: utf-8 -*-
"""
위어드섹터 회사소개서 v2 (빌드 v3)
- 이모지 완전 제거
- 가격/요금제 슬라이드 삭제
- 인포그래픽 4장 삽입 (프로세스, 역량, 성과 수치, 레퍼런스)
- 15슬라이드 풀블리드 야경 감성
"""
import sys, io
if sys.platform == 'win32':
    sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree
from pathlib import Path
from PIL import Image
import os

ASSETS = Path("C:/Agent/pepper/output/wirdsector_assets")
OUT    = Path("C:/Agent/pepper/output/wirdsector_company_profile_v2.pptx")

# ── 컬러 ──────────────────────────────────────────
C_NAVY    = RGBColor(0x0D, 0x1B, 0x2A)
C_NAVY2   = RGBColor(0x06, 0x0E, 0x18)
C_TEAL    = RGBColor(0x00, 0xD4, 0xAA)
C_WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
C_GRAY    = RGBColor(0xB0, 0xBC, 0xCC)
C_GOLD    = RGBColor(0xFF, 0xC8, 0x4B)
C_PURPLE  = RGBColor(0x9B, 0x7A, 0xFF)
C_DARK    = RGBColor(0x00, 0x00, 0x00)
C_RED     = RGBColor(0xFF, 0x6B, 0x6B)

SW = Inches(13.33)
SH = Inches(7.5)
TOTAL = 15

def new_prs():
    prs = Presentation()
    prs.slide_width  = SW
    prs.slide_height = SH
    return prs

def blank(prs):
    return prs.slide_layouts[6]

# ── 이미지 헬퍼 ────────────────────────────────────
def bg_img(slide, img_path):
    p = Path(img_path)
    if not p.exists():
        return
    try:
        slide.shapes.add_picture(str(p), Inches(0), Inches(0), SW, SH)
    except Exception as e:
        print(f"  [WARN] bg_img: {e}")

def overlay(slide, x, y, w, h, r=0, g=0, b=0, alpha=65):
    sp = slide.shapes.add_shape(1, x, y, w, h)
    sp.line.fill.background()
    sp.fill.solid()
    sp.fill.fore_color.rgb = RGBColor(r, g, b)
    spPr = sp.element.spPr if hasattr(sp.element, 'spPr') else sp.element.find(qn('p:spPr'))
    sf = spPr.find('.//' + qn('a:solidFill'))
    if sf is not None:
        clr = sf.find(qn('a:srgbClr'))
        if clr is not None:
            a = etree.SubElement(clr, qn('a:alpha'))
            a.set('val', str(int(alpha * 1000)))
    return sp

def rrect(slide, x, y, w, h, fill=None, alpha=100):
    sp = slide.shapes.add_shape(5, x, y, w, h)
    sp.line.fill.background()
    if fill:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    else:
        sp.fill.background()
    if alpha < 100 and fill:
        spPr = sp.element.spPr if hasattr(sp.element, 'spPr') else sp.element.find(qn('p:spPr'))
        sf = spPr.find('.//' + qn('a:solidFill'))
        if sf is not None:
            clr = sf.find(qn('a:srgbClr'))
            if clr is not None:
                a = etree.SubElement(clr, qn('a:alpha'))
                a.set('val', str(int(alpha * 1000)))
    return sp

def rect(slide, x, y, w, h, fill=None):
    sp = slide.shapes.add_shape(1, x, y, w, h)
    sp.line.fill.background()
    if fill:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    else:
        sp.fill.background()
    return sp

def txt(slide, text, x, y, w, h, sz=20, bold=False, color=None,
        align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p  = tf.paragraphs[0]
    p.alignment = align
    r  = p.add_run()
    r.text  = text
    r.font.size   = Pt(sz)
    r.font.bold   = bold
    r.font.italic = italic
    r.font.name   = "Malgun Gothic"
    if color:
        r.font.color.rgb = color
    return tb

def slide_num(slide, n):
    overlay(slide, Inches(12.5), Inches(7.1), Inches(0.83), Inches(0.35), 0, 0, 0, 40)
    txt(slide, f"{n:02d} / {TOTAL:02d}",
        Inches(12.5), Inches(7.12), Inches(0.8), Inches(0.3),
        sz=11, color=C_GRAY, align=PP_ALIGN.CENTER)

def logo_badge(slide):
    rrect(slide, Inches(0.4), Inches(0.3), Inches(2.5), Inches(0.55), fill=C_TEAL, alpha=90)
    txt(slide, "Weirdsector",
        Inches(0.45), Inches(0.32), Inches(2.4), Inches(0.5),
        sz=16, bold=True, color=C_NAVY2, align=PP_ALIGN.CENTER)

def teal_line(slide, y=Inches(1.0)):
    rect(slide, Inches(0.5), y, Inches(1.5), Inches(0.05), fill=C_TEAL)

def section_header(slide, en_label, ko_title, subtitle=None):
    """공통 섹션 헤더 (영문 레이블 + 한글 제목 + 밑줄)"""
    txt(slide, en_label,
        Inches(0.5), Inches(0.85), Inches(6.0), Inches(0.4),
        sz=14, color=C_TEAL, bold=True)
    txt(slide, ko_title,
        Inches(0.5), Inches(1.25), Inches(9.5), Inches(0.85),
        sz=44, bold=True, color=C_WHITE)
    teal_line(slide, Inches(2.2))
    if subtitle:
        txt(slide, subtitle,
            Inches(0.5), Inches(2.4), Inches(9.0), Inches(0.5),
            sz=20, color=C_TEAL)

def num_badge(slide, num_str, x, y, c=None):
    """번호 배지 (이모지 대체)"""
    c = c or C_TEAL
    rrect(slide, x, y, Inches(0.6), Inches(0.6), fill=c, alpha=90)
    txt(slide, num_str, x, y+Inches(0.04), Inches(0.6), Inches(0.52),
        sz=16, bold=True, color=C_NAVY2, align=PP_ALIGN.CENTER)

def infographic_slide(prs, img_path, en_label, ko_title, bg_file, n):
    """인포그래픽 이미지를 풀스크린에 가깝게 삽입하는 슬라이드"""
    sl = prs.slides.add_slide(blank(prs))
    bg_img(sl, ASSETS / bg_file)
    overlay(sl, 0, 0, SW, SH, 4, 10, 20, 82)
    logo_badge(sl)
    slide_num(sl, n)

    txt(sl, en_label,
        Inches(0.5), Inches(0.85), Inches(6.0), Inches(0.4),
        sz=14, color=C_TEAL, bold=True)
    txt(sl, ko_title,
        Inches(0.5), Inches(1.25), Inches(9.5), Inches(0.65),
        sz=38, bold=True, color=C_WHITE)
    teal_line(sl, Inches(2.0))

    # 인포그래픽 이미지 삽입
    p = Path(img_path)
    if p.exists():
        # 이미지 크기 확인해서 비율 맞춤
        try:
            im = Image.open(str(p))
            img_w, img_h = im.size
            aspect = img_w / img_h
            # 사용 가능 영역: 12.4" × 5.0" (y=2.2부터 7.1까지)
            avail_w = Inches(12.4)
            avail_h = Inches(5.0)
            if aspect > avail_w / avail_h:
                w = avail_w
                h = Emu(int(avail_w / aspect))
            else:
                h = avail_h
                w = Emu(int(avail_h * aspect))
            left = Emu(int((SW - w) / 2))
            top  = Inches(2.25) + Emu(int((avail_h - h) / 2))
            sl.shapes.add_picture(str(p), left, top, w, h)
        except Exception as e:
            print(f"  [WARN] infographic insert: {e}")
    return sl


# ══════════════════════════════════════════════════
def build():
    prs = new_prs()

    # ─────────────────────────────────────────────
    # Slide 1: 표지
    # ─────────────────────────────────────────────
    sl = prs.slides.add_slide(blank(prs))
    bg_img(sl, ASSETS/"city_city_aerial.jpg")
    overlay(sl, 0, 0, SW, SH, 6, 14, 26, 72)
    overlay(sl, 0, Inches(4.5), SW, Inches(3.0), 0, 0, 0, 55)

    rect(sl, Inches(0.5), Inches(1.5), Inches(0.06), Inches(4.0), fill=C_TEAL)

    txt(sl, "Weirdsector",
        Inches(0.8), Inches(1.5), Inches(9.0), Inches(1.3),
        sz=62, bold=True, color=C_WHITE)
    txt(sl, "위어드섹터",
        Inches(0.8), Inches(2.8), Inches(9.0), Inches(0.8),
        sz=36, bold=False, color=C_TEAL)
    txt(sl, "데이터로 성장을 만드는\n마케팅 테크 파트너",
        Inches(0.8), Inches(3.7), Inches(8.0), Inches(1.3),
        sz=26, color=C_GRAY)

    txt(sl, '"Be yourself, Embrace your weirdness!"',
        Inches(0.8), Inches(5.2), Inches(9.0), Inches(0.55),
        sz=17, italic=True, color=C_TEAL)

    rect(sl, Inches(0.8), Inches(5.0), Inches(5.0), Inches(0.03), fill=C_TEAL)

    for i, tag in enumerate(["Data & Code", "MarTech", "CRM", "Labbit", "DataNugget"]):
        rrect(sl, Inches(0.8 + i*2.35), Inches(6.4), Inches(2.2), Inches(0.45),
              fill=C_NAVY2, alpha=80)
        txt(sl, tag, Inches(0.8 + i*2.35), Inches(6.42), Inches(2.2), Inches(0.42),
            sz=14, color=C_TEAL, align=PP_ALIGN.CENTER)

    txt(sl, "2026", Inches(12.0), Inches(7.0), Inches(1.2), Inches(0.4),
        sz=14, color=C_GRAY, align=PP_ALIGN.RIGHT)
    slide_num(sl, 1)
    print("  [1] 표지")

    # ─────────────────────────────────────────────
    # Slide 2: 회사 소개
    # ─────────────────────────────────────────────
    sl = prs.slides.add_slide(blank(prs))
    bg_img(sl, ASSETS/"city_city_night.jpg")
    overlay(sl, 0, 0, SW, SH, 6, 14, 26, 75)
    logo_badge(sl)
    slide_num(sl, 2)

    section_header(sl, "Company Overview", "위어드섹터 소개",
                   "데이터 기반 마케팅의 실행 파트너")

    txt(sl, "개발력, 마케팅 감각, 데이터 분석 — 세 가지를 한 팀에서 제공합니다.\n단순한 대행사가 아닌, 함께 성장하는 테크 파트너를 경험하세요.",
        Inches(0.5), Inches(3.05), Inches(7.5), Inches(1.1),
        sz=18, color=C_GRAY)

    # 핵심 팩트 카드 (이모지 → 번호 배지)
    facts = [
        ("01", "서울 & 창원", "노원 본사 · 도봉 연구소\n창원 지사 3거점 운영"),
        ("02", "하이브리드 팀", "개발 + 마케팅 + 데이터\n한 팀에서 원스톱 제공"),
        ("03", "자체 서비스", "Labbit · DataNugget\n자체 SaaS 보유 에이전시"),
    ]
    for i, (num, ttl, body) in enumerate(facts):
        bx = Inches(0.5 + i*4.15)
        rrect(sl, bx, Inches(4.5), Inches(3.95), Inches(2.6), fill=C_NAVY2, alpha=80)
        rect(sl, bx, Inches(4.5), Inches(3.95), Inches(0.05), fill=C_TEAL)
        num_badge(sl, num, bx+Inches(0.2), Inches(4.58), C_TEAL)
        txt(sl, ttl, bx+Inches(0.95), Inches(4.65), Inches(2.7), Inches(0.5),
            sz=18, bold=True, color=C_WHITE)
        txt(sl, body, bx+Inches(0.2), Inches(5.22), Inches(3.6), Inches(1.1),
            sz=14, color=C_GRAY)
    print("  [2] 회사 소개")

    # ─────────────────────────────────────────────
    # Slide 3: 문제 정의
    # ─────────────────────────────────────────────
    sl = prs.slides.add_slide(blank(prs))
    bg_img(sl, ASSETS/"city_data_viz.jpg")
    overlay(sl, 0, 0, SW, SH, 4, 10, 20, 80)
    logo_badge(sl)
    slide_num(sl, 3)

    section_header(sl, "Pain Points", "마케팅팀이 겪는 고통점")

    pains = [
        ("01", "데이터가 쌓이는데 쓸 수가 없다", "GA4 이벤트 설계 오류, 트래킹 누락, 데이터 사일로"),
        ("02", "채널별 마케팅이 연결되지 않는다", "Push · Email · SMS 각각 관리, 통합 뷰 없음"),
        ("03", "광고비는 나가는데 ROI가 불분명하다", "성과 측정 체계 부재, 기여도 분석 불가"),
        ("04", "경쟁사는 AI 개인화인데 우린 일괄발송", "사용자 이탈 가속, 전환율 정체"),
        ("05", "개발팀·마케팅팀 커뮤니케이션 비용 폭발", "태그 추가 하나에 2주, 캠페인 지연 반복"),
        ("06", "리텐션 전략 없이 신규 유입만 의존", "CAC 상승, LTV 하락, 성장 정체"),
    ]
    for i, (num, ttl, sub) in enumerate(pains):
        row, col = divmod(i, 3)
        bx = Inches(0.4 + col*4.3)
        ry = Inches(2.55 + row*2.2)
        c  = C_TEAL if row == 0 else C_GOLD
        rrect(sl, bx, ry, Inches(4.1), Inches(2.0), fill=C_NAVY2, alpha=75)
        rect(sl, bx, ry, Inches(0.06), Inches(2.0), fill=c)
        num_badge(sl, num, bx+Inches(0.15), ry+Inches(0.1), c)
        txt(sl, ttl, bx+Inches(0.88), ry+Inches(0.15), Inches(3.0), Inches(0.5),
            sz=15, bold=True, color=C_WHITE)
        txt(sl, sub, bx+Inches(0.2), ry+Inches(0.72), Inches(3.7), Inches(1.1),
            sz=13, color=C_GRAY)
    print("  [3] 문제 정의")

    # ─────────────────────────────────────────────
    # Slide 4: 솔루션 개요
    # ─────────────────────────────────────────────
    sl = prs.slides.add_slide(blank(prs))
    bg_img(sl, ASSETS/"bg_technology.jpg")
    overlay(sl, 0, 0, SW, SH, 4, 10, 20, 82)
    logo_badge(sl)
    slide_num(sl, 4)

    section_header(sl, "Our Solutions", "위어드섹터의 5가지 솔루션")

    solutions = [
        ("01", "데이터 &\n코드 관리", "GA4 이벤트 설계\n로그정의서 GitHub 관리\n트래킹 오류 제거"),
        ("02", "마테크\n(MarTech)", "CleverTap 도입/운영\n마케팅 자동화 플랫폼\n캠페인 플로우 설계"),
        ("03", "CRM 운영", "고객 세그멘테이션\n리텐션 캠페인 설계\n퍼널 분석 & 최적화"),
        ("04", "Labbit", "그로스해킹 파트너\nA/B 테스트 플랫폼\n퍼포먼스 마케팅"),
        ("05", "DataNugget", "데이터 인사이트 자동화\n브랜드 성장 리포트\n커스텀 대시보드"),
    ]
    for i, (num, ttl, body) in enumerate(solutions):
        bx = Inches(0.35 + i*2.6)
        rrect(sl, bx, Inches(2.5), Inches(2.45), Inches(4.65), fill=C_NAVY2, alpha=80)
        rect(sl, bx, Inches(2.5), Inches(2.45), Inches(0.06), fill=C_TEAL)
        num_badge(sl, num, bx+Inches(0.9), Inches(2.6), C_TEAL)
        txt(sl, ttl, bx+Inches(0.1), Inches(3.35), Inches(2.25), Inches(0.75),
            sz=16, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        txt(sl, body, bx+Inches(0.1), Inches(4.15), Inches(2.25), Inches(2.0),
            sz=13, color=C_GRAY, align=PP_ALIGN.CENTER)
    print("  [4] 솔루션 개요")

    # ─────────────────────────────────────────────
    # Slide 5: 서비스 프로세스 인포그래픽 (NEW)
    # ─────────────────────────────────────────────
    infographic_slide(prs,
        ASSETS/"infographic_process.png",
        "Service Process",
        "위어드섹터 서비스 프로세스",
        "bg_technology.jpg", 5)
    print("  [5] 서비스 프로세스 인포그래픽")

    # ─────────────────────────────────────────────
    # Slide 6: 서비스 1 — 데이터 & 코드 관리
    # ─────────────────────────────────────────────
    sl = prs.slides.add_slide(blank(prs))
    bg_img(sl, ASSETS/"bg_data.jpg")
    overlay(sl, 0, 0, SW, SH, 4, 10, 20, 80)
    logo_badge(sl)
    slide_num(sl, 6)

    txt(sl, "Service 01",
        Inches(0.5), Inches(0.85), Inches(4.0), Inches(0.4),
        sz=14, color=C_TEAL, bold=True)
    txt(sl, "데이터 & 코드 관리",
        Inches(0.5), Inches(1.25), Inches(8.0), Inches(0.85),
        sz=44, bold=True, color=C_WHITE)
    teal_line(sl, Inches(2.2))
    txt(sl, "정확한 데이터가 모든 마케팅 의사결정의 기반입니다.",
        Inches(0.5), Inches(2.4), Inches(8.5), Inches(0.55),
        sz=20, color=C_TEAL)

    data_items = [
        ("01", "GA4 이벤트 택소노미 설계",
         "표준화된 이벤트 구조로 분석 정확도를 높이고\n불필요한 데이터 수집을 제거합니다."),
        ("02", "로그정의서 GitHub 기반 관리",
         "마케팅팀과 개발팀이 공유하는 단일 소스로\n커뮤니케이션 오류를 90% 줄입니다."),
        ("03", "GTM & SDK 최적화",
         "Google Tag Manager, AppsFlyer, Firebase 등\n스크립트 통합 설치 및 최적화"),
        ("04", "데이터 품질 감사 (Audit)",
         "기존 트래킹 오류 발견 및 수정\n데이터 신뢰도 복원"),
    ]
    for i, (num, ttl, body) in enumerate(data_items):
        ry = Inches(3.0 + i*1.05)
        rrect(sl, Inches(0.5), ry, Inches(6.8), Inches(0.92), fill=C_NAVY2, alpha=75)
        num_badge(sl, num, Inches(0.65), ry+Inches(0.14), C_TEAL)
        txt(sl, ttl, Inches(1.4), ry+Inches(0.08), Inches(5.6), Inches(0.38),
            sz=16, bold=True, color=C_WHITE)
        txt(sl, body, Inches(1.4), ry+Inches(0.48), Inches(5.6), Inches(0.4),
            sz=12, color=C_GRAY)

    # 우측 핵심 지표
    rrect(sl, Inches(7.8), Inches(2.5), Inches(5.1), Inches(4.6), fill=C_NAVY2, alpha=75)
    rect(sl, Inches(7.8), Inches(2.5), Inches(5.1), Inches(0.05), fill=C_TEAL)
    txt(sl, "도입 효과", Inches(8.0), Inches(2.65), Inches(4.6), Inches(0.5),
        sz=20, bold=True, color=C_TEAL)
    for i, (v, l) in enumerate([
        ("90%+", "커뮤니케이션 오류 감소"),
        ("100%", "데이터 정확도 달성"),
        ("즉시", "시작 가능"),
    ]):
        ry = Inches(3.3 + i*1.2)
        txt(sl, v, Inches(8.0), ry, Inches(2.5), Inches(0.65),
            sz=40, bold=True, color=C_TEAL)
        txt(sl, l, Inches(8.0), ry+Inches(0.62), Inches(4.6), Inches(0.42),
            sz=15, color=C_GRAY)
    print("  [6] 데이터 & 코드")

    # ─────────────────────────────────────────────
    # Slide 7: 서비스 2 — 마테크 & CRM
    # ─────────────────────────────────────────────
    sl = prs.slides.add_slide(blank(prs))
    bg_img(sl, ASSETS/"bg_dashboard.jpg")
    overlay(sl, 0, 0, SW, SH, 4, 10, 20, 80)
    logo_badge(sl)
    slide_num(sl, 7)

    txt(sl, "Service 02 · 03",
        Inches(0.5), Inches(0.85), Inches(6.0), Inches(0.4),
        sz=14, color=C_TEAL, bold=True)
    txt(sl, "마테크 & CRM",
        Inches(0.5), Inches(1.25), Inches(8.0), Inches(0.85),
        sz=44, bold=True, color=C_WHITE)
    teal_line(sl, Inches(2.2))

    # 마테크 카드
    rrect(sl, Inches(0.5), Inches(2.5), Inches(5.9), Inches(4.65), fill=C_NAVY2, alpha=78)
    rect(sl, Inches(0.5), Inches(2.5), Inches(5.9), Inches(0.06), fill=C_TEAL)
    txt(sl, "MarTech 운영",
        Inches(0.7), Inches(2.65), Inches(5.4), Inches(0.55),
        sz=20, bold=True, color=C_TEAL)
    martech = [
        "CleverTap · Braze 등 인게이지먼트 플랫폼 도입",
        "Push · 인앱 · 이메일 · SMS 자동화 플로우 구축",
        "온보딩 / 이탈방지 / Win-back 캠페인 설계",
        "A/B 테스트 기반 메시지 최적화",
        "전환 퍼널 분석 및 개선 사이클 운영",
    ]
    for i, item in enumerate(martech):
        txt(sl, f"▸  {item}",
            Inches(0.7), Inches(3.3 + i*0.68), Inches(5.5), Inches(0.6),
            sz=15, color=C_WHITE if i == 0 else C_GRAY)

    # CRM 카드
    rrect(sl, Inches(6.9), Inches(2.5), Inches(6.0), Inches(4.65), fill=C_NAVY2, alpha=78)
    rect(sl, Inches(6.9), Inches(2.5), Inches(6.0), Inches(0.06), fill=C_GOLD)
    txt(sl, "CRM 운영",
        Inches(7.1), Inches(2.65), Inches(5.5), Inches(0.55),
        sz=20, bold=True, color=C_GOLD)
    crm = [
        "행동 기반 사용자 세그멘테이션 (RFM, 예측)",
        "고객 생애주기(LTV) 분석 및 전략 수립",
        "이탈 위험 고객 조기 감지 & 재활성화",
        "구독 전환율 / 재구독율 향상 캠페인",
        "주간·월간 CRM 성과 리포팅",
    ]
    for i, item in enumerate(crm):
        txt(sl, f"▸  {item}",
            Inches(7.1), Inches(3.3 + i*0.68), Inches(5.6), Inches(0.6),
            sz=15, color=C_WHITE if i == 0 else C_GRAY)
    print("  [7] 마테크 & CRM")

    # ─────────────────────────────────────────────
    # Slide 8: Labbit
    # ─────────────────────────────────────────────
    sl = prs.slides.add_slide(blank(prs))
    bg_img(sl, ASSETS/"city_neon_tech.jpg")
    overlay(sl, 0, 0, SW, SH, 4, 8, 18, 78)
    logo_badge(sl)
    slide_num(sl, 8)

    txt(sl, "Own Service 01",
        Inches(0.5), Inches(0.85), Inches(6.0), Inches(0.4),
        sz=14, color=C_TEAL, bold=True)
    txt(sl, "Labbit",
        Inches(0.5), Inches(1.25), Inches(5.0), Inches(1.0),
        sz=54, bold=True, color=C_WHITE)
    txt(sl, "래빗  —  그로스해킹 파트너",
        Inches(0.5), Inches(2.25), Inches(7.0), Inches(0.55),
        sz=24, color=C_TEAL)
    teal_line(sl, Inches(2.92))

    txt(sl, "고객 행동 데이터를 기반으로 전환율과 리텐션을 개선하는\n데이터 드리븐 그로스해킹 서비스입니다.",
        Inches(0.5), Inches(3.1), Inches(7.5), Inches(1.0),
        sz=18, color=C_GRAY)

    labbit_svcs = [
        ("01", "구매 단계 트래킹",    "전환 퍼널 분석"),
        ("02", "광고 성과 추적",       "채널별 ROI 측정"),
        ("03", "커스텀 타겟 발굴",     "행동 기반 세그먼트"),
        ("04", "사이트 UX 개선",       "A/B 테스트 운영"),
        ("05", "데이터 시각화 리포트", "대시보드 커스텀 설계"),
        ("06", "커스텀 솔루션",        "1:1 그로스해커 배정"),
    ]
    for i, (num, ttl, sub) in enumerate(labbit_svcs):
        row, col = divmod(i, 3)
        bx = Inches(0.5 + col*4.2)
        ry = Inches(4.25 + row*1.42)
        rrect(sl, bx, ry, Inches(4.0), Inches(1.28), fill=C_NAVY2, alpha=78)
        rect(sl, bx, ry, Inches(0.06), Inches(1.28), fill=C_TEAL)
        num_badge(sl, num, bx+Inches(0.15), ry+Inches(0.12), C_TEAL)
        txt(sl, ttl, bx+Inches(0.88), ry+Inches(0.14), Inches(2.9), Inches(0.42),
            sz=15, bold=True, color=C_WHITE)
        txt(sl, sub, bx+Inches(0.88), ry+Inches(0.6), Inches(2.9), Inches(0.4),
            sz=13, color=C_GRAY)

    # URL 배지 (이모지 제거)
    rrect(sl, Inches(8.5), Inches(1.3), Inches(4.4), Inches(0.65), fill=C_TEAL, alpha=90)
    txt(sl, "labbit.kr",
        Inches(8.5), Inches(1.32), Inches(4.4), Inches(0.6),
        sz=18, bold=True, color=C_NAVY2, align=PP_ALIGN.CENTER)
    print("  [8] Labbit")

    # ─────────────────────────────────────────────
    # Slide 9: DataNugget
    # ─────────────────────────────────────────────
    sl = prs.slides.add_slide(blank(prs))
    bg_img(sl, ASSETS/"city_data_viz.jpg")
    overlay(sl, 0, 0, SW, SH, 2, 8, 18, 78)
    logo_badge(sl)
    slide_num(sl, 9)

    txt(sl, "Own Service 02",
        Inches(0.5), Inches(0.85), Inches(6.0), Inches(0.4),
        sz=14, color=C_GOLD, bold=True)
    txt(sl, "DataNugget",
        Inches(0.5), Inches(1.25), Inches(7.0), Inches(1.0),
        sz=54, bold=True, color=C_WHITE)
    txt(sl, "데이터너겟  —  브랜드 성장박스",
        Inches(0.5), Inches(2.25), Inches(7.0), Inches(0.55),
        sz=24, color=C_GOLD)
    rect(sl, Inches(0.5), Inches(2.92), Inches(1.5), Inches(0.05), fill=C_GOLD)

    txt(sl, "광고 성과 조회부터 홈페이지 리뉴얼까지\n데이터에 대한 모든 질문을 해결하는 데이터 서비스입니다.",
        Inches(0.5), Inches(3.1), Inches(7.5), Inches(1.0),
        sz=18, color=C_GRAY)

    nugget_features = [
        ("01", "광고 성과 통합 분석", "모든 매체 성과를 하나의 대시보드로"),
        ("02", "홈페이지 데이터 진단", "방문자 행동 패턴 및 UX 개선점 도출"),
        ("03", "데이터 시각화 자동화", "커스텀 리포트 자동 생성 및 공유"),
        ("04", "인사이트 리포팅",      "데이터에서 실행 가능한 인사이트 추출"),
    ]
    for i, (num, ttl, sub) in enumerate(nugget_features):
        row, col = divmod(i, 2)
        bx = Inches(0.5 + col*6.3)
        ry = Inches(4.3 + row*1.5)
        rrect(sl, bx, ry, Inches(6.1), Inches(1.32), fill=C_NAVY2, alpha=78)
        rect(sl, bx, ry, Inches(6.1), Inches(0.06), fill=C_GOLD)
        num_badge(sl, num, bx+Inches(0.2), ry+Inches(0.12), C_GOLD)
        txt(sl, ttl, bx+Inches(0.95), ry+Inches(0.14), Inches(4.8), Inches(0.45),
            sz=17, bold=True, color=C_WHITE)
        txt(sl, sub, bx+Inches(0.95), ry+Inches(0.62), Inches(4.8), Inches(0.42),
            sz=14, color=C_GRAY)

    # URL 배지 (이모지 제거)
    rrect(sl, Inches(8.5), Inches(1.3), Inches(4.4), Inches(0.65), fill=C_GOLD, alpha=90)
    txt(sl, "datanugget.io",
        Inches(8.5), Inches(1.32), Inches(4.4), Inches(0.6),
        sz=18, bold=True, color=C_NAVY2, align=PP_ALIGN.CENTER)
    print("  [9] DataNugget")

    # ─────────────────────────────────────────────
    # Slide 10: 핵심 역량 인포그래픽 (NEW)
    # ─────────────────────────────────────────────
    infographic_slide(prs,
        ASSETS/"infographic_capabilities.png",
        "Core Capabilities",
        "핵심 역량 그리드",
        "bg_technology.jpg", 10)
    print("  [10] 핵심 역량 인포그래픽")

    # ─────────────────────────────────────────────
    # Slide 11: 협업 방식 (M/M)
    # ─────────────────────────────────────────────
    sl = prs.slides.add_slide(blank(prs))
    bg_img(sl, ASSETS/"city_office_night.jpg")
    overlay(sl, 0, 0, SW, SH, 4, 10, 22, 80)
    logo_badge(sl)
    slide_num(sl, 11)

    section_header(sl, "How We Work", "M/M 기반 월정액 협업",
                   "매월 우선순위를 합의하고, 전담 팀이 실행합니다.")

    steps = [
        ("01", "월초 미팅", "우선순위 작업 목록 공유\n및 이번 달 목표 설정"),
        ("02", "실행",      "전담 팀이 합의된 작업\n순차적으로 수행"),
        ("03", "주간 보고", "매주 진행 상황 슬랙\n리포팅 & 이슈 공유"),
        ("04", "월말 리뷰", "성과 측정 + 다음 달\n우선순위 재조정"),
    ]
    for i, (num, ttl, body) in enumerate(steps):
        bx = Inches(0.5 + i*3.18)
        rrect(sl, bx, Inches(3.1), Inches(3.0), Inches(3.95), fill=C_NAVY2, alpha=80)
        txt(sl, num, bx+Inches(0.2), Inches(3.2), Inches(1.2), Inches(0.7),
            sz=38, bold=True, color=C_TEAL)
        txt(sl, ttl, bx+Inches(0.2), Inches(3.95), Inches(2.6), Inches(0.5),
            sz=20, bold=True, color=C_WHITE)
        txt(sl, body, bx+Inches(0.2), Inches(4.52), Inches(2.7), Inches(1.0),
            sz=14, color=C_GRAY)
        if i < 3:
            txt(sl, "->", bx+Inches(2.8), Inches(3.6), Inches(0.5), Inches(0.55),
                sz=24, color=C_TEAL, align=PP_ALIGN.CENTER)

    for i, principle in enumerate(["전담 슬랙 채널 운영", "격주 정기 미팅", "월간 성과 리포트", "유연한 서비스 조정"]):
        rrect(sl, Inches(0.5 + i*3.18), Inches(7.1), Inches(3.0), Inches(0.35),
              fill=C_TEAL, alpha=85)
        txt(sl, f"✓  {principle}",
            Inches(0.5 + i*3.18 + 0.1), Inches(7.12), Inches(2.8), Inches(0.3),
            sz=13, color=C_NAVY2, bold=True)
    print("  [11] 협업 방식")

    # ─────────────────────────────────────────────
    # Slide 12: 핵심 성과 수치 인포그래픽 (NEW)
    # ─────────────────────────────────────────────
    infographic_slide(prs,
        ASSETS/"infographic_stats.png",
        "Key Metrics",
        "주요 성과 지표",
        "bg_growth.jpg", 12)
    print("  [12] 핵심 성과 수치 인포그래픽")

    # ─────────────────────────────────────────────
    # Slide 13: 케이스 스터디
    # ─────────────────────────────────────────────
    sl = prs.slides.add_slide(blank(prs))
    bg_img(sl, ASSETS/"bg_growth.jpg")
    overlay(sl, 0, 0, SW, SH, 4, 10, 20, 82)
    logo_badge(sl)
    slide_num(sl, 13)

    section_header(sl, "Case Study", "실적 & 케이스 스터디")

    cases = [
        {
            "industry": "OTT · 미디어 플랫폼",
            "challenge": "신규 가입자 이탈률 60% 이상, CRM 자동화 부재",
            "solution": "온보딩 플로우 재설계 + CleverTap 세그먼트 구축 + 이탈방지 Push 자동화",
            "results": [("-38%", "이탈률"), ("+52%", "D30 리텐션"), ("2.4x", "캠페인 ROI")],
            "color": C_TEAL,
        },
        {
            "industry": "이커머스 · 쇼핑몰",
            "challenge": "GA4 이전 후 데이터 신뢰도 0%, 광고 최적화 불가",
            "solution": "GA4 전면 재설계 + GTM 정리 + DataNugget 대시보드 구축",
            "results": [("100%", "데이터 정확도"), ("-45%", "광고 비효율"), ("+28%", "전환율")],
            "color": C_GOLD,
        },
    ]
    for i, c in enumerate(cases):
        bx = Inches(0.4 + i*6.5)
        rrect(sl, bx, Inches(2.5), Inches(6.2), Inches(4.7), fill=C_NAVY2, alpha=80)
        rect(sl, bx, Inches(2.5), Inches(6.2), Inches(0.07), fill=c["color"])
        txt(sl, c["industry"], bx+Inches(0.25), Inches(2.62), Inches(5.7), Inches(0.5),
            sz=18, bold=True, color=c["color"])
        txt(sl, "과제: " + c["challenge"],
            bx+Inches(0.25), Inches(3.18), Inches(5.7), Inches(0.75),
            sz=14, color=C_GRAY)
        txt(sl, "솔루션: " + c["solution"],
            bx+Inches(0.25), Inches(4.0), Inches(5.7), Inches(0.85),
            sz=14, color=C_WHITE)
        rect(sl, bx+Inches(0.25), Inches(4.95), Inches(5.5), Inches(0.04), fill=c["color"])
        for j, (v, l) in enumerate(c["results"]):
            kx = bx + Inches(0.25 + j*1.95)
            txt(sl, v, kx, Inches(5.1), Inches(1.8), Inches(0.7),
                sz=34, bold=True, color=c["color"])
            txt(sl, l, kx, Inches(5.82), Inches(1.8), Inches(0.38),
                sz=13, color=C_GRAY)
    print("  [13] 케이스 스터디")

    # ─────────────────────────────────────────────
    # Slide 14: 레퍼런스 파트너 인포그래픽 (NEW)
    # ─────────────────────────────────────────────
    infographic_slide(prs,
        ASSETS/"infographic_references.png",
        "Our References",
        "주요 파트너 & 레퍼런스",
        "city_skyline_blue.jpg", 14)
    print("  [14] 레퍼런스 파트너 인포그래픽")

    # ─────────────────────────────────────────────
    # Slide 15: CTA
    # ─────────────────────────────────────────────
    sl = prs.slides.add_slide(blank(prs))
    bg_img(sl, ASSETS/"city_city_dawn.jpg")
    overlay(sl, 0, 0, SW, SH, 4, 10, 22, 72)
    overlay(sl, 0, Inches(3.5), SW, Inches(4.0), 2, 6, 14, 65)
    slide_num(sl, 15)

    txt(sl, "지금 시작하세요",
        Inches(1.5), Inches(1.2), Inches(10.5), Inches(1.2),
        sz=54, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    txt(sl, "위어드섹터와 함께 데이터로 성장하는 마케팅을 경험하세요",
        Inches(1.5), Inches(2.5), Inches(10.5), Inches(0.65),
        sz=22, color=C_TEAL, align=PP_ALIGN.CENTER)

    rrect(sl, Inches(4.7), Inches(3.4), Inches(4.0), Inches(0.8), fill=C_TEAL, alpha=95)
    txt(sl, "미팅 신청하기  ->",
        Inches(4.7), Inches(3.44), Inches(4.0), Inches(0.72),
        sz=20, bold=True, color=C_NAVY2, align=PP_ALIGN.CENTER)

    # 연락처 카드 (이모지 대신 레이블 배지)
    contacts = [
        ("W", "웹사이트",   "weirdsector.co.kr"),
        ("E", "이메일",     "info@weirdsector.co.kr"),
        ("T", "전화",       "02-458-0601"),
    ]
    for i, (badge, lb, vl) in enumerate(contacts):
        bx = Inches(1.0 + i*3.8)
        rrect(sl, bx, Inches(4.6), Inches(3.5), Inches(2.2), fill=C_NAVY2, alpha=80)
        rect(sl, bx, Inches(4.6), Inches(3.5), Inches(0.05), fill=C_TEAL)
        # 배지 대신 글자 배지
        rrect(sl, bx+Inches(1.35), Inches(4.72), Inches(0.78), Inches(0.62), fill=C_TEAL, alpha=90)
        txt(sl, badge, bx+Inches(1.35), Inches(4.74), Inches(0.78), Inches(0.58),
            sz=22, bold=True, color=C_NAVY2, align=PP_ALIGN.CENTER)
        txt(sl, lb, bx+Inches(0.2), Inches(5.4), Inches(3.1), Inches(0.38),
            sz=13, color=C_GRAY, align=PP_ALIGN.CENTER)
        txt(sl, vl, bx+Inches(0.2), Inches(5.82), Inches(3.1), Inches(0.42),
            sz=15, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    txt(sl, '"Be yourself, Embrace your weirdness!"',
        Inches(1.5), Inches(7.0), Inches(10.5), Inches(0.38),
        sz=15, italic=True, color=C_TEAL, align=PP_ALIGN.CENTER)
    print("  [15] CTA")

    # ── 저장
    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    sz = OUT.stat().st_size // 1024
    print(f"\n✅ 저장: {OUT}")
    print(f"   슬라이드: {len(prs.slides)}장 / {sz:,}KB")
    return len(prs.slides)


if __name__ == "__main__":
    print("위어드섹터 회사소개서 v2 빌드 시작\n")
    n = build()
    print(f"\n완료: {n}장")
