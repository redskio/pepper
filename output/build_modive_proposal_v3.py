# -*- coding: utf-8 -*-
"""
모다이브 × 위어드섹터 제안서 v3
- Slides 1~8: v2와 동일 (댓글 8개 반영)
- Slides 9~12: 완전 재설계
  [9]  솔루션 상세 — 3컬럼 서비스 카드 + UI 크롭
  [10] 기대 도입 효과 — infographic_ott_kpi 강조
  [11] 3개월 도입 로드맵 — infographic_implementation_v2
  [12] CTA — 임팩트 강화
"""
import sys, io
if sys.platform == 'win32':
    sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree
from pathlib import Path
from PIL import Image

ASSETS = Path('C:/Agent/pepper/output/modive_assets')
WS_AST = Path('C:/Agent/pepper/output/wirdsector_assets')
SCRAPE = Path('C:/Agent/pepper/output/brand_scrape')
OUT    = Path('C:/Agent/pepper/output/modive_martech_v2.pptx')

C_OCEAN  = RGBColor(0x04, 0x15, 0x20)
C_DEEP   = RGBColor(0x06, 0x1E, 0x35)
C_CARD   = RGBColor(0x0A, 0x22, 0x36)
C_CARD2  = RGBColor(0x0D, 0x2D, 0x48)
C_TEAL   = RGBColor(0x00, 0xD4, 0xAA)
C_CYAN   = RGBColor(0x00, 0xB4, 0xFF)
C_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
C_GRAY   = RGBColor(0x8A, 0xAD, 0xCC)
C_GOLD   = RGBColor(0xFF, 0xC8, 0x4B)
C_PURPLE = RGBColor(0xA7, 0x8B, 0xFA)
C_NAVY2  = RGBColor(0x06, 0x0E, 0x18)
C_GREEN  = RGBColor(0x34, 0xD3, 0x99)
C_RED    = RGBColor(0xE8, 0x23, 0x0B)
C_BLUE2  = RGBColor(0x17, 0x69, 0xFF)

SW = Inches(13.33)
SH = Inches(7.5)
TOTAL = 12

def new_prs():
    prs = Presentation()
    prs.slide_width = SW; prs.slide_height = SH
    return prs

def blank(prs): return prs.slide_layouts[6]

def bg_img(slide, path):
    p = Path(path)
    if not p.exists(): return
    try: slide.shapes.add_picture(str(p), Inches(0), Inches(0), SW, SH)
    except Exception as e: print(f'  [W] bg: {e}')

def overlay(slide, x, y, w, h, r=0, g=0, b=0, alpha=65):
    sp = slide.shapes.add_shape(1, x, y, w, h)
    sp.line.fill.background(); sp.fill.solid()
    sp.fill.fore_color.rgb = RGBColor(r, g, b)
    spPr = sp.element.spPr if hasattr(sp.element, 'spPr') else sp.element.find(qn('p:spPr'))
    sf = spPr.find('.//' + qn('a:solidFill'))
    if sf is not None:
        clr = sf.find(qn('a:srgbClr'))
        if clr is not None:
            a = etree.SubElement(clr, qn('a:alpha'))
            a.set('val', str(int(alpha * 1000)))
    return sp

def rrect(slide, x, y, w, h, fill=None, alpha=100):
    sp = slide.shapes.add_shape(5, x, y, w, h)
    sp.line.fill.background()
    if fill: sp.fill.solid(); sp.fill.fore_color.rgb = fill
    else: sp.fill.background()
    if alpha < 100 and fill:
        spPr = sp.element.spPr if hasattr(sp.element, 'spPr') else sp.element.find(qn('p:spPr'))
        sf = spPr.find('.//' + qn('a:solidFill'))
        if sf is not None:
            clr = sf.find(qn('a:srgbClr'))
            if clr is not None:
                a = etree.SubElement(clr, qn('a:alpha'))
                a.set('val', str(int(alpha * 1000)))
    return sp

def rect(slide, x, y, w, h, fill=None):
    sp = slide.shapes.add_shape(1, x, y, w, h)
    sp.line.fill.background()
    if fill: sp.fill.solid(); sp.fill.fore_color.rgb = fill
    else: sp.fill.background()
    return sp

def txt(slide, text, x, y, w, h, sz=18, bold=False, color=None,
        align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    para = tf.paragraphs[0]; para.alignment = align
    run = para.add_run()
    run.text = text; run.font.size = Pt(sz)
    run.font.bold = bold; run.font.italic = italic
    run.font.name = 'Malgun Gothic'
    if color: run.font.color.rgb = color
    return tb

def txt_multi(slide, lines, x, y, w, h, sz=14, color=None, line_color_pairs=None):
    """Multi-line textbox with optional per-line color."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            para = tf.paragraphs[0]
        else:
            para = tf.add_paragraph()
        para.space_before = Pt(2)
        run = para.add_run()
        run.text = line
        run.font.size = Pt(sz)
        run.font.name = 'Malgun Gothic'
        c = (line_color_pairs[i] if line_color_pairs and i < len(line_color_pairs)
             else color)
        if c: run.font.color.rgb = c
    return tb

def pic(slide, path, x, y, w, h):
    p = Path(path)
    if not p.exists(): return
    try: slide.shapes.add_picture(str(p), x, y, w, h)
    except Exception as e: print(f'  [W] pic: {e}')

def infographic_embed(slide, img_path, top=Inches(2.15), avail_h=Inches(5.05)):
    p = Path(img_path)
    if not p.exists(): return
    try:
        im = Image.open(str(p))
        iw, ih = im.size
        aspect = iw / ih
        avail_w = Inches(12.4)
        if aspect > avail_w / avail_h:
            w = avail_w; h = Emu(int(avail_w / aspect))
        else:
            h = avail_h; w = Emu(int(avail_h * aspect))
        left = Emu(int((SW - w) / 2))
        t = top + Emu(int((avail_h - h) / 2))
        slide.shapes.add_picture(str(p), left, t, w, h)
    except Exception as e:
        print(f'  [W] infographic: {e}')

def snum(slide, n):
    overlay(slide, Inches(12.5), Inches(7.12), Inches(0.83), Inches(0.33), 0, 0, 0, 40)
    txt(slide, f'{n:02d} / {TOTAL:02d}', Inches(12.5), Inches(7.13),
        Inches(0.8), Inches(0.3), sz=10, color=C_GRAY, align=PP_ALIGN.CENTER)

def ws_badge(slide):
    rrect(slide, Inches(0.4), Inches(0.28), Inches(2.5), Inches(0.5),
          fill=C_TEAL, alpha=90)
    txt(slide, 'Weirdsector', Inches(0.45), Inches(0.3), Inches(2.4), Inches(0.46),
        sz=14, bold=True, color=C_NAVY2, align=PP_ALIGN.CENTER)

def teal_line(slide, y=Inches(1.0)):
    rect(slide, Inches(0.5), y, Inches(1.5), Inches(0.05), fill=C_TEAL)

def sec_hdr(slide, en_lbl, ko_title, sub=None):
    txt(slide, en_lbl, Inches(0.5), Inches(0.85), Inches(6.0), Inches(0.38),
        sz=13, color=C_TEAL, bold=True)
    txt(slide, ko_title, Inches(0.5), Inches(1.25), Inches(9.5), Inches(0.82),
        sz=40, bold=True, color=C_WHITE)
    teal_line(slide, Inches(2.17))
    if sub:
        txt(slide, sub, Inches(0.5), Inches(2.35), Inches(9.0), Inches(0.48),
            sz=18, color=C_TEAL)

def num_b(slide, n, x, y, c=None):
    c = c or C_TEAL
    rrect(slide, x, y, Inches(0.56), Inches(0.56), fill=c, alpha=90)
    txt(slide, n, x, y+Inches(0.03), Inches(0.56), Inches(0.5),
        sz=14, bold=True, color=C_NAVY2, align=PP_ALIGN.CENTER)

def infographic_page(prs, img_path, en_lbl, ko_title, bg_file, n,
                     bg_dir=None, top=Inches(2.15), avail_h=Inches(5.05)):
    sl = prs.slides.add_slide(blank(prs))
    bg_dir = bg_dir or ASSETS
    bg_img(sl, Path(bg_dir) / bg_file)
    overlay(sl, 0, 0, SW, SH, 4, 16, 28, 82)
    ws_badge(sl); snum(sl, n)
    txt(sl, en_lbl, Inches(0.5), Inches(0.85), Inches(7.0), Inches(0.38),
        sz=13, color=C_TEAL, bold=True)
    txt(sl, ko_title, Inches(0.5), Inches(1.25), Inches(9.5), Inches(0.62),
        sz=34, bold=True, color=C_WHITE)
    teal_line(sl, Inches(2.0))
    infographic_embed(sl, img_path, top=top, avail_h=avail_h)
    return sl


# ══════════════════════════════════════════════════
def build():
    prs = new_prs()

    # ── 1. 커버 ──────────────────────────────────
    sl = prs.slides.add_slide(blank(prs))
    bg_img(sl, ASSETS/'ocean_deep.jpg')
    overlay(sl, 0, 0, SW, SH, 4, 16, 28, 68)
    overlay(sl, 0, Inches(4.0), SW, Inches(3.5), 0, 0, 0, 55)
    rect(sl, Inches(0.5), Inches(1.3), Inches(0.06), Inches(4.5), fill=C_TEAL)
    txt(sl, '모다이브', Inches(0.8), Inches(1.3), Inches(11.0), Inches(1.0),
        sz=54, bold=True, color=C_WHITE)
    txt(sl, 'Weirdsector MarTech 제안서',
        Inches(0.8), Inches(2.35), Inches(11.0), Inches(0.72),
        sz=30, color=C_TEAL)
    txt(sl, 'Data Engineering · MarTech · CRM · Growth',
        Inches(0.8), Inches(3.15), Inches(9.5), Inches(0.55),
        sz=20, color=C_GRAY)
    rect(sl, Inches(0.8), Inches(3.08), Inches(5.5), Inches(0.03), fill=C_TEAL)
    txt(sl, '"데이터로 구독자를 만들고, 마케팅으로 팬을 만듭니다"',
        Inches(0.8), Inches(3.75), Inches(9.5), Inches(0.5),
        sz=16, italic=True, color=C_TEAL)
    for i, tag in enumerate(['Data & Code', 'MarTech', 'CRM', 'Labbit', 'DataNugget']):
        rrect(sl, Inches(0.8 + i*2.35), Inches(6.4), Inches(2.2), Inches(0.43),
              fill=C_OCEAN, alpha=80)
        txt(sl, tag, Inches(0.8 + i*2.35), Inches(6.42), Inches(2.2), Inches(0.4),
            sz=13, color=C_TEAL, align=PP_ALIGN.CENTER)
    txt(sl, '2026', Inches(12.0), Inches(7.02), Inches(1.2), Inches(0.36),
        sz=13, color=C_GRAY, align=PP_ALIGN.RIGHT)
    snum(sl, 1)
    print('  [1] 커버')

    # ── 2. 위어드섹터 소개 (파트너십 배지) ───────
    sl = prs.slides.add_slide(blank(prs))
    bg_img(sl, ASSETS/'ocean_blue.jpg')
    overlay(sl, 0, 0, SW, SH, 4, 16, 28, 78)
    ws_badge(sl); snum(sl, 2)
    sec_hdr(sl, 'About Weirdsector', '위어드섹터 소개', '데이터 기반 마케팅의 실행 파트너')
    ws_img = SCRAPE / 'ws_index1.jpg'
    if ws_img.exists():
        pic(sl, ws_img, Inches(9.2), Inches(1.0), Inches(3.8), Inches(6.2))
        overlay(sl, Inches(9.2), Inches(1.0), Inches(3.8), Inches(6.2), 4, 16, 28, 55)
    txt(sl, '개발력 · 마케팅 감각 · 데이터 분석\n세 가지를 한 팀에서 원스톱으로 제공합니다.',
        Inches(0.5), Inches(2.95), Inches(8.4), Inches(0.72), sz=17, color=C_GRAY)
    for i, (label, c) in enumerate([('Adobe\n공식 파트너', C_RED),
                                     ('Amplitude\n공식 파트너', C_BLUE2),
                                     ('GA4\n10년 파트너', C_GOLD)]):
        bx = Inches(0.5 + i * 2.65)
        rrect(sl, bx, Inches(3.75), Inches(2.55), Inches(0.52), fill=c, alpha=15)
        rect(sl, bx, Inches(3.75), Inches(0.05), Inches(0.52), fill=c)
        txt(sl, label, bx+Inches(0.12), Inches(3.77), Inches(2.43), Inches(0.48),
            sz=11, bold=True, color=c, align=PP_ALIGN.CENTER)
    for i, (num, ttl, body) in enumerate([
        ('01', '하이브리드 팀', '개발 + 마케팅 + 데이터 전문가\n한 팀에서 원스톱 실행'),
        ('02', '자체 서비스', 'Labbit(그로스해킹) + DataNugget(데이터)\n자체 SaaS 보유 에이전시'),
        ('03', '공식 파트너십', 'Adobe · Amplitude · GA4\n글로벌 솔루션 공식 파트너 에이전시'),
    ]):
        by = Inches(4.38 + i * 0.95)
        rrect(sl, Inches(0.5), by, Inches(8.4), Inches(0.82), fill=C_CARD, alpha=82)
        rect(sl, Inches(0.5), by, Inches(0.05), Inches(0.82), fill=C_TEAL)
        num_b(sl, num, Inches(0.62), by+Inches(0.12), C_TEAL)
        txt(sl, ttl, Inches(1.32), by+Inches(0.08), Inches(2.6), Inches(0.36),
            sz=15, bold=True, color=C_WHITE)
        txt(sl, body, Inches(4.1), by+Inches(0.05), Inches(4.6), Inches(0.72),
            sz=12, color=C_GRAY)
    print('  [2] 위어드섹터 소개')

    # ── 3. Pain Points v2 ──────────────────────
    infographic_page(prs, ASSETS/'infographic_ott_problems_v2.png',
                     'Pain Points', '제안 배경 및 문제 인식',
                     'ocean_surface.jpg', 3, top=Inches(2.1), avail_h=Inches(5.1))
    print('  [3] 문제 인식')

    # ── 4. 솔루션 개요 ─────────────────────────
    sl = prs.slides.add_slide(blank(prs))
    bg_img(sl, ASSETS/'ocean_dark.jpg')
    overlay(sl, 0, 0, SW, SH, 4, 14, 25, 84)
    ws_badge(sl); snum(sl, 4)
    sec_hdr(sl, 'Solution Overview', '핵심 솔루션 3가지')
    for i, (num, ttl, body, c) in enumerate([
        ('01', 'Data\nEngineering',
         'GA4 이벤트 설계\nGTM 재구조화\nBigQuery 파이프라인\n데이터 품질 감사', C_TEAL),
        ('02', 'MarTech\n자동화',
         '인게이지먼트 플랫폼 도입/운영\n온보딩·이탈방지 플로우\nPush · 인앱 · 이메일\nA/B 테스트 최적화', C_CYAN),
        ('03', 'CRM &\n성과 측정',
         '구독자 세그멘테이션\nLTV 분석 & 예측\n리텐션 캠페인 설계\nROI 리포팅 대시보드', C_GOLD),
    ]):
        bx = Inches(0.4 + i*4.3)
        rrect(sl, bx, Inches(2.48), Inches(4.1), Inches(4.65), fill=C_CARD, alpha=85)
        rect(sl, bx, Inches(2.48), Inches(4.1), Inches(0.06), fill=c)
        num_b(sl, num, bx+Inches(1.6), Inches(2.58), c)
        txt(sl, ttl, bx+Inches(0.15), Inches(3.35), Inches(3.8), Inches(0.75),
            sz=18, bold=True, color=c, align=PP_ALIGN.CENTER)
        txt(sl, body, bx+Inches(0.15), Inches(4.15), Inches(3.8), Inches(2.7),
            sz=14, color=C_GRAY, align=PP_ALIGN.CENTER)
    print('  [4] 솔루션 개요')

    # ── 5. 기능 1 — 데이터 분석 ───────────────
    sl = prs.slides.add_slide(blank(prs))
    bg_img(sl, ASSETS/'ocean_night.jpg')
    overlay(sl, 0, 0, SW, SH, 4, 14, 25, 82)
    ws_badge(sl); snum(sl, 5)
    sec_hdr(sl, 'Feature 01', '데이터 분석 & 인사이트',
            '정확한 데이터가 모든 의사결정의 기반입니다.')
    dash_img = ASSETS / 'labbit_svc1.png'
    if (ASSETS / 'datanugget_crop.jpg').exists() and not dash_img.exists():
        dash_img = ASSETS / 'datanugget_crop.jpg'
    if dash_img.exists():
        pic(sl, dash_img, Inches(9.0), Inches(2.5), Inches(4.0), Inches(4.6))
        overlay(sl, Inches(9.0), Inches(2.5), Inches(4.0), Inches(4.6), 4, 14, 25, 45)
    for i, (num, ttl, body) in enumerate([
        ('01', 'GA4 이벤트 택소노미 설계',
         '재생 시작·완료·구독 이벤트 표준화\n콘텐츠 카테고리 × 사용자 행동 분류'),
        ('02', 'GTM & SDK 최적화',
         'Google Tag Manager 재구조화\n모바일 SDK · Firebase 통합 설치'),
        ('03', 'BigQuery 데이터 파이프라인',
         'GA4 → BigQuery 자동화\n콘텐츠별 시청 지표 실시간 집계'),
        ('04', '대시보드 구축',
         '커스텀 대시보드 솔루션\n구독·이탈·LTV 통합 리포팅'),
    ]):
        ry = Inches(2.78 + i*1.1)
        rrect(sl, Inches(0.5), ry, Inches(8.2), Inches(0.98), fill=C_CARD, alpha=78)
        rect(sl, Inches(0.5), ry, Inches(0.05), Inches(0.98), fill=C_TEAL)
        num_b(sl, num, Inches(0.62), ry+Inches(0.18), C_TEAL)
        txt(sl, ttl, Inches(1.32), ry+Inches(0.1), Inches(6.6), Inches(0.38),
            sz=15, bold=True, color=C_WHITE)
        txt(sl, body, Inches(1.32), ry+Inches(0.52), Inches(6.6), Inches(0.42),
            sz=12, color=C_GRAY)
    print('  [5] 데이터 분석')

    # ── 6. 기능 2 — 마케팅 자동화 ─────────────
    sl = prs.slides.add_slide(blank(prs))
    bg_img(sl, ASSETS/'ocean_media.jpg')
    overlay(sl, 0, 0, SW, SH, 4, 14, 25, 82)
    ws_badge(sl); snum(sl, 6)
    sec_hdr(sl, 'Feature 02', '마케팅 자동화 & CRM',
            '개인화 캠페인으로 구독자를 팬으로 만듭니다.')
    svc_img = ASSETS / 'labbit_svc0.png'
    if not svc_img.exists():
        svc_img = ASSETS / 'labbit_crop_hero.jpg'
    if svc_img.exists():
        pic(sl, svc_img, Inches(9.0), Inches(2.5), Inches(4.0), Inches(4.6))
        overlay(sl, Inches(9.0), Inches(2.5), Inches(4.0), Inches(4.6), 4, 14, 25, 45)
    for i, (ttl, body) in enumerate([
        ('인게이지먼트 플랫폼 기반 자동화',
         'Push · 인앱 · 이메일 · SMS\n통합 멀티채널 캠페인 운영'),
        ('온보딩 플로우 구축',
         '가입 후 D1/D3/D7 단계별\n콘텐츠 추천 + 시청 유도 시퀀스'),
        ('이탈 방지 캠페인',
         '30일 비활성 감지 자동 트리거\nWin-back 개인화 메시지'),
        ('A/B 테스트 최적화',
         '메시지 카피 · 발송 시간 · CTA\n통계적 유의성 기반 의사결정'),
    ]):
        ry = Inches(2.78 + i*1.1)
        rrect(sl, Inches(0.5), ry, Inches(8.2), Inches(0.98), fill=C_CARD, alpha=78)
        rect(sl, Inches(0.5), ry, Inches(0.05), Inches(0.98), fill=C_CYAN)
        txt(sl, ttl, Inches(0.7), ry+Inches(0.1), Inches(7.8), Inches(0.38),
            sz=15, bold=True, color=C_WHITE)
        txt(sl, body, Inches(0.7), ry+Inches(0.52), Inches(7.8), Inches(0.42),
            sz=12, color=C_GRAY)
    print('  [6] 마케팅 자동화')

    # ── 7. 기능 3 — 성과 측정 ─────────────────
    sl = prs.slides.add_slide(blank(prs))
    bg_img(sl, ASSETS/'ocean_stream.jpg')
    overlay(sl, 0, 0, SW, SH, 4, 14, 25, 82)
    ws_badge(sl); snum(sl, 7)
    sec_hdr(sl, 'Feature 03', '성과 측정 & 리포팅',
            '모든 캠페인 성과를 실시간으로 측정합니다.')
    svc2_img = ASSETS / 'labbit_svc2.png'
    if not svc2_img.exists():
        svc2_img = ASSETS / 'labbit_crop_portfolio.jpg'
    if svc2_img.exists():
        pic(sl, svc2_img, Inches(8.8), Inches(2.5), Inches(4.2), Inches(4.6))
        overlay(sl, Inches(8.8), Inches(2.5), Inches(4.2), Inches(4.6), 4, 14, 25, 40)
    for i, (ttl, body) in enumerate([
        ('실시간 대시보드',
         '대시보드 솔루션 + GA4 연동\n콘텐츠별 시청률/이탈점 실시간 모니터링'),
        ('캠페인 기여도 분석',
         '채널별 구독 기여도 (Multi-touch)\n광고 → 구독 전환 경로 추적'),
        ('구독자 LTV 분석',
         'D7 / D30 / D90 코호트 분석\n콘텐츠 장르별 리텐션 비교'),
        ('주간·월간 리포팅',
         '성과 현황 + 다음 주 실행 계획\n슬랙/이메일 자동 발송'),
    ]):
        ry = Inches(2.78 + i*1.1)
        rrect(sl, Inches(0.5), ry, Inches(8.0), Inches(0.98), fill=C_CARD, alpha=78)
        rect(sl, Inches(0.5), ry, Inches(0.05), Inches(0.98), fill=C_GOLD)
        num_b(sl, f'{i+1:02d}', Inches(0.62), ry+Inches(0.18), C_GOLD)
        txt(sl, ttl, Inches(1.32), ry+Inches(0.1), Inches(6.4), Inches(0.38),
            sz=15, bold=True, color=C_WHITE)
        txt(sl, body, Inches(1.32), ry+Inches(0.52), Inches(6.4), Inches(0.42),
            sz=12, color=C_GRAY)
    print('  [7] 성과 측정')

    # ── 8. 구독자 라이프사이클 퍼널 v2 ──────────
    infographic_page(prs, ASSETS/'infographic_subscriber_funnel_v2.png',
                     'Subscriber Lifecycle', '구독자 라이프사이클 & 개입 포인트',
                     'ocean_deep.jpg', 8, top=Inches(2.1), avail_h=Inches(5.1))
    print('  [8] 구독자 라이프사이클')

    # ════════════════════════════════════════════
    # SLIDES 9–12: 완전 재설계
    # ════════════════════════════════════════════

    # ── 9. 솔루션 & 서비스 상세 ─────────────────
    # 3컬럼 서비스 카드 + 우측 UI 크롭 패널
    sl = prs.slides.add_slide(blank(prs))
    bg_img(sl, ASSETS/'ocean_dark.jpg')
    overlay(sl, 0, 0, SW, SH, 4, 14, 25, 88)
    ws_badge(sl); snum(sl, 9)

    txt(sl, 'Service Detail', Inches(0.5), Inches(0.85), Inches(6.0), Inches(0.38),
        sz=13, color=C_TEAL, bold=True)
    txt(sl, '제공 서비스 상세', Inches(0.5), Inches(1.22), Inches(9.5), Inches(0.68),
        sz=36, bold=True, color=C_WHITE)
    teal_line(sl, Inches(2.0))

    services = [
        ('01', 'Data\nEngineering', C_TEAL,
         ['GA4 이벤트 택소노미 설계', 'GTM 재구조화 & 배포',
          'BigQuery 파이프라인 구축', '데이터 품질 감사 & 정제',
          '실시간 대시보드 연동']),
        ('02', 'MarTech\n자동화', C_CYAN,
         ['온보딩 자동화 플로우 설계', 'Push · 인앱 · 이메일 운영',
          '행동 기반 세그멘테이션', 'A/B 테스트 최적화',
          '이탈 방지 Win-back 캠페인']),
        ('03', 'CRM &\nGrowth', C_GOLD,
         ['구독자 LTV 모델링', 'D7/D30/D90 코호트 분석',
          'ROI 리포팅 대시보드', '그로스 전략 수립 & 실행',
          '월간 성과 리포팅']),
    ]

    card_top  = Inches(2.12)
    card_h    = Inches(5.1)
    card_w    = Inches(4.05)
    card_gap  = Inches(0.09)

    for i, (num, title, c, feats) in enumerate(services):
        bx = Inches(0.38) + i * (card_w + card_gap)

        # 카드 베이스
        rrect(sl, bx, card_top, card_w, card_h, fill=C_CARD, alpha=88)

        # 상단 컬러 바
        rect(sl, bx, card_top, card_w, Inches(0.06), fill=c)

        # 아이콘 원형 (상단 중앙)
        cx_emu = bx + card_w / 2
        circ_r = Inches(0.42)
        circ_x = cx_emu - circ_r
        circ_y = card_top + Inches(0.18)
        sp_circ = slide_shapes_add_circle(sl, circ_x, circ_y, circ_r * 2, circ_r * 2, c, alpha=25)

        # 번호 텍스트 대신 간결한 숫자
        txt(sl, num, circ_x, circ_y + Inches(0.06), circ_r * 2, circ_r * 1.8,
            sz=18, bold=True, color=c, align=PP_ALIGN.CENTER)

        # 서비스 제목
        txt(sl, title, bx + Inches(0.12), card_top + Inches(1.12),
            card_w - Inches(0.24), Inches(0.82),
            sz=17, bold=True, color=c, align=PP_ALIGN.CENTER)

        # 구분선
        rect(sl, bx + Inches(0.3), card_top + Inches(2.0),
             card_w - Inches(0.6), Inches(0.03), fill=c)

        # 기능 목록
        for j, feat in enumerate(feats):
            fy = card_top + Inches(2.1) + Inches(j * 0.57)
            rect(sl, bx + Inches(0.22), fy + Inches(0.18),
                 Inches(0.08), Inches(0.08), fill=c)
            txt(sl, feat, bx + Inches(0.38), fy,
                card_w - Inches(0.5), Inches(0.5),
                sz=12, color=C_GRAY)

    print('  [9] 솔루션 & 서비스 상세')

    # ── 10. 기대 도입 효과 (KPI 인포그래픽) ─────
    sl = infographic_page(prs, ASSETS/'infographic_ott_kpi.png',
                          'Expected Results', '기대 도입 효과 & KPI',
                          'ocean_surface.jpg', 10,
                          top=Inches(2.08), avail_h=Inches(5.12))

    # 하단 보충 텍스트
    txt(sl, '모다이브 도입 후 OTT 플랫폼 평균 성과 기준 (Weirdsector 레퍼런스 데이터)',
        Inches(0.5), Inches(7.15), Inches(12.0), Inches(0.28),
        sz=10, italic=True, color=C_GRAY, align=PP_ALIGN.CENTER)
    print('  [10] 기대 도입 효과 / KPI')

    # ── 11. 3개월 도입 로드맵 ─────────────────
    sl = infographic_page(prs, ASSETS/'infographic_implementation_v2.png',
                          'Implementation Roadmap', '3개월 도입 로드맵',
                          'ocean_night.jpg', 11,
                          top=Inches(2.05), avail_h=Inches(5.15))

    # 하단 단계 요약 뱃지
    phases = [('Month 1', 'Foundation', C_TEAL),
              ('Month 2', 'Automation', C_CYAN),
              ('Month 3', 'Optimization', C_GOLD)]
    for i, (m, th, c) in enumerate(phases):
        bx = Inches(0.5 + i * 4.25)
        rrect(sl, bx, Inches(7.1), Inches(3.8), Inches(0.3), fill=c, alpha=20)
        rect(sl, bx, Inches(7.1), Inches(0.05), Inches(0.3), fill=c)
        txt(sl, f'{m}  {th}', bx + Inches(0.1), Inches(7.1),
            Inches(3.6), Inches(0.3), sz=11, bold=True, color=c)
    print('  [11] 3개월 도입 로드맵')

    # ── 12. CTA (임팩트 강화) ─────────────────
    sl = prs.slides.add_slide(blank(prs))
    bg_img(sl, ASSETS/'ocean_deep.jpg')
    # 깊은 오버레이로 강렬한 분위기
    overlay(sl, 0, 0, SW, SH, 2, 8, 16, 80)

    # 좌측 TEAL 수직 액센트 바
    rect(sl, Inches(0.5), Inches(0.6), Inches(0.1), Inches(6.3), fill=C_TEAL)

    # 상단 소형 배지
    rrect(sl, Inches(0.8), Inches(0.65), Inches(3.2), Inches(0.42),
          fill=C_TEAL, alpha=18)
    rect(sl, Inches(0.8), Inches(0.65), Inches(3.2), Inches(0.03), fill=C_TEAL)
    txt(sl, 'Weirdsector × 모다이브', Inches(0.88), Inches(0.67),
        Inches(3.0), Inches(0.38), sz=12, bold=True, color=C_TEAL)

    # 메인 헤드라인
    txt(sl, '지금 바로', Inches(0.8), Inches(1.25), Inches(12.0), Inches(1.0),
        sz=62, bold=True, color=C_WHITE)
    txt(sl, '시작하세요.',
        Inches(0.8), Inches(2.2), Inches(12.0), Inches(1.0),
        sz=62, bold=True, color=C_TEAL)

    # 슬로건
    txt(sl, '"데이터로 구독자를 만들고, 마케팅으로 팬을 만듭니다"',
        Inches(0.8), Inches(3.28), Inches(11.0), Inches(0.48),
        sz=16, italic=True, color=C_GRAY)

    # 미팅 CTA 버튼 (강조)
    rrect(sl, Inches(0.8), Inches(3.95), Inches(5.0), Inches(0.8),
          fill=C_TEAL, alpha=100)
    txt(sl, '무료 도입 상담 신청  →', Inches(0.8), Inches(3.99),
        Inches(5.0), Inches(0.72), sz=20, bold=True,
        color=C_NAVY2, align=PP_ALIGN.CENTER)

    # 구분선
    rect(sl, Inches(0.8), Inches(4.95), Inches(11.6), Inches(0.02), fill=C_TEAL)

    # 연락처 3종 (세로 배치 → 가로 3칸)
    contacts = [
        ('WEB', 'weirdsector.co.kr', C_TEAL),
        ('MAIL', 'info@weirdsector.co.kr', C_CYAN),
        ('TEL', '02-458-0601', C_GOLD),
    ]
    for i, (badge, val, c) in enumerate(contacts):
        bx = Inches(0.8 + i * 4.15)
        rrect(sl, bx, Inches(5.1), Inches(3.9), Inches(1.8),
              fill=C_CARD, alpha=75)
        rect(sl, bx, Inches(5.1), Inches(3.9), Inches(0.05), fill=c)
        rrect(sl, bx + Inches(0.15), Inches(5.22), Inches(0.8), Inches(0.5),
              fill=c, alpha=22)
        txt(sl, badge, bx + Inches(0.15), Inches(5.24), Inches(0.8), Inches(0.48),
            sz=12, bold=True, color=c, align=PP_ALIGN.CENTER)
        txt(sl, val, bx + Inches(1.1), Inches(5.35), Inches(2.65), Inches(0.36),
            sz=13, bold=True, color=C_WHITE)

    snum(sl, 12)
    print('  [12] CTA (개선)')

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    sz = OUT.stat().st_size // 1024
    print(f'\n✅ Saved: {OUT}')
    print(f'   Slides: {len(prs.slides)} / {sz:,}KB')
    return len(prs.slides)


def slide_shapes_add_circle(slide, x, y, w, h, color, alpha=100):
    """Add an oval/ellipse shape (MSO_SHAPE_TYPE 9 = oval)."""
    sp = slide.shapes.add_shape(9, x, y, w, h)
    sp.line.fill.background()
    sp.fill.solid()
    sp.fill.fore_color.rgb = color
    if alpha < 100:
        spPr = sp.element.spPr if hasattr(sp.element, 'spPr') else sp.element.find(qn('p:spPr'))
        sf = spPr.find('.//' + qn('a:solidFill'))
        if sf is not None:
            clr = sf.find(qn('a:srgbClr'))
            if clr is not None:
                a = etree.SubElement(clr, qn('a:alpha'))
                a.set('val', str(int(alpha * 1000)))
    return sp


if __name__ == '__main__':
    print('모다이브 제안서 v3 빌드 시작 (슬라이드 9~12 재설계)\n')
    n = build()
    print(f'\nDone: {n} slides')
