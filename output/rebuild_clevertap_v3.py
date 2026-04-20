# -*- coding: utf-8 -*-
"""
CleverTap 한국어 제안서 v3 — 이미지 레이아웃 전면 수정
핵심 원칙:
- 텍스트 영역(좌): 0~7인치 → 절대 이미지 침범 금지
- 이미지 영역(우): 7~13.33인치, 높이 1.7~7인치 → 깔끔하게 배치
- 이미지는 하나씩, 오버레이 없이 투명하게
- 스크린샷에 subtle 라운드 박스 테두리 적용
"""
import sys, io
if sys.platform == 'win32':
    sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')
    sys.stderr = io.TextIOWrapper(sys.stderr.buffer, encoding='utf-8', errors='replace')

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pathlib import Path
from PIL import Image
import os, copy
from lxml import etree

ASSETS = Path("C:/Agent/pepper/output/clevertap_assets")
SRC    = Path("C:/Agent/pepper/output/clevertap_proposal_kr.pptx")
OUTPUT = Path("C:/Agent/pepper/output/clevertap_proposal_kr_v3.pptx")

# ── 컬러 ──────────────────────────────────────────────────────
C_PURPLE      = RGBColor(0x6B, 0x4E, 0xFF)
C_PURPLE_DARK = RGBColor(0x3D, 0x1F, 0xCC)
C_PURPLE_LIGHT= RGBColor(0x9B, 0x8A, 0xFF)
C_DARK        = RGBColor(0x1A, 0x1A, 0x2E)
C_DARK2       = RGBColor(0x16, 0x21, 0x3E)
C_MID         = RGBColor(0x2D, 0x2D, 0x4E)
C_WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
C_GRAY        = RGBColor(0x6C, 0x75, 0x8D)
C_LIGHT_GRAY  = RGBColor(0xF4, 0xF4, 0xF8)
C_ACCENT      = RGBColor(0xFF, 0x6B, 0x35)
C_GREEN       = RGBColor(0x00, 0xC9, 0x85)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

# 텍스트 존 / 이미지 존 경계
TEXT_RIGHT = Inches(6.8)   # 텍스트는 여기까지
IMG_LEFT   = Inches(7.0)   # 이미지는 여기서부터
IMG_TOP    = Inches(1.65)  # 이미지 상단 (헤더 아래)
IMG_BOTTOM = Inches(7.2)   # 이미지 하단 (여백)
IMG_RIGHT  = Inches(13.1)  # 이미지 우측 여백
IMG_W      = IMG_RIGHT - IMG_LEFT    # ~6.1인치
IMG_H      = IMG_BOTTOM - IMG_TOP    # ~5.55인치


def blank(prs):
    return prs.slide_layouts[6]


def add_rect(slide, x, y, w, h, fill=None, line=None, line_w=None):
    s = slide.shapes.add_shape(1, x, y, w, h)
    if fill:
        s.fill.solid(); s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if line:
        s.line.color.rgb = line
        if line_w: s.line.width = line_w
    else:
        s.line.fill.background()
    return s


def add_text(slide, text, x, y, w, h, size=18, bold=False,
             color=None, align=PP_ALIGN.LEFT, italic=False, wrap=True):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.name = "Malgun Gothic"
    if color: r.font.color.rgb = color
    return tb


def add_image_clean(slide, img_path, x, y, w, h, caption=None):
    """이미지를 깔끔하게 배치 — 경계박스 + 캡션 포함"""
    img_path = str(img_path)
    if not os.path.exists(img_path):
        print(f"  [이미지 없음] {img_path}")
        return

    # 비율 유지 계산
    try:
        img = Image.open(img_path)
        iw, ih = img.size
        ratio = iw / ih
        # 가용 영역에 맞게 조정
        if w / h > ratio:
            # 높이 기준
            actual_h = h
            actual_w = h * ratio
        else:
            actual_w = w
            actual_h = w / ratio
        # 중앙 정렬
        ox = x + (w - actual_w) / 2
        oy = y + (h - actual_h) / 2
    except:
        actual_w, actual_h, ox, oy = w, h, x, y

    # 이미지 배경 박스 (약간 더 크게 — 프레임 효과)
    pad = Inches(0.06)
    add_rect(slide, ox - pad, oy - pad,
             actual_w + pad*2, actual_h + pad*2,
             fill=RGBColor(0x22, 0x22, 0x44),
             line=C_PURPLE_DARK, line_w=Pt(1))

    # 실제 이미지
    slide.shapes.add_picture(img_path, ox, oy, actual_w, actual_h)

    # 캡션
    if caption:
        add_text(slide, caption,
                 x, oy + actual_h + Inches(0.05),
                 w, Inches(0.3),
                 size=9, color=C_GRAY, align=PP_ALIGN.CENTER, italic=True)


# ══════════════════════════════════════════════════════════════
def build_all():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    # ── 공통 배경 빌더 ─────────────────────────────────────────
    def header(slide, num, title, subtitle, img_path=None, caption=None):
        """다크 배경 + 사이드바 + 헤더 + 선택적 이미지"""
        add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill=C_DARK)
        add_rect(slide, 0, 0, Inches(0.3), SLIDE_H, fill=C_PURPLE)
        add_rect(slide, 0, 0, SLIDE_W, Inches(1.55), fill=C_PURPLE_DARK)
        add_text(slide, f"{num}  {title}",
                 Inches(0.5), Inches(0.3), Inches(TEXT_RIGHT - Inches(0.5)), Inches(0.85),
                 size=30, bold=True, color=C_WHITE)
        add_text(slide, subtitle,
                 Inches(0.5), Inches(1.08), Inches(TEXT_RIGHT - Inches(0.5)), Inches(0.42),
                 size=13, color=C_PURPLE_LIGHT)
        # 이미지 (우측)
        if img_path:
            add_image_clean(slide, img_path,
                            IMG_LEFT, IMG_TOP, IMG_W, IMG_H, caption=caption)

    # ──────────────────────────────────────────────────────────
    # SLIDE 1: 표지
    # ──────────────────────────────────────────────────────────
    slide = prs.slides.add_slide(blank(prs))
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill=C_DARK)
    add_rect(slide, 0, 0, Inches(0.3), SLIDE_H, fill=C_PURPLE)

    # 배경 이미지 — 우측 절반에만, 조용하게
    ss = ASSETS / "screenshot_main.png"
    if ss.exists():
        add_image_clean(slide, ss, Inches(7.2), Inches(0.8), Inches(5.8), Inches(5.5),
                        caption="clevertap.com")

    # 로고 배지
    add_rect(slide, Inches(0.5), Inches(0.5), Inches(2.5), Inches(0.55), fill=C_PURPLE)
    add_text(slide, "CleverTap", Inches(0.55), Inches(0.52), Inches(2.4), Inches(0.5),
             size=18, bold=True, color=C_WHITE, align=PP_ALIGN.LEFT)

    add_text(slide, "CleverTap", Inches(0.5), Inches(1.7), Inches(6.3), Inches(1.0),
             size=54, bold=True, color=C_WHITE)
    add_text(slide, "솔루션 제안서", Inches(0.5), Inches(2.7), Inches(6.3), Inches(0.9),
             size=44, bold=True, color=C_PURPLE_LIGHT)

    add_rect(slide, Inches(0.5), Inches(3.8), Inches(3.8), Inches(0.05), fill=C_PURPLE)

    add_text(slide, "모바일 마케팅 자동화 & 실시간 고객 인게이지먼트 플랫폼",
             Inches(0.5), Inches(4.0), Inches(6.3), Inches(0.65),
             size=15, color=C_LIGHT_GRAY)
    add_text(slide, "2026년 4월", Inches(0.5), Inches(6.7), Inches(3.0), Inches(0.5),
             size=14, color=C_GRAY)

    # 통계 카드
    stats = [("10,000+", "글로벌 앱"), ("12조+", "월간 이벤트"), ("25개국+", "서비스")]
    for i, (num_s, label) in enumerate(stats):
        x = Inches(0.5 + i * 2.05)
        add_rect(slide, x, Inches(5.2), Inches(1.85), Inches(0.95), fill=C_MID)
        add_text(slide, num_s, x+Inches(0.1), Inches(5.28), Inches(1.65), Inches(0.45),
                 size=19, bold=True, color=C_PURPLE_LIGHT, align=PP_ALIGN.CENTER)
        add_text(slide, label, x+Inches(0.1), Inches(5.72), Inches(1.65), Inches(0.35),
                 size=11, color=C_GRAY, align=PP_ALIGN.CENTER)

    print("  [1] 표지")

    # ──────────────────────────────────────────────────────────
    # SLIDE 2: 목차
    # ──────────────────────────────────────────────────────────
    slide = prs.slides.add_slide(blank(prs))
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill=C_DARK)
    add_rect(slide, 0, 0, Inches(0.3), SLIDE_H, fill=C_PURPLE)
    add_rect(slide, 0, 0, SLIDE_W, Inches(1.55), fill=C_DARK2)
    add_text(slide, "목 차", Inches(0.5), Inches(0.3), Inches(6.0), Inches(0.8),
             size=34, bold=True, color=C_WHITE)
    add_text(slide, "CONTENTS", Inches(0.5), Inches(1.08), Inches(4.0), Inches(0.42),
             size=13, color=C_PURPLE_LIGHT)

    toc = [
        ("01", "CleverTap 소개"), ("02", "주요 글로벌 고객사"),
        ("03", "핵심 가치 제안"), ("04", "제품 아키텍처 개요"),
        ("05", "실시간 분석 (Analytics)"), ("06", "사용자 세그멘테이션"),
        ("07", "멀티채널 인게이지먼트"), ("08", "AI/ML 기반 개인화"),
        ("09", "A/B 테스트 & 최적화"), ("10", "핀테크/금융 활용 사례"),
        ("11", "도입 효과 & ROI"), ("12", "도입 프로세스 & 타임라인"),
        ("13", "가격 정책 & 다음 단계"),
    ]
    for i, (num_s, title) in enumerate(toc):
        col = i % 2; row = i // 2
        x = Inches(0.6 + col * 6.4)
        y = Inches(1.7 + row * 0.72)
        add_rect(slide, x, y+Inches(0.05), Inches(0.52), Inches(0.52), fill=C_PURPLE)
        add_text(slide, num_s, x+Inches(0.02), y+Inches(0.09), Inches(0.48), Inches(0.38),
                 size=11, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        add_text(slide, title, x+Inches(0.62), y, Inches(5.6), Inches(0.58),
                 size=16, color=C_WHITE)

    # 우측 장식 이미지 — 목차 하단, 텍스트 없는 영역
    add_image_clean(slide, ASSETS / "viewport_product.png",
                    Inches(9.5), Inches(5.8), Inches(3.5), Inches(1.5))
    print("  [2] 목차")

    # ──────────────────────────────────────────────────────────
    # SLIDE 3: CleverTap 소개
    # ──────────────────────────────────────────────────────────
    slide = prs.slides.add_slide(blank(prs))
    header(slide, "01", "CleverTap 소개",
           "모바일 마케팅 자동화의 글로벌 리더",
           ASSETS / "viewport_main.png", "clevertap.com 메인 화면")

    desc = ("CleverTap은 2013년 설립된 글로벌 모바일 마케팅 자동화 & 분석 플랫폼으로,\n"
            "앱 비즈니스가 사용자를 유지하고 성장시킬 수 있도록 돕는 올인원 솔루션입니다.\n\n"
            "실시간 행동 데이터를 기반으로 개인화된 메시지를 적시에 전달하여\n"
            "고객 참여도와 매출을 극대화합니다.")
    add_text(slide, desc, Inches(0.5), Inches(1.8), Inches(6.1), Inches(2.2),
             size=14, color=C_LIGHT_GRAY, wrap=True)

    key_stats = [("2013년", "설립"), ("10,000+", "글로벌 고객사"), ("50+", "국가 서비스"),
                 ("12조+", "월간 이벤트"), ("600+", "임직원"), ("$105M+", "누적 투자")]
    for i, (val, lbl) in enumerate(key_stats):
        col = i % 3; row = i // 3
        x = Inches(0.5 + col * 2.1); y = Inches(4.1 + row * 1.3)
        add_rect(slide, x, y, Inches(1.92), Inches(1.12), fill=C_MID)
        add_rect(slide, x, y, Inches(1.92), Inches(0.06), fill=C_PURPLE)
        add_text(slide, val, x+Inches(0.1), y+Inches(0.1), Inches(1.72), Inches(0.55),
                 size=21, bold=True, color=C_PURPLE_LIGHT, align=PP_ALIGN.CENTER)
        add_text(slide, lbl, x+Inches(0.1), y+Inches(0.65), Inches(1.72), Inches(0.38),
                 size=11, color=C_GRAY, align=PP_ALIGN.CENTER)
    print("  [3] 소개")

    # ──────────────────────────────────────────────────────────
    # SLIDE 4: 주요 고객사 — 이미지 제거, 텍스트 카드만
    # ──────────────────────────────────────────────────────────
    slide = prs.slides.add_slide(blank(prs))
    header(slide, "02", "주요 글로벌 고객사",
           "전 세계 10,000+ 앱이 CleverTap을 신뢰합니다")
    # 우측에 고객사 스크린샷 (적절한 크기)
    add_image_clean(slide, ASSETS / "screenshot_customers.png",
                    IMG_LEFT, IMG_TOP, IMG_W, IMG_H,
                    "clevertap.com/customers/ 화면")

    customers = [
        ("Jio", "인도 최대 통신사", "통신"),
        ("Swiggy", "푸드 딜리버리 1위", "이커머스"),
        ("PhonePe", "UPI 결제 플랫폼", "핀테크"),
        ("Axis Bank", "인도 3대 민간은행", "금융"),
        ("Disney+ Hotstar", "OTT 스트리밍", "미디어"),
        ("Vodafone", "글로벌 통신사", "통신"),
        ("Puma", "스포츠 브랜드", "리테일"),
        ("OLX", "중고거래 플랫폼", "마켓"),
        ("Carousell", "동남아 커머스", "이커머스"),
    ]
    cols = 3
    for i, (name, desc, cat) in enumerate(customers):
        col = i % cols; row = i // cols
        x = Inches(0.5 + col * 2.05); y = Inches(1.9 + row * 1.75)
        add_rect(slide, x, y, Inches(1.85), Inches(1.55), fill=C_MID)
        add_rect(slide, x, y, Inches(0.1), Inches(1.55), fill=C_PURPLE)
        add_rect(slide, x+Inches(0.18), y+Inches(0.12), Inches(0.8), Inches(0.26), fill=C_PURPLE)
        add_text(slide, cat, x+Inches(0.2), y+Inches(0.13), Inches(0.76), Inches(0.23),
                 size=8, bold=True, color=C_WHITE)
        add_text(slide, name, x+Inches(0.18), y+Inches(0.45), Inches(1.55), Inches(0.42),
                 size=15, bold=True, color=C_WHITE)
        add_text(slide, desc, x+Inches(0.18), y+Inches(0.9), Inches(1.55), Inches(0.55),
                 size=10, color=C_GRAY, wrap=True)
    print("  [4] 고객사")

    # ──────────────────────────────────────────────────────────
    # SLIDE 5: 핵심 가치 제안
    # ──────────────────────────────────────────────────────────
    slide = prs.slides.add_slide(blank(prs))
    header(slide, "03", "핵심 가치 제안", "왜 CleverTap인가?")

    values = [
        ("🔄", "올인원 플랫폼", "분석→세그멘테이션→캠페인→최적화\n하나의 플랫폼에서 완결"),
        ("⚡", "실시간 처리", "초당 수백만 이벤트 처리\n즉각적인 의사결정 지원"),
        ("🎯", "AI 기반 개인화", "ML 알고리즘으로 최적 메시지\n사용자별 자동 전달"),
        ("📱", "멀티채널 통합", "푸시·이메일·SMS·인앱·WhatsApp\n단일 플랫폼에서 통합 관리"),
        ("📊", "강력한 분석", "퍼널·코호트·플로우 분석\n깊이 있는 행동 인사이트"),
        ("🔒", "엔터프라이즈 보안", "SOC2·GDPR·ISO 27001 인증\n금융/헬스케어 수준 보안"),
    ]
    # 2열 × 3행으로 전체 슬라이드 활용 (이미지 없이 콘텐츠 풀)
    for i, (icon, title, desc) in enumerate(values):
        col = i % 2; row = i // 2
        x = Inches(0.5 + col * 6.4); y = Inches(1.75 + row * 1.85)
        add_rect(slide, x, y, Inches(6.0), Inches(1.65), fill=C_MID)
        add_rect(slide, x, y, Inches(6.0), Inches(0.05), fill=C_PURPLE)
        add_rect(slide, x+Inches(0.18), y+Inches(0.22), Inches(0.68), Inches(0.68), fill=C_PURPLE_DARK)
        add_text(slide, icon, x+Inches(0.18), y+Inches(0.22), Inches(0.68), Inches(0.68),
                 size=22, align=PP_ALIGN.CENTER)
        add_text(slide, title, x+Inches(1.0), y+Inches(0.18), Inches(4.8), Inches(0.42),
                 size=15, bold=True, color=C_WHITE)
        add_text(slide, desc, x+Inches(1.0), y+Inches(0.62), Inches(4.8), Inches(0.9),
                 size=12, color=C_LIGHT_GRAY, wrap=True)
    print("  [5] 가치제안")

    # ──────────────────────────────────────────────────────────
    # SLIDE 6: 제품 아키텍처
    # ──────────────────────────────────────────────────────────
    slide = prs.slides.add_slide(blank(prs))
    header(slide, "04", "제품 아키텍처 개요",
           "데이터 수집부터 실행까지 — 완전 통합 플랫폼",
           ASSETS / "viewport_product.png", "clevertap.com/product/ 화면")

    layers = [
        ("데이터 수집 레이어", "SDK (iOS/Android/Web) · REST API · S3 커넥터 · 서드파티 통합", C_PURPLE),
        ("데이터 처리 레이어", "실시간 스트림 처리 · 사용자 프로필 병합 · 이벤트 저장소", C_PURPLE_DARK),
        ("분석 & AI 레이어", "행동 분석 · ML 모델 · 예측 세그멘테이션 · 최적 전송 시간", RGBColor(0x3D, 0x5A, 0xCC)),
        ("실행 레이어", "캠페인 관리 · 멀티채널 전송 · A/B 테스트 · 자동화 워크플로우", C_MID),
    ]
    for i, (title, desc, color) in enumerate(layers):
        y = Inches(1.9 + i * 1.12)
        add_rect(slide, Inches(0.5), y, Inches(6.1), Inches(0.98), fill=color)
        add_rect(slide, Inches(0.5), y, Inches(0.38), Inches(0.98), fill=C_PURPLE_DARK)
        add_text(slide, str(i+1), Inches(0.52), y+Inches(0.22), Inches(0.34), Inches(0.5),
                 size=13, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        add_text(slide, title, Inches(1.05), y+Inches(0.05), Inches(5.4), Inches(0.38),
                 size=13, bold=True, color=C_WHITE)
        add_text(slide, desc, Inches(1.05), y+Inches(0.47), Inches(5.4), Inches(0.44),
                 size=11, color=C_LIGHT_GRAY)
    print("  [6] 아키텍처")

    # ──────────────────────────────────────────────────────────
    # SLIDE 7: Analytics
    # ──────────────────────────────────────────────────────────
    slide = prs.slides.add_slide(blank(prs))
    header(slide, "05", "실시간 분석 (Analytics)",
           "모든 데이터를 실시간으로 — 즉각적인 비즈니스 인사이트",
           ASSETS / "screenshot_analytics.png", "Analytics 대시보드")

    features = [
        ("퍼널 분석", "전환율 추적, 이탈 지점 파악, 퍼널 최적화"),
        ("코호트 분석", "사용자 그룹별 리텐션 트렌드 분석"),
        ("플로우 분석", "사용자 탐색 패턴 시각화"),
        ("이벤트 분석", "커스텀 이벤트 실시간 트래킹"),
        ("피벗 분석", "다차원 데이터 분석 및 리포팅"),
    ]
    for i, (title, desc) in enumerate(features):
        y = Inches(1.85 + i * 0.95)
        add_rect(slide, Inches(0.5), y, Inches(6.1), Inches(0.82), fill=C_MID)
        add_rect(slide, Inches(0.5), y, Inches(0.05), Inches(0.82), fill=C_GREEN)
        add_text(slide, title, Inches(0.7), y+Inches(0.05), Inches(2.5), Inches(0.36),
                 size=13, bold=True, color=C_WHITE)
        add_text(slide, desc, Inches(0.7), y+Inches(0.42), Inches(5.7), Inches(0.33),
                 size=11, color=C_GRAY)
    print("  [7] Analytics")

    # ──────────────────────────────────────────────────────────
    # SLIDE 8: 세그멘테이션
    # ──────────────────────────────────────────────────────────
    slide = prs.slides.add_slide(blank(prs))
    header(slide, "06", "사용자 세그멘테이션",
           "정밀한 타겟팅으로 관련성 높은 경험 제공",
           ASSETS / "screenshot_segmentation.png", "Segmentation 화면")

    segs = [
        ("정적 세그먼트", "고정된 기준으로 사용자 그룹화\n예: 연령대별, 지역별, 가입일별"),
        ("동적 세그먼트", "실시간 행동에 따라 자동 업데이트\n예: 최근 7일 미접속자"),
        ("RFM 세그먼트", "Recency·Frequency·Monetary 기반\n고가치 사용자 자동 식별"),
        ("예측 세그먼트", "ML로 이탈/구매 의향 사전 예측\n선제적 대응 가능"),
    ]
    for i, (title, desc) in enumerate(segs):
        col = i % 2; row = i // 2
        x = Inches(0.5 + col * 3.2); y = Inches(1.9 + row * 2.2)
        add_rect(slide, x, y, Inches(2.9), Inches(1.95), fill=C_MID)
        add_rect(slide, x, y, Inches(2.9), Inches(0.06), fill=C_PURPLE)
        add_text(slide, title, x+Inches(0.15), y+Inches(0.15), Inches(2.6), Inches(0.42),
                 size=13, bold=True, color=C_WHITE)
        add_text(slide, desc, x+Inches(0.15), y+Inches(0.65), Inches(2.6), Inches(1.1),
                 size=12, color=C_LIGHT_GRAY, wrap=True)
    print("  [8] 세그멘테이션")

    # ──────────────────────────────────────────────────────────
    # SLIDE 9: 멀티채널
    # ──────────────────────────────────────────────────────────
    slide = prs.slides.add_slide(blank(prs))
    header(slide, "07", "멀티채널 인게이지먼트",
           "고객이 있는 모든 채널에서 일관된 경험 제공",
           ASSETS / "screenshot_engagement.png", "Engagement 화면")

    channels = [
        ("📲", "푸시 알림", "iOS/Android\n리치 푸시, 딥링크"),
        ("📧", "이메일", "드래그&드롭 에디터\nAMP 이메일 지원"),
        ("💬", "SMS/MMS", "글로벌 통신사 연동\n개인화 메시지"),
        ("📱", "인앱 메시지", "풀스크린/배너/팝업\n실시간 노출 조건"),
        ("🌐", "웹 푸시", "브라우저 알림\n세션 외 재참여"),
        ("💚", "WhatsApp", "공식 비즈니스 API\n템플릿 메시지"),
    ]
    for i, (icon, name, desc) in enumerate(channels):
        col = i % 3; row = i // 3
        x = Inches(0.5 + col * 2.1); y = Inches(1.9 + row * 2.45)
        add_rect(slide, x, y, Inches(1.9), Inches(2.2), fill=C_MID)
        add_rect(slide, x, y, Inches(1.9), Inches(0.06), fill=C_PURPLE)
        add_text(slide, icon, x+Inches(0.1), y+Inches(0.12), Inches(0.65), Inches(0.65),
                 size=26, align=PP_ALIGN.CENTER)
        add_text(slide, name, x+Inches(0.1), y+Inches(0.82), Inches(1.7), Inches(0.42),
                 size=13, bold=True, color=C_WHITE)
        add_text(slide, desc, x+Inches(0.1), y+Inches(1.28), Inches(1.7), Inches(0.78),
                 size=11, color=C_GRAY, wrap=True)
    print("  [9] 멀티채널")

    # ──────────────────────────────────────────────────────────
    # SLIDE 10: AI/ML
    # ──────────────────────────────────────────────────────────
    slide = prs.slides.add_slide(blank(prs))
    header(slide, "08", "AI/ML 기반 개인화",
           "머신러닝으로 구현하는 1:1 초개인화 마케팅",
           ASSETS / "screenshot_ai_ml.png", "AI 기능 화면")

    ai_feats = [
        ("🤖 Clever.AI", "CleverTap 독자 AI 엔진 — 고객 행동 학습 및 예측"),
        ("⏰ 최적 전송 시간", "각 사용자별 최적 수신 시간 자동 선택 (오픈율 47% 향상)"),
        ("🎯 예측 세그먼트", "이탈 예측·구매 예측·구독 해지 예측, LTV 기반 자동 분류"),
        ("💡 동적 콘텐츠", "개인 프로필 기반 실시간 콘텐츠 변환 & 제품 추천 자동화"),
        ("📈 자동 최적화", "A/B 테스트 자동 승자 선택, 트래픽 동적 배분"),
        ("🔮 이탈 방지", "이탈 위험 사용자 선제 감지, 재참여 캠페인 자동 트리거"),
    ]
    for i, (title, desc) in enumerate(ai_feats):
        col = i % 2; row = i // 2  # 정확히 3행 2열
        x = Inches(0.5 + col * 3.3); y = Inches(1.85 + row * 1.65)
        add_rect(slide, x, y, Inches(3.05), Inches(1.48), fill=C_MID)
        add_rect(slide, x, y, Inches(0.06), Inches(1.48), fill=C_PURPLE)
        add_text(slide, title, x+Inches(0.2), y+Inches(0.1), Inches(2.7), Inches(0.42),
                 size=13, bold=True, color=C_WHITE)
        add_text(slide, desc, x+Inches(0.2), y+Inches(0.58), Inches(2.7), Inches(0.78),
                 size=11, color=C_LIGHT_GRAY, wrap=True)
    print("  [10] AI/ML")

    # ──────────────────────────────────────────────────────────
    # SLIDE 11: A/B 테스트
    # ──────────────────────────────────────────────────────────
    slide = prs.slides.add_slide(blank(prs))
    header(slide, "09", "A/B 테스트 & 최적화",
           "데이터 기반 의사결정으로 마케팅 성과 극대화",
           ASSETS / "screenshot_ab_testing.png", "A/B Testing 화면")

    ab_items = [
        ("메시지 A/B 테스트", "제목·본문·CTA 버튼·이미지·이모지 테스트"),
        ("전송 시간 테스트", "발송 시간대별 성과 비교, 요일/시간 최적화"),
        ("채널 A/B 테스트", "푸시 vs 이메일 vs SMS 채널별 효과 측정"),
        ("자동 승자 선택", "통계적 유의성 기반 자동 선택, 트래픽 손실 최소화"),
    ]
    for i, (title, desc) in enumerate(ab_items):
        y = Inches(1.85 + i * 1.12)
        add_rect(slide, Inches(0.5), y, Inches(6.1), Inches(0.98), fill=C_MID)
        add_rect(slide, Inches(0.5), y, Inches(0.05), Inches(0.98), fill=C_ACCENT)
        add_text(slide, title, Inches(0.7), y+Inches(0.08), Inches(3.0), Inches(0.4),
                 size=14, bold=True, color=C_WHITE)
        add_text(slide, desc, Inches(0.7), y+Inches(0.52), Inches(5.7), Inches(0.38),
                 size=12, color=C_LIGHT_GRAY)

    # 성과 수치 (하단 좌측)
    for i, (val, lbl) in enumerate([("47%", "오픈율 향상"), ("32%", "전환율 개선"), ("2.3x", "ROI 증가")]):
        x = Inches(0.5 + i * 2.1)
        add_rect(slide, x, Inches(6.3), Inches(1.9), Inches(1.0), fill=C_PURPLE_DARK)
        add_text(slide, val, x+Inches(0.1), Inches(6.35), Inches(1.7), Inches(0.55),
                 size=26, bold=True, color=C_GREEN, align=PP_ALIGN.CENTER)
        add_text(slide, lbl, x+Inches(0.1), Inches(6.9), Inches(1.7), Inches(0.35),
                 size=11, color=C_GRAY, align=PP_ALIGN.CENTER)
    print("  [11] A/B 테스트")

    # ──────────────────────────────────────────────────────────
    # SLIDE 12: 핀테크 활용 사례
    # ──────────────────────────────────────────────────────────
    slide = prs.slides.add_slide(blank(prs))
    header(slide, "10", "핀테크/금융 업계 활용 사례",
           "금융 서비스에 최적화된 마케팅 자동화",
           ASSETS / "screenshot_customers.png", "글로벌 금융 고객사")

    use_cases = [
        ("신규 가입 온보딩", "가입 후 7일 온보딩 시나리오 — 단계별 안내로 첫 거래 유도"),
        ("KYC 완료 독려", "미완료 사용자 자동 리마인드 — 채널별 순차 발송"),
        ("투자 상품 추천", "거래 이력 기반 AI 상품 추천 — 맞춤형 리스크 등급별 제안"),
        ("이탈 방지 캠페인", "30일 미거래 사용자 자동 감지 — 특별 혜택 오퍼 발송"),
        ("거래 알림 & 크로스셀", "실시간 거래 알림 + 크로스셀 — 수수료 면제 프로모션"),
    ]
    for i, (title, desc) in enumerate(use_cases):
        y = Inches(1.85 + i * 0.97)
        add_rect(slide, Inches(0.5), y, Inches(6.1), Inches(0.84), fill=C_MID)
        add_rect(slide, Inches(0.5), y, Inches(0.05), Inches(0.84), fill=C_ACCENT)
        add_text(slide, title, Inches(0.7), y+Inches(0.05), Inches(2.6), Inches(0.36),
                 size=13, bold=True, color=C_WHITE)
        add_text(slide, desc, Inches(0.7), y+Inches(0.44), Inches(5.7), Inches(0.34),
                 size=11, color=C_GRAY)
    print("  [12] 핀테크")

    # ──────────────────────────────────────────────────────────
    # SLIDE 13: ROI
    # ──────────────────────────────────────────────────────────
    slide = prs.slides.add_slide(blank(prs))
    header(slide, "11", "도입 효과 & ROI",
           "측정된 성과 — 실제 고객사 기준 평균 수치")

    roi = [
        ("43%", "리텐션율 향상", "이탈 방지 캠페인 적용 후"),
        ("2.8배", "마케팅 ROI", "전통 채널 대비"),
        ("62%", "인게이지먼트 증가", "개인화 메시지 적용 후"),
        ("30%", "마케팅 비용 절감", "자동화로 운영 효율 개선"),
        ("5배", "캠페인 실행 속도", "수동 대비 자동화 효과"),
        ("93%", "고객 만족도", "NPS 기준 CleverTap 도입사"),
    ]
    for i, (val, title, desc) in enumerate(roi):
        col = i % 3; row = i // 3
        x = Inches(0.5 + col * 4.2); y = Inches(1.85 + row * 2.4)
        add_rect(slide, x, y, Inches(3.9), Inches(2.15), fill=C_MID)
        add_rect(slide, x, y, Inches(3.9), Inches(0.06), fill=C_GREEN)
        add_text(slide, val, x+Inches(0.15), y+Inches(0.15), Inches(3.6), Inches(0.88),
                 size=44, bold=True, color=C_GREEN, align=PP_ALIGN.CENTER)
        add_text(slide, title, x+Inches(0.15), y+Inches(1.05), Inches(3.6), Inches(0.42),
                 size=14, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        add_text(slide, desc, x+Inches(0.15), y+Inches(1.52), Inches(3.6), Inches(0.52),
                 size=11, color=C_GRAY, align=PP_ALIGN.CENTER)
    print("  [13] ROI")

    # ──────────────────────────────────────────────────────────
    # SLIDE 14: 도입 프로세스
    # ──────────────────────────────────────────────────────────
    slide = prs.slides.add_slide(blank(prs))
    header(slide, "12", "도입 프로세스 & 타임라인",
           "빠르고 안전한 구축 — 평균 4주 내 라이브")

    phases = [
        ("Week 1", "킥오프 & 설계", "요구사항 정의\nSDK 통합 계획\n이벤트 텍소노미 설계"),
        ("Week 2", "SDK 통합", "iOS/Android SDK\n이벤트 태깅\n데이터 검증"),
        ("Week 3", "캠페인 구성", "세그멘테이션 설정\n캠페인 템플릿\nA/B 테스트 설계"),
        ("Week 4", "라이브 & 최적화", "첫 캠페인 발송\n성과 모니터링\n피드백 반영"),
    ]
    add_rect(slide, Inches(0.8), Inches(3.5), Inches(11.7), Inches(0.06), fill=C_PURPLE)
    for i, (week, title, desc) in enumerate(phases):
        x = Inches(0.8 + i * 3.0)
        circle = slide.shapes.add_shape(9, x+Inches(1.1), Inches(3.28), Inches(0.48), Inches(0.48))
        circle.fill.solid(); circle.fill.fore_color.rgb = C_PURPLE
        circle.line.fill.background()
        add_rect(slide, x, Inches(1.85), Inches(2.7), Inches(1.38), fill=C_MID)
        add_text(slide, week, x+Inches(0.1), Inches(1.9), Inches(2.5), Inches(0.38),
                 size=12, color=C_PURPLE_LIGHT, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, title, x+Inches(0.1), Inches(2.3), Inches(2.5), Inches(0.42),
                 size=14, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        add_rect(slide, x, Inches(4.0), Inches(2.7), Inches(2.1), fill=C_MID)
        add_text(slide, desc, x+Inches(0.18), Inches(4.1), Inches(2.4), Inches(1.85),
                 size=12, color=C_LIGHT_GRAY, wrap=True)

    add_rect(slide, Inches(0.5), Inches(6.9), Inches(12.3), Inches(0.42), fill=C_PURPLE_DARK)
    add_text(slide, "전담 CSM 배정  |  한국어 기술 지원  |  24/7 모니터링  |  전용 온보딩 플레이북",
             Inches(0.6), Inches(6.95), Inches(12.1), Inches(0.35),
             size=12, color=C_WHITE, align=PP_ALIGN.CENTER)
    print("  [14] 타임라인")

    # ──────────────────────────────────────────────────────────
    # SLIDE 15: 가격 정책
    # ──────────────────────────────────────────────────────────
    slide = prs.slides.add_slide(blank(prs))
    header(slide, "13", "가격 정책 & 다음 단계",
           "비즈니스 규모에 맞는 유연한 플랜")

    plans = [
        ("Essentials", "스타트업 & 성장기\nMAU 50만 이하",
         ["핵심 분석 기능", "푸시 + 이메일", "기본 세그멘테이션", "이메일 지원"], False),
        ("Advanced", "Mid-Market\nMAU 500만 이하",
         ["Essentials 전체", "AI 예측 세그먼트", "멀티채널 전체", "전담 CSM"], True),
        ("Enterprise", "대기업 & 금융기관\nMAU 무제한",
         ["Advanced 전체", "커스텀 SLA", "온프레미스 옵션", "전용 인프라"], False),
    ]
    for i, (plan_name, desc, features, is_rec) in enumerate(plans):
        x = Inches(0.5 + i * 4.25)
        bg = C_PURPLE if is_rec else C_MID
        add_rect(slide, x, Inches(1.85), Inches(3.9), Inches(5.1), fill=bg)
        if is_rec:
            add_rect(slide, x, Inches(1.85), Inches(3.9), Inches(0.4), fill=C_PURPLE_DARK)
            add_text(slide, "★ 추천 플랜", x+Inches(0.1), Inches(1.88), Inches(3.7), Inches(0.33),
                     size=12, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        add_text(slide, plan_name, x+Inches(0.2), Inches(2.38), Inches(3.5), Inches(0.52),
                 size=22, bold=True, color=C_WHITE)
        add_text(slide, desc, x+Inches(0.2), Inches(2.92), Inches(3.5), Inches(0.6),
                 size=12, color=C_LIGHT_GRAY, wrap=True)
        add_rect(slide, x+Inches(0.12), Inches(3.58), Inches(3.66), Inches(0.04), fill=C_PURPLE_LIGHT)
        for j, feat in enumerate(features):
            fy = Inches(3.72 + j * 0.6)
            add_text(slide, "✓  " + feat, x+Inches(0.2), fy, Inches(3.5), Inches(0.48),
                     size=13, color=C_WHITE)

    add_text(slide, "다음 단계:", Inches(0.5), Inches(7.0), Inches(2.5), Inches(0.38),
             size=14, bold=True, color=C_PURPLE_LIGHT)
    add_text(slide, "① POC 범위 확정  →  ② 무료 파일럿 (2주)  →  ③ 맞춤 견적  →  ④ 계약 & 온보딩",
             Inches(0.5), Inches(7.25), Inches(12.3), Inches(0.38),
             size=13, color=C_WHITE)
    print("  [15] 가격")

    # ──────────────────────────────────────────────────────────
    # SLIDE 16: CTA
    # ──────────────────────────────────────────────────────────
    slide = prs.slides.add_slide(blank(prs))
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill=C_DARK)
    add_rect(slide, 0, 0, Inches(0.3), SLIDE_H, fill=C_PURPLE)

    # 중앙 카드 (이미지 없이 깔끔하게)
    add_rect(slide, Inches(1.8), Inches(1.4), Inches(9.7), Inches(4.8), fill=C_MID)
    add_rect(slide, Inches(1.8), Inches(1.4), Inches(9.7), Inches(0.07), fill=C_PURPLE)

    add_text(slide, "지금 바로 시작하세요",
             Inches(2.0), Inches(1.95), Inches(9.3), Inches(1.0),
             size=40, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, "2주 무료 파일럿 프로그램 신청",
             Inches(2.0), Inches(2.98), Inches(9.3), Inches(0.62),
             size=24, color=C_PURPLE_LIGHT, align=PP_ALIGN.CENTER)
    add_text(slide, "실제 데이터로 CleverTap의 가치를 직접 확인하세요",
             Inches(2.0), Inches(3.64), Inches(9.3), Inches(0.45),
             size=16, color=C_LIGHT_GRAY, align=PP_ALIGN.CENTER)

    add_rect(slide, Inches(4.7), Inches(4.42), Inches(3.9), Inches(0.7), fill=C_PURPLE)
    add_text(slide, "데모 신청하기  →",
             Inches(4.7), Inches(4.47), Inches(3.9), Inches(0.6),
             size=18, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    contacts = [("웹사이트", "clevertap.com"), ("이메일", "korea@clevertap.com"), ("전화", "+82-2-XXXX-XXXX")]
    for i, (lbl, val) in enumerate(contacts):
        x = Inches(2.2 + i * 3.5)
        add_text(slide, lbl, x, Inches(6.3), Inches(3.0), Inches(0.32),
                 size=12, color=C_GRAY, align=PP_ALIGN.CENTER)
        add_text(slide, val, x, Inches(6.65), Inches(3.0), Inches(0.42),
                 size=14, bold=True, color=C_PURPLE_LIGHT, align=PP_ALIGN.CENTER)

    add_text(slide, "CleverTap", Inches(5.7), Inches(0.25), Inches(2.0), Inches(0.6),
             size=26, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    print("  [16] CTA")

    # ── 저장 ──────────────────────────────────────────────────
    prs.save(str(OUTPUT))
    size = OUTPUT.stat().st_size // 1024
    print(f"\n저장 완료: {OUTPUT}")
    print(f"슬라이드: {len(prs.slides)}장 / 크기: {size}KB")
    return str(OUTPUT)


if __name__ == "__main__":
    print("CleverTap 제안서 v3 빌드 시작")
    build_all()
    print("완료")
