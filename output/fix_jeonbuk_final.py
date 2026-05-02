# -*- coding: utf-8 -*-
"""전북은행 Amplitude 제안서 — 최종본 (확인 수치 복원 + 미확인 수치 완화)"""
import sys, io
if sys.platform == 'win32':
    sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')

from pptx import Presentation
from pathlib import Path

SRC = r'C:\Agent\pepper\outputs\전북은행_Amplitude_도입제안서_v2.pptx'
DST = r'C:\Agent\pepper\전북은행_Amplitude_도입제안서_최종.pptx'

prs = Presentation(SRC)
changes = []

def set_text(shape, para_idx, new_text):
    """Set text of a specific paragraph, keeping first run's formatting."""
    para = shape.text_frame.paragraphs[para_idx]
    if para.runs:
        para.runs[0].text = new_text
        for r in para.runs[1:]:
            r.text = ''

# ================================================================
# SLIDE 3 — ✅ MAU 71만 복원 (JB다이렉트, 출처: 그린이코노미)
# ================================================================
s3 = prs.slides[2]

# S4: big number → 71만명 복원
set_text(s3.shapes[4], 0, '71만명')
changes.append('S3: 모바일 확대 -> 71만명 (확인된 수치 복원)')

# S5: subtitle
set_text(s3.shapes[5], 0, 'JB다이렉트 MAU (2024 Q1)')
changes.append('S3: MAU subtitle -> JB다이렉트 MAU')

# S6: growth description
set_text(s3.shapes[6], 0, '2023년 초 39만 → 63만 → 71만 성장')
changes.append('S3: 성장 수치 복원 (그린이코노미 출처)')

# S8~S10: 고령 고객 → 완화 (출처 미확인이므로 34.8% 복원하지 않음)
set_text(s3.shapes[8], 0, '고령 고객')
set_text(s3.shapes[9], 0, '디지털 접근성 과제')
set_text(s3.shapes[10], 0, '고령 고객 비중과 디지털 접근성 이슈')
changes.append('S3: 고령 고객 카드 -> 디지털 접근성 과제 (출처 미확인)')

# S20 P2: bullet — 고령 고객 완화
set_text(s3.shapes[20], 2,
    '  지역 기반 은행 특성상 고령 고객 비중과 디지털 접근성 이슈가 중요한 과제')
changes.append('S3: 고령 고객 bullet -> 태스크 지시 문구 적용')

# S20 — add source note (append to last para or modify existing)
# P4 is last bullet. Add source as additional context
# Let's check if there's room to add source
# We'll modify P0 to include source at bottom-level
# Actually, we need to add source to the card subtitle area
# Best approach: modify S6 to include source
set_text(s3.shapes[6], 0, '2023년 초 39만 → 63만 → 71만 성장 (출처: 그린이코노미)')
changes.append('S3: 출처 그린이코노미 표기')

# ================================================================
# SLIDE 4 — ✅ MAU 71만 참조 복원
# ================================================================
s4 = prs.slides[3]
# S6 P0: 기회 영역 첫 bullet
set_text(s4.shapes[6], 0,
    "  JB다이렉트 MAU 71만명(2024 Q1) — 앱 행동 데이터 분석 시 즉각 효과 확인 가능 (출처: 그린이코노미)")
changes.append('S4: MAU 71만 복원 + 출처 명시')

# ================================================================
# SLIDE 9 — 시나리오 퍼널 수치는 시뮬레이션이므로 유지하되 명확 표기
# ================================================================
s9 = prs.slides[8]
# S26: 이미 '적용 가능 시나리오'로 표기됨 - OK
# S27 P1: "5~8%p 개선 가능" → 완화
set_text(s9.shapes[27], 1,
    '  첫 거래 전환율 개선 — 온보딩 가이드 최적화(Guides 기능) 활용')
changes.append('S9: 5~8%p -> 전환율 개선 (수치 제거)')
# S27 P2: "D7 리텐션 20%+" → 완화
set_text(s9.shapes[27], 2,
    '  재방문율 개선 — 코호트 분석 기반 리텐션 캠페인 설계')
changes.append('S9: D7 리텐션 20%+ -> 리텐션 캠페인 설계')

# ================================================================
# SLIDE 12 — ROI: 수치 완화 + 출처 정비
# ================================================================
s12 = prs.slides[11]

# S5 (big number "퍼널 개선") — OK 유지
# S6 (subtitle) — 이미 패치됨. 좀 더 명확하게
set_text(s12.shapes[6], 0, 'Amplitude 도입 기업들의 사례에서 전환율 개선이')
# S6이 2개 run이므로 직접 수정
p6 = s12.shapes[6].text_frame.paragraphs[0]
if len(p6.runs) >= 2:
    p6.runs[0].text = 'Amplitude 도입 기업들의 사례에서 전환율 개선이'
    p6.runs[1].text = '공통적으로 보고됩니다.'
changes.append('S12: 전환율 subtitle -> Amplitude 사례 기반')

# S9 "+25~30%" (리텐션) → 완화
set_text(s12.shapes[9], 0, '개선 기대')
changes.append('S12: +25~30% -> 개선 기대')

# S10 (리텐션 subtitle)
p10 = s12.shapes[10].text_frame.paragraphs[0]
if len(p10.runs) >= 2:
    p10.runs[0].text = '코호트 기반 리텐션 분석으로'
    p10.runs[1].text = '맞춤 캠페인 전략 수립'
else:
    p10.runs[0].text = '코호트 기반 리텐션 분석으로 맞춤 캠페인 전략 수립'
changes.append('S12: 리텐션 subtitle -> 코호트 기반 분석')

# S13 ("셀프서비스") — OK
# S14 — 이미 패치됨. 출처 추가
p14 = s12.shapes[14].text_frame.paragraphs[0]
if len(p14.runs) >= 2:
    p14.runs[0].text = 'Amplitude 도입 사례에서 분석 속도 향상이'
    p14.runs[1].text = '공통적으로 보고됩니다.'
else:
    p14.runs[0].text = 'Amplitude 도입 사례에서 분석 속도 향상이 공통적으로 보고됩니다.'
changes.append('S12: 분석시간 subtitle -> Amplitude 사례 기반')

# S17 "+35~40%" (마케팅 ROI) → 완화
set_text(s12.shapes[17], 0, '효율 개선')
changes.append('S12: +35~40% -> 효율 개선')

# S18 (마케팅 ROI subtitle)
p18 = s12.shapes[18].text_frame.paragraphs[0]
if len(p18.runs) >= 2:
    p18.runs[0].text = '데이터 기반 의사결정 고도화를 통해'
    p18.runs[1].text = '마케팅 효율 및 비용 절감 기대'
else:
    p18.runs[0].text = '데이터 기반 의사결정 고도화를 통해 마케팅 효율 및 비용 절감 기대'
changes.append('S12: 마케팅 ROI subtitle -> 비용 절감 기대')

# S20 ROI 시뮬레이션 — 추가 정리
set_text(s12.shapes[20], 4,
    '  예상 ROI: Forrester TEI 기준 217% ROI 달성 사례 (출처: Forrester, 2023)')
changes.append('S12: ROI 출처 정리')
set_text(s12.shapes[20], 5,
    '  (출처: Forrester Total Economic Impact of Amplitude, 2023 / Amplitude 공식 케이스 스터디)')
changes.append('S12: 출처 라인 보강')

# ================================================================
# SLIDE 15 — References: 글로벌 일반화 + 하나증권 완화
# ================================================================
s15 = prs.slides[14]

# S5: title
set_text(s15.shapes[5], 0, '국내 주요 금융사 (하나증권·키움증권 등)')
changes.append('S15: title 수정')

# S6 P0: 첫 줄
set_text(s15.shapes[6], 0,
    '  국내 주요 금융사들이 점진적으로 Amplitude를 도입하는 추세')
changes.append('S15: 국내 최초 -> 점진적 도입 추세')

# S6 P4: 키움증권 — 유지 (확인된 사실)

# S6 P7~P9: 인사이트 인용문 완화
set_text(s15.shapes[6], 7, '  핵심 인사이트:')
set_text(s15.shapes[6], 8, '  "핵심 퍼널 데이터를 현업이 직접 확인하면서')
set_text(s15.shapes[6], 9, '  조직 내 데이터 기반 의사결정 문화가 확산되고 있습니다"')
changes.append('S15: 인용문 완화')

# S9: "글로벌 금융사 사례" → "글로벌 핀테크 기업들의 도입 사례"
set_text(s15.shapes[9], 0, '글로벌 핀테크 기업들의 도입 사례')
changes.append('S15: 글로벌 금융사 -> 글로벌 핀테크 도입 사례')

# S10 P0: PayPal, Square, Capital One, Intuit → 일반화
set_text(s15.shapes[10], 0,
    '  글로벌 핀테크 및 금융 기업들이 Amplitude를 활용하여 사용자 행동 분석 고도화')
changes.append('S15: 개별 기업명 -> 일반화')

# S10 P1: Forrester TEI — 유지 (공식 출처)

# S10 P2: Square 개별 사례 → 일반화
set_text(s15.shapes[10], 2,
    '  A/B 테스트, 퍼널 분석, 코호트 분석 등 핵심 기능을 금융 도메인에 적용')
changes.append('S15: Square 사례 -> 일반 기능 설명')

# S10 P3: 보안 컴플라이언스 — 유지 (사실)

# S10 P4: N26 → 일반화
set_text(s15.shapes[10], 4,
    '  Pathfinder, Session Replay 등으로 이탈 구간 발견 및 UX 개선 사례 다수')
changes.append('S15: N26 -> 일반 사례')

# S10 P6~P9: 트렌드 문구 완화
set_text(s15.shapes[10], 7,
    '  국내 주요 금융사들이 점진적으로')
set_text(s15.shapes[10], 8,
    '  이벤트 기반 분석 플랫폼을 도입하는 추세이며,')
set_text(s15.shapes[10], 9,
    '  데이터 기반 의사결정이 업계 표준으로 자리잡는 중')
changes.append('S15: 트렌드 문구 완화')

# ================================================================
# SLIDE 16 — Before/After: 일부 문구 완화
# ================================================================
s16 = prs.slides[15]
# S6 P1: "파악할 수 없음" → 완화
set_text(s16.shapes[6], 1, '  사용자 이탈 원인을 체계적으로 파악하기 어려움')
changes.append('S16: 파악할 수 없음 -> 파악하기 어려움')

# ================================================================
# Save
# ================================================================
prs.save(DST)
print(f'Saved: {DST}')
print(f'\nTotal changes: {len(changes)}')
for c in changes:
    print(f'  - {c}')
