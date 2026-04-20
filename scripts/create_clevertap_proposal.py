"""
CleverTap 한국어 제안서 생성 스크립트
Pepper Potts Design Agent
"""

import sys
import io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

import os
import requests
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import urllib.request

# ─── 경로 설정 ───────────────────────────────────────────────────────────────
OUTPUT_DIR = Path(r"C:\Agent\pepper\output")
ASSETS_DIR = OUTPUT_DIR / "clevertap_assets"
OUTPUT_FILE = OUTPUT_DIR / "clevertap-proposal-kr.pptx"
OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
ASSETS_DIR.mkdir(parents=True, exist_ok=True)

# ─── 색상 팔레트 (CleverTap 브랜드) ──────────────────────────────────────────
DARK_BLUE   = RGBColor(0x1A, 0x23, 0x7E)   # 진한 네이비
MID_BLUE    = RGBColor(0x00, 0x6F, 0xC6)   # 미드 블루
ACCENT_CYAN = RGBColor(0x00, 0xBC, 0xD4)   # 청록
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY  = RGBColor(0xF5, 0xF7, 0xFA)
DARK_GRAY   = RGBColor(0x37, 0x47, 0x4F)
ORANGE      = RGBColor(0xFF, 0x6F, 0x00)   # 강조 포인트
LIGHT_BLUE  = RGBColor(0xE3, 0xF2, 0xFD)

# ─── 이미지 다운로드 ─────────────────────────────────────────────────────────
IMAGE_URLS = {
    "all_in_one":    "https://clevertap.com/wp-content/uploads/2025/09/all-in-one_76562c.png",
    "ai_fold":       "https://clevertap.com/wp-content/uploads/2025/09/Home-page_Ai-fold-1.webp",
    "decisioning":   "https://clevertap.com/wp-content/uploads/2025/08/decisioning_2da01e.webp",
    "platform_1":    "https://clevertap.com/wp-content/uploads/2022/11/Banners-04.jpg",
    "platform_2":    "https://clevertap.com/wp-content/uploads/2022/09/Artboard-15-copy-5402x-1.png",
    "platform_3":    "https://clevertap.com/wp-content/uploads/2022/09/Artboard-8-copy-5402x.png",
    "badge":         "https://clevertap.com/wp-content/uploads/2026/02/Badge.webp",
}

downloaded = {}
headers = {
    "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36"
}

print("🖼️  이미지 다운로드 중...")
for name, url in IMAGE_URLS.items():
    try:
        resp = requests.get(url, headers=headers, timeout=15)
        if resp.status_code == 200 and len(resp.content) > 1000:
            ext = url.split(".")[-1].split("?")[0]
            if ext in ("webp", "svg"):
                ext = "png"  # webp/svg → png로 저장 시도
            fpath = ASSETS_DIR / f"{name}.{ext}"
            fpath.write_bytes(resp.content)
            downloaded[name] = str(fpath)
            print(f"  ✅ {name}: {len(resp.content)//1024}KB")
        else:
            print(f"  ⚠️  {name}: HTTP {resp.status_code}")
    except Exception as e:
        print(f"  ❌ {name}: {e}")

print(f"\n📦 다운로드 완료: {len(downloaded)}/{len(IMAGE_URLS)}개\n")

# ─── 헬퍼 함수 ──────────────────────────────────────────────────────────────
def add_slide(prs, layout_idx=6):
    """빈 슬라이드 추가"""
    layout = prs.slide_layouts[layout_idx]
    return prs.slides.add_slide(layout)

def fill_bg(slide, color: RGBColor):
    """슬라이드 배경 단색 채우기"""
    from pptx.oxml.ns import qn
    from lxml import etree
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, left, top, width, height, fill_color, alpha=None):
    """사각형 도형 추가"""
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape

def add_text_box(slide, text, left, top, width, height,
                 font_size=18, bold=False, color=WHITE,
                 align=PP_ALIGN.LEFT, word_wrap=True):
    """텍스트 박스 추가"""
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "맑은 고딕"
    return txBox

def add_multi_para(slide, lines, left, top, width, height,
                   font_size=14, color=DARK_GRAY, bold_first=False):
    """여러 줄 텍스트 (리스트)를 텍스트박스로 추가"""
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = line
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
        run.font.name = "맑은 고딕"
        if bold_first and i == 0:
            run.font.bold = True
    return txBox

def add_image_safe(slide, img_key, left, top, width, height):
    """이미지 안전하게 추가 (실패 시 placeholder 사각형)"""
    if img_key in downloaded:
        try:
            slide.shapes.add_picture(
                downloaded[img_key],
                Inches(left), Inches(top),
                Inches(width), Inches(height)
            )
            return True
        except Exception as e:
            print(f"  ⚠️  이미지 삽입 실패 ({img_key}): {e}")
    # fallback: 회색 박스
    rect = add_rect(slide, left, top, width, height, LIGHT_BLUE)
    add_text_box(slide, "📊 CleverTap Dashboard", left + 0.1, top + height/2 - 0.2,
                 width - 0.2, 0.4, font_size=11, color=MID_BLUE, align=PP_ALIGN.CENTER)
    return False

def add_logo_watermark(slide):
    """우하단 CleverTap 로고 워터마크"""
    add_text_box(slide, "CleverTap", 8.8, 6.9, 1.5, 0.3,
                 font_size=9, color=RGBColor(0xBB, 0xBB, 0xBB),
                 bold=True, align=PP_ALIGN.RIGHT)

def add_header_bar(slide, color=DARK_BLUE, height=1.15):
    """상단 헤더 바"""
    add_rect(slide, 0, 0, 10, height, color)

def add_accent_line(slide, color=ACCENT_CYAN, top=1.15):
    """헤더 아래 액센트 라인"""
    add_rect(slide, 0, top, 10, 0.05, color)

def slide_header(slide, title, subtitle=None, header_color=DARK_BLUE):
    fill_bg(slide, WHITE)
    add_header_bar(slide, header_color, 1.2)
    add_accent_line(slide, ACCENT_CYAN, 1.2)
    add_text_box(slide, title, 0.4, 0.15, 9.0, 0.9,
                 font_size=28, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        add_text_box(slide, subtitle, 0.4, 0.85, 9.0, 0.4,
                     font_size=13, color=ACCENT_CYAN, align=PP_ALIGN.LEFT)
    add_logo_watermark(slide)

# ─── 프레젠테이션 생성 ───────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(10)
prs.slide_height = Inches(7.5)

# ════════════════════════════════════════════════════════════════
# 슬라이드 1: 표지
# ════════════════════════════════════════════════════════════════
slide = add_slide(prs)
fill_bg(slide, DARK_BLUE)

# 배경 그라데이션 효과 (사각형 오버레이)
add_rect(slide, 0, 0, 10, 7.5, DARK_BLUE)
add_rect(slide, 5.5, 0, 4.5, 7.5, MID_BLUE)
add_rect(slide, 0, 6.8, 10, 0.7, ACCENT_CYAN)

# 메인 텍스트
add_text_box(slide, "CleverTap", 0.5, 1.2, 8, 1.0,
             font_size=52, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
add_text_box(slide, "한국어 제안서", 0.5, 2.15, 8, 0.8,
             font_size=36, bold=False, color=ACCENT_CYAN, align=PP_ALIGN.LEFT)
add_text_box(slide, "차세대 고객 인게이지먼트 플랫폼", 0.5, 2.9, 8.5, 0.6,
             font_size=20, bold=False, color=RGBColor(0xBB, 0xDE, 0xFB), align=PP_ALIGN.LEFT)

# 구분선
add_rect(slide, 0.5, 3.65, 5.5, 0.06, ACCENT_CYAN)

# 서브 정보
add_text_box(slide, "2,000+ 글로벌 브랜드의 선택", 0.5, 3.85, 9, 0.45,
             font_size=15, color=RGBColor(0xBB, 0xDE, 0xFB))
add_text_box(slide, "Powered by CleverAI™ — AI 기반 고객 생애가치 극대화", 0.5, 4.3, 9, 0.45,
             font_size=13, color=RGBColor(0x90, 0xCA, 0xF9))

# 이미지
add_image_safe(slide, "all_in_one", 6.0, 1.5, 3.7, 4.5)

add_text_box(slide, "Confidential | © 2025 CleverTap Inc.", 0.5, 7.1, 9, 0.3,
             font_size=9, color=RGBColor(0x78, 0x90, 0x9C), align=PP_ALIGN.LEFT)

# ════════════════════════════════════════════════════════════════
# 슬라이드 2: 목차
# ════════════════════════════════════════════════════════════════
slide = add_slide(prs)
slide_header(slide, "목차", "제안서 구성 안내")

toc_items = [
    ("01", "CleverTap 소개", "회사 개요 및 글로벌 현황"),
    ("02", "시장 현황", "고객 인게이지먼트의 중요성"),
    ("03", "핵심 기능 — 데이터 & 분석", "실시간 데이터 분석 & 세그멘테이션"),
    ("04", "핵심 기능 — 옴니채널", "캠페인 자동화 & 채널 연동"),
    ("05", "핵심 기능 — CleverAI™", "AI/ML 기반 개인화 & 예측"),
    ("06", "핵심 기능 — 리텐션 최적화", "퍼널 분석 & 생애주기 관리"),
    ("07", "핀테크 & 금융 적용 사례", "도입 성과 및 ROI"),
    ("08", "글로벌 고객사", "2,000+ 브랜드 레퍼런스"),
    ("09", "경쟁사 비교", "CleverTap vs 경쟁 솔루션"),
    ("10", "도입 프로세스", "온보딩 & 기술 통합"),
    ("11", "가격 & 패키지", "도입 옵션 안내"),
    ("12", "다음 단계 & CTA", "데모 신청 및 문의"),
]

col1 = toc_items[:6]
col2 = toc_items[6:]

for i, (num, title, sub) in enumerate(col1):
    y = 1.45 + i * 0.83
    add_rect(slide, 0.35, y, 0.55, 0.55, DARK_BLUE)
    add_text_box(slide, num, 0.35, y + 0.07, 0.55, 0.4,
                 font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text_box(slide, title, 1.0, y + 0.02, 3.8, 0.3,
                 font_size=13, bold=True, color=DARK_BLUE)
    add_text_box(slide, sub, 1.0, y + 0.3, 3.8, 0.25,
                 font_size=10, color=DARK_GRAY)

for i, (num, title, sub) in enumerate(col2):
    y = 1.45 + i * 0.83
    add_rect(slide, 5.2, y, 0.55, 0.55, ACCENT_CYAN)
    add_text_box(slide, num, 5.2, y + 0.07, 0.55, 0.4,
                 font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text_box(slide, title, 5.85, y + 0.02, 3.8, 0.3,
                 font_size=13, bold=True, color=DARK_BLUE)
    add_text_box(slide, sub, 5.85, y + 0.3, 3.8, 0.25,
                 font_size=10, color=DARK_GRAY)

# ════════════════════════════════════════════════════════════════
# 슬라이드 3: CleverTap 소개
# ════════════════════════════════════════════════════════════════
slide = add_slide(prs)
slide_header(slide, "CleverTap 소개", "회사 개요 및 글로벌 현황")

# 좌측 텍스트
add_text_box(slide, "올인원 고객 인게이지먼트 플랫폼", 0.4, 1.45, 5.2, 0.5,
             font_size=16, bold=True, color=DARK_BLUE)
add_multi_para(slide, [
    "CleverTap은 '고객 생애가치(CLV)의 한계 없는 성장'을 미션으로",
    "설립된 세계 최고의 올인원 고객 인게이지먼트 플랫폼입니다.",
    "",
    "데이터 분석, AI 기반 개인화, 옴니채널 캠페인 자동화를",
    "하나의 플랫폼에서 통합 제공합니다.",
], 0.4, 1.95, 5.0, 1.8, font_size=13, color=DARK_GRAY)

# 수치 카드
stats = [
    ("2,000+", "글로벌 브랜드"),
    ("50+", "국가 & 시장"),
    ("10년+", "업력"),
    ("G2 리더", "고객 만족 1위"),
]
for i, (val, label) in enumerate(stats):
    x = 0.3 + i * 2.4
    add_rect(slide, x, 4.0, 2.2, 1.5, DARK_BLUE)
    add_text_box(slide, val, x, 4.15, 2.2, 0.7,
                 font_size=26, bold=True, color=ACCENT_CYAN, align=PP_ALIGN.CENTER)
    add_text_box(slide, label, x, 4.85, 2.2, 0.45,
                 font_size=11, color=WHITE, align=PP_ALIGN.CENTER)

# 우측 이미지
add_image_safe(slide, "ai_fold", 5.4, 1.4, 4.3, 2.4)

# 핵심 미션 박스
add_rect(slide, 0.3, 5.8, 9.4, 0.9, LIGHT_BLUE)
add_text_box(slide,
    '🎯  미션: "사람, 프로세스, 기술을 통합하여 모든 고객 접점에서 최적화된 인게이지먼트를 실현한다"',
    0.5, 5.9, 9.0, 0.7, font_size=13, color=DARK_BLUE, align=PP_ALIGN.LEFT)

# ════════════════════════════════════════════════════════════════
# 슬라이드 4: 시장 현황
# ════════════════════════════════════════════════════════════════
slide = add_slide(prs)
slide_header(slide, "시장 현황 & 고객 인게이지먼트의 중요성",
             "왜 지금 CleverTap이 필요한가")

challenges = [
    ("📉", "고객 이탈 급증", "앱 설치 후 30일 내 90%의 사용자가 이탈.\n고객 획득 비용(CAC) 급등으로 리텐션 전략이 생존 과제."),
    ("📊", "데이터 사일로", "마케팅, CRM, 분석 데이터가 분산되어\n통합 고객 뷰 구축 불가능."),
    ("🤖", "AI 시대의 개인화", "고객은 '나를 아는' 경험을 기대.\n일괄 메시지는 더 이상 효과 없음."),
    ("📱", "채널 파편화", "앱, 웹, 이메일, 카카오, SMS 등\n복잡한 채널을 통합 운영해야 하는 부담."),
]

for i, (icon, title, desc) in enumerate(challenges):
    col = i % 2
    row = i // 2
    x = 0.3 + col * 4.9
    y = 1.5 + row * 2.55
    add_rect(slide, x, y, 4.5, 2.3, LIGHT_GRAY)
    add_rect(slide, x, y, 0.12, 2.3, ACCENT_CYAN)
    add_text_box(slide, icon + " " + title, x + 0.25, y + 0.15, 4.1, 0.5,
                 font_size=15, bold=True, color=DARK_BLUE)
    add_text_box(slide, desc, x + 0.25, y + 0.65, 4.1, 1.5,
                 font_size=12, color=DARK_GRAY)

# 결론 박스
add_rect(slide, 0.3, 6.7, 9.4, 0.55, DARK_BLUE)
add_text_box(slide, "✅  CleverTap은 이 모든 문제를 하나의 플랫폼으로 해결합니다",
             0.5, 6.75, 9.0, 0.45, font_size=13, bold=True, color=WHITE)

# ════════════════════════════════════════════════════════════════
# 슬라이드 5: 핵심 기능 1 — 데이터 & 분석
# ════════════════════════════════════════════════════════════════
slide = add_slide(prs)
slide_header(slide, "핵심 기능 ① — 고객 데이터 & 분석",
             "실시간 빅데이터 기반 세그멘테이션 & 액션 가능한 인사이트")

# 좌측 텍스트
features_data = [
    ("🔍 자동화된 세그멘테이션",
     "수백만 고객을 실시간으로 행동/속성/RFM 기반으로 자동 분류"),
    ("📊 퍼널 & 코호트 분석",
     "전환율 병목 구간 즉각 파악, 이탈 지점 정밀 분석"),
    ("⚡ 실시간 이벤트 처리",
     "초당 수십억 이벤트를 실시간 처리, 즉각적인 인게이지먼트 트리거"),
    ("🎯 액션 가능한 대시보드",
     "KPI 모니터링부터 캠페인 성과까지 단일 뷰에서 확인"),
]
for i, (title, desc) in enumerate(features_data):
    y = 1.5 + i * 1.28
    add_rect(slide, 0.3, y, 5.0, 1.15, LIGHT_GRAY)
    add_rect(slide, 0.3, y, 0.1, 1.15, MID_BLUE)
    add_text_box(slide, title, 0.5, y + 0.08, 4.7, 0.4,
                 font_size=13, bold=True, color=DARK_BLUE)
    add_text_box(slide, desc, 0.5, y + 0.5, 4.7, 0.6,
                 font_size=11, color=DARK_GRAY)

# 우측 이미지
add_image_safe(slide, "decisioning", 5.6, 1.4, 4.1, 3.5)

# 성과 지표
add_rect(slide, 0.3, 6.55, 9.4, 0.75, DARK_BLUE)
add_text_box(slide, "📈  도입 성과: 5x 리텐션 향상 | 40% 캠페인 CTR 증가 | 90% 세그멘테이션 시간 단축",
             0.5, 6.65, 9.0, 0.55, font_size=12, bold=True, color=ACCENT_CYAN, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════
# 슬라이드 6: 핵심 기능 2 — 옴니채널
# ════════════════════════════════════════════════════════════════
slide = add_slide(prs)
slide_header(slide, "핵심 기능 ② — 옴니채널 캠페인 자동화",
             "고객이 있는 모든 채널에서 일관된 메시지 전달")

# 채널 아이콘 그리드
channels = [
    ("📱", "푸시 알림", "RenderMax™ 기술로\n전달율 극대화"),
    ("📧", "이메일", "AI 기반 최적 발송\n시간 자동 결정"),
    ("💬", "WhatsApp", "대화형 메시지로\n높은 오픈율"),
    ("💻", "인앱 메시지", "앱 내 타겟 팝업\n& 배너"),
    ("🌐", "웹 메시징", "웹사이트 방문자\n실시간 인게이지"),
    ("📩", "SMS / RCS", "도달 보장 채널\n풍부한 미디어 지원"),
]

for i, (icon, ch, desc) in enumerate(channels):
    col = i % 3
    row = i // 3
    x = 0.3 + col * 3.2
    y = 1.5 + row * 2.2
    add_rect(slide, x, y, 2.9, 2.0, LIGHT_GRAY)
    add_text_box(slide, icon, x + 0.1, y + 0.15, 2.7, 0.55,
                 font_size=30, align=PP_ALIGN.CENTER, color=MID_BLUE)
    add_text_box(slide, ch, x + 0.1, y + 0.7, 2.7, 0.4,
                 font_size=14, bold=True, color=DARK_BLUE, align=PP_ALIGN.CENTER)
    add_text_box(slide, desc, x + 0.1, y + 1.1, 2.7, 0.75,
                 font_size=10, color=DARK_GRAY, align=PP_ALIGN.CENTER)

# IntelliNODE 소개
add_rect(slide, 0.3, 6.55, 9.4, 0.75, MID_BLUE)
add_text_box(slide,
    "🧠  IntelliNODE™: AI가 각 고객에게 가장 효과적인 채널 & 시간 & 메시지를 자동 선택",
    0.5, 6.65, 9.0, 0.55, font_size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_logo_watermark(slide)

# ════════════════════════════════════════════════════════════════
# 슬라이드 7: 핵심 기능 3 — CleverAI™
# ════════════════════════════════════════════════════════════════
slide = add_slide(prs)
fill_bg(slide, DARK_BLUE)
add_rect(slide, 0, 0, 10, 1.2, RGBColor(0x0D, 0x1A, 0x6E))
add_accent_line(slide, ACCENT_CYAN, 1.2)
add_text_box(slide, "핵심 기능 ③ — CleverAI™", 0.4, 0.15, 9, 0.75,
             font_size=28, bold=True, color=WHITE)
add_text_box(slide, "AI · ML · Generative AI · Agentic AI — 4세대 인공지능 인게이지먼트",
             0.4, 0.82, 9, 0.45, font_size=13, color=ACCENT_CYAN)

ai_features = [
    ("🔮", "Predictive AI", "고객 이탈 예측, LTV 예측,\n최적 발송 시간 예측"),
    ("✍️", "Generative AI (Scribe™)", "클릭 한 번으로 캠페인 카피\n자동 생성 — A/B 변형 포함"),
    ("🤖", "Agentic AI", "자율 실행 에이전트가\n24/7 캠페인 최적화"),
    ("🎯", "AI 세그멘테이션", "머신러닝이 고가치 고객군을\n자동 발굴 & 업데이트"),
    ("📈", "예측 분석", "다음 구매 가능성, 이탈 위험도,\n업셀 기회를 실시간 스코어링"),
    ("⚙️", "자동 최적화", "캠페인 성과를 학습하여\n지속적으로 전략 개선"),
]

for i, (icon, title, desc) in enumerate(ai_features):
    col = i % 3
    row = i // 3
    x = 0.3 + col * 3.2
    y = 1.45 + row * 2.25
    add_rect(slide, x, y, 2.9, 2.1, RGBColor(0x1E, 0x35, 0x9E))
    add_rect(slide, x, y, 2.9, 0.08, ACCENT_CYAN)
    add_text_box(slide, icon, x + 0.1, y + 0.2, 2.7, 0.55,
                 font_size=28, align=PP_ALIGN.CENTER, color=ACCENT_CYAN)
    add_text_box(slide, title, x + 0.1, y + 0.72, 2.7, 0.4,
                 font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text_box(slide, desc, x + 0.1, y + 1.15, 2.7, 0.85,
                 font_size=10, color=RGBColor(0xBB, 0xDE, 0xFB), align=PP_ALIGN.CENTER)

add_logo_watermark(slide)

# ════════════════════════════════════════════════════════════════
# 슬라이드 8: 핵심 기능 4 — 리텐션 & 퍼널 최적화
# ════════════════════════════════════════════════════════════════
slide = add_slide(prs)
slide_header(slide, "핵심 기능 ④ — 리텐션 & 생애주기 최적화",
             "고객을 오래, 더 깊이 — 전환율과 리텐션을 동시에 개선")

# 생애주기 단계
stages = [
    ("획득\nAcquire", "신규 가입자\n온보딩 자동화", MID_BLUE),
    ("활성화\nActivate", "첫 핵심행동\n유도 & 가이드", ACCENT_CYAN),
    ("유지\nRetain", "재방문 유도\n이탈 방지 캠페인", RGBColor(0x00, 0x8B, 0x45)),
    ("성장\nGrow", "업셀 & 크로스셀\n수익 극대화", ORANGE),
    ("복구\nWin-back", "이탈 고객\n재활성화", RGBColor(0xD3, 0x2F, 0x2F)),
]

for i, (stage, desc, color) in enumerate(stages):
    x = 0.2 + i * 1.92
    add_rect(slide, x, 1.5, 1.7, 2.8, color)
    add_text_box(slide, stage, x + 0.05, 1.65, 1.6, 0.9,
                 font_size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text_box(slide, desc, x + 0.05, 2.6, 1.6, 1.5,
                 font_size=10, color=WHITE, align=PP_ALIGN.CENTER)
    if i < 4:
        add_text_box(slide, "→", x + 1.7, 2.7, 0.25, 0.4,
                     font_size=18, bold=True, color=DARK_GRAY, align=PP_ALIGN.CENTER)

# 실험/최적화
add_text_box(slide, "실험 & 최적화 (Experimentation)", 0.3, 4.5, 9.4, 0.4,
             font_size=15, bold=True, color=DARK_BLUE)

exp_items = [
    "A/B 테스트 & 멀티배리에이트 테스트 내장",
    "실시간 실험 결과 모니터링",
    "통계적 유의성 자동 판별",
    "위닝 배리언트 자동 적용",
]
for i, item in enumerate(exp_items):
    x = 0.3 + (i % 2) * 4.8
    y = 5.0 + (i // 2) * 0.55
    add_text_box(slide, "✅  " + item, x, y, 4.6, 0.5,
                 font_size=12, color=DARK_GRAY)

add_image_safe(slide, "platform_1", 7.0, 4.4, 2.7, 1.8)

# ════════════════════════════════════════════════════════════════
# 슬라이드 9: 핀테크 적용 사례
# ════════════════════════════════════════════════════════════════
slide = add_slide(prs)
slide_header(slide, "핀테크 & 금융 업계 적용 사례",
             "실제 도입 성과 — 수익 성장부터 운영 효율화까지")

cases = [
    {
        "company": "APLAZO",
        "type": "🏦 핀테크 — BNPL",
        "country": "🇲🇽 멕시코",
        "result1": "50%+",
        "label1": "장바구니 이탈 복구율",
        "result2": "94%",
        "label2": "운영 업무 자동화",
        "desc": "CleverTap의 자동화 캠페인으로 이탈 고객을 적시에 재활성화. 마케팅 팀의 운영 부담 94% 감소.",
        "color": MID_BLUE,
    },
    {
        "company": "M-KOPA",
        "type": "💳 핀테크 — 금융포용",
        "country": "🌍 아프리카",
        "result1": "3x",
        "label1": "전환율 스케일",
        "result2": "60%",
        "label2": "캠페인 실행시간 단축",
        "desc": "개인화 인게이지먼트로 비은행권 고객 전환을 3배 확대. 캠페인 실행 속도 대폭 향상.",
        "color": ACCENT_CYAN,
    },
    {
        "company": "Banco Promerica",
        "type": "🏛️ 은행 — 디지털",
        "country": "🇬🇹 과테말라",
        "result1": "105%",
        "label1": "대출 실행 증가",
        "result2": "2x",
        "label2": "디지털 전환율",
        "desc": "개인화된 금융 상품 추천과 적시 알림으로 대출 실행 105% 증가 달성.",
        "color": RGBColor(0x00, 0x8B, 0x45),
    },
]

for i, case in enumerate(cases):
    x = 0.25 + i * 3.25
    color = case["color"]
    add_rect(slide, x, 1.5, 3.05, 5.3, LIGHT_GRAY)
    add_rect(slide, x, 1.5, 3.05, 0.5, color)
    add_text_box(slide, case["company"], x + 0.1, 1.55, 2.85, 0.4,
                 font_size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text_box(slide, case["type"], x + 0.05, 2.1, 2.95, 0.3,
                 font_size=10, bold=True, color=color, align=PP_ALIGN.CENTER)
    add_text_box(slide, case["country"], x + 0.05, 2.4, 2.95, 0.3,
                 font_size=10, color=DARK_GRAY, align=PP_ALIGN.CENTER)
    # 수치
    add_rect(slide, x + 0.15, 2.8, 2.75, 0.85, color)
    add_text_box(slide, case["result1"], x + 0.15, 2.82, 2.75, 0.55,
                 font_size=28, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text_box(slide, case["label1"], x + 0.15, 3.35, 2.75, 0.3,
                 font_size=9, color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(slide, x + 0.15, 3.75, 2.75, 0.85, RGBColor(0xE0, 0xE0, 0xE0))
    add_text_box(slide, case["result2"], x + 0.15, 3.77, 2.75, 0.55,
                 font_size=28, bold=True, color=color, align=PP_ALIGN.CENTER)
    add_text_box(slide, case["label2"], x + 0.15, 4.32, 2.75, 0.3,
                 font_size=9, color=DARK_GRAY, align=PP_ALIGN.CENTER)
    add_text_box(slide, case["desc"], x + 0.1, 4.7, 2.85, 2.0,
                 font_size=10, color=DARK_GRAY, align=PP_ALIGN.LEFT)

# ════════════════════════════════════════════════════════════════
# 슬라이드 10: 글로벌 고객사
# ════════════════════════════════════════════════════════════════
slide = add_slide(prs)
slide_header(slide, "2,000+ 글로벌 브랜드의 선택",
             "이커머스 · 핀테크 · 미디어 · 통신 · 리테일 전 업종")

customers = [
    "Picsart", "Tata CLIQ Luxury", "Times of India", "Zee5 Global",
    "Vodafone", "Circles.Life", "Banco Promerica", "APLAZO",
    "M-KOPA", "Ooredoo Qatar", "HUFT", "WalaPlus",
    "TMOV", "Eden Red", "Mobile Premier League", "Chakkizza",
]

cols = 4
rows = 4
for i, name in enumerate(customers[:16]):
    col = i % cols
    row = i // cols
    x = 0.3 + col * 2.4
    y = 1.5 + row * 1.38
    add_rect(slide, x, y, 2.2, 1.2, WHITE)
    add_rect(slide, x, y, 2.2, 0.06, MID_BLUE)
    add_text_box(slide, name, x + 0.1, y + 0.3, 2.0, 0.6,
                 font_size=11, bold=True, color=DARK_BLUE, align=PP_ALIGN.CENTER)

# 산업군
industries = ["이커머스", "핀테크", "미디어/엔터", "통신", "리테일", "여행/물류", "HR테크"]
add_rect(slide, 0.3, 7.05, 9.4, 0.35, LIGHT_GRAY)
add_text_box(slide, "주요 산업군:  " + "  |  ".join(industries),
             0.5, 7.08, 9.0, 0.3, font_size=11, color=DARK_GRAY)

add_logo_watermark(slide)

# ════════════════════════════════════════════════════════════════
# 슬라이드 11: 경쟁사 비교
# ════════════════════════════════════════════════════════════════
slide = add_slide(prs)
slide_header(slide, "경쟁사 비교",
             "CleverTap vs 주요 마케팅 자동화 솔루션")

# 테이블 헤더
headers_row = ["기능 / 역량", "CleverTap", "Braze", "MoEngage", "Adobe"]
col_widths = [2.8, 1.7, 1.7, 1.7, 1.7]
x_positions = [0.2, 3.1, 4.85, 6.6, 8.35]

# 헤더 행
add_rect(slide, 0.2, 1.45, 9.6, 0.5, DARK_BLUE)
for j, (h, x, w) in enumerate(zip(headers_row, x_positions, col_widths)):
    add_text_box(slide, h, x + 0.05, 1.5, w - 0.1, 0.4,
                 font_size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# 비교 데이터
comparison = [
    ("올인원 플랫폼",        "✅ 완전 통합",  "⚡ 채널 강세",  "✅ 통합형",  "⚠️ 복잡"),
    ("AI/ML 기능",          "✅ CleverAI™", "✅ 있음",       "✅ 있음",    "✅ Sensei"),
    ("실시간 분석",          "✅ 초당 처리",  "✅ 실시간",    "✅ 실시간",  "⚠️ 지연"),
    ("핀테크 특화",          "✅ 전문 사례", "⚠️ 일반형",    "⚠️ 일반형",  "❌ 엔터급"),
    ("구현 난이도",          "✅ 쉬움",      "⚠️ 중간",      "✅ 쉬움",    "❌ 어려움"),
    ("가격 경쟁력",          "✅ 합리적",    "⚠️ 고가",      "✅ 합리적",  "❌ 고가"),
    ("한국어/아시아 지원",   "✅ 강력",      "⚠️ 제한적",    "✅ 있음",    "⚠️ 제한적"),
]

row_colors = [LIGHT_GRAY, WHITE, LIGHT_GRAY, WHITE, LIGHT_GRAY, WHITE, LIGHT_GRAY]
for r_idx, (row_data) in enumerate(comparison):
    y = 2.0 + r_idx * 0.67
    add_rect(slide, 0.2, y, 9.6, 0.62, row_colors[r_idx])
    for j, (cell, x, w) in enumerate(zip(row_data, x_positions, col_widths)):
        is_clevertap = j == 1
        color = MID_BLUE if is_clevertap else DARK_GRAY
        bold = is_clevertap
        add_text_box(slide, cell, x + 0.05, y + 0.1, w - 0.1, 0.45,
                     font_size=11, bold=bold, color=color, align=PP_ALIGN.CENTER)

# 결론
add_rect(slide, 0.2, 6.7, 9.6, 0.6, ACCENT_CYAN)
add_text_box(slide,
    "🏆  CleverTap은 올인원 통합 · AI 역량 · 합리적 가격 · 아시아 특화 지원에서 경쟁 우위",
    0.4, 6.78, 9.2, 0.45, font_size=12, bold=True, color=DARK_BLUE)

add_logo_watermark(slide)

# ════════════════════════════════════════════════════════════════
# 슬라이드 12: 도입 프로세스
# ════════════════════════════════════════════════════════════════
slide = add_slide(prs)
slide_header(slide, "도입 프로세스 & 기술 통합",
             "빠르고 안전한 온보딩 — 평균 4~6주 완성")

steps = [
    ("01", "킥오프 & 요구사항 분석", "비즈니스 목표 정의\n현재 시스템 분석\nKPI 설정", MID_BLUE),
    ("02", "SDK & 데이터 연동", "iOS/Android/Web SDK 설치\n이벤트 태깅 & 검증\nCRM 데이터 마이그레이션", ACCENT_CYAN),
    ("03", "세그멘테이션 & 캠페인 설계", "초기 세그먼트 정의\n첫 캠페인 구성\nA/B 테스트 설계", RGBColor(0x00, 0x8B, 0x45)),
    ("04", "파일럿 런치 & 검증", "소규모 파일럿 실행\n성과 측정 & 피드백\n최적화 반영", ORANGE),
    ("05", "전체 롤아웃 & 운영", "전사 확대 적용\n팀 교육 & 전문가 지원\n지속적 최적화", DARK_BLUE),
]

for i, (num, title, desc, color) in enumerate(steps):
    x = 0.15 + i * 1.95
    add_rect(slide, x, 1.5, 1.75, 4.0, color)
    add_text_box(slide, num, x + 0.05, 1.6, 1.65, 0.6,
                 font_size=24, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text_box(slide, title, x + 0.05, 2.25, 1.65, 0.85,
                 font_size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text_box(slide, desc, x + 0.05, 3.15, 1.65, 2.2,
                 font_size=10, color=WHITE, align=PP_ALIGN.CENTER)
    if i < 4:
        add_text_box(slide, "→", x + 1.74, 3.2, 0.22, 0.4,
                     font_size=20, bold=True, color=DARK_GRAY)

# 기술 통합 지원
add_rect(slide, 0.2, 5.7, 9.6, 1.55, LIGHT_GRAY)
add_text_box(slide, "🔗 기술 통합 지원", 0.4, 5.8, 9.2, 0.4,
             font_size=14, bold=True, color=DARK_BLUE)
integrations = [
    "✅ REST API & Webhook 완전 지원",
    "✅ iOS / Android / Web / React Native / Flutter SDK",
    "✅ Salesforce, HubSpot, Segment, Amplitude 연동",
    "✅ Google Analytics, BigQuery, Snowflake 연동",
    "✅ AWS, GCP, Azure 클라우드 호환",
    "✅ 전담 기술 지원팀 (한국어 지원)",
]
for i, item in enumerate(integrations):
    x = 0.4 + (i % 2) * 4.8
    y = 6.25 + (i // 2) * 0.33
    add_text_box(slide, item, x, y, 4.7, 0.3,
                 font_size=10, color=DARK_GRAY)

add_logo_watermark(slide)

# ════════════════════════════════════════════════════════════════
# 슬라이드 13: 가격 & 패키지
# ════════════════════════════════════════════════════════════════
slide = add_slide(prs)
slide_header(slide, "가격 정책 & 패키지",
             "비즈니스 규모에 맞는 유연한 플랜")

packages = [
    {
        "name": "Essentials",
        "price": "맞춤 견적",
        "target": "스타트업 · 초기 단계",
        "color": MID_BLUE,
        "features": [
            "핵심 분석 & 세그멘테이션",
            "푸시 · 인앱 · 이메일",
            "기본 A/B 테스트",
            "표준 대시보드",
            "이메일 지원",
        ],
    },
    {
        "name": "Advanced",
        "price": "맞춤 견적",
        "target": "성장 단계 기업",
        "color": DARK_BLUE,
        "features": [
            "모든 Essentials 포함",
            "옴니채널 자동화",
            "CleverAI™ 기본 기능",
            "고급 실험 & 최적화",
            "WhatsApp · SMS · RCS",
            "전담 CSM 지원",
        ],
    },
    {
        "name": "Enterprise",
        "price": "맞춤 견적",
        "target": "대기업 · 금융 · 이커머스",
        "color": ACCENT_CYAN,
        "features": [
            "모든 Advanced 포함",
            "CleverAI™ 풀 기능",
            "자체 데이터센터 옵션",
            "무제한 이벤트 처리",
            "전담 기술지원 & SLA",
            "맞춤형 통합 개발",
        ],
    },
]

for i, pkg in enumerate(packages):
    x = 0.3 + i * 3.25
    add_rect(slide, x, 1.4, 3.05, 5.3, pkg["color"])
    add_text_box(slide, pkg["name"], x + 0.1, 1.5, 2.85, 0.55,
                 font_size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text_box(slide, pkg["target"], x + 0.1, 2.0, 2.85, 0.35,
                 font_size=10, color=RGBColor(0xBB, 0xDE, 0xFB), align=PP_ALIGN.CENTER)
    add_rect(slide, x + 0.15, 2.4, 2.75, 0.7, RGBColor(0xFF, 0xFF, 0xFF))
    add_text_box(slide, pkg["price"], x + 0.15, 2.45, 2.75, 0.55,
                 font_size=16, bold=True, color=pkg["color"], align=PP_ALIGN.CENTER)
    for j, feat in enumerate(pkg["features"]):
        add_text_box(slide, "✓  " + feat, x + 0.15, 3.2 + j * 0.53, 2.75, 0.45,
                     font_size=10, color=WHITE)

# 하단 메모
add_rect(slide, 0.2, 6.9, 9.6, 0.45, LIGHT_GRAY)
add_text_box(slide,
    "💡  모든 플랜은 사용량 기반 맞춤 견적 | 연간 계약 할인 제공 | 30일 POC(파일럿) 옵션 가능",
    0.4, 6.95, 9.2, 0.35, font_size=10, color=DARK_GRAY)

add_logo_watermark(slide)

# ════════════════════════════════════════════════════════════════
# 슬라이드 14: 다음 단계 & CTA
# ════════════════════════════════════════════════════════════════
slide = add_slide(prs)
fill_bg(slide, DARK_BLUE)
add_rect(slide, 0, 0, 10, 0.08, ACCENT_CYAN)

add_text_box(slide, "지금 시작하세요", 0.5, 0.5, 9, 0.8,
             font_size=38, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text_box(slide, "CleverTap으로 고객 생애가치를 무한히 성장시키세요",
             0.5, 1.25, 9, 0.5, font_size=16, color=ACCENT_CYAN, align=PP_ALIGN.CENTER)

ctas = [
    ("🚀", "무료 데모 신청", "30분 맞춤 데모로\n플랫폼을 직접 확인하세요"),
    ("📊", "POC 파일럿", "30일 무료 파일럿으로\n실제 성과를 검증하세요"),
    ("💬", "기술 컨설팅", "아키텍처 & 통합 방안을\n전문가와 논의하세요"),
    ("📋", "맞춤 견적", "비즈니스 규모에 맞는\n최적 플랜을 제안받으세요"),
]

for i, (icon, title, desc) in enumerate(ctas):
    x = 0.25 + i * 2.4
    add_rect(slide, x, 2.0, 2.2, 2.5, RGBColor(0x1E, 0x35, 0x9E))
    add_rect(slide, x, 2.0, 2.2, 0.07, ACCENT_CYAN)
    add_text_box(slide, icon, x + 0.05, 2.15, 2.1, 0.6,
                 font_size=30, align=PP_ALIGN.CENTER, color=ACCENT_CYAN)
    add_text_box(slide, title, x + 0.05, 2.78, 2.1, 0.4,
                 font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text_box(slide, desc, x + 0.05, 3.22, 2.1, 1.1,
                 font_size=10, color=RGBColor(0xBB, 0xDE, 0xFB), align=PP_ALIGN.CENTER)

# 연락처
add_rect(slide, 0.3, 4.75, 9.4, 1.5, RGBColor(0x0D, 0x1A, 0x6E))
add_text_box(slide, "📞 문의하기", 0.5, 4.9, 9.0, 0.4,
             font_size=15, bold=True, color=ACCENT_CYAN, align=PP_ALIGN.CENTER)
add_text_box(slide, "🌐 clevertap.com  |  📧 contact@clevertap.com  |  데모 신청: clevertap.com/request-demo",
             0.5, 5.35, 9.0, 0.4, font_size=12, color=WHITE, align=PP_ALIGN.CENTER)
add_text_box(slide, "한국 파트너: 도입 문의 및 기술 지원은 현지 파트너를 통해 제공됩니다",
             0.5, 5.82, 9.0, 0.35, font_size=10, color=RGBColor(0x90, 0xCA, 0xF9), align=PP_ALIGN.CENTER)

add_logo_watermark(slide)

# ════════════════════════════════════════════════════════════════
# 슬라이드 15: Q&A / 마무리
# ════════════════════════════════════════════════════════════════
slide = add_slide(prs)
fill_bg(slide, WHITE)
add_rect(slide, 0, 0, 10, 7.5, WHITE)
add_rect(slide, 0, 0, 0.4, 7.5, DARK_BLUE)
add_rect(slide, 0.4, 0, 9.6, 0.06, ACCENT_CYAN)

add_text_box(slide, "Q & A", 1.0, 1.5, 8, 1.2,
             font_size=68, bold=True, color=DARK_BLUE, align=PP_ALIGN.CENTER)
add_text_box(slide, "질문이 있으신가요?", 1.0, 2.8, 8, 0.6,
             font_size=24, color=ACCENT_CYAN, align=PP_ALIGN.CENTER)
add_text_box(slide, "CleverTap과 함께라면, 고객 인게이지먼트의 한계는 없습니다.",
             1.0, 3.5, 8, 0.5, font_size=14, color=DARK_GRAY, align=PP_ALIGN.CENTER)

# 요약 수치
summary_stats = [
    ("2,000+", "글로벌 브랜드"),
    ("5x", "리텐션 향상"),
    ("60%", "CTR 증가"),
    ("94%", "운영 자동화"),
]
for i, (val, lbl) in enumerate(summary_stats):
    x = 0.8 + i * 2.2
    add_rect(slide, x, 4.3, 1.9, 1.3, DARK_BLUE)
    add_text_box(slide, val, x, 4.35, 1.9, 0.7,
                 font_size=26, bold=True, color=ACCENT_CYAN, align=PP_ALIGN.CENTER)
    add_text_box(slide, lbl, x, 5.05, 1.9, 0.4,
                 font_size=10, color=WHITE, align=PP_ALIGN.CENTER)

add_text_box(slide, "clevertap.com", 1.0, 5.9, 8, 0.4,
             font_size=13, color=MID_BLUE, align=PP_ALIGN.CENTER)
add_text_box(slide, "© 2025 CleverTap Inc. All rights reserved.",
             1.0, 7.1, 8, 0.3, font_size=9, color=RGBColor(0xBB, 0xBB, 0xBB), align=PP_ALIGN.CENTER)

add_logo_watermark(slide)

# ─── 저장 ────────────────────────────────────────────────────────────────────
prs.save(str(OUTPUT_FILE))
print(f"\n✅ PPTX 저장 완료!")
print(f"   경로: {OUTPUT_FILE}")
print(f"   슬라이드 수: {len(prs.slides)}장")
print(f"   활용 이미지: {len(downloaded)}개")
